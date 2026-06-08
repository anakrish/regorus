#!/usr/bin/env python3
"""
Generate an editable PowerPoint deck: AWS Policy Analysis — 8-Year Comparison.
All content is native PowerPoint shapes/tables — fully editable, no images.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from lxml import etree
import os

# ── Color palette ──────────────────────────────────────────────────────────
AWS_ORANGE     = RGBColor(0xFF, 0x99, 0x00)
AWS_ORANGE_LT  = RGBColor(0xFF, 0xE8, 0xCC)
AWS_ORANGE_BG  = RGBColor(0xFF, 0xF5, 0xE6)
MS_BLUE        = RGBColor(0x00, 0x78, 0xD4)
MS_BLUE_LT     = RGBColor(0xCC, 0xE4, 0xF7)
MS_BLUE_BG     = RGBColor(0xE6, 0xF2, 0xFB)
DARK_BLUE      = RGBColor(0x0D, 0x47, 0xA1)
GREEN          = RGBColor(0x38, 0x8E, 0x3C)
GREEN_LT       = RGBColor(0xC8, 0xE6, 0xC9)
GREEN_BG       = RGBColor(0xE8, 0xF5, 0xE9)
RED            = RGBColor(0xE5, 0x39, 0x35)
RED_LT         = RGBColor(0xFF, 0xCD, 0xD2)
RED_BG         = RGBColor(0xFF, 0xEB, 0xEE)
PURPLE         = RGBColor(0x7B, 0x1F, 0xA2)
PURPLE_LT      = RGBColor(0xE1, 0xBE, 0xE7)
PURPLE_BG      = RGBColor(0xF3, 0xE5, 0xF5)
YELLOW         = RGBColor(0xF9, 0xA8, 0x25)
YELLOW_LT      = RGBColor(0xFF, 0xF9, 0xC4)
TEAL           = RGBColor(0x00, 0x89, 0x7B)
TEAL_LT        = RGBColor(0xB2, 0xDF, 0xDB)
GRAY           = RGBColor(0x60, 0x60, 0x60)
GRAY_LT        = RGBColor(0xE0, 0xE0, 0xE0)
GRAY_BG        = RGBColor(0xF5, 0xF5, 0xF5)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
DARK           = RGBColor(0x1A, 0x1A, 0x2E)
SLIDE_BG       = RGBColor(0xF8, 0xF9, 0xFA)
AMBER_BG       = RGBColor(0xFF, 0xF8, 0xE1)
AMBER          = RGBColor(0xFF, 0x8F, 0x00)

# ── Helpers ────────────────────────────────────────────────────────────────

def set_slide_bg(slide, color=SLIDE_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, title, subtitle=None, bg_color=DARK_BLUE):
    shapes = slide.shapes
    bar = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                           Inches(13.333), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = bg_color
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
        p2.alignment = PP_ALIGN.LEFT


def add_box(slide, left, top, width, height, text, fill_color, border_color,
            font_size=12, bold=False, font_color=DARK,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE, align=PP_ALIGN.CENTER):
    s = slide.shapes.add_shape(shape, Inches(left), Inches(top),
                               Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    s.line.color.rgb = border_color
    s.line.width = Pt(2)
    s.shadow.inherit = False
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = font_color
        p.alignment = align
    tf.auto_size = None
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
    return s


def add_label(slide, left, top, width, text, font_size=10, color=GRAY,
              bold=False, align=PP_ALIGN.CENTER, height=0.3):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox


def add_multiline_label(slide, left, top, width, height, lines, font_size=10,
                        color=DARK, bold=False, align=PP_ALIGN.LEFT):
    """Add a textbox with multiple lines."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
        p.space_after = Pt(3)
    return txBox


def add_arrow(slide, start_left, start_top, end_left, end_top, color=GRAY):
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(start_left), Inches(start_top),
        Inches(end_left), Inches(end_top)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    ln = connector._element.find('.//a:ln', nsmap)
    if ln is None:
        spPr = connector._element.find('.//a:spPr', nsmap) or connector._element
        ln = etree.SubElement(spPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}ln')
    tail = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd')
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('len', 'med')
    return connector


def add_group_box(slide, left, top, width, height, title, fill_color,
                  border_color, title_color=DARK):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
                               Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    s.line.color.rgb = border_color
    s.line.width = Pt(2)
    add_label(slide, left + 0.15, top + 0.08, width - 0.3, title,
              font_size=11, color=title_color, bold=True, align=PP_ALIGN.LEFT)
    return s


def add_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def add_table(slide, left, top, width, row_height, headers, rows,
              header_fill=DARK_BLUE, header_font_color=WHITE,
              row_fill_even=WHITE, row_fill_odd=GRAY_BG,
              font_size=10, col_widths=None):
    """Add a styled table with alternating row colors."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols,
                                         Inches(left), Inches(top),
                                         Inches(width), Inches(row_height * n_rows))
    table = table_shape.table

    # Column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    else:
        cw = width / n_cols
        for i in range(n_cols):
            table.columns[i].width = Inches(cw)

    # Style header
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.bold = True
            paragraph.font.color.rgb = header_font_color
            paragraph.alignment = PP_ALIGN.LEFT
        _set_cell_fill(cell, header_fill)
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.05)
        cell.margin_top = Inches(0.03)
        cell.margin_bottom = Inches(0.03)

    # Style rows
    for i, row in enumerate(rows):
        fill = row_fill_even if i % 2 == 0 else row_fill_odd
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.color.rgb = DARK
                paragraph.alignment = PP_ALIGN.LEFT
            _set_cell_fill(cell, fill)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)

    return table_shape


def _set_cell_fill(cell, color):
    """Set table cell fill color via XML."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = etree.SubElement(tcPr,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    srgb = etree.SubElement(solidFill,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
    srgb.set('val', str(color))


def add_big_number(slide, left, top, number, label, number_color=DARK_BLUE,
                   number_size=36, label_size=11):
    """Add a large stat number with descriptive label below."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(2.6), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(number_size)
    p.font.bold = True
    p.font.color.rgb = number_color
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(label_size)
    p2.font.color.rgb = GRAY
    p2.alignment = PP_ALIGN.CENTER
    return txBox


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    """Slide 0: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BLUE)

    add_label(slide, 0.8, 1.2, 11.7, "AWS Formal Policy Analysis",
              font_size=40, color=WHITE, bold=True, align=PP_ALIGN.LEFT)
    add_label(slide, 0.8, 2.0, 11.7, "8-Year Investment Timeline & Comparison with Policy Intelligence",
              font_size=22, color=RGBColor(0xBB, 0xDE, 0xFB), bold=False, align=PP_ALIGN.LEFT)

    # Divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.8), Inches(2.8), Inches(4.0), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = AWS_ORANGE
    line.line.fill.background()

    add_label(slide, 0.8, 3.2, 8.0,
              "Competitive Intelligence  •  February 2026  •  Confidential",
              font_size=14, color=RGBColor(0x90, 0xCA, 0xF9), align=PP_ALIGN.LEFT)

    # Side boxes with key stats
    y = 4.2
    add_box(slide, 0.8, y, 3.5, 0.7, "~1 Billion SMT Queries/Day", AWS_ORANGE_LT, AWS_ORANGE,
            font_size=14, bold=True, font_color=DARK)
    add_box(slide, 4.6, y, 3.5, 0.7, "8+ Years of Production Use", AWS_ORANGE_LT, AWS_ORANGE,
            font_size=14, bold=True, font_color=DARK)
    add_box(slide, 8.4, y, 3.5, 0.7, "8+ Named Enterprise Customers", AWS_ORANGE_LT, AWS_ORANGE,
            font_size=14, bold=True, font_color=DARK)

    add_label(slide, 0.8, 5.5, 6.0,
              "Data from publicly available AWS sources only.",
              font_size=11, color=RGBColor(0x90, 0xCA, 0xF9), align=PP_ALIGN.LEFT)

    add_notes(slide, """TITLE SLIDE — AWS Formal Policy Analysis: 8-Year Investment & Comparison
• All data sourced from public AWS blogs, product pages, and Amazon Science publications.
• This deck maps AWS's investment timeline, products, impact, and where MS Policy Intelligence differs.
• Key framing: AWS has 8 years of production-hardened formal analysis; MS has novel prototype capabilities.""")


def slide_exec_summary(prs):
    """Slide 1: Executive Summary."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Executive Summary",
                  "AWS's ~8-Year Investment in Formal Policy Analysis")

    # AWS investment summary — left column
    add_group_box(slide, 0.4, 1.4, 6.2, 3.3, "AWS Investment (Production)",
                  AWS_ORANGE_BG, AWS_ORANGE, title_color=DARK)

    items_aws = [
        "• Automated Reasoning Group (ARG) — dedicated research org",
        "• Zelkova engine: SMT-based policy analysis since ~2017",
        "• IAM Access Analyzer: GA since Dec 2019 (5 features)",
        "• Cedar language: open source May 2023 (Apache 2.0)",
        "• Verified Permissions: managed Cedar-as-a-service",
        "• \"Provable Security\" brand with customer testimonials",
        "• ~1 billion SMT queries/day across AWS services",
        "• Academic: CAV 2022 keynote, OOPSLA 2023 paper"
    ]
    add_multiline_label(slide, 0.6, 1.85, 5.8, 2.7, items_aws,
                        font_size=11, color=DARK)

    # MS status — right column
    add_group_box(slide, 6.9, 1.4, 6.0, 3.3, "Microsoft Policy Intelligence (Prototype)",
                  MS_BLUE_BG, MS_BLUE, title_color=DARK)

    items_ms = [
        "• Regorus Rego engine: production-grade, open source (MIT)",
        "• Z3 analysis: ~10K lines, working CLI (feature branch)",
        "• Novel: policy diff, subsumption, test gen, MC/DC",
        "• Novel: dead rule detection, line-targeted coverage",
        "• Supports Rego + Cedar + Azure Policy (multi-language)",
        "• Verus verification: design doc only (not yet started)",
        "• No production deployment; no public customers",
        "• No academic publications yet"
    ]
    add_multiline_label(slide, 7.1, 1.85, 5.6, 2.7, items_ms,
                        font_size=11, color=DARK)

    # Maturity callout
    add_box(slide, 0.4, 5.0, 12.5, 0.7,
            "⚠  Maturity gap: AWS features are production-grade (GA, 8 years, 1B queries/day). "
            "MS Z3 features are working prototypes on a feature branch — not production-shipped.",
            AMBER_BG, AMBER, font_size=11, font_color=DARK,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE, align=PP_ALIGN.LEFT)

    add_notes(slide, """EXECUTIVE SUMMARY
• AWS has invested ~8 years building formal analysis — from internal Zelkova to production services.
• Microsoft has novel capabilities (diff, subsumes, gen-tests) but they're prototypes.
• The Regorus Rego engine itself is production-grade and already used in Azure.
• Key honest framing: maturity gap is real but MS has capabilities AWS doesn't have at all.""")


def slide_timeline(prs):
    """Slide 2: AWS Timeline — visual horizontal timeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "AWS Timeline — 8 Years of Formal Policy Analysis",
                  "Major milestones from ~2017 to 2026")

    # Horizontal timeline axis
    axis_y = 3.5
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.6), Inches(axis_y), Inches(12.1), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = GRAY
    line.line.fill.background()

    # Timeline events — alternating above/below
    events = [
        ("~2017",    "Zelkova\n(internal)", True),
        ("2018",     "S3 Block\nPublic Access", False),
        ("2019",     "IAM Access\nAnalyzer GA", True),
        ("2021",     "Policy Validation\n& Generation", False),
        ("2022",     "1B SMT\nQueries/Day", True),
        ("2023",     "Cedar Open\nSourced", False),
        ("2023",     "Verified\nPermissions GA", True),
        ("2023",     "OOPSLA\nPaper", False),
        ("2025-26",  "Cedar 4.x\n+ cvc5 collab", True),
    ]

    n = len(events)
    start_x = 0.8
    spacing = 12.0 / n

    for i, (year, label, above) in enumerate(events):
        cx = start_x + i * spacing
        # Dot on axis
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(cx + 0.35), Inches(axis_y - 0.07),
                                     Inches(0.18), Inches(0.18))
        dot.fill.solid()
        dot.fill.fore_color.rgb = AWS_ORANGE
        dot.line.fill.background()

        if above:
            # Event box above
            box_y = 1.55
            box_h = 1.0
            add_box(slide, cx, box_y, 1.2, box_h, label,
                    AWS_ORANGE_LT, AWS_ORANGE, font_size=9, font_color=DARK)
            # Year label
            add_label(slide, cx, box_y - 0.3, 1.2, year,
                      font_size=9, color=AWS_ORANGE, bold=True)
            # Connector line down to axis
            vline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                           Inches(cx + 0.58), Inches(box_y + box_h),
                                           Inches(0.02), Inches(axis_y - box_y - box_h))
            vline.fill.solid()
            vline.fill.fore_color.rgb = AWS_ORANGE
            vline.line.fill.background()
        else:
            # Event box below
            box_y = 3.9
            box_h = 1.0
            add_box(slide, cx, box_y, 1.2, box_h, label,
                    AWS_ORANGE_BG, AWS_ORANGE, font_size=9, font_color=DARK)
            # Year label
            add_label(slide, cx, box_y + box_h + 0.02, 1.2, year,
                      font_size=9, color=AWS_ORANGE, bold=True)
            # Connector line from axis down
            vline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                           Inches(cx + 0.58), Inches(axis_y + 0.04),
                                           Inches(0.02), Inches(box_y - axis_y - 0.04))
            vline.fill.solid()
            vline.fill.fore_color.rgb = AWS_ORANGE
            vline.line.fill.background()

    add_label(slide, 0.4, 5.8, 12.5,
              "Source: AWS Security Blog, Amazon Science, AWS product pages  •  All dates from public sources",
              font_size=9, color=GRAY, align=PP_ALIGN.LEFT)

    add_notes(slide, """TIMELINE
• Zelkova was first disclosed in a June 2018 blog; believed to have started ~2017.
• S3 Block Public Access (Nov 2018) was the first customer-facing use — it runs in the critical path.
• IAM Access Analyzer launched at re:Invent 2019 — the first service customers interact with directly.
• 2021: Policy validation (100+ checks) and policy generation (from CloudTrail).
• 2022: Neha Rungta's CAV keynote revealed ~1B SMT queries/day — that was 4 years ago.
• 2023: Cedar open-sourced + Verified Permissions GA + OOPSLA paper — a big year.
• Ongoing: Active Cedar 4.x releases, Stanford cvc5 collaboration.""")


def slide_aws_products_zelkova(prs):
    """Slide 3: Zelkova — the engine."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "AWS Product: Zelkova",
                  "The internal SMT engine powering provable security")

    # Central architecture diagram
    # IAM/S3 Policy → Zelkova → Portfolio Solver → Result
    add_box(slide, 0.6, 1.8, 2.4, 1.2,
            "IAM / S3 / Resource\nPolicies",
            GRAY_BG, GRAY, font_size=13, font_color=DARK)

    add_arrow(slide, 3.1, 2.4, 3.9, 2.4, AWS_ORANGE)

    add_box(slide, 4.0, 1.6, 3.0, 1.6,
            "Zelkova Engine\n\nPolicy → Math Formula\n→ SMT Constraints",
            AWS_ORANGE_LT, AWS_ORANGE, font_size=12, bold=False, font_color=DARK)

    add_arrow(slide, 7.1, 2.4, 7.9, 2.4, AWS_ORANGE)

    # Portfolio solver box
    add_group_box(slide, 8.0, 1.5, 4.6, 1.8, "Portfolio Solver",
                  AWS_ORANGE_BG, AWS_ORANGE)
    solver_y = 2.05
    add_box(slide, 8.2, solver_y, 1.0, 0.55, "Z3", WHITE, GRAY, font_size=11, bold=True)
    add_box(slide, 9.3, solver_y, 1.0, 0.55, "CVC4", WHITE, GRAY, font_size=11, bold=True)
    add_box(slide, 10.4, solver_y, 1.0, 0.55, "cvc5", WHITE, GRAY, font_size=11, bold=True)
    add_box(slide, 11.5, solver_y, 1.0, 0.55, "Custom\nAutomata", WHITE, GRAY, font_size=9, bold=True)
    add_label(slide, 8.2, 2.7, 4.2,
              "Winner-take-all: first solver to answer wins",
              font_size=9, color=GRAY, align=PP_ALIGN.LEFT)

    # Services powered by Zelkova
    add_label(slide, 0.6, 3.65, 12.0, "Services Powered by Zelkova",
              font_size=14, color=DARK, bold=True, align=PP_ALIGN.LEFT)

    svc_y = 4.1
    services = [
        ("S3 Block\nPublic Access", "Critical path\nof every S3 put"),
        ("IAM Access\nAnalyzer", "External access\nfinding + validation"),
        ("AWS Config\nRules", "s3-bucket-*\nmanaged rules"),
        ("VPC Network\nAnalyzer", "Network\nreachability proofs"),
        ("Amazon\nMacie", "S3 accessibility\nclassification"),
    ]
    sx = 0.6
    for name, desc in services:
        add_box(slide, sx, svc_y, 2.2, 0.7, name, WHITE, AWS_ORANGE,
                font_size=11, bold=True, font_color=DARK)
        add_label(slide, sx, svc_y + 0.75, 2.2, desc,
                  font_size=9, color=GRAY, align=PP_ALIGN.CENTER)
        sx += 2.5

    # Big stat
    add_box(slide, 9.0, 4.1, 3.8, 1.3,
            "~1 Billion\nSMT Queries / Day\n(CAV 2022 keynote)",
            AWS_ORANGE_LT, AWS_ORANGE, font_size=16, bold=True, font_color=DARK)

    add_notes(slide, """ZELKOVA
• Internal engine — not exposed directly to customers.
• Key insight: "services ask questions on behalf of customers" — users don't run SMT queries.
• Portfolio solver approach: run Z3, CVC4, cvc5, and custom automata solver in parallel; first answer wins.
• Upgrade challenge (CAV talk): moving from CVC4→cvc5 caused regressions; they kept both.
• The ~1B figure is from 2022 — likely higher now.
• Services powered: S3 Block Public Access runs Zelkova synchronously in the request path.""")


def slide_aws_products_iam_aa(prs):
    """Slide 4: IAM Access Analyzer features."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "AWS Product: IAM Access Analyzer",
                  "Customer-facing service since Dec 2019 — 5 major features")

    # 5 feature boxes
    features = [
        ("External Access\nFindings", "2019", "S3, IAM roles, KMS,\nLambda, SQS shared\noutside account/org"),
        ("Policy\nValidation", "2021", "100+ authoring-time\nchecks for security\nanti-patterns"),
        ("Policy\nGeneration", "2021", "CloudTrail-based\nleast-privilege\npolicy suggestions"),
        ("Custom Policy\nChecks", "2023", "CI/CD integration —\nautomated reasoning\nvalidates custom rules"),
        ("Unused Access\nAnalyzer", "2023", "Org-wide detection\nof unused roles,\npermissions, keys"),
    ]

    fx = 0.4
    for name, year, desc in features:
        add_box(slide, fx, 1.5, 2.35, 0.8, name, AWS_ORANGE_LT, AWS_ORANGE,
                font_size=12, bold=True, font_color=DARK)
        add_label(slide, fx, 2.35, 2.35, f"Launched {year}",
                  font_size=9, color=AWS_ORANGE, bold=True)
        add_box(slide, fx, 2.7, 2.35, 1.1, desc, WHITE, GRAY_LT,
                font_size=10, font_color=DARK, align=PP_ALIGN.LEFT)
        fx += 2.55

    # Customer quotes
    add_label(slide, 0.4, 4.2, 12.0, "Customer Impact (Public Quotes)",
              font_size=14, color=DARK, bold=True, align=PP_ALIGN.LEFT)

    quotes = [
        ("USAA", "Uses Access Analyzer for\ndata perimeter strategy"),
        ("GoTo Technologies", "\"Reduced processing time\nfrom days to minutes\"\n(custom policy checks in CI/CD)"),
        ("Attentive", "Cleaned up unused roles\nand permissions at scale"),
    ]
    qx = 0.4
    for company, quote in quotes:
        add_box(slide, qx, 4.65, 3.8, 1.0, f"{company}\n{quote}",
                GREEN_BG, GREEN, font_size=10, font_color=DARK, align=PP_ALIGN.LEFT)
        qx += 4.1

    add_notes(slide, """IAM ACCESS ANALYZER
• This is the primary customer-facing automated-reasoning service.
• External access findings: FREE (included in IAM) — finds resources shared outside org boundary.
• Policy validation: 100+ checks at authoring time — free.
• Policy generation: Analyzes CloudTrail to suggest least-privilege policies — free.
• Custom policy checks: Paid feature for CI/CD pipelines; GoTo's quote is very strong.
• Unused access analyzer: Paid per analyzer; finds unused roles/permissions across the org.""")


def slide_aws_products_cedar(prs):
    """Slide 5: Cedar + Verified Permissions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "AWS Products: Cedar + Amazon Verified Permissions",
                  "Purpose-built policy language + managed authorization service")

    # Left: Cedar
    add_group_box(slide, 0.4, 1.4, 6.0, 2.8, "Cedar Language (Open Source, May 2023)",
                  GREEN_BG, GREEN)

    cedar_items = [
        "• Purpose-built for authorization (RBAC + ABAC)",
        "• Designed for analysis: no loops, no UDFs, decidable",
        "• Apache 2.0 on GitHub — ~5,700+ stars",
        "• Verification-guided development:",
        "    — Dafny formal model of authorizer",
        "    — 100M differential random tests/night",
        "    — Proven: explicit-permit, forbid-overrides-permit",
        "• OOPSLA 2023 publication"
    ]
    add_multiline_label(slide, 0.6, 1.9, 5.6, 2.2, cedar_items,
                        font_size=11, color=DARK)

    # Right: Verified Permissions
    add_group_box(slide, 6.7, 1.4, 6.2, 2.8, "Amazon Verified Permissions (GA, May 2023)",
                  PURPLE_BG, PURPLE)

    avp_items = [
        "• Managed Cedar-as-a-service (authorization)",
        "• Central policy store, real-time eval, audit logging",
        "• Schema validation for entities & actions",
        "• Target: app developers needing fine-grained authz",
        "• Not just cloud IAM — application-level authorization"
    ]
    add_multiline_label(slide, 6.9, 1.9, 5.8, 2.0, avp_items,
                        font_size=11, color=DARK)

    # AVP Customers
    add_label(slide, 0.4, 4.5, 12.0, "Named Customers (Public Testimonials)",
              font_size=14, color=DARK, bold=True, align=PP_ALIGN.LEFT)

    customers = [
        ("TELUS", "IoT/smart home\nauthorization"),
        ("Stedi", "700M B2B EDI\ntransactions/mo"),
        ("Twilio Flex", "Contact center\nauthorization"),
        ("FIS", "$50T annual\ntransactions"),
        ("Grosvenor Eng", "1.5B building\nassets managed"),
    ]
    cx = 0.4
    for name, desc in customers:
        add_box(slide, cx, 4.95, 2.2, 0.75, f"{name}\n{desc}",
                WHITE, PURPLE, font_size=10, font_color=DARK)
        cx += 2.5

    add_notes(slide, """CEDAR + VERIFIED PERMISSIONS
• Cedar was intentionally designed to be analyzable — no loops, no user-defined functions.
• This makes Cedar's analysis decidable, but limits expressiveness vs Rego.
• Verification-guided development: they write Dafny formal models FIRST, then implement in Rust.
• 100M differential random tests per night compare Dafny model vs Rust implementation.
• Verified Permissions is the managed service — customers write Cedar, AWS manages infra.
• Customer names are from AWS product pages — all public.""")


def slide_impact_numbers(prs):
    """Slide 6: Impact — By the Numbers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "AWS Impact — By the Numbers",
                  "Scale, testing, and enterprise adoption")

    # Big number grid — 2 rows of 4
    stats = [
        ("~1B", "SMT Queries/Day", AWS_ORANGE),
        ("4", "Solvers in Portfolio", AWS_ORANGE),
        ("100M+", "DRT Tests/Night", GREEN),
        ("8+", "Named Customers", PURPLE),
        ("100+", "Policy Validation\nChecks", DARK_BLUE),
        ("700M", "Stedi req/month", DARK_BLUE),
        ("$50T", "FIS Annual\nTransactions", DARK_BLUE),
        ("All", "Commercial AWS\nRegions", DARK_BLUE),
    ]

    positions = [
        (0.6, 1.5), (3.4, 1.5), (6.2, 1.5), (9.8, 1.5),
        (0.6, 3.4), (3.4, 3.4), (6.2, 3.4), (9.8, 3.4),
    ]

    for (num, label, color), (x, y) in zip(stats, positions):
        add_box(slide, x, y, 2.8, 1.4, "", WHITE, GRAY_LT, font_size=10)
        add_label(slide, x, y + 0.1, 2.8, num,
                  font_size=32, color=color, bold=True)
        add_label(slide, x, y + 0.75, 2.8, label,
                  font_size=11, color=GRAY, height=0.5)

    # Source note
    add_label(slide, 0.4, 5.5, 12.0,
              "Sources: CAV 2022 keynote, Amazon Science blog, AWS product pages, GitHub",
              font_size=9, color=GRAY, align=PP_ALIGN.LEFT)

    add_notes(slide, """IMPACT BY THE NUMBERS
• The 1B queries/day figure is from Neha Rungta's CAV 2022 keynote — likely higher now.
• Portfolio solver: Z3 + CVC4 + cvc5 + custom automata solver; winner-take-all.
• 100M DRT tests/night: differential random testing comparing Dafny model vs Rust Cedar implementation.
• Named customers from product pages: USAA, Bridgewater, GoTo, TELUS, Stedi, Twilio, FIS, Grosvenor, Attentive.""")


def slide_comparison_aws_leads(prs):
    """Slide 7: Head-to-Head — Where AWS Leads."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Head-to-Head: Where AWS Leads",
                  "Production-grade capabilities with years of hardening")

    headers = ["Capability", "AWS", "Microsoft PI", "Status"]
    rows = [
        ["Public access detection", "✅ S3 Block Public Access", "N/A (different use case)", "AWS production"],
        ["Cross-account access detection", "✅ IAM Access Analyzer", "N/A (different use case)", "AWS production"],
        ["Policy validation (authoring)", "✅ 100+ built-in checks", "Partial (via analysis)", "AWS production"],
        ["Policy generation from usage", "✅ CloudTrail-based", "❌ Not available", "AWS production"],
        ["Unused access detection", "✅ Org-wide detection", "❌ Not available", "AWS production"],
        ["Custom CI/CD policy checks", "✅ Integrated in IAM AA", "Possible via CLI", "AWS integrated"],
        ["Managed authz service", "✅ Verified Permissions", "❌ Not a service", "AWS GA"],
        ["SMT solver diversity", "Portfolio (4 solvers)", "Z3 only", "AWS production"],
        ["Engine verification", "✅ Dafny + 100M DRT/night", "Planned (Verus design doc)", "AWS shipping"],
        ["Production scale", "~1B queries/day, 8 years", "Feature branch prototype", "AWS production"],
        ["Enterprise adoption", "8+ named public customers", "Internal Azure use only", "AWS ahead"],
        ["Academic publications", "CAV, OOPSLA, Amazon Sci", "None yet", "AWS ahead"],
        ["Marketing/branding", "\"Provable Security\" brand", "Not yet branded", "AWS ahead"],
    ]

    add_table(slide, 0.3, 1.3, 12.7, 0.35, headers, rows,
              header_fill=AWS_ORANGE, col_widths=[2.8, 3.5, 3.2, 1.5],
              font_size=9)

    add_notes(slide, """WHERE AWS LEADS
• AWS's lead is structural: 8 years, billions of queries, production-hardened.
• IAM posture management (access analyzer, policy validation, generation) is AWS's home turf.
• The "Provable Security" brand is established — customers search for it.
• Engine verification: AWS ships Dafny proofs of Cedar today; MS Verus is a design doc.
• Key message: don't compete directly here — compete on the capabilities AWS LACKS.""")


def slide_comparison_ms_leads(prs):
    """Slide 8: Head-to-Head — Where MS Policy Intelligence Is Novel."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Head-to-Head: Where MS Has Novel Capabilities",
                  "Capabilities AWS doesn't offer — all currently prototype stage")

    headers = ["Capability", "AWS", "Microsoft PI", "MS Maturity"]
    rows = [
        ["Policy diff (semantic)", "❌ Not available", "regorus diff — proves equivalence\nor finds distinguishing input", "Prototype\n(working CLI)"],
        ["Policy subsumption", "❌ Not available", "regorus subsumes — proves\nnew ⊇ old or counterexample", "Prototype\n(working CLI)"],
        ["Automatic test generation", "❌ Not available", "regorus gen-tests — full\npath-coverage test suite", "Prototype\n(working CLI)"],
        ["MC/DC condition coverage", "❌ Not available", "Aviation-grade condition\ncoverage (--condition-coverage)", "Prototype\n(working CLI)"],
        ["Line-targeted coverage", "❌ Not available", "--cover-line / --avoid-line\nconstraints", "Prototype\n(working CLI)"],
        ["Dead rule detection", "❌ Not available", "UNSAT path condition =\nprovably unreachable code", "Prototype"],
        ["Input synthesis", "Limited\n(simulation)", "Full Z3-driven counterexample\ngeneration with JSON output", "Prototype\n(working CLI)"],
        ["SMT dump/inspection", "Not exposed", "--dump-smt for full\ntransparency", "Prototype\n(working CLI)"],
        ["Rego support", "❌ (Cedar only)", "Full Rego: loops, UDFs,\npartial rules (bounded)", "Prototype\n(Z3 analysis)"],
        ["Why denied (root cause)", "Limited\n(simulation)", "MAX-SAT approach designed", "Planned\n(not impl.)"],
    ]

    add_table(slide, 0.3, 1.3, 12.7, 0.38, headers, rows,
              header_fill=MS_BLUE, col_widths=[2.5, 2.0, 3.8, 1.6],
              font_size=9)

    # Caveat bar at bottom
    add_box(slide, 0.3, 6.0, 12.7, 0.45,
            "⚠  All MS capabilities above are prototypes on a feature branch (z3-redux). "
            "Working CLI demos exist but no production deployment.",
            AMBER_BG, AMBER, font_size=10, font_color=DARK,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE, align=PP_ALIGN.LEFT)

    add_notes(slide, """WHERE MS HAS NOVEL CAPABILITIES
• These are capabilities that AWS simply does not offer today — real differentiation.
• BUT: they are all prototypes, not production. Honest framing is critical for credibility.
• Policy diff & subsumption: most compelling for safe policy migration scenarios.
• MC/DC condition coverage: unique — borrowed from aviation safety standards (DO-178C).
• Dead rule detection: emerges naturally from UNSAT path analysis — no extra work.
• Rego support is the strategic wedge — OPA/Rego is the industry standard; Cedar is growing but smaller.
• Why denied (MAX-SAT): designed but not yet implemented — don't claim it.""")


def slide_comparison_language(prs):
    """Slide 9: Language & Expressiveness Comparison."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Language & Expressiveness",
                  "Restricted (Cedar) vs. general-purpose (Rego) — tradeoffs")

    # Left: Cedar
    add_group_box(slide, 0.4, 1.4, 5.8, 3.6, "Cedar (AWS)",
                  AWS_ORANGE_BG, AWS_ORANGE)

    add_box(slide, 0.6, 1.9, 5.4, 0.5, "Restricted by Design",
            WHITE, AWS_ORANGE, font_size=13, bold=True)

    cedar_pros = [
        "✅  RBAC + ABAC model",
        "✅  No loops, no user-defined functions",
        "✅  Analysis is decidable (guaranteed termination)",
        "✅  Entity schema enforced at authoring time",
        "✅  Dafny-verified authorizer (provably correct)",
    ]
    cedar_cons = [
        "❌  Cannot express iteration over collections",
        "❌  No custom functions (copy-paste patterns)",
        "❌  No partial rules or incremental definitions",
        "❌  Smaller ecosystem (newer, fewer adopters)",
    ]
    add_multiline_label(slide, 0.6, 2.55, 5.4, 1.2, cedar_pros,
                        font_size=10, color=GREEN)
    add_multiline_label(slide, 0.6, 3.7, 5.4, 1.0, cedar_cons,
                        font_size=10, color=RED)

    # Right: Rego
    add_group_box(slide, 6.5, 1.4, 6.4, 3.6, "Rego / OPA (Microsoft PI)",
                  MS_BLUE_BG, MS_BLUE)

    add_box(slide, 6.7, 1.9, 6.0, 0.5, "General-Purpose with Bounded Analysis",
            WHITE, MS_BLUE, font_size=13, bold=True)

    rego_pros = [
        "✅  Industry standard (CNCF graduated project)",
        "✅  Loops, comprehensions, user-defined functions",
        "✅  Partial rules + incremental definitions",
        "✅  JSON Schema support for analysis constraints",
        "✅  Also supports Cedar + Azure Policy",
    ]
    rego_cons = [
        "❌  Analysis is bounded (loop unrolling depth limit)",
        "❌  More expressive = harder to analyze",
        "❌  Z3 analysis is prototype, not production",
        "❌  No Dafny/Verus proofs of engine (yet)",
    ]
    add_multiline_label(slide, 6.7, 2.55, 6.0, 1.2, rego_pros,
                        font_size=10, color=GREEN)
    add_multiline_label(slide, 6.7, 3.7, 6.0, 1.0, rego_cons,
                        font_size=10, color=RED)

    # Bottom takeaway
    add_box(slide, 0.4, 5.3, 12.5, 0.6,
            "Key insight: Cedar is easy to analyze because it's restricted. "
            "Rego is harder to analyze but covers more real-world policies. "
            "MS needs bounded analysis to work reliably at scale.",
            GRAY_BG, GRAY, font_size=11, font_color=DARK,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE, align=PP_ALIGN.LEFT)

    add_notes(slide, """LANGUAGE COMPARISON
• Cedar: intentionally restricted — no loops, no UDFs. This makes analysis decidable.
• Rego: general-purpose, Turing-powerful over finite data. Industry standard via OPA/CNCF.
• Tradeoff: Cedar is easier to analyze but less expressive. Rego handles more real policies.
• MS bounded analysis: loops unrolled to a depth limit (default 5). This is sound but incomplete.
• The strategic bet: most enterprises already use Rego/OPA. Cedar adoption is growing but smaller.""")


def slide_aws_investment(prs):
    """Slide 10: AWS Organizational Investment."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "AWS Organizational Investment",
                  "Dedicated team, academic credibility, brand")

    headers = ["Dimension", "AWS Investment", "Evidence"]
    rows = [
        ["Dedicated org", "Automated Reasoning Group (ARG)", "Internal advanced innovation team"],
        ["Leadership", "Neha Rungta (Dir. Applied Science);\nMike Hicks (Sr. Principal Scientist)", "CAV keynote; Cedar publications"],
        ["Team composition", "Multiple PhD scientists + engineering\nteams across 5+ services", "Job postings, conference papers"],
        ["Publications", "CAV 2022 keynote, OOPSLA 2023 paper,\nre:Invent talks yearly, Amazon Science blog", "Public record"],
        ["Open source", "Cedar SDK (GitHub, Apache 2.0)\n~5,700+ stars", "github.com/cedar-policy"],
        ["Branding", "\"Provable Security\" — dedicated page,\ncustomer testimonials, re:Inforce talks", "aws.amazon.com/security/provable-security"],
        ["Academic collab", "Stanford (cvc5 solver)\nformally verified solver integration", "Amazon Science (Feb 2026)"],
    ]

    add_table(slide, 0.3, 1.3, 12.7, 0.5, headers, rows,
              header_fill=AWS_ORANGE, col_widths=[2.3, 5.0, 3.5],
              font_size=10)

    add_notes(slide, """AWS ORGANIZATIONAL INVESTMENT
• The Automated Reasoning Group (ARG) is a distinct org within AWS — not a side project.
• Leadership: Neha Rungta gave the CAV 2022 keynote (top formal methods conference).
• Mike Hicks: senior principal scientist, formerly CS professor at UMD.
• The "Provable Security" brand is embedded in AWS's security marketing.
• Academic collaborations (Stanford/cvc5) give additional credibility.
• This level of organizational commitment is hard to replicate quickly.""")


def slide_takeaways_aws(prs):
    """Slide 11: Key Takeaways — AWS's Unassailable Leads."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Key Takeaways: Where AWS Leads",
                  "Structural advantages that are hard to replicate", bg_color=AWS_ORANGE)

    items = [
        ("8 Years of Production Hardening",
         "Zelkova has been in the S3 critical path since 2018.\n"
         "Every S3 policy attachment goes through SMT analysis."),
        ("~1 Billion SMT Queries/Day",
         "This scale is a moat. The portfolio solver approach\n"
         "(Z3 + CVC4 + cvc5 + custom) is battle-tested."),
        ("Named Enterprise Customers",
         "USAA, Bridgewater, GoTo, TELUS, Stedi, Twilio, FIS,\n"
         "Grosvenor, Attentive — all with public testimonials."),
        ("\"Provable Security\" Brand",
         "Established marketing category on AWS. Customers\n"
         "search for it. Dedicated page, re:Inforce presence."),
        ("Managed Service (GA)",
         "Amazon Verified Permissions is a pay-as-you-go\n"
         "managed Cedar authorization service."),
    ]

    y = 1.5
    for title, desc in items:
        add_box(slide, 0.5, y, 3.5, 0.85, title, AWS_ORANGE_LT, AWS_ORANGE,
                font_size=13, bold=True, font_color=DARK)
        add_label(slide, 4.2, y + 0.08, 8.5, desc,
                  font_size=11, color=DARK, align=PP_ALIGN.LEFT, height=0.7)
        y += 0.95

    add_notes(slide, """KEY TAKEAWAYS — AWS LEADS
• These are structural advantages built over 8 years — not quick to replicate.
• The moat is real: production scale, customer proof points, academic credibility, brand.
• Recommendation: acknowledge these and compete on different dimensions (the novel capabilities).""")


def slide_takeaways_ms(prs):
    """Slide 12: Key Takeaways — MS Differentiation (Prototype Stage)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Key Takeaways: Where MS Is Differentiated",
                  "Novel prototype capabilities AWS doesn't offer — not yet in production",
                  bg_color=MS_BLUE)

    items = [
        ("Policy Diff & Subsumption",
         "Proves two policies equivalent or finds distinguishing input.\n"
         "Proves new policy ⊇ old. No AWS equivalent.",
         "Prototype (working CLI)"),
        ("Automatic Test Generation",
         "Auto-generates test suite covering all paths.\n"
         "MC/DC condition coverage (aviation-grade).",
         "Prototype (working CLI)"),
        ("Multi-Language (Rego/Cedar/Azure)",
         "Rego (OPA) is the industry standard. Cedar-only\n"
         "limits AWS's reach. Regorus engine is production-grade.",
         "Engine: Production\nZ3: Prototype"),
        ("Dead Rule Detection",
         "UNSAT path conditions provably identify unreachable\n"
         "code. Novel — no AWS counterpart.",
         "Prototype"),
        ("Full Rego Expressiveness",
         "Loops, user-defined functions, partial rules.\n"
         "Cedar can't express these by design.",
         "Prototype (Z3 analysis)"),
    ]

    y = 1.5
    for title, desc, maturity in items:
        add_box(slide, 0.5, y, 3.2, 0.85, title, MS_BLUE_LT, MS_BLUE,
                font_size=12, bold=True, font_color=DARK)
        add_label(slide, 3.9, y + 0.08, 6.5, desc,
                  font_size=11, color=DARK, align=PP_ALIGN.LEFT, height=0.7)
        add_box(slide, 10.7, y + 0.1, 2.2, 0.65, maturity,
                AMBER_BG, AMBER, font_size=9, font_color=DARK)
        y += 0.95

    add_notes(slide, """KEY TAKEAWAYS — MS DIFFERENTIATION
• The diff/subsumes/gen-tests capabilities are genuinely novel — AWS has nothing like them.
• BUT: all prototype stage. Honest framing builds credibility.
• Rego multi-language support is the strategic wedge — most enterprises use OPA/Rego.
• The Regorus Rego engine itself IS production-grade and open source — that's real.
• Path to value: harden prototypes → ship → get customer proof points → brand.""")


def slide_recommendations(prs):
    """Slide 13: Strategic Recommendations."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Strategic Recommendations",
                  "Where to focus given the competitive landscape")

    recs = [
        ("1", "Don't Compete on\nIAM Posture Mgmt",
         "AWS Access Analyzer owns this space (5 features, production since 2019). "
         "Focus on analysis capabilities AWS lacks: diff, subsumes, test gen.",
         RED_LT, RED),
        ("2", "Harden the\nPrototype",
         "diff / subsumes / gen-tests are novel but need production-quality testing, "
         "performance, and documentation before head-to-head comparison.",
         AMBER_BG, AMBER),
        ("3", "Rego Is the\nWedge",
         "Most enterprises use OPA/Rego. Cedar adoption is growing but smaller. "
         "Regorus engine is already production-grade — leverage it.",
         GREEN_LT, GREEN),
        ("4", "Publish the\nRVM-to-Z3 Work",
         "AWS has CAV keynote + OOPSLA paper. Publish the symbolic translation "
         "to establish academic credibility. The approach is novel.",
         MS_BLUE_LT, MS_BLUE),
        ("5", "Ship Then Brand",
         "AWS \"Provable Security\" is built on years of production use. "
         "Ship features first, then create the marketing narrative.",
         PURPLE_LT, PURPLE),
        ("6", "Get Customer\nProof Points",
         "The #1 gap vs AWS is named customers willing to speak publicly. "
         "Start with internal Azure teams, then ISVs.",
         TEAL_LT, TEAL),
    ]

    col1_x, col2_x = 0.4, 6.7
    y1 = 1.4
    for i, (num, title, desc, fill, border) in enumerate(recs):
        x = col1_x if i < 3 else col2_x
        y = y1 + (i % 3) * 1.55
        # Number badge
        add_box(slide, x, y, 0.45, 0.45, num, border, border,
                font_size=16, bold=True, font_color=WHITE,
                shape=MSO_SHAPE.OVAL)
        # Title
        add_box(slide, x + 0.55, y, 2.3, 0.9, title, fill, border,
                font_size=11, bold=True, font_color=DARK)
        # Description
        add_label(slide, x + 3.0, y + 0.05, 3.2, desc,
                  font_size=10, color=DARK, align=PP_ALIGN.LEFT, height=0.85)

    add_notes(slide, """STRATEGIC RECOMMENDATIONS
1. Don't compete where AWS is strongest (IAM posture) — focus on novel analysis.
2. Invest in hardening: the prototypes work but need production polish.
3. Rego is the strategic wedge — it's where the customers are.
4. Publish: academic credibility matters in formal methods. The approach is novel enough for a paper.
5. Ship then brand: don't create a "Provable Security" equivalent until features are production-grade.
6. Customer proof points: start internal (Azure teams), expand to design partners.""")


def slide_closing(prs):
    """Slide 14: Closing summary."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BLUE)

    add_label(slide, 0.8, 1.0, 11.7, "Summary",
              font_size=36, color=WHITE, bold=True, align=PP_ALIGN.LEFT)

    # Three column summary
    cols = [
        ("AWS Has Built", AWS_ORANGE_LT, AWS_ORANGE, [
            "8 years of production use",
            "~1B SMT queries/day",
            "Zelkova + IAM AA + Cedar + AVP",
            "\"Provable Security\" brand",
            "8+ named enterprise customers",
            "CAV keynote + OOPSLA paper",
        ]),
        ("MS Has Started", MS_BLUE_LT, MS_BLUE, [
            "Novel: diff, subsumes, gen-tests",
            "MC/DC condition coverage",
            "Dead rule detection",
            "Multi-language (Rego + Cedar)",
            "Regorus engine: production-grade",
            "~10K lines of analysis code",
        ]),
        ("MS Needs To", GREEN_LT, GREEN, [
            "Harden prototypes to production",
            "Publish the RVM-to-Z3 research",
            "Get customer proof points",
            "Implement Verus verification",
            "Ship features before branding",
            "Close the credibility gap",
        ]),
    ]

    cx = 0.6
    for title, fill, border, items in cols:
        add_box(slide, cx, 1.9, 3.8, 0.6, title, fill, border,
                font_size=16, bold=True, font_color=DARK)
        iy = 2.6
        for item in items:
            add_label(slide, cx + 0.1, iy, 3.6, f"•  {item}",
                      font_size=12, color=WHITE, align=PP_ALIGN.LEFT, height=0.35)
            iy += 0.38

        cx += 4.2

    add_label(slide, 0.8, 5.5, 11.7,
              "All data from public sources  •  Feb 2026  •  Confidential",
              font_size=11, color=RGBColor(0x90, 0xCA, 0xF9), align=PP_ALIGN.LEFT)

    add_notes(slide, """CLOSING SUMMARY
• Three columns: what AWS has built (their moat), what MS has started (prototypes), what MS needs to do.
• The honest framing builds credibility — don't oversell the prototypes.
• The novel capabilities (diff, subsumes, gen-tests, MC/DC) are real differentiation.
• Path forward: harden → publish → ship → get customers → brand.""")


def slide_sources(prs):
    """Slide 15: Sources."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Sources",
                  "All data from publicly available sources")

    sources = [
        "1.  AWS Security Blog — Zelkova (Jun 2018)",
        "2.  AWS IAM Access Analyzer product page",
        "3.  AWS IAM Access Analyzer features page",
        "4.  Amazon Science — \"A Billion SMT Queries a Day\" (CAV 2022 keynote)",
        "5.  Amazon Science — Cedar verification-guided development (May 2023)",
        "6.  AWS Security Blog — IAM Access Analyzer policy generation (Apr 2021)",
        "7.  Amazon Verified Permissions product page",
        "8.  AWS Provable Security page",
        "9.  Cedar policy language website (cedarpolicy.com)",
        "10. Cedar GitHub repository (github.com/cedar-policy)",
        "11. Amazon Science — Academic collaboration on cvc5 (Feb 2026)",
        "12. Internal: rvm-to-z3.md — comparison table (Section 17)",
    ]

    add_multiline_label(slide, 0.6, 1.5, 12.0, 4.5, sources,
                        font_size=12, color=DARK)

    add_label(slide, 0.6, 5.5, 12.0,
              "Full comparison document: docs/diagrams/AWS_Policy_Analysis_Comparison.md",
              font_size=11, color=GRAY, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Build slides in order
    slide_title(prs)                  #  0: Title
    slide_exec_summary(prs)           #  1: Executive Summary
    slide_timeline(prs)               #  2: AWS Timeline (visual)
    slide_aws_products_zelkova(prs)   #  3: Zelkova architecture
    slide_aws_products_iam_aa(prs)    #  4: IAM Access Analyzer features
    slide_aws_products_cedar(prs)     #  5: Cedar + Verified Permissions
    slide_impact_numbers(prs)         #  6: Impact stats
    slide_aws_investment(prs)         #  7: Organizational investment table
    slide_comparison_aws_leads(prs)   #  8: H2H — AWS leads
    slide_comparison_ms_leads(prs)    #  9: H2H — MS novel capabilities
    slide_comparison_language(prs)    # 10: Language comparison
    slide_takeaways_aws(prs)          # 11: Takeaways — AWS leads
    slide_takeaways_ms(prs)           # 12: Takeaways — MS differentiation
    slide_recommendations(prs)        # 13: Strategic recommendations
    slide_closing(prs)                # 14: Closing summary
    slide_sources(prs)                # 15: Sources

    outpath = os.path.join(os.path.dirname(__file__),
                           "AWS_Policy_Comparison.pptx")
    prs.save(outpath)
    print(f"Saved {len(prs.slides)} slides → {outpath}")


if __name__ == "__main__":
    main()
