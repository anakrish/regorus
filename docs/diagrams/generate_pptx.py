#!/usr/bin/env python3
"""
Generate an editable PowerPoint presentation showcasing Policy Intelligence.
All diagrams are native PowerPoint shapes — fully editable, no images.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
import os

# ── Color palette ──────────────────────────────────────────────────────────
BLUE       = RGBColor(0x1A, 0x73, 0xE8)
BLUE_LIGHT = RGBColor(0xBB, 0xDE, 0xFB)
BLUE_BG    = RGBColor(0xE8, 0xF0, 0xFE)
ORANGE     = RGBColor(0xFF, 0x98, 0x00)
ORANGE_LT  = RGBColor(0xFF, 0xE0, 0xB2)
ORANGE_BG  = RGBColor(0xFF, 0xF3, 0xE0)
GREEN      = RGBColor(0x38, 0x8E, 0x3C)
GREEN_LT   = RGBColor(0xC8, 0xE6, 0xC9)
GREEN_BG   = RGBColor(0xE8, 0xF5, 0xE9)
GREEN_DK   = RGBColor(0xA5, 0xD6, 0xA7)
RED        = RGBColor(0xE5, 0x39, 0x35)
RED_LT     = RGBColor(0xFF, 0xCD, 0xD2)
RED_BG     = RGBColor(0xFF, 0xEB, 0xEE)
YELLOW_LT  = RGBColor(0xFF, 0xF9, 0xC4)
YELLOW     = RGBColor(0xF9, 0xA8, 0x25)
PURPLE     = RGBColor(0x8E, 0x24, 0xAA)
PURPLE_LT  = RGBColor(0xE1, 0xBE, 0xE7)
PURPLE_BG  = RGBColor(0xF3, 0xE5, 0xF5)
INDIGO     = RGBColor(0x3F, 0x51, 0xB5)
INDIGO_LT  = RGBColor(0xC5, 0xCA, 0xE9)
INDIGO_BG  = RGBColor(0xE8, 0xEA, 0xF6)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK       = RGBColor(0x1A, 0x1A, 0x2E)
GRAY       = RGBColor(0x60, 0x60, 0x60)
SLIDE_BG   = RGBColor(0xF8, 0xF9, 0xFA)

# ── Helpers ────────────────────────────────────────────────────────────────

def set_slide_bg(slide, color=SLIDE_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, title, subtitle=None):
    """Dark blue title bar across the top."""
    shapes = slide.shapes
    # Title bar background
    bar = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x0D, 0x47, 0xA1)
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
        p2.alignment = PP_ALIGN.LEFT


def add_box(slide, left, top, width, height, text, fill_color, border_color,
            font_size=12, bold=False, font_color=DARK, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
            align=PP_ALIGN.CENTER):
    """Add a styled shape with centered text."""
    s = slide.shapes.add_shape(shape, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    s.line.color.rgb = border_color
    s.line.width = Pt(2)
    s.shadow.inherit = False
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = font_color
        p.alignment = align
    tf.auto_size = None
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except:
        pass
    return s


def add_arrow(slide, start_left, start_top, end_left, end_top, color=GRAY):
    """Add a connector arrow between two points with arrowhead."""
    from lxml import etree
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(start_left), Inches(start_top),
        Inches(end_left), Inches(end_top)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    # Arrowhead via direct XML manipulation
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    ln = connector._element.find('.//a:ln', nsmap)
    if ln is None:
        spPr = connector._element.find('.//a:spPr', nsmap) or connector._element
        ln = etree.SubElement(spPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}ln')
    tail = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd')
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('len', 'med')
    return connector


def add_label(slide, left, top, width, text, font_size=10, color=GRAY, bold=False,
              align=PP_ALIGN.CENTER):
    """Add a text label (no border)."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox


def add_group_box(slide, left, top, width, height, title, fill_color, border_color,
                  title_color=DARK):
    """Add a group/section box with a title label."""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
                                Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    s.line.color.rgb = border_color
    s.line.width = Pt(2)
    # Title label
    add_label(slide, left + 0.1, top + 0.05, width - 0.2, title,
              font_size=11, color=title_color, bold=True, align=PP_ALIGN.LEFT)
    return s


def add_notes(slide, text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


# ── Slide builders ─────────────────────────────────────────────────────────

def slide_what_is_z3(prs):
    """What is Z3 / Theorem Proving — PM-friendly intro."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "What is Theorem Proving?",
                  "From Spot-Checking to Mathematical Proof")

    # ── LEFT: Traditional Testing ──
    add_group_box(slide, 0.4, 1.5, 4.0, 3.6,
                  "❌  Traditional Testing", RED_BG, RED)
    add_box(slide, 0.7, 2.2, 3.4, 0.7,
            "Pick some test inputs by hand", RED_LT, RED, font_size=13)
    add_arrow(slide, 2.4, 2.9, 2.4, 3.15, RED)
    add_box(slide, 0.7, 3.15, 3.4, 0.7,
            "Run them against the policy", RED_LT, RED, font_size=13)
    add_arrow(slide, 2.4, 3.85, 2.4, 4.1, RED)
    add_box(slide, 0.7, 4.1, 3.4, 0.7,
            "Hope you didn't miss a case",
            RED_LT, RED, font_size=13, bold=True)

    add_label(slide, 0.7, 4.9, 3.4,
              "Covers a few points in the input space",
              font_size=10, color=RED)

    # ── CENTER: VS ──
    add_label(slide, 4.7, 3.0, 1.0, "VS", font_size=26, color=INDIGO, bold=True)

    # ── RIGHT: Theorem Proving ──
    add_group_box(slide, 5.8, 1.5, 7.2, 3.6,
                  "✅  Theorem Proving (Z3)", GREEN_BG, GREEN)
    add_box(slide, 6.1, 2.2, 3.2, 0.7,
            "Translate policy to\nmathematical constraints",
            GREEN_LT, GREEN, font_size=13)
    add_arrow(slide, 9.3, 2.55, 9.8, 2.55, ORANGE)
    add_box(slide, 9.8, 2.2, 3.0, 0.7,
            "🧠 Z3 Solver\n(Microsoft Research)",
            ORANGE_LT, ORANGE, font_size=12, bold=True)

    add_arrow(slide, 11.3, 2.9, 11.3, 3.3, GREEN)
    add_box(slide, 6.1, 3.3, 6.7, 1.5,
            "Z3 checks ALL possible inputs simultaneously:\n\n"
            "  • SAT  →  found a concrete input (here it is!)\n"
            "  • UNSAT  →  no such input exists (mathematical proof)",
            GREEN_LT, GREEN, font_size=12, align=PP_ALIGN.LEFT)

    add_label(slide, 6.1, 4.9, 6.7,
              "Covers the ENTIRE input space — nothing is missed",
              font_size=10, color=GREEN)

    # ── Bottom: what Z3 is ──
    add_group_box(slide, 0.4, 5.5, 12.5, 1.7,
                  "About Z3", BLUE_BG, BLUE)

    facts = [
        ("Built by", "Microsoft Research — one of the most widely used theorem provers in the world"),
        ("Used in", "Windows driver verification, Azure security, AWS, compilers, cryptography"),
        ("Key idea", "Instead of running tests, Z3 reasons about ALL inputs at once using math"),
    ]
    for i, (label, desc) in enumerate(facts):
        y = 5.9 + i * 0.4
        add_label(slide, 0.8, y, 1.5, label + ":",
                  font_size=11, color=BLUE, bold=True, align=PP_ALIGN.LEFT)
        add_label(slide, 2.3, y, 10.0, desc,
                  font_size=11, color=DARK, align=PP_ALIGN.LEFT)


def slide_policy_to_smt(prs):
    """Show a Cedar policy and its Z3/SMT translation side by side."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Policy → SMT:  Under the Hood",
                  "How a Cedar Policy Becomes Mathematical Constraints")

    # ── LEFT: Cedar policy ──
    add_group_box(slide, 0.3, 1.4, 5.0, 4.2,
                  "📋  Cedar Policy  (IAM Zero Trust)", BLUE_BG, BLUE)

    cedar_rule1 = (
        '// Allow admin login with MFA\n'
        '// from internal IPs\n'
        'permit(\n'
        '  principal in User::"admins",\n'
        '  action == Action::"login",\n'
        '  resource == App::"portal"\n'
        ')\n'
        'when {\n'
        '  context.mfa == true &&\n'
        '  context.ip like "10.*"\n'
        '};'
    )
    add_box(slide, 0.5, 2.1, 4.6, 2.0,
            cedar_rule1,
            GREEN_LT, GREEN, font_size=10, align=PP_ALIGN.LEFT,
            font_color=RGBColor(0x1B, 0x5E, 0x20))

    cedar_rule2 = (
        '// Block if account suspended\n'
        'forbid(\n'
        '  principal in User::"admins",\n'
        '  action == Action::"login",\n'
        '  resource == App::"portal"\n'
        ')\n'
        'when { context.suspended == true };'
    )
    add_box(slide, 0.5, 4.25, 4.6, 1.2,
            cedar_rule2,
            RED_LT, RED, font_size=10, align=PP_ALIGN.LEFT,
            font_color=RGBColor(0xB7, 0x1C, 0x1C))

    # ── CENTER: Arrow ──
    add_arrow(slide, 5.3, 3.3, 5.9, 3.3, ORANGE)
    add_label(slide, 5.1, 3.6, 1.2, "compiles\nto SMT",
              font_size=10, color=ORANGE, bold=True)

    # ── RIGHT: SMT constraints ──
    add_group_box(slide, 5.9, 1.4, 7.1, 4.2,
                  "🧮  Z3 / SMT Constraints  (auto-generated)", ORANGE_BG, ORANGE)

    smt_decls = (
        '; Variables for each input field\n'
        '(declare-fun input.principal () String)\n'
        '(declare-fun input.action () String)\n'
        '(declare-fun input.resource () String)\n'
        '(declare-fun input.context.mfa () Bool)\n'
        '(declare-fun input.context.ip () String)\n'
        '(declare-fun input.context.suspended () Bool)'
    )
    add_box(slide, 6.1, 2.1, 6.7, 1.3,
            smt_decls,
            INDIGO_LT, INDIGO, font_size=9, align=PP_ALIGN.LEFT,
            font_color=RGBColor(0x1A, 0x23, 0x7E))

    smt_logic = (
        '; Policy logic: permit ∧ ¬ forbid\n'
        '(assert\n'
        '  (and\n'
        '    (= input.principal "User::admins")\n'
        '    (= input.action "Action::login")\n'
        '    (= input.resource "App::portal")\n'
        '    (= input.context.mfa true)\n'
        '    (str.in_re input.context.ip\n'
        '      (re.++ (str.to_re "10.") re.all))\n'
        '    (not (= input.context.suspended true))))'
    )
    add_box(slide, 6.1, 3.5, 6.7, 2.0,
            smt_logic,
            ORANGE_LT, ORANGE, font_size=9, align=PP_ALIGN.LEFT,
            font_color=RGBColor(0xE6, 0x51, 0x00))

    # ── Bottom: Z3 answer ──
    add_group_box(slide, 0.3, 5.85, 12.7, 1.5,
                  "⚡ Z3 Result:  SAT — found a concrete input that satisfies all constraints", GREEN_BG, GREEN)

    result_json = (
        '{ "principal": "User::admins",  "action": "Action::login",  "resource": "App::portal",\n'
        '  "context": { "mfa": true,  "ip": "10.",  "suspended": false } }'
    )
    add_box(slide, 0.5, 6.35, 12.3, 0.75,
            result_json,
            GREEN_LT, GREEN, font_size=11, align=PP_ALIGN.LEFT,
            font_color=RGBColor(0x1B, 0x5E, 0x20), bold=True)


def slide_title(prs):
    """Slide 0: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, RGBColor(0x0D, 0x47, 0xA1))

    # Big title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Policy Intelligence"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Formal Analysis for Authorization & Governance Policies"
    p2.font.size = Pt(24)
    p2.font.color.rgb = BLUE_LIGHT
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.space_before = Pt(30)
    p3.text = "Automated Testing  •  Change Validation  •  Compliance Proofs  •  Impact Analysis"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(0x90, 0xCA, 0xF9)
    p3.alignment = PP_ALIGN.CENTER

    # Bottom bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.8), Inches(13.333), Inches(0.7))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x0A, 0x36, 0x7A)
    bar.line.fill.background()
    tf2 = bar.text_frame
    p4 = tf2.paragraphs[0]
    p4.text = "Regorus — Microsoft Policy Framework"
    p4.font.size = Pt(14)
    p4.font.color.rgb = RGBColor(0x64, 0xB5, 0xF6)
    p4.alignment = PP_ALIGN.CENTER


def slide_01_overview(prs):
    """Slide 1: Policy Intelligence Overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Policy Intelligence — From Policies to Proven Guarantees")

    # ── Input group ──
    add_group_box(slide, 0.5, 1.5, 3.0, 4.2, "📋  What You Have", BLUE_BG, BLUE)
    add_box(slide, 0.8, 2.1, 2.4, 0.8, "Rego · Cedar ·\nAzure Policies", BLUE_LIGHT, BLUE, font_size=13, bold=True)
    add_box(slide, 0.8, 3.1, 2.4, 0.8, "JSON Schema\n(optional)", BLUE_LIGHT, BLUE, font_size=13)
    add_box(slide, 0.8, 4.1, 2.4, 0.8, "Reference Data", BLUE_LIGHT, BLUE, font_size=13)

    # ── Engine group ──
    add_group_box(slide, 4.5, 1.5, 3.5, 4.2, "🧠  Policy Intelligence Engine", ORANGE_BG, ORANGE)
    add_box(slide, 5.0, 2.3, 2.5, 1.0, "Compile &\nAnalyze", ORANGE_LT, ORANGE, font_size=14, bold=True)
    add_arrow(slide, 6.25, 3.3, 6.25, 3.7, ORANGE)
    add_box(slide, 5.0, 3.7, 2.5, 1.2, "Mathematical\nReasoning\n(Z3 Solver)", ORANGE_LT, ORANGE, font_size=13)

    # ── Output group ──
    add_group_box(slide, 9.0, 1.5, 3.8, 4.2, "✅  What You Get", GREEN_BG, GREEN)
    add_box(slide, 9.3, 2.1, 3.2, 0.7, "Auto-Generated Test Suites", GREEN_LT, GREEN, font_size=13, bold=True)
    add_box(slide, 9.3, 3.0, 3.2, 0.7, "Change Impact Analysis", GREEN_LT, GREEN, font_size=13)
    add_box(slide, 9.3, 3.9, 3.2, 0.7, "Root Cause Explanations", GREEN_LT, GREEN, font_size=13)
    add_box(slide, 9.3, 4.8, 3.2, 0.7, "Compliance Proofs", GREEN_LT, GREEN, font_size=13)

    # ── Arrows between groups ──
    add_arrow(slide, 3.5, 3.6, 4.5, 3.6, GRAY)
    add_arrow(slide, 8.0, 3.6, 9.0, 3.6, GRAY)

    # ── Bottom tagline ──
    add_label(slide, 0.5, 6.2, 12, "Policies go in → Mathematical guarantees come out",
              font_size=14, color=GRAY, bold=True)


def slide_02_before_after(prs):
    """Slide 2: Manual Testing vs Automated."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Automated Test Generation",
                  "Zero Manual Effort, Full Coverage")

    # ── LEFT: Today (red) ──
    add_group_box(slide, 0.5, 1.5, 4.5, 4.5, "❌  Today: Manual Testing", RED_BG, RED)
    add_box(slide, 0.9, 2.2, 3.7, 0.9, "Engineer writes\ntest inputs by hand", RED_LT, RED, font_size=14)
    add_arrow(slide, 2.75, 3.1, 2.75, 3.4, RED)
    add_box(slide, 0.9, 3.4, 3.7, 0.9, "Misses edge cases\n& corner scenarios", RED_LT, RED, font_size=14)
    add_arrow(slide, 2.75, 4.3, 2.75, 4.6, RED)
    add_box(slide, 0.9, 4.6, 3.7, 0.9, "Incomplete coverage\n= security gaps", RED_LT, RED, font_size=14, bold=True)

    # ── CENTER: VS ──
    add_label(slide, 5.5, 3.3, 1.5, "VS", font_size=28, color=INDIGO, bold=True)

    # ── RIGHT: Future (green) ──
    add_group_box(slide, 7.3, 1.5, 5.5, 4.5, "✅  With Policy Intelligence", GREEN_BG, GREEN)
    add_box(slide, 7.7, 2.2, 4.7, 0.9, "Engine reads policy & schema\nautomatically", GREEN_LT, GREEN, font_size=14)
    add_arrow(slide, 10.05, 3.1, 10.05, 3.4, GREEN)
    add_box(slide, 7.7, 3.4, 4.7, 0.9, "Mathematically finds\nALL decision paths", GREEN_LT, GREEN, font_size=14, bold=True)
    add_arrow(slide, 10.05, 4.3, 10.05, 4.6, GREEN)
    add_box(slide, 7.7, 4.6, 4.7, 0.9, "Generates minimal test suite\nwith 100% coverage", GREEN_LT, GREEN, font_size=14, bold=True)


def slide_03_policy_migration(prs):
    """Slide 3: Safe Policy Migration."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Safe Policy Migration",
                  "Prove Equivalence Before You Deploy")

    # Policy v1
    add_box(slide, 0.8, 2.2, 2.5, 1.2, "Policy v1\n(Current)", BLUE_LIGHT, BLUE, font_size=15, bold=True)
    # Policy v2
    add_box(slide, 0.8, 4.2, 2.5, 1.2, "Policy v2\n(Proposed)", BLUE_LIGHT, BLUE, font_size=15, bold=True)

    # Arrows to diff
    add_arrow(slide, 3.3, 2.8, 4.3, 3.5, GRAY)
    add_arrow(slide, 3.3, 4.8, 4.3, 3.9, GRAY)

    # Diff engine
    add_box(slide, 4.3, 2.8, 2.8, 1.8, "🔍 Policy Diff\nAnalysis\n\nCompares ALL\npossible inputs",
            ORANGE_LT, ORANGE, font_size=13, bold=False)

    # Arrow to safe
    add_arrow(slide, 7.1, 3.2, 8.3, 2.5, GREEN)
    add_label(slide, 7.2, 2.5, 1.5, "Identical\nbehavior", font_size=10, color=GREEN)

    # Arrow to diff found
    add_arrow(slide, 7.1, 4.0, 8.3, 4.8, YELLOW)
    add_label(slide, 7.2, 4.3, 1.5, "Behavior\ndiffers", font_size=10, color=YELLOW)

    # Safe box
    add_box(slide, 8.3, 1.8, 4.2, 1.3, "✅ Safe to Deploy\nMathematically Equivalent\nfor ALL inputs",
            GREEN_LT, GREEN, font_size=14, bold=True)

    # Diff found box
    add_box(slide, 8.3, 4.2, 4.2, 1.5, "⚠️ Found Difference!\nHere's the exact input\nthat causes divergent behavior",
            YELLOW_LT, YELLOW, font_size=13, bold=True)

    # Fix loop arrow
    add_arrow(slide, 9.2, 5.7, 2.0, 5.7, PURPLE)
    add_arrow(slide, 2.0, 5.7, 2.0, 5.4, PURPLE)
    add_label(slide, 4.5, 5.8, 3, "Fix & Re-verify", font_size=11, color=PURPLE, bold=True)


def slide_04_why_denied(prs):
    """Slide 4: Why Denied root cause analysis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Why Denied?",
                  "Pinpoint the Exact Failing Condition")

    # Denied request
    add_box(slide, 0.5, 2.5, 2.8, 1.5, "🚫 Request Denied\n\n\"Why was I blocked?\"",
            RED_LT, RED, font_size=14, bold=True)

    # Arrow to analysis
    add_arrow(slide, 3.3, 3.25, 4.3, 3.25, GRAY)

    # Analysis engine
    add_box(slide, 4.3, 2.2, 2.5, 2.1, "🧠 Root Cause\nAnalysis\n\n(MAX-SAT\nSolver)", ORANGE_LT, ORANGE, font_size=13)

    # Arrows to conditions
    y_positions = [1.8, 2.8, 3.8, 4.8]
    labels_data = [
        ("Condition 1 ✅\nUser is authenticated", GREEN_LT, GREEN, False),
        ("Condition 2 ✅\nValid subscription", GREEN_LT, GREEN, False),
        ("Condition 3 ❌\nGroup membership\nnot enforced", RED_LT, RED, True),
        ("Condition 4 ✅\nRegion allowed", GREEN_LT, GREEN, False),
    ]

    for i, (text, fill, border, is_fail) in enumerate(labels_data):
        y = y_positions[i]
        add_arrow(slide, 6.8, 3.25, 7.5, y + 0.4, border)
        bw = Pt(3) if is_fail else Pt(2)
        s = add_box(slide, 7.5, y, 2.7, 0.9, text, fill, border,
                    font_size=12, bold=is_fail)
        if is_fail:
            s.line.width = Pt(3)

    # Arrow from failing condition to action
    add_arrow(slide, 10.2, 4.25, 10.7, 4.25, BLUE)

    # Actionable insight
    add_box(slide, 10.7, 3.5, 2.2, 1.5,
            "💡 Actionable Insight\n\n\"managedState\"\nmust be \"enforced\"\n— contact\ngroup owner",
            BLUE_BG, BLUE, font_size=11, bold=False)


def slide_05_compliance(prs):
    """Slide 5: Compliance Proofs."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Compliance Proofs",
                  "Mathematical Guarantees, Not Just Testing")

    # Questions
    add_group_box(slide, 0.5, 1.5, 3.5, 4.5, "❓  Compliance Questions", INDIGO_BG, INDIGO)
    add_box(slide, 0.8, 2.2, 2.9, 1.0, "Can any request\nbypass this rule?", INDIGO_LT, INDIGO, font_size=13)
    add_box(slide, 0.8, 3.4, 2.9, 1.0, "Is the new policy at\nleast as restrictive?", INDIGO_LT, INDIGO, font_size=13)
    add_box(slide, 0.8, 4.6, 2.9, 1.0, "Are there unreachable\n(dead) rules?", INDIGO_LT, INDIGO, font_size=13)

    # Arrow to engine
    add_arrow(slide, 4.0, 3.75, 5.0, 3.75, GRAY)

    # Engine
    add_box(slide, 5.0, 2.5, 2.5, 2.5, "🧠 Policy\nIntelligence\nEngine\n\n(Z3 Theorem\nProver)",
            ORANGE_LT, ORANGE, font_size=13, bold=True)

    # Arrows to answers
    add_arrow(slide, 7.5, 3.0, 8.5, 2.3, GREEN)
    add_arrow(slide, 7.5, 3.75, 8.5, 3.75, GREEN)
    add_arrow(slide, 7.5, 4.5, 8.5, 5.0, YELLOW)

    # Answer boxes
    add_box(slide, 8.5, 1.7, 4.3, 1.2,
            "✅ Mathematically Impossible to Bypass\nNo input exists that circumvents this rule\n(UNSAT proof)",
            GREEN_LT, GREEN, font_size=12, bold=True)
    add_box(slide, 8.5, 3.1, 4.3, 1.2,
            "✅ New Policy is Strictly Tighter\nEvery input allowed by old is also allowed by new\n(Subsumption proof)",
            GREEN_LT, GREEN, font_size=12, bold=True)
    add_box(slide, 8.5, 4.5, 4.3, 1.2,
            "🧹 3 Rules Are Unreachable\nNo input can ever trigger them\n— safe to remove",
            YELLOW_LT, YELLOW, font_size=12)


def slide_06_mcdc(prs):
    """Slide 6: MC/DC Condition Coverage."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "MC/DC Condition Coverage",
                  "Aviation-Grade Testing Standard (DO-178C) — Applied to Policies")

    # Policy rule box
    add_group_box(slide, 0.3, 1.5, 3.3, 2.8, "Policy Rule: applicable", INDIGO_BG, INDIGO)
    add_box(slide, 0.5, 2.2, 2.9, 1.8,
            "applicable if {\n  appId NOT in allowlist\n  AND isFetchSuccess\n  AND managedState\n      == \"enforced\"\n}",
            INDIGO_LT, INDIGO, font_size=12, align=PP_ALIGN.LEFT)

    # Arrow to MC/DC
    add_arrow(slide, 3.6, 3.0, 4.3, 3.0, ORANGE)

    # MC/DC box
    add_box(slide, 4.3, 2.2, 2.0, 1.6, "🧠 MC/DC\nAnalysis\n\nFlip each\ncondition\nindependently",
            ORANGE_LT, ORANGE, font_size=12)

    # Arrow to tests
    add_arrow(slide, 6.3, 3.0, 6.8, 3.0, ORANGE)

    # Test pairs
    add_group_box(slide, 6.8, 1.3, 6.2, 5.2,
                  "Generated Test Pairs — Each Condition Independently Determines Outcome",
                  PURPLE_BG, PURPLE)

    # Condition A
    add_label(slide, 7.0, 2.0, 5.8, "Condition A: appId in allowlist", font_size=11, color=INDIGO, bold=True)
    add_box(slide, 7.0, 2.4, 2.8, 0.8, "✅ TRUE\nappId = \"ff177ae3...\"\n(allowed app)", GREEN_LT, GREEN, font_size=10)
    add_box(slide, 10.0, 2.4, 2.8, 0.8, "❌ FALSE\nappId = \"aaaaaaaa...\"\n(unknown app)", RED_LT, RED, font_size=10)

    # Condition B
    add_label(slide, 7.0, 3.4, 5.8, "Condition B: isFetchSuccess", font_size=11, color=INDIGO, bold=True)
    add_box(slide, 7.0, 3.7, 2.8, 0.8, "✅ TRUE\nfetchResponse.code = 200", GREEN_LT, GREEN, font_size=10)
    add_box(slide, 10.0, 3.7, 2.8, 0.8, "❌ FALSE\nfetchResponse.code = 500", RED_LT, RED, font_size=10)

    # Condition C
    add_label(slide, 7.0, 4.7, 5.8, "Condition C: managedState == \"enforced\"", font_size=11, color=INDIGO, bold=True)
    add_box(slide, 7.0, 5.0, 2.8, 0.8, "✅ TRUE\nmanagedState = \"enforced\"", GREEN_LT, GREEN, font_size=10)
    add_box(slide, 10.0, 5.0, 2.8, 0.8, "❌ FALSE\nmanagedState = \"reportOnly\"", RED_LT, RED, font_size=10)

    # Callout
    add_label(slide, 0.3, 5.8, 12.5,
              "MC/DC = Modified Condition/Decision Coverage — the same standard required for flight-critical avionics software",
              font_size=12, color=GRAY)


def slide_07_decision_tree(prs):
    """Slide 7: Real Policy Decision Tree."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Real Policy — Azure Group Governance",
                  "5 Decision Paths, All Automatically Explored by Z3")

    # Request entry
    add_box(slide, 0.3, 3.0, 2.2, 1.2, "Incoming\nRequest\n(appId +\nfetchResponse)",
            INDIGO_LT, INDIGO, font_size=12, bold=True)

    # Decision 1: allowlist
    add_arrow(slide, 2.5, 3.6, 3.2, 3.6, GRAY)
    d1 = add_box(slide, 3.2, 2.8, 2.0, 1.5, "appId in\nallowlist?",
                 ORANGE_LT, ORANGE, font_size=13, bold=True, shape=MSO_SHAPE.DIAMOND)

    # Yes → Allow (path 1)
    add_arrow(slide, 4.2, 2.8, 4.2, 1.7, GREEN)
    add_label(slide, 3.5, 2.1, 0.6, "Yes", font_size=10, color=GREEN, bold=True)
    add_box(slide, 3.2, 1.2, 2.0, 0.6, "✅ ALLOW\n(Allowed App)", GREEN_LT, GREEN, font_size=11, bold=True)

    # No → HTTP code check
    add_arrow(slide, 5.2, 3.6, 6.0, 3.6, GRAY)
    add_label(slide, 5.3, 3.2, 0.5, "No", font_size=10, color=RED, bold=True)

    d2 = add_box(slide, 6.0, 2.8, 2.0, 1.5, "HTTP\nstatus\ncode?",
                 ORANGE_LT, ORANGE, font_size=13, bold=True, shape=MSO_SHAPE.DIAMOND)

    # 200 → managedState check
    add_arrow(slide, 8.0, 3.6, 8.8, 3.6, GRAY)
    add_label(slide, 8.0, 3.2, 0.6, "200", font_size=10, color=BLUE, bold=True)

    d3 = add_box(slide, 8.8, 2.8, 2.0, 1.5, "managed\nState?",
                 ORANGE_LT, ORANGE, font_size=13, bold=True, shape=MSO_SHAPE.DIAMOND)

    # enforced → DENY (path 2)
    add_arrow(slide, 9.8, 2.8, 9.8, 1.7, RED)
    add_label(slide, 9.2, 2.1, 1.0, "enforced", font_size=10, color=RED, bold=True)
    add_box(slide, 9.0, 1.2, 1.8, 0.6, "🚫 DENY\n(Enforced)", RED_LT, RED, font_size=11, bold=True)

    # reportOnly → AUDIT (path 3)
    add_arrow(slide, 10.8, 3.6, 11.5, 3.0, YELLOW)
    add_label(slide, 11.0, 2.7, 1.2, "reportOnly", font_size=9, color=YELLOW, bold=True)
    add_box(slide, 11.3, 1.2, 1.8, 0.6, "📋 AUDIT\n(Report Only)", YELLOW_LT, YELLOW, font_size=11, bold=True)

    # 403/404 → ALLOW (path 4)
    add_arrow(slide, 7.0, 4.3, 7.0, 5.1, GREEN)
    add_label(slide, 6.1, 4.5, 1.2, "403/404", font_size=10, color=GREEN, bold=True)
    add_box(slide, 5.8, 5.1, 2.5, 0.7, "✅ ALLOW\n(Expected Error)", GREEN_LT, GREEN, font_size=11, bold=True)

    # null/other → DENY (path 5)
    add_arrow(slide, 7.0, 2.8, 7.0, 1.7, RED)
    add_label(slide, 7.3, 2.0, 1.5, "null/other", font_size=10, color=RED, bold=True)
    add_box(slide, 6.0, 1.2, 2.0, 0.6, "🚫 DENY\n(Fail-Closed)", RED_LT, RED, font_size=11, bold=True)

    # Legend
    add_label(slide, 0.3, 6.2, 12,
              "Z3 automatically finds an input for each of the 5 paths — no manual test authoring needed",
              font_size=12, color=GRAY, bold=True)


def slide_08_real_results(prs):
    """Slide 8: Real Results with numbers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Real Results — AGS Group Governance Policy",
                  "23 Tests, 100% Coverage, Zero Manual Effort")

    # Input section
    add_group_box(slide, 0.3, 1.5, 3.5, 3.0, "Input", INDIGO_BG, INDIGO)
    add_box(slide, 0.5, 2.2, 3.1, 0.9,
            "group_governance.rego\n139 lines · 5 decision paths", INDIGO_LT, INDIGO, font_size=12)
    add_box(slide, 0.5, 3.3, 3.1, 0.9,
            "JSON Schema\ntype constraints · required fields\nenum values", INDIGO_LT, INDIGO, font_size=11)

    # Arrow
    add_arrow(slide, 3.8, 3.1, 4.5, 3.1, ORANGE)

    # Command
    add_box(slide, 4.5, 2.3, 2.5, 1.5,
            "⚡ gen-tests\n\n--condition-\ncoverage",
            ORANGE_LT, ORANGE, font_size=14, bold=True)

    # Arrow
    add_arrow(slide, 7.0, 3.1, 7.5, 3.1, ORANGE)

    # Results section
    add_group_box(slide, 7.5, 1.3, 5.5, 5.3, "Automatic Results", GREEN_BG, GREEN)

    # Stats - big numbers
    # Line coverage
    add_box(slide, 7.8, 2.0, 2.5, 0.9,
            "34 / 34 lines\n100% Line Coverage", GREEN_DK, GREEN, font_size=13, bold=True)
    # MC/DC
    add_box(slide, 10.5, 2.0, 2.3, 0.9,
            "46 / 46 conditions\n100% MC/DC", GREEN_DK, GREEN, font_size=13, bold=True)
    # Test count
    add_box(slide, 7.8, 3.1, 5.0, 0.6,
            "23 test cases auto-generated", GREEN_LT, GREEN, font_size=15, bold=True)

    # Sample tests
    add_label(slide, 7.8, 3.9, 5.0, "Sample Generated Inputs:", font_size=11, color=DARK, bold=True)

    add_box(slide, 7.8, 4.3, 5.0, 0.5,
            "Test 2: DENY — code: 200, managedState: \"enforced\", appid: \"aaaaaaaa-...\"",
            RED_LT, RED, font_size=10)
    add_box(slide, 7.8, 4.9, 5.0, 0.5,
            "Test 3: ALLOW — code: 200, managedState: \"reportOnly\", appid: \"aaaaaaaa-...\"",
            GREEN_LT, GREEN, font_size=10)
    add_box(slide, 7.8, 5.5, 5.0, 0.5,
            "Test 5: ALLOW — code: 403, appid: \"aaaaaaaa-...\"",
            GREEN_LT, GREEN, font_size=10)


def slide_section_vision(prs):
    """Section divider: What's Next — Vision slides."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(0x0D, 0x47, 0xA1))

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "What's Next"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "High-Impact Scenarios Enabled by Policy Intelligence"
    p2.font.size = Pt(20)
    p2.font.color.rgb = BLUE_LIGHT
    p2.alignment = PP_ALIGN.CENTER

    # Bullet list of what's coming
    items = [
        "Copilot + PI verification loop",
        "PR review bot with formal diff",
        "Shadow mode prediction — zero traffic",
        "Regulatory impact scanning",
        "Separation of duty proofs",
        "Gap detection for undefined decisions",
    ]
    y = 4.2
    for item in items:
        add_label(slide, 3.5, y, 6.0, "→  " + item,
                  font_size=16, color=RGBColor(0x90, 0xCA, 0xF9), bold=False)
        y += 0.4


def slide_09_copilot_pi_loop(prs):
    """Slide 9: Copilot + Policy Intelligence verification loop."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Copilot + Policy Intelligence",
                  "AI Writes the Policy — Math Proves It's Correct")

    # ── Step 1: Developer prompt ──
    add_box(slide, 0.4, 1.8, 2.5, 1.3,
            "👤 Developer\n\n\"Only allow reads\nfor interns\"",
            INDIGO_LT, INDIGO, font_size=13, bold=True)

    # Arrow to Copilot
    add_arrow(slide, 2.9, 2.45, 3.7, 2.45, GRAY)

    # ── Step 2: Copilot generates ──
    add_box(slide, 3.7, 1.8, 2.5, 1.3,
            "🤖 Copilot\n\nGenerates Rego\npolicy from intent",
            PURPLE_LT, PURPLE, font_size=13, bold=True)

    # Arrow to PI
    add_arrow(slide, 6.2, 2.45, 7.0, 2.45, GRAY)

    # ── Step 3: PI verifies ──
    add_box(slide, 7.0, 1.6, 2.8, 1.7,
            "🧠 Policy Intelligence\n\nVerifies against\nformal properties:\n• No privilege escalation\n• No bypasses possible",
            ORANGE_LT, ORANGE, font_size=12, bold=False)

    # Arrow to result PASS
    add_arrow(slide, 9.8, 2.1, 10.8, 2.1, GREEN)
    add_box(slide, 10.8, 1.6, 2.2, 1.0,
            "✅ VERIFIED\n\nSafe to deploy",
            GREEN_LT, GREEN, font_size=13, bold=True)

    # Arrow to result FAIL → feedback loop
    add_arrow(slide, 8.4, 3.3, 8.4, 4.2, RED)

    add_box(slide, 7.0, 4.2, 2.8, 1.3,
            "❌ ISSUE FOUND\n\n\"Intern with role=manager\ncan still write\"",
            RED_LT, RED, font_size=12, bold=True)

    # Feedback arrow back to Copilot
    add_arrow(slide, 7.0, 4.85, 5.0, 4.85, PURPLE)
    add_arrow(slide, 5.0, 4.85, 5.0, 3.1, PURPLE)
    add_label(slide, 5.2, 4.3, 2.5, "Auto-fix prompt\nwith counterexample",
              font_size=10, color=PURPLE, bold=True)

    # ── Bottom: the value prop ──
    add_group_box(slide, 0.4, 5.8, 12.5, 1.2,
                  "Key Insight", GREEN_BG, GREEN)
    add_label(slide, 0.6, 6.2, 12.0,
              "Every AI-generated policy gets a mathematical proof before it reaches production. "
              "Copilot iterates automatically using counterexamples until the proof passes.",
              font_size=13, color=DARK)


def slide_10_pr_review_bot(prs):
    """Slide 10: GitHub PR policy review bot."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "PR Review Bot — Policy Diff in CI",
                  "Every Policy Change Gets a Formal Impact Report")

    # PR box
    add_box(slide, 0.4, 1.8, 2.5, 1.4,
            "📝 Pull Request\n\nModifies\npolicy.rego",
            INDIGO_LT, INDIGO, font_size=13, bold=True)

    # Arrow to CI
    add_arrow(slide, 2.9, 2.5, 3.7, 2.5, GRAY)

    # CI pipeline
    add_box(slide, 3.7, 1.8, 2.3, 1.4,
            "⚙️ CI Pipeline\n\nruns regorus\ndiff v1 vs v2",
            ORANGE_LT, ORANGE, font_size=13, bold=True)

    # Arrow to analysis
    add_arrow(slide, 6.0, 2.5, 6.8, 2.5, ORANGE)

    # Diff results
    add_group_box(slide, 6.8, 1.5, 6.2, 3.5,
                  "📊  PR Comment — Formal Diff Report", BLUE_BG, BLUE)

    add_box(slide, 7.0, 2.2, 5.8, 0.7,
            "🔬 Behavioral Changes Found:  2 inputs now produce different outcomes",
            YELLOW_LT, YELLOW, font_size=12)
    add_box(slide, 7.0, 3.1, 5.8, 0.7,
            "Input #1: role=contractor, resource=internal-docs → was ALLOW, now DENY",
            RED_LT, RED, font_size=11)
    add_box(slide, 7.0, 4.0, 5.8, 0.7,
            "Input #2: role=vendor, action=read → was DENY, now ALLOW  ⚠️ wider access",
            YELLOW_LT, YELLOW, font_size=11)

    # Bottom: reviewer approval
    add_group_box(slide, 0.4, 5.4, 12.5, 1.8,
                  "What Reviewers See", GREEN_BG, GREEN)
    items = [
        ("No behavioral change → ✅ auto-approve", GREEN),
        ("Changes found → reviewer sees exact inputs with before/after outcomes", BLUE),
        ("Wider access detected → ⚠️ mandatory security review", RED),
    ]
    for i, (txt, color) in enumerate(items):
        add_label(slide, 0.8, 5.85 + i * 0.45, 11.5, txt,
                  font_size=12, color=color, bold=True, align=PP_ALIGN.LEFT)


def slide_11_shadow_prediction(prs):
    """Slide 11: Shadow mode prediction without live traffic."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Shadow Mode — Without the Shadow",
                  "Predict Policy Impact Mathematically, Not by Routing Traffic")

    # ── LEFT: Old way ──
    add_group_box(slide, 0.4, 1.5, 4.0, 3.5,
                  "❌  Shadow Mode (Today)", RED_BG, RED)
    add_box(slide, 0.7, 2.2, 3.4, 0.7,
            "Deploy policy to shadow fleet", RED_LT, RED, font_size=12)
    add_arrow(slide, 2.4, 2.9, 2.4, 3.1, RED)
    add_box(slide, 0.7, 3.1, 3.4, 0.7,
            "Wait days/weeks for traffic", RED_LT, RED, font_size=12)
    add_arrow(slide, 2.4, 3.8, 2.4, 4.0, RED)
    add_box(slide, 0.7, 4.0, 3.4, 0.7,
            "Hope you saw enough variety", RED_LT, RED, font_size=12, bold=True)

    # ── CENTER: VS ──
    add_label(slide, 4.7, 3.0, 1.0, "VS", font_size=24, color=INDIGO, bold=True)

    # ── RIGHT: PI way ──
    add_group_box(slide, 5.8, 1.5, 7.2, 3.5,
                  "✅  Policy Intelligence (Instant)", GREEN_BG, GREEN)
    add_box(slide, 6.1, 2.2, 3.0, 0.7,
            "Old Policy + New Policy", BLUE_LIGHT, BLUE, font_size=12, bold=True)
    add_arrow(slide, 9.1, 2.55, 9.6, 2.55, ORANGE)
    add_box(slide, 9.6, 2.2, 3.2, 0.7,
            "🧠 Diff Engine (Z3)", ORANGE_LT, ORANGE, font_size=12, bold=True)

    add_arrow(slide, 11.2, 2.9, 11.2, 3.3, GREEN)
    add_box(slide, 6.1, 3.3, 6.7, 1.3,
            "Instant Report:\n"
            "• 99.7% of input space: identical behavior\n"
            "• 0.3% affected: contractors accessing internal-docs\n"
            "• Zero traffic needed — mathematically exhaustive",
            GREEN_LT, GREEN, font_size=12, align=PP_ALIGN.LEFT)

    # ── Bottom comparison ──
    add_group_box(slide, 0.4, 5.5, 12.5, 1.5,
                  "Comparison", INDIGO_BG, INDIGO)
    comparisons = [
        ("Shadow Mode", "Days to weeks", "Only traffic you see", RED),
        ("Policy Intelligence", "Seconds", "ALL possible inputs", GREEN),
    ]
    for i, (label, time, coverage, color) in enumerate(comparisons):
        x = 0.8 + i * 6.3
        add_label(slide, x, 5.9, 2.0, label, font_size=13, color=color, bold=True)
        add_label(slide, x + 2.0, 5.9, 1.5, time, font_size=12, color=DARK)
        add_label(slide, x + 3.5, 5.9, 2.5, coverage, font_size=12, color=DARK)


def slide_12_regulatory_impact(prs):
    """Slide 12: Regulatory impact analysis — scan for compliance gaps."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Regulatory Impact Analysis",
                  "Scan All Policies for Compliance Gaps — Automatically")

    # Regulation input
    add_box(slide, 0.4, 2.0, 2.5, 1.5,
            "📜 Regulation\n\n\"MFA required\nfor all PII\naccess\"",
            INDIGO_LT, INDIGO, font_size=13, bold=True)

    # Arrow to scanner
    add_arrow(slide, 2.9, 2.75, 3.8, 2.75, GRAY)

    # Scanner / formalization
    add_box(slide, 3.8, 1.8, 2.5, 1.8,
            "🧠 Formalize &\nScan\n\n∀ request:\naccesses_pii(r) →\nhas_mfa(r)",
            ORANGE_LT, ORANGE, font_size=12)

    # Arrow to policy set
    add_arrow(slide, 6.3, 2.75, 7.2, 2.75, ORANGE)

    # Policy fleet scan
    add_group_box(slide, 7.2, 1.5, 5.8, 2.7,
                  "🏢  Policy Fleet (42 services)", BLUE_BG, BLUE)
    # Results
    add_box(slide, 7.5, 2.3, 2.5, 0.7,
            "✅ 39 services\nCompliant", GREEN_LT, GREEN, font_size=12, bold=True)
    add_box(slide, 10.2, 2.3, 2.5, 0.7,
            "❌ 3 services\nGap Found", RED_LT, RED, font_size=12, bold=True)
    add_box(slide, 7.5, 3.2, 5.2, 0.7,
            "Each result backed by a mathematical proof — not sampling",
            GREEN_LT, GREEN, font_size=11)

    # Arrow to details
    add_arrow(slide, 11.45, 3.0, 11.45, 4.5, RED)

    # Detailed findings
    add_group_box(slide, 0.4, 4.5, 12.5, 2.5,
                  "Findings Detail — Counterexamples", RED_BG, RED)

    findings = [
        ("Service: billing-api", "PII endpoint /invoices has no MFA check", "→ Add MFA gate"),
        ("Service: hr-portal", "Admin override bypasses PII check", "→ Remove override"),
        ("Service: analytics", "Logging endpoint exposes email without MFA", "→ Restrict field"),
    ]
    for i, (svc, issue, fix) in enumerate(findings):
        y = 5.1 + i * 0.6
        add_box(slide, 0.7, y, 2.8, 0.45, svc, INDIGO_LT, INDIGO, font_size=10, bold=True)
        add_box(slide, 3.6, y, 5.0, 0.45, issue, RED_LT, RED, font_size=10)
        add_box(slide, 8.7, y, 4.0, 0.45, fix, GREEN_LT, GREEN, font_size=10, bold=True)


def slide_13_separation_of_duty(prs):
    """Slide 13: Separation of Duty — SOX compliance."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Separation of Duty — Formal Proof",
                  "Prove No Role Can Both Submit AND Approve (SOX / SOC-2)")

    # ── Property to prove ──
    add_group_box(slide, 0.4, 1.5, 4.0, 2.0,
                  "🎯  Security Property", INDIGO_BG, INDIGO)
    add_box(slide, 0.7, 2.2, 3.4, 1.0,
            "∀ user, request:\ncan_submit(user, request)\n  → ¬ can_approve(user, request)",
            INDIGO_LT, INDIGO, font_size=13, align=PP_ALIGN.LEFT)

    # Arrow to engine
    add_arrow(slide, 4.4, 2.5, 5.3, 2.5, GRAY)

    # Engine
    add_box(slide, 5.3, 1.8, 2.5, 1.5,
            "🧠 Z3 Solver\n\nSearches ALL\nrole combinations\nexhaustively",
            ORANGE_LT, ORANGE, font_size=12, bold=False)

    # Path A: PROVEN
    add_arrow(slide, 7.8, 2.2, 8.8, 1.8, GREEN)
    add_box(slide, 8.8, 1.5, 4.2, 1.0,
            "✅ PROVEN: No Violation Possible\n\nMathematical certificate for auditors",
            GREEN_LT, GREEN, font_size=12, bold=True)

    # Path B: VIOLATION
    add_arrow(slide, 7.8, 2.8, 8.8, 3.5, RED)
    add_box(slide, 8.8, 3.0, 4.2, 1.3,
            "❌ VIOLATION FOUND\n\nUser: finance_lead\nRole: [submitter, approver]\nAmount: $50,000",
            RED_LT, RED, font_size=12, bold=True)

    # ── Bottom: real-world compliance impact ──
    add_group_box(slide, 0.4, 4.8, 12.5, 2.2,
                  "Compliance Impact", GREEN_BG, GREEN)

    audits = [
        ("🏦", "SOX Compliance",
         "Prove financial controls enforce separation of duty across all roles"),
        ("🔐", "SOC-2 Type II",
         "Continuous proof that access controls meet trust service criteria"),
        ("🏥", "HIPAA",
         "Verify no single user can access AND authorize patient data release"),
    ]
    for i, (icon, title, desc) in enumerate(audits):
        y = 5.2 + i * 0.55
        add_box(slide, 0.7, y, 0.6, 0.4, icon, GREEN_LT, GREEN, font_size=14)
        add_label(slide, 1.4, y, 2.0, title, font_size=12, color=GREEN, bold=True, align=PP_ALIGN.LEFT)
        add_label(slide, 3.4, y, 9.3, desc, font_size=11, color=DARK, align=PP_ALIGN.LEFT)


def slide_14_gap_detection(prs):
    """Slide 14: Gap detection — find undefined policy paths."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Gap Detection — Find Undefined Decisions",
                  "Discover Inputs Where Policy Returns UNDEFINED (Falls Through All Rules)")

    # ── Policy ──
    add_box(slide, 0.4, 2.0, 2.5, 1.3,
            "📋 Your Policy\n\n3 explicit rules\n(allow, deny, audit)",
            BLUE_LIGHT, BLUE, font_size=13, bold=True)

    # Arrow to gap finder
    add_arrow(slide, 2.9, 2.65, 3.8, 2.65, GRAY)

    # Gap finder
    add_box(slide, 3.8, 1.8, 2.5, 1.8,
            "🔍 Gap Finder\n\n\"Is there any input\nthat matches NO\nrule at all?\"",
            ORANGE_LT, ORANGE, font_size=12)

    # Arrow to results
    add_arrow(slide, 6.3, 2.65, 7.2, 2.65, ORANGE)

    # Results
    add_group_box(slide, 7.2, 1.5, 5.8, 3.0,
                  "🕳️  Gaps Found: 2 Unhandled Input Regions", YELLOW_LT, YELLOW)

    add_box(slide, 7.5, 2.3, 5.3, 0.9,
            "Gap 1:  role = \"contractor\" AND resource = \"internal-wiki\"\n"
            "→ No rule matches — default behavior is UNDEFINED",
            RED_LT, RED, font_size=11)
    add_box(slide, 7.5, 3.4, 5.3, 0.9,
            "Gap 2:  action = \"delete\" AND resource_type = \"audit-log\"\n"
            "→ Falls through all rules — should this be explicitly denied?",
            RED_LT, RED, font_size=11)

    # ── Bottom: Why this matters ──
    add_group_box(slide, 0.4, 5.0, 12.5, 2.1,
                  "Why Gaps Are Dangerous", RED_BG, RED)
    dangers = [
        ("🔓", "Security Risk",
         "Undefined decisions may default to ALLOW in some engines — silent backdoors"),
        ("🐛", "Hard to Find",
         "Manual testing rarely covers combinations that match no rule — invisible by definition"),
        ("✅", "The Fix",
         "PI finds ALL gaps mathematically — add explicit deny-by-default or new rules"),
    ]
    for i, (icon, title, desc) in enumerate(dangers):
        y = 5.4 + i * 0.55
        add_box(slide, 0.7, y, 0.6, 0.4, icon, RED_LT, RED, font_size=14)
        add_label(slide, 1.4, y, 2.0, title, font_size=12, color=RED, bold=True, align=PP_ALIGN.LEFT)
        add_label(slide, 3.4, y, 9.3, desc, font_size=11, color=DARK, align=PP_ALIGN.LEFT)


def slide_15_verus(prs):
    """Slide 15: Verified Engine — Verus formal verification of RVM."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Verified Engine — Proving the Engine Itself Correct",
                  "Verus Formal Verification of the RVM (Rego Virtual Machine)")

    # ── LEFT: The problem ──
    add_group_box(slide, 0.4, 1.5, 4.0, 2.5,
                  "❓  The Question", INDIGO_BG, INDIGO)
    add_box(slide, 0.7, 2.2, 3.4, 1.5,
            "Z3 proves properties about\nyour POLICIES...\n\nBut who proves the ENGINE\nthat runs them is correct?",
            INDIGO_LT, INDIGO, font_size=13)

    # Arrow to answer
    add_arrow(slide, 4.4, 2.75, 5.3, 2.75, GRAY)

    # ── CENTER: Verus approach ──
    add_group_box(slide, 5.3, 1.5, 3.5, 2.5,
                  "🛡️  Verus — Rust Verification", ORANGE_BG, ORANGE)
    add_box(slide, 5.6, 2.2, 2.9, 1.5,
            "Verus verifies Rust code\nat compile time\n\nMathematical proofs that\nthe RVM implementation\nmatches its specification",
            ORANGE_LT, ORANGE, font_size=12)

    # Arrow to what we verify
    add_arrow(slide, 8.8, 2.75, 9.5, 2.75, GREEN)

    # ── RIGHT: What's verified ──
    add_group_box(slide, 9.5, 1.5, 3.5, 2.5,
                  "✅  Verified Properties", GREEN_BG, GREEN)
    add_box(slide, 9.8, 2.1, 2.9, 0.5, "No register out-of-bounds", GREEN_LT, GREEN, font_size=11, bold=True)
    add_box(slide, 9.8, 2.7, 2.9, 0.5, "Instruction limits enforced", GREEN_LT, GREEN, font_size=11, bold=True)
    add_box(slide, 9.8, 3.3, 2.9, 0.5, "Value invariants preserved", GREEN_LT, GREEN, font_size=11, bold=True)

    # ── Bottom: The verification stack ──
    add_group_box(slide, 0.4, 4.5, 12.5, 2.2,
                  "Verification Stack — Two Layers of Formal Guarantees", BLUE_BG, BLUE)

    # Layer 1: Z3
    add_box(slide, 0.8, 5.1, 5.5, 1.0,
            "Layer 1: Z3 Analysis\n"
            "Proves properties about YOUR POLICIES\n"
            "(test gen, diff, compliance, gap detection)",
            ORANGE_LT, ORANGE, font_size=12, bold=False)
    add_label(slide, 3.0, 6.2, 1.5, "External → Z3 SMT solver",
              font_size=9, color=GRAY)

    # Plus
    add_label(slide, 6.5, 5.3, 0.5, "+", font_size=24, color=BLUE, bold=True)

    # Layer 2: Verus
    add_box(slide, 7.2, 5.1, 5.5, 1.0,
            "Layer 2: Verus Verification\n"
            "Proves the ENGINE ITSELF is correct\n"
            "(no crashes, no memory errors, spec compliance)",
            GREEN_LT, GREEN, font_size=12, bold=False)
    add_label(slide, 9.5, 6.2, 1.5, "Internal → Rust compile-time proofs",
              font_size=9, color=GRAY)

    # Tagline
    add_label(slide, 0.4, 6.9, 12.5,
              "Verified policies running on a verified engine — end-to-end mathematical trust",
              font_size=13, color=DARK, bold=True)


def slide_16_value_prop(prs):
    """Slide 16: Business value summary."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Business Impact",
                  "Why This Matters for Your Service")

    items = [
        ("🔬", "Eliminate Coverage Gaps",
         "Automatically test every reachable decision path — including edge cases engineers would miss",
         GREEN_LT, GREEN),
        ("🔄", "Safe Policy Changes",
         "Mathematically prove that a policy refactor doesn't change behavior before deploying to production",
         BLUE_LIGHT, BLUE),
        ("🔍", "Faster Incident Response",
         "When a request is denied, instantly identify which specific condition failed and why",
         ORANGE_LT, ORANGE),
        ("🛡️", "Compliance Evidence",
         "Produce mathematical proofs that security rules cannot be bypassed — not just test results",
         INDIGO_LT, INDIGO),
        ("⏱️", "Reduce QA Effort",
         "Replace weeks of manual test authoring with a single command that generates complete test suites",
         PURPLE_LT, PURPLE),
        ("🔒", "Verified Engine",
         "The policy engine itself is formally verified with Verus — proven correct at compile time, not just tested",
         RED_LT, RED),
    ]

    for i, (icon, title, desc, fill, border) in enumerate(items):
        y = 1.5 + i * 0.9
        # Icon
        add_box(slide, 0.5, y, 0.7, 0.7, icon, fill, border, font_size=18)
        # Title
        add_label(slide, 1.4, y, 3.0, title, font_size=14, color=DARK, bold=True, align=PP_ALIGN.LEFT)
        # Description
        add_label(slide, 1.4, y + 0.32, 11.0, desc, font_size=11, color=GRAY, align=PP_ALIGN.LEFT)


def slide_17_closing(prs):
    """Slide 17: Closing / next steps."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(0x0D, 0x47, 0xA1))

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(11), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Policy Intelligence"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Key points
    points = [
        "Not just testing — mathematical proof",
        "Aviation-grade MC/DC coverage for any policy",
        "Copilot integration — AI writes, math verifies",
        "Gap detection, SoD proofs, regulatory scanning",
        "Verified engine — Verus proves the RVM itself correct",
        "Works with Rego (OPA), Cedar (AWS), and Azure Policy",
    ]

    y = 2.5
    for pt in points:
        add_box(slide, 2.5, y, 8.3, 0.55, "→  " + pt,
                RGBColor(0x0A, 0x36, 0x7A), RGBColor(0x1A, 0x73, 0xE8),
                font_size=15, font_color=WHITE, align=PP_ALIGN.LEFT)
        y += 0.7

    # Bottom
    add_label(slide, 1, 6.8, 11,
              "Regorus — microsoft/regorus — Policy Framework Team",
              font_size=14, color=RGBColor(0x64, 0xB5, 0xF6))


# ── Main ────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    # Widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_title(prs)                # 0: Title
    slide_what_is_z3(prs)           # 1: What is Z3 / Theorem Proving?
    slide_policy_to_smt(prs)        # 2: Policy → SMT Under the Hood
    slide_01_overview(prs)          # 3: Overview
    slide_02_before_after(prs)      # 4: Manual vs Auto
    slide_03_policy_migration(prs)  # 5: Safe Migration
    slide_04_why_denied(prs)        # 6: Why Denied
    slide_05_compliance(prs)        # 7: Compliance Proofs
    slide_06_mcdc(prs)              # 8: MC/DC
    slide_07_decision_tree(prs)     # 9: Decision Tree
    slide_08_real_results(prs)      # 10: Real Results
    slide_section_vision(prs)       # 11: Section Divider — What's Next
    slide_09_copilot_pi_loop(prs)   # 12: Copilot + PI
    slide_10_pr_review_bot(prs)     # 13: PR Review Bot
    slide_11_shadow_prediction(prs) # 14: Shadow Mode Prediction
    slide_12_regulatory_impact(prs) # 15: Regulatory Impact
    slide_13_separation_of_duty(prs)# 16: Separation of Duty
    slide_14_gap_detection(prs)     # 17: Gap Detection
    slide_15_verus(prs)             # 18: Verified Engine (Verus)
    slide_16_value_prop(prs)        # 19: Business Impact
    slide_17_closing(prs)           # 20: Closing

    # ── Speaker Notes ──────────────────────────────────────────────────────
    # Add presenter notes with talking points and live-demo commands.
    # BIN is the regorus binary path — the presenter should set this.
    # ───────────────────────────────────────────────────────────────────────
    BIN = "regorus"  # or: target/debug/examples/regorus
    D = "examples/demos"

    notes = {}

    notes[0] = """\
TITLE SLIDE — Policy Intelligence

Talking points:
• Regorus is Microsoft's Rego policy engine (open source).
• "Policy Intelligence" = the Z3 symbolic-analysis layer on top.
• Works with Rego (OPA), Cedar (AWS), and Azure Policy definitions.
• One engine — multiple policy languages, all get formal guarantees.

SETUP (run once before the presentation):
  cd /path/to/regorus
  cargo build --example regorus --features z3-analysis,cedar,azure_policy
  export BIN="target/debug/examples/regorus"
"""

    notes[1] = """\
WHAT IS THEOREM PROVING? — Z3 Introduction

Talking points:
• Traditional testing: pick some inputs, run them, hope for the best.
• Theorem proving: check ALL inputs at once using mathematical reasoning.
• Z3 is an SMT solver built by Microsoft Research — one of the most
  widely used theorem provers in the world.
• Key concepts for the audience:
  - SAT = "satisfiable" — Z3 found a concrete input (and returns it).
  - UNSAT = "unsatisfiable" — Z3 mathematically proved no such input exists.
• Used in: Windows driver verification, Azure, AWS, compilers, cryptography.
• Our contribution: translate policy languages into Z3 constraints automatically.

Analogy for PM audience:
  "Unit tests are flashlights that illuminate spots in a dark room.
   Z3 is a floodlight that illuminates the entire room at once."
"""

    notes[2] = f"""\
POLICY → SMT: UNDER THE HOOD

Talking points:
• Show a Cedar policy (2 rules: permit + forbid) on the left.
• Show the auto-generated SMT constraints on the right.
• Walk through the correspondence:
  - Each input field becomes a Z3 variable (declare-fun).
  - The permit condition becomes an AND of constraints.
  - The forbid condition becomes a negated AND.
  - String patterns like "10.*" become regex constraints (str.in_re).
• The bottom shows Z3's answer: a concrete JSON input that satisfies all constraints.
• This is fully automatic — the presenter never writes SMT by hand.

DEMO — Generate the SMT yourself:
  {BIN} analyze \\
    -d examples/cedar/examples/iam_zero_trust/policy.cedar \\
    -d examples/cedar/examples/iam_zero_trust/entities.json \\
    -e cedar.authorize -o 1 \\
    --dump-smt /tmp/cedar_iam.smt2

  cat /tmp/cedar_iam.smt2
  → Shows the exact SMT encoding from the slide.
"""

    notes[3] = f"""\
OVERVIEW — From Policies to Proven Guarantees

Talking points:
• Three inputs: policy files, optional JSON schema, reference data.
• Engine compiles to RVM bytecode, translates to Z3 SMT constraints.
• Four outputs: test suites, change impact, root cause, compliance proofs.

DEMO — Show basic analysis (synthesize an input):
  {BIN} analyze \\
    -d {D}/container_admission.rego \\
    -e data.container_admission.allow \\
    -o false \\
    -i {D}/container_admission_input.json \\
    -s {D}/container_admission_schema.json \\
    --max-loops 3

  → Z3 instantly produces a concrete JSON input that violates the policy.
"""

    notes[4] = f"""\
AUTOMATED TEST GENERATION — Zero Manual Effort, Full Coverage

Talking points:
• Today: engineers write tests by hand, miss edge cases, leave gaps.
• With PI: one command → mathematically complete test suite.
• Every reachable decision path gets at least one test.

DEMO — Generate tests for the server policy:
  {BIN} gen-tests \\
    -d examples/server/allowed_server.rego \\
    -e data.example.allow \\
    -i examples/server/input.json \\
    -s examples/server/input_schema.json \\
    --max-loops 3 --max-tests 10

  → Shows JSON array of auto-generated test cases.
"""

    notes[5] = f"""\
SAFE POLICY MIGRATION — Prove Equivalence Before You Deploy

Talking points:
• Policy refactors are risky — how do you know behavior didn't change?
• PI compares ALL possible inputs, not just a sample.
• If identical: UNSAT proof. If different: exact counterexample.

DEMO — Diff two network segmentation versions:
  {BIN} diff \\
    --policy1 {D}/network_segmentation.rego \\
    --policy2 {D}/network_segmentation_v2.rego \\
    -e data.network_segmentation.compliant \\
    -i {D}/network_segmentation_input.json \\
    -s {D}/network_segmentation_schema.json \\
    --max-loops 3

  → Z3 finds a PII-related input where v1 and v2 disagree.

DEMO — Subsumption check (is v2 at least as permissive as v1?):
  {BIN} subsumes \\
    --old {D}/network_segmentation.rego \\
    --new {D}/network_segmentation_v2.rego \\
    -e data.network_segmentation.compliant \\
    -i {D}/network_segmentation_input.json \\
    -s {D}/network_segmentation_schema.json \\
    --max-loops 3
"""

    notes[6] = f"""\
WHY DENIED — Pinpoint the Exact Failing Condition

Talking points:
• Support tickets: "Why was my request blocked?"
• PI uses MAX-SAT to find the minimal set of failing conditions.
• Shows which specific condition failed AND how to fix it.

DEMO — Targeted analysis (cover a specific deny line):
  {BIN} analyze \\
    -d {D}/ags_group_governance.rego \\
    -e data.graph.elm_governance_group_membership.deny_result \\
    -o true \\
    -i {D}/ags_group_governance_input.json \\
    -s {D}/ags_group_governance_schema.json

  → Z3 produces an input that triggers the deny path,
    showing which condition (e.g., managedState != "enforced") caused it.
"""

    notes[7] = f"""\
COMPLIANCE PROOFS — Mathematical Guarantees

Talking points:
• "Can any request bypass this rule?" → UNSAT = impossible.
• "Is the new policy at least as restrictive?" → Subsumption proof.
• "Are there unreachable (dead) rules?" → Z3 checks reachability.

DEMO — Prove v2 subsumes v1 (stricter deny behavior):
  {BIN} subsumes \\
    --old {D}/azure_storage_https_v1_definition.json \\
    --new {D}/azure_storage_https_v2_definition.json \\
    -e main -o '"deny"' \\
    --azure-aliases {D}/azure_policy_aliases.json \\
    -i {D}/azure_storage_input.json \\
    -s {D}/azure_storage_schema.json \\
    --max-loops 3

  → "New subsumes old" = v2 denies everything v1 denies (and more).
"""

    notes[8] = f"""\
MC/DC CONDITION COVERAGE — Aviation-Grade Testing

Talking points:
• MC/DC (Modified Condition/Decision Coverage) is required by DO-178C.
• Each condition is independently shown to affect the outcome.
• Applied here to policy rules, not avionics — same rigor.

DEMO — Generate MC/DC tests for the AGS policy:
  {BIN} gen-tests \\
    -d {D}/ags_group_governance.rego \\
    -e data.graph.elm_governance_group_membership.deny_result \\
    -i {D}/ags_group_governance_input.json \\
    -s {D}/ags_group_governance_schema.json \\
    --condition-coverage \\
    --format annotated

  → 23 tests, 46/46 condition goals, 100% MC/DC.
  → The "annotated" format shows each condition's T/F per test.
"""

    notes[9] = f"""\
DECISION TREE — AGS Group Governance Policy

Talking points:
• Real Azure policy with 5 decision paths.
• Z3 automatically finds an input for each path.
• Show the annotated output — each test covers specific conditions.

DEMO — Same command as MC/DC slide (with annotated output):
  {BIN} gen-tests \\
    -d {D}/ags_group_governance.rego \\
    -e data.graph.elm_governance_group_membership.deny_result \\
    -i {D}/ags_group_governance_input.json \\
    -s {D}/ags_group_governance_schema.json \\
    --condition-coverage --format annotated

  → Walk through the annotated source listing: each line shows
    which tests cover it and what condition values they use.
"""

    notes[10] = f"""\
REAL RESULTS — AGS Group Governance Policy

Talking points:
• 139-line policy, 5 decision paths.
• 23 auto-generated tests, 34/34 lines, 46/46 MC/DC conditions.
• Zero manual effort — all from one command.

DEMO — Show JSON output format:
  {BIN} gen-tests \\
    -d {D}/ags_group_governance.rego \\
    -e data.graph.elm_governance_group_membership.deny_result \\
    -i {D}/ags_group_governance_input.json \\
    -s {D}/ags_group_governance_schema.json \\
    --condition-coverage --format json | head -80

  → Scroll through the JSON — each test has input fields and
    the lines/conditions it covers.
"""

    notes[11] = """\
SECTION DIVIDER — What's Next

Talking points:
• "Everything you've seen so far works today — one command."
• "Now let's look at high-impact scenarios this unlocks."
• Use this as a natural pause / audience check-in.
"""

    notes[12] = """\
COPILOT + POLICY INTELLIGENCE

Talking points:
• Vision: Copilot generates policy from natural language intent.
• PI verifies the generated policy against formal properties.
• If a violation is found, the counterexample is fed back as an
  auto-fix prompt — Copilot iterates until the proof passes.
• Every AI-generated policy gets a math proof before production.

(No live demo — this is a future integration concept.)
"""

    notes[13] = f"""\
PR REVIEW BOT — Policy Diff in CI

Talking points:
• Every PR that modifies a .rego file triggers a formal diff.
• If behavior is identical: auto-approve.
• If behavior changes: the PR comment shows exact inputs with
  before/after outcomes.
• Wider access detected → mandatory security review.

DEMO — Simulate a PR diff (v1 → v2):
  {BIN} diff \\
    --policy1 examples/server/allowed_server.rego \\
    --policy2 examples/server/allowed_server_v2.rego \\
    -e data.example.allow \\
    -i examples/server/input.json \\
    -s examples/server/input_schema.json \\
    --max-loops 3

  → Z3 finds: telnet server allowed by v2, denied by v1.
  → This is exactly what the PR bot would post as a comment.
"""

    notes[14] = f"""\
SHADOW MODE — WITHOUT THE SHADOW

Talking points:
• Shadow mode: deploy to separate fleet, wait days/weeks.
• PI: compare old vs new in seconds, covering ALL inputs.
• Example: "99.7% identical, 0.3% affected: contractors on internal-docs."
• No infrastructure cost, no traffic dependency.

DEMO — Diff as an instant "shadow mode":
  {BIN} diff \\
    --policy1 {D}/network_segmentation.rego \\
    --policy2 {D}/network_segmentation_v2.rego \\
    -e data.network_segmentation.compliant \\
    -i {D}/network_segmentation_input.json \\
    -s {D}/network_segmentation_schema.json \\
    --max-loops 3

  → "This is what shadow mode does in weeks — we do it in seconds."
"""

    notes[15] = f"""\
REGULATORY IMPACT ANALYSIS

Talking points:
• Formalize a regulation as a property (e.g., MFA for PII).
• Scan all policies fleet-wide — each result is a math proof.
• Counterexamples show exactly HOW the gap can be exploited.

DEMO — Azure Policy: find resources that violate HTTPS-only:
  {BIN} analyze \\
    -d {D}/azure_storage_https_v1_definition.json \\
    -e main -o '"deny"' \\
    --azure-aliases {D}/azure_policy_aliases.json \\
    -i {D}/azure_storage_input.json \\
    -s {D}/azure_storage_schema.json \\
    --max-loops 3

  → Z3 finds: supportsHttpsTrafficOnly=false triggers deny.
  → "Now imagine running this across 42 service policies."
"""

    notes[16] = """\
SEPARATION OF DUTY — SOX / SOC-2 / HIPAA

Talking points:
• Security property: no user can both submit AND approve.
• Z3 exhaustively checks ALL role combinations.
• Proven → mathematical certificate for auditors.
• Violation → exact counterexample (user, roles, amount).
• Works for any property expressible as ∀ user: P(user) → ¬Q(user).

(No standalone demo command for SoD — compose with analyze
 using cover-line / avoid-line to target specific rule paths.)
"""

    notes[17] = f"""\
GAP DETECTION — Find Undefined Decisions

Talking points:
• "Gap" = an input that matches NO rule → UNDEFINED result.
• In some engines, undefined defaults to ALLOW — silent backdoor.
• PI finds ALL gaps mathematically — manual testing can't do this.
• Fix: add explicit deny-by-default rules.

DEMO — Find an input where result is defined but also try:
  {BIN} analyze \\
    -d {D}/container_admission.rego \\
    -e data.container_admission.allow \\
    -i {D}/container_admission_input.json \\
    -s {D}/container_admission_schema.json \\
    --max-loops 3

  (Without -o flag: Z3 finds any input where the result IS defined.
   Contrast with the error when no satisfying input exists.)
"""

    notes[18] = """\
VERIFIED ENGINE — Verus Formal Verification of the RVM

Talking points:
• Z3 proves properties about policies. But who ensures the engine is correct?
• Verus is a Rust verification framework — proves code correct at compile time.
• We're using Verus to formally verify the RVM (Rego Virtual Machine):
  - Register bounds: no out-of-bounds access, ever (proven, not just tested).
  - Instruction limits: DoS protection is mathematically enforced.
  - Value invariants: arrays, objects, sets maintain well-formedness.
• This gives TWO layers of formal guarantees:
  Layer 1: Z3 proves your policies are correct.
  Layer 2: Verus proves the engine running them is correct.
• The verification is staged — starting with core safety properties
  and expanding to full functional correctness.

Why this matters for PMs:
  "Even if your policy is perfect, a bug in the engine could undermine everything.
   Verus eliminates that risk — proven at compile time, not discovered in prod."

Reference: verus.md in the repo for the detailed verification plan.
"""

    notes[19] = """\
BUSINESS IMPACT — Why This Matters

Talking points:
• Coverage gaps → security incidents. PI eliminates them.
• Safe policy changes → fewer rollbacks, faster deployments.
• Faster incident response → lower MTTR.
• Compliance evidence → auditors get proofs, not test reports.
• QA effort → weeks of manual testing → one command.
• Verified engine → the engine itself is proven crash-free.

Transition: "Let's wrap up with what's next."
"""

    notes[20] = """\
CLOSING — Policy Intelligence

Talking points:
• Not just testing — mathematical proof.
• Works today with Rego, Cedar, Azure Policy.
• Verified engine — Verus proves the RVM correct.
• Open source: microsoft/regorus on GitHub.
• Contact: Policy Framework Team.

Call to action:
  "Try it on YOUR policies — one command, zero manual effort."
"""

    # Apply notes to each slide
    for idx, text in notes.items():
        add_notes(prs.slides[idx], text)

    # ── Hidden Demo Cheat Sheet slide ──────────────────────────────────────
    # This slide is hidden (not shown during slideshow) and contains
    # ALL demo commands in the speaker notes for quick reference.
    cheat = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(cheat, RGBColor(0x2C, 0x2C, 0x2C))
    add_label(cheat, 1, 0.5, 11,
              "DEMO CHEAT SHEET (hidden slide — presenter reference only)",
              font_size=24, color=WHITE, bold=True)

    # Mark slide as hidden via direct XML manipulation
    from lxml import etree
    # In OOXML, a slide is hidden by adding show="0" to the <p:sld> element
    cheat_element = cheat._element
    cheat_element.set('show', '0')

    # Visual command list on the slide itself (for quick glance)
    commands_on_slide = [
        ("Analyze (find violation)", f"{BIN} analyze -d {D}/container_admission.rego -e data.container_admission.allow -o false -i ... -s ... --max-loops 3"),
        ("Gen-Tests (MC/DC)", f"{BIN} gen-tests -d {D}/ags_group_governance.rego -e data.graph.elm_governance_group_membership.deny_result -i ... -s ... --condition-coverage --format annotated"),
        ("Diff (two versions)", f"{BIN} diff --policy1 {D}/network_segmentation.rego --policy2 {D}/network_segmentation_v2.rego -e data.network_segmentation.compliant -i ... -s ..."),
        ("Subsumes (strictness)", f"{BIN} subsumes --old {D}/azure_storage_https_v1_definition.json --new {D}/azure_storage_https_v2_definition.json -e main -o '\"deny\"' --azure-aliases ..."),
        ("Azure Policy", f"{BIN} analyze -d {D}/azure_storage_https_v1_definition.json -e main -o '\"deny\"' --azure-aliases {D}/azure_policy_aliases.json -i ... -s ..."),
        ("Cedar (IAM)", f"{BIN} analyze -d examples/cedar/examples/iam_zero_trust/policy.cedar -d .../entities.json -e cedar.authorize -o 1"),
    ]
    y = 1.3
    for label, cmd in commands_on_slide:
        add_box(cheat, 0.5, y, 2.5, 0.7, label, ORANGE_LT, ORANGE, font_size=11, bold=True)
        add_box(cheat, 3.2, y, 9.5, 0.7, cmd, RGBColor(0x3C, 0x3C, 0x3C), GRAY, font_size=9, font_color=WHITE, align=PP_ALIGN.LEFT)
        y += 0.85

    # Full commands in notes
    cheat_notes = f"""\
═══════════════════════════════════════════════════════════
  DEMO CHEAT SHEET — COPY-PASTE COMMANDS
  (all paths relative to the regorus repo root)
═══════════════════════════════════════════════════════════

PREREQUISITE — Build once:
  BINDGEN_EXTRA_CLANG_ARGS="-I/opt/homebrew/include" \\
  LIBRARY_PATH="/opt/homebrew/lib" \\
  cargo build --example regorus --features z3-analysis,cedar,azure_policy

  BIN=target/debug/examples/regorus

───────────────────────────────────────────────────────────
1) ANALYZE — Synthesize a violating input
───────────────────────────────────────────────────────────
  $BIN analyze \\
    -d {D}/container_admission.rego \\
    -e data.container_admission.allow \\
    -o false \\
    -i {D}/container_admission_input.json \\
    -s {D}/container_admission_schema.json \\
    --max-loops 3

───────────────────────────────────────────────────────────
2) ANALYZE — Compliant input
───────────────────────────────────────────────────────────
  $BIN analyze \\
    -d {D}/container_admission.rego \\
    -e data.container_admission.allow \\
    -o true \\
    -i {D}/container_admission_input.json \\
    -s {D}/container_admission_schema.json \\
    --max-loops 3

───────────────────────────────────────────────────────────
3) ANALYZE — Targeted line coverage
───────────────────────────────────────────────────────────
  $BIN analyze \\
    -d {D}/container_admission.rego \\
    -e data.container_admission.allow \\
    -o false \\
    -l container_admission.rego:101 \\
    --avoid-line container_admission.rego:75 \\
    -i {D}/container_admission_input.json \\
    -s {D}/container_admission_schema.json \\
    --max-loops 3

───────────────────────────────────────────────────────────
4) GEN-TESTS — MC/DC with annotated output ★ BEST DEMO ★
───────────────────────────────────────────────────────────
  $BIN gen-tests \\
    -d {D}/ags_group_governance.rego \\
    -e data.graph.elm_governance_group_membership.deny_result \\
    -i {D}/ags_group_governance_input.json \\
    -s {D}/ags_group_governance_schema.json \\
    --condition-coverage \\
    --format annotated

───────────────────────────────────────────────────────────
5) GEN-TESTS — JSON output (for piping/automation)
───────────────────────────────────────────────────────────
  $BIN gen-tests \\
    -d {D}/ags_group_governance.rego \\
    -e data.graph.elm_governance_group_membership.deny_result \\
    -i {D}/ags_group_governance_input.json \\
    -s {D}/ags_group_governance_schema.json \\
    --condition-coverage --format json

───────────────────────────────────────────────────────────
6) DIFF — Network segmentation v1 vs v2
───────────────────────────────────────────────────────────
  $BIN diff \\
    --policy1 {D}/network_segmentation.rego \\
    --policy2 {D}/network_segmentation_v2.rego \\
    -e data.network_segmentation.compliant \\
    -i {D}/network_segmentation_input.json \\
    -s {D}/network_segmentation_schema.json \\
    --max-loops 3

───────────────────────────────────────────────────────────
7) DIFF — Server allowed v1 vs v2 (telnet removal)
───────────────────────────────────────────────────────────
  $BIN diff \\
    --policy1 examples/server/allowed_server.rego \\
    --policy2 examples/server/allowed_server_v2.rego \\
    -e data.example.allow \\
    -i examples/server/input.json \\
    -s examples/server/input_schema.json \\
    --max-loops 3

───────────────────────────────────────────────────────────
8) SUBSUMES — Azure Policy v2 ⊇ v1
───────────────────────────────────────────────────────────
  $BIN subsumes \\
    --old {D}/azure_storage_https_v1_definition.json \\
    --new {D}/azure_storage_https_v2_definition.json \\
    -e main -o '"deny"' \\
    --azure-aliases {D}/azure_policy_aliases.json \\
    -i {D}/azure_storage_input.json \\
    -s {D}/azure_storage_schema.json \\
    --max-loops 3

───────────────────────────────────────────────────────────
9) SUBSUMES — Reverse (expect failure + counterexample)
───────────────────────────────────────────────────────────
  $BIN subsumes \\
    --old {D}/azure_storage_https_v2_definition.json \\
    --new {D}/azure_storage_https_v1_definition.json \\
    -e main -o '"deny"' \\
    --azure-aliases {D}/azure_policy_aliases.json \\
    -i {D}/azure_storage_input.json \\
    -s {D}/azure_storage_schema.json \\
    --max-loops 3

───────────────────────────────────────────────────────────
10) AZURE POLICY — Find deny-triggering resource
───────────────────────────────────────────────────────────
  $BIN analyze \\
    -d {D}/azure_storage_https_v1_definition.json \\
    -e main -o '"deny"' \\
    --azure-aliases {D}/azure_policy_aliases.json \\
    -i {D}/azure_storage_input.json \\
    -s {D}/azure_storage_schema.json \\
    --max-loops 3

───────────────────────────────────────────────────────────
11) CEDAR — IAM Zero Trust (permitted request)
───────────────────────────────────────────────────────────
  $BIN analyze \\
    -d examples/cedar/examples/iam_zero_trust/policy.cedar \\
    -d examples/cedar/examples/iam_zero_trust/entities.json \\
    -e cedar.authorize -o 1

───────────────────────────────────────────────────────────
12) CEDAR — Healthcare HIPAA (permitted request)
───────────────────────────────────────────────────────────
  $BIN analyze \\
    -d examples/cedar/examples/hipaa_healthcare/policy.cedar \\
    -d examples/cedar/examples/hipaa_healthcare/entities.json \\
    -e cedar.authorize -o 1

───────────────────────────────────────────────────────────
13) RUN ALL DEMOS (the full script):
───────────────────────────────────────────────────────────
  bash {D}/run_demos.sh

═══════════════════════════════════════════════════════════
TIP: Use Presenter View (Cmd+Shift+Enter on Mac) to see
these notes alongside the slides. Keep a terminal in
split-screen for live demos.
═══════════════════════════════════════════════════════════
"""
    add_notes(cheat, cheat_notes)

    out_path = os.path.join(os.path.dirname(__file__), "Policy_Intelligence.pptx")
    prs.save(out_path)
    print(f"✅ Saved: {out_path}")
    print(f"   {len(prs.slides)} slides, all shapes editable")


if __name__ == "__main__":
    main()
