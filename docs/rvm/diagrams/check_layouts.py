#!/usr/bin/env python3
"""Check available slide layouts in SlidesToReuse.pptx."""
from pptx import Presentation
import os
base = os.path.dirname(os.path.abspath(__file__))
prs = Presentation(os.path.join(base, "SlidesToReuse.pptx"))
for i, layout in enumerate(prs.slide_layouts):
    print(f"  [{i}] {layout.name}  (placeholders: {len(layout.placeholders)})")
