#!/usr/bin/env python3
"""Extract images from SlidesToReuse.pptx and list SmartArt/diagram info."""
from pptx import Presentation
from pptx.oxml.ns import qn
import os, hashlib

base = os.path.dirname(os.path.abspath(__file__))
path = os.path.join(base, "SlidesToReuse.pptx")
out_dir = os.path.join(base, "extracted_images")
os.makedirs(out_dir, exist_ok=True)

prs = Presentation(path)

img_idx = 0
for si, slide in enumerate(prs.slides, 1):
    print(f"=== SLIDE {si} ===")
    for sh in slide.shapes:
        # Images
        if sh.shape_type == 13:  # PICTURE
            img_idx += 1
            img = sh.image
            ext = img.content_type.split("/")[-1]
            if ext == "jpeg":
                ext = "jpg"
            fname = f"slide{si}_img{img_idx}.{ext}"
            fpath = os.path.join(out_dir, fname)
            with open(fpath, "wb") as f:
                f.write(img.blob)
            size_kb = len(img.blob) / 1024
            print(f"  Image: {sh.name}")
            print(f"    Saved: {fname} ({size_kb:.0f} KB, {img.content_type})")
            print(f"    Position: ({sh.left/914400:.2f}\", {sh.top/914400:.2f}\")")
            print(f"    Size: {sh.width/914400:.2f}\" x {sh.height/914400:.2f}\"")
            # Try to get alt text
            sp = sh._element
            nvp = sp.find(qn("p:nvPicPr"))
            if nvp is not None:
                cnvpr = nvp.find(qn("p:cNvPr"))
                if cnvpr is None:
                    cnvpr = nvp.find(qn("p:nvPr"))
            # try nvSpPr path
            for tag in ["p:nvPicPr/p:cNvPr", "p:nvSpPr/p:cNvPr"]:
                pass
            descr = sp.attrib.get("descr", "")
            if not descr:
                # Look in cNvPr
                for el in sp.iter():
                    if el.tag.endswith("}cNvPr") or el.tag == "cNvPr":
                        descr = el.attrib.get("descr", "")
                        if descr:
                            break
            if descr:
                print(f"    Alt text: {descr}")
            print()

        # Diagrams / SmartArt
        if sh.shape_type is None:
            print(f"  SmartArt/Diagram: {sh.name}")
            print(f"    Position: ({sh.left/914400:.2f}\", {sh.top/914400:.2f}\")")
            print(f"    Size: {sh.width/914400:.2f}\" x {sh.height/914400:.2f}\"")
            # Try to extract text from diagram
            texts = []
            for el in sh._element.iter():
                if el.tag.endswith("}t"):
                    t = el.text
                    if t and t.strip():
                        texts.append(t.strip())
            if texts:
                print(f"    Diagram text: {texts}")
            print()

        # Text content summary
        if sh.has_text_frame and sh.text_frame.text.strip():
            txt = sh.text_frame.text.strip().replace("\n", " | ")
            if len(txt) > 120:
                txt = txt[:120] + "..."
            if sh.shape_type != 13:   # skip if already printed as image
                pass  # already printed in previous script

print(f"\nImages extracted to: {out_dir}")
print(f"Total images: {img_idx}")
