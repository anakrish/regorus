#!/usr/bin/env python3
"""Generate the final RVM presentation — demo-driven, minimal slides.

Structure (11 slides):
  1    Problem — engine proliferation (no solution revealed yet)
  2    CTO Memo — preserved from SlidesToReuse.pptx (SmartArt, gradient)
  3    Regorus Overview — what it is + adoption (open source & internal)
  4    JVM / RVM Analogy — solution reveal
  5    DEMO 1 — RVM Playground (Rego + Cedar + Azure Policy)
  6    Cedar Contrast — A Comparison with Amazon's Cedar
  7    Policy Intelligence — Z3 capabilities + MSR quote + AWS context
  8    DEMO 2 — Formal Verification
  9    Copilot + Policy Intelligence — AI writes, math verifies
  10   PR Review Bot + Regulatory Impact — combined vision
  11   Closing

Output: RVM-Final.pptx
"""

import os, copy
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

try:
    from pptx.enum.text import MSO_ANCHOR
except ImportError:
    MSO_ANCHOR = None

BASE = os.path.dirname(os.path.abspath(__file__))
REUSE = os.path.join(BASE, "SlidesToReuse.pptx")

# ── Dimensions ───────────────────────────────────────────────────────────
SW, SH = 13.333, 7.5
M = 0.6
U = SW - 2 * M

# ── Palette ──────────────────────────────────────────────────────────────
CHARCOAL    = RGBColor(0x2D, 0x2D, 0x2D)
DARK_TEXT   = RGBColor(0x33, 0x33, 0x33)
MID_GRAY    = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT_GRAY  = RGBColor(0xE0, 0xE0, 0xE0)
FAINT_GRAY  = RGBColor(0xF5, 0xF5, 0xF5)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

GREEN_D   = RGBColor(0x1B, 0x5E, 0x20)
GREEN_M   = RGBColor(0x2E, 0x7D, 0x32)
GREEN_L   = RGBColor(0xC8, 0xE6, 0xC9)
GREEN_XL  = RGBColor(0xE8, 0xF5, 0xE9)
BLUE_D    = RGBColor(0x15, 0x65, 0xC0)
BLUE_L    = RGBColor(0xE3, 0xF2, 0xFD)
ORANGE_D  = RGBColor(0xEF, 0x6C, 0x00)
ORANGE_L  = RGBColor(0xFF, 0xF3, 0xE0)
RED_D     = RGBColor(0xC6, 0x28, 0x28)
RED_L     = RGBColor(0xFF, 0xCD, 0xD2)
PURPLE_D  = RGBColor(0x7B, 0x1F, 0xA2)
PURPLE_L  = RGBColor(0xF3, 0xE5, 0xF5)
TEAL_D    = RGBColor(0x00, 0x69, 0x5C)
AZURE     = RGBColor(0x00, 0x78, 0xD4)
DEEP_PURPLE = RGBColor(0x4A, 0x14, 0x8C)


# ── Helpers ──────────────────────────────────────────────────────────────

def I(v): return Inches(v)

def dist(n, w, total=U, start=M):
    if n <= 1: return [start + (total - w) / 2]
    gap = (total - n * w) / (n - 1)
    return [start + i * (w + gap) for i in range(n)]

def setbg(slide, color):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def shape(slide, l, t, w, h, fill, border=None, txt="",
          sz=14, color=DARK_TEXT, bold=False,
          align=PP_ALIGN.CENTER, rr=True):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rr else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(st, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border; sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    tf = sh.text_frame; tf.word_wrap = True; tf.auto_size = None
    tf.margin_left = Pt(8); tf.margin_right = Pt(8)
    tf.margin_top = Pt(4); tf.margin_bottom = Pt(4)
    if MSO_ANCHOR:
        try: tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except: pass
    if txt:
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = txt
        r.font.size = Pt(sz); r.font.color.rgb = color; r.font.bold = bold
    return sh

def text(slide, l, t, w, h, txt, sz=14, color=DARK_TEXT,
         bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = txt
    r.font.size = Pt(sz); r.font.color.rgb = color; r.font.bold = bold
    return tb

def text_italic(slide, l, t, w, h, txt, sz=14, color=DARK_TEXT,
                bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = txt
    r.font.size = Pt(sz); r.font.color.rgb = color; r.font.bold = bold
    r.font.italic = True
    return tb

def hyperlink_text(slide, l, t, w, h, txt, url, sz=14, color=AZURE,
                   bold=False, align=PP_ALIGN.LEFT, underline=True):
    """Add a textbox whose text is a clickable hyperlink."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = txt
    r.font.size = Pt(sz); r.font.color.rgb = color
    r.font.bold = bold; r.font.underline = underline
    r.hyperlink.address = url
    return tb

def bullets(slide, l, t, w, h, items, sz=14, color=DARK_TEXT,
            spacing=Pt(8), bold=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        r = p.add_run(); r.text = item
        r.font.size = Pt(sz); r.font.color.rgb = color; r.font.bold = bold
    return tb

def delete_slide(prs, index):
    rIdList = prs.slides._sldIdLst
    sldId = list(rIdList)[index]
    rId = sldId.get(qn("r:id"))
    prs.part.drop_rel(rId)
    rIdList.remove(sldId)

def move_slide_to_front(prs, src_index):
    """Move slide at src_index to position 0."""
    rIdList = prs.slides._sldIdLst
    items = list(rIdList)
    target = items[src_index]
    rIdList.remove(target)
    rIdList.insert(0, target)


# ═════════════════════════════════════════════════════════════════════════
#   SLIDE 1 — Problem (no solution revealed yet)
# ═════════════════════════════════════════════════════════════════════════

def slide_problem(prs, BL):
    s = prs.slides.add_slide(BL)
    setbg(s, WHITE)

    text(s, I(M), I(0.4), I(U), I(0.8),
         "The Problem", sz=44, color=RED_D, bold=True)

    text(s, I(M), I(1.3), I(U), I(0.6),
         "Azure Policy must support multiple policy languages.\n"
         "Across Microsoft, teams maintain separate engines.",
         sz=20, color=MID_GRAY)

    engines = ["OPA / Rego", "Cedar", "Azure Policy", "AKS Admission", "Custom"]
    bw = 2.1
    for name, x in zip(engines, dist(5, bw)):
        shape(s, I(x), I(2.3), I(bw), I(1.0),
              RED_L, RED_D, name, sz=16, color=RED_D, bold=True)

    shape(s, I(M), I(3.7), I(U), I(0.65), FAINT_GRAY, LIGHT_GRAY,
          "N engines  =  N \u00D7 audits  +  N \u00D7 test suites  +  "
          "N \u00D7 pipelines  +  N \u00D7 teams",
          sz=18, color=RED_D, bold=True)

    text(s, I(M), I(4.8), I(U), I(0.6),
         "Multiplicative cost. Inconsistent behavior. Fragmented tooling.",
         sz=20, color=DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)

    text(s, I(M), I(5.7), I(U), I(0.6),
         "Every new policy language multiplies the problem.",
         sz=18, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════
#   SLIDE 2 — CTO Memo (preserved from SlidesToReuse.pptx)
#   (handled via file merge — not generated here)
# ═════════════════════════════════════════════════════════════════════════


# ═════════════════════════════════════════════════════════════════════════
#   SLIDE 3 — Regorus Overview (diagrams from SlidesToReuse + adoption)
# ═════════════════════════════════════════════════════════════════════════

def slide_regorus_overview(prs, BL, image_blobs):
    """Regorus overview: CC diagram + perf diagram + Rust logo + adoption."""
    s = prs.slides.add_slide(BL)
    setbg(s, WHITE)

    CEDAR_PAPER = "https://arxiv.org/abs/2403.04651"

    # ── Title + Rust logo ──
    text(s, I(M), I(0.15), I(9.0), I(0.45),
         "Regorus Engine", sz=36, color=GREEN_D, bold=True)

    if "Picture 8" in image_blobs:
        s.shapes.add_picture(BytesIO(image_blobs["Picture 8"]),
                             I(11.4), I(0.1), I(1.4), I(0.93))

    # ── Subtitle ──
    text(s, I(M), I(0.58), I(10.5), I(0.3),
         "Rego-Rus(t) \u2014 A fast, lightweight, open-source Rego "
         "interpreter written in Rust.  "
         "Rigorous enforcer of well-defined Rego semantics.",
         sz=14, color=MID_GRAY)

    # ── Two diagrams side by side ──
    diag_y = 0.95

    # Left: CC diagram
    text(s, I(0.3), I(diag_y), I(5.5), I(0.22),
         "Origin: Confidential Containers", sz=12, color=DARK_TEXT, bold=True)
    if "Picture 2" in image_blobs:
        s.shapes.add_picture(BytesIO(image_blobs["Picture 2"]),
                             I(0.3), I(diag_y + 0.25), I(5.3), I(3.25))

    # Right: Performance diagram
    text(s, I(7.2), I(diag_y), I(5.8), I(0.22),
         "Regorus 6.5\u00D7 faster than Cedar, 80\u00D7 faster than OPA",
         sz=12, color=DARK_TEXT, bold=True)
    if "Picture 10" in image_blobs:
        s.shapes.add_picture(BytesIO(image_blobs["Picture 10"]),
                             I(7.2), I(diag_y + 0.25), I(4.0), I(2.95))

    # Cedar paper link under perf diagram
    hyperlink_text(s, I(7.2), I(diag_y + 3.25), I(5.0), I(0.22),
                   "Cedar paper: arxiv.org/abs/2403.04651", CEDAR_PAPER,
                   sz=10, color=AZURE, align=PP_ALIGN.LEFT)

    # ── Divider ──
    div_y = 4.55
    shape(s, I(M), I(div_y), I(U), I(0.03), LIGHT_GRAY, None, "", rr=False)

    # ── Open Source Adoption ──
    oss_y = div_y + 0.1
    text(s, I(M), I(oss_y), I(5.0), I(0.25),
         "Open Source Adoption", sz=17, color=GREEN_D, bold=True)

    oss = [
        ("Kata Containers",
         "Container runtime for confidential containers \u2014 "
         "uses Regorus for policy evaluation"),
        ("NVIDIA Attestation SDK",
         "GPU attestation and verification \u2014 "
         "Regorus enforces attestation policies"),
    ]
    ow = 5.8
    for (name, desc), x in zip(oss, dist(2, ow)):
        shape(s, I(x), I(oss_y + 0.3), I(ow), I(0.28),
              GREEN_D, GREEN_D, name, sz=11, color=WHITE, bold=True)
        shape(s, I(x), I(oss_y + 0.61), I(ow), I(0.32),
              GREEN_XL, GREEN_D, desc, sz=10, color=DARK_TEXT, align=PP_ALIGN.LEFT)

    # ── Internal Microsoft Adoption ──
    int_y = oss_y + 1.1
    text(s, I(M), I(int_y), I(5.0), I(0.25),
         "Internal Microsoft Adoption", sz=17, color=BLUE_D, bold=True)

    teams = [
        "AKS Confidential\nContainers",
        "Policy Control\nPoint Service",
        "Azure Graph\nService",
        "JIT\n(Just-In-Time)",
        "Network Config\nValidator",
    ]
    tw = 2.2
    for name_t, x in zip(teams, dist(5, tw, total=U, start=M)):
        shape(s, I(x), I(int_y + 0.28), I(tw), I(0.45),
              BLUE_L, BLUE_D, name_t, sz=10, color=BLUE_D, bold=True)

    # ── Footer ──
    text(s, I(M), I(int_y + 0.82), I(U), I(0.2),
         "Open source (MIT) \u2014 github.com/microsoft/regorus  \u00B7  "
         "and more across Azure, Windows, and M365",
         sz=10, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════
#   SLIDE 4 — JVM / RVM Analogy (solution reveal)
# ═════════════════════════════════════════════════════════════════════════

def flow_row(slide, y, label, accent, boxes):
    text(slide, I(M), I(y - 0.45), I(5), I(0.4),
         label, sz=18, color=accent, bold=True)
    bw, gap = 2.4, 0.55
    total = 4 * bw + 3 * gap
    sx = (SW - total) / 2
    for i, (tv, bc) in enumerate(boxes):
        x = sx + i * (bw + gap)
        shape(slide, I(x), I(y), I(bw), I(0.95),
              bc, accent, tv, sz=14, color=DARK_TEXT)
        if i < 3:
            ax = x + bw + (gap - 0.3) / 2
            text(slide, I(ax), I(y + 0.2), I(0.3), I(0.5),
                 "\u2192", sz=24, color=accent, bold=True,
                 align=PP_ALIGN.CENTER)

def slide_jvm_analogy(prs, BL):
    s = prs.slides.add_slide(BL)
    setbg(s, WHITE)

    text(s, I(M), I(0.3), I(U), I(0.7),
         "What Java Did for Apps, RVM Does for Policy",
         sz=34, color=GREEN_D, bold=True)

    flow_row(s, 1.5, "JVM  (Applications)", ORANGE_D, [
        ("Java \u00B7 Kotlin\nScala \u00B7 Groovy", ORANGE_L),
        ("Compilers", ORANGE_L),
        ("Java Bytecode", RGBColor(0xFF, 0xE0, 0xB2)),
        ("JVM: Any Platform", ORANGE_L),
    ])

    flow_row(s, 3.3, "RVM  (Policy)", GREEN_D, [
        ("Rego \u00B7 Cedar\nAzure Policy \u00B7 more", GREEN_XL),
        ("Compilers", GREEN_XL),
        ("RVM Bytecode", GREEN_L),
        ("RVM: Any Platform", GREEN_XL),
    ])

    text(s, I(M), I(4.8), I(U), I(0.4),
         "What this gives you", sz=20, color=BLUE_D, bold=True)

    bullets(s, I(M + 0.2), I(5.3), I(U), I(2.0), [
        "\u2013  One security surface to audit",
        "\u2013  One engine to optimize",
        "\u2013  One test framework to maintain",
        "\u2013  One artifact to deploy",
        "\u2013  Any new language gets the entire tooling ecosystem for free",
    ], sz=15, color=DARK_TEXT)


# ═════════════════════════════════════════════════════════════════════════
#   SLIDE 4 — DEMO 1: RVM Playground (Rego + Cedar + Azure Policy)
# ═════════════════════════════════════════════════════════════════════════

def slide_demo_playground(prs, BL):
    s = prs.slides.add_slide(BL)
    setbg(s, CHARCOAL)

    text(s, I(1.5), I(1.2), I(10.3), I(1.5),
         "LIVE DEMO", sz=72, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    text(s, I(1.5), I(3.0), I(10.3), I(0.6),
         "RVM Playground \u2014 Cross-language policy in your browser",
         sz=22, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    text(s, I(1.5), I(4.0), I(10.3), I(0.5),
         "Rego  \u2192  Compile to bytecode  \u2192  Run",
         sz=18, color=GREEN_L, align=PP_ALIGN.CENTER)

    text(s, I(1.5), I(4.6), I(10.3), I(0.5),
         "Cedar  \u2192  Same bytecode target  \u2192  Same engine",
         sz=18, color=GREEN_L, align=PP_ALIGN.CENTER)

    text(s, I(1.5), I(5.2), I(10.3), I(0.5),
         "Azure Policy  \u2192  Same bytecode target  \u2192  Same engine",
         sz=18, color=GREEN_L, align=PP_ALIGN.CENTER)

    hyperlink_text(s, I(1.5), I(6.2), I(10.3), I(0.5),
                   "https://anakrish.github.io/rego-virtual-machine-playground/",
                   "https://anakrish.github.io/rego-virtual-machine-playground/",
                   sz=16, color=AZURE, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════
#   SLIDE 5 — Cedar Contrast (independent efforts)
# ═════════════════════════════════════════════════════════════════════════

def slide_cedar_contrast(prs, BL):
    s = prs.slides.add_slide(BL)
    setbg(s, WHITE)

    text(s, I(M), I(0.25), I(U), I(0.6),
         "A Comparison with Amazon\u2019s Cedar", sz=36, color=GREEN_D, bold=True)

    text(s, I(M), I(0.82), I(U), I(0.4),
         "Cedar and RVM are independent efforts with different architectures "
         "and different design goals.",
         sz=16, color=MID_GRAY)

    # ── Left column: Cedar SDK ──
    lx, lw = 0.5, 5.7
    shape(s, I(lx), I(1.45), I(lw), I(0.55),
          ORANGE_D, ORANGE_D, "Cedar SDK  (Amazon)", sz=18, color=WHITE, bold=True)

    cedar_items = [
        ("\u2713  Strong formal foundations \u2014 Dafny proofs, 100M DRT/night",    DARK_TEXT),
        ("\u2713  Purpose-built for authorization (RBAC + ABAC)",                     DARK_TEXT),
        ("\u2713  Amazon Verified Permissions \u2014 managed cloud service",          DARK_TEXT),
        ("\u2713  Open source (Apache 2.0), ~5.7K GitHub stars",                     DARK_TEXT),
        ("",                                                                          None),
        ("\u2013  Single language \u2014 Cedar only, no Rego or Azure Policy",       RED_D),
        ("\u2013  AST tree-walking evaluator \u2014 no bytecode compilation",        RED_D),
        ("\u2013  No mid-evaluation host callbacks",                                  RED_D),
        ("\u2013  No suspendable execution, no step-through debugging",              RED_D),
        ("\u2013  Limited expressiveness \u2014 no loops, no user-defined functions", RED_D),
    ]
    cy = 2.15
    for (txt_val, clr) in cedar_items:
        if txt_val == "":
            cy += 0.12
            continue
        text(s, I(lx + 0.15), I(cy), I(lw - 0.2), I(0.28),
             txt_val, sz=12, color=clr)
        cy += 0.3

    # ── Right column: RVM ──
    rx, rw = 6.8, 6.0
    shape(s, I(rx), I(1.45), I(rw), I(0.55),
          GREEN_D, GREEN_D, "Regorus / RVM  (Microsoft)", sz=18, color=WHITE, bold=True)

    rvm_items = [
        "\u2713  Multi-language: Rego + Cedar + Azure Policy + extensible",
        "\u2713  Compiled register-based bytecode \u2014 high performance",
        "\u2713  HostAwait: pause mid-eval for external host data",
        "\u2713  Suspendable execution \u2014 breakpoints, step-through debugging",
        "\u2713  8 language bindings (C, C++, C#, Go, Java, Python, Ruby, JS)",
        "\u2713  Compile once, run anywhere \u2014 WASM, cloud, edge, bare metal",
        "\u2713  Instruction budget + memory limits per evaluation",
        "\u2713  Open source (MIT) \u2014 github.com/microsoft/regorus",
        "\u2713  Full Rego expressiveness \u2014 loops, functions, partial rules",
        "\u2713  Z3-powered formal analysis \u2014 diff, subsumption, test gen",
    ]
    ry = 2.15
    for i, txt_val in enumerate(rvm_items):
        bg = GREEN_XL if i % 2 == 0 else FAINT_GRAY
        shape(s, I(rx), I(ry), I(rw), I(0.33),
              bg, None, txt_val, sz=12, color=DARK_TEXT, align=PP_ALIGN.LEFT)
        ry += 0.37

    # ── Bottom banner ──
    shape(s, I(M), I(6.05), I(U), I(0.6), GREEN_D, GREEN_D,
          "RVM can also compile Cedar to bytecode \u2014 "
          "Cedar users get RVM performance, tooling, "
          "and formal analysis for free",
          sz=15, color=WHITE, bold=True)


# ═════════════════════════════════════════════════════════════════════════
#   SLIDE 6 — Policy Intelligence (Z3 capabilities + MSR quote)
# ═════════════════════════════════════════════════════════════════════════

def slide_policy_intelligence(prs, BL):
    s = prs.slides.add_slide(BL)
    setbg(s, WHITE)

    text(s, I(M), I(0.2), I(U), I(0.55),
         "Policy Intelligence", sz=36, color=DEEP_PURPLE, bold=True)

    text(s, I(M), I(0.72), I(U), I(0.4),
         "Bytecode + Z3 unlocks formal analysis "
         "across every language RVM supports.",
         sz=17, color=MID_GRAY)

    # ── AWS context bar ──
    shape(s, I(M), I(1.25), I(U), I(0.6), ORANGE_L, ORANGE_D,
          "AWS invested ~8 years building formal analysis "
          "(Zelkova, IAM Access Analyzer, Cedar proofs) "
          "\u2014 only for IAM/Cedar.  "
          "None of the capabilities below exist in their stack.",
          sz=12, color=DARK_TEXT, align=PP_ALIGN.LEFT)

    # ── Capability cards (2 rows x 3) ──
    caps = [
        ("Policy Diff",
         "Prove two policies are equivalent\nor find a distinguishing input.",
         "regorus diff"),
        ("Policy Subsumption",
         "Prove new policy is a superset\nof old \u2014 safe migration guarantee.",
         "regorus subsumes"),
        ("Auto Test Generation",
         "Z3-driven \u2014 generate inputs\nthat cover every code path.",
         "regorus gen-tests"),
        ("MC/DC Coverage",
         "Modified condition / decision\ncoverage \u2014 each condition toggled.",
         "--condition-coverage"),
        ("Dead Rule Detection",
         "UNSAT path = provably unreachable.\nNo inputs can ever trigger it.",
         "UNSAT analysis"),
        ("Input Synthesis",
         "Z3-driven input generation\nwith JSON output. Target any line.",
         "--cover-line / --avoid-line"),
    ]

    cw = 3.8
    rows = [caps[:3], caps[3:]]
    for ri, row in enumerate(rows):
        y_top = 2.1 + ri * 1.95
        for ci, (title, desc, cli) in enumerate(row):
            x = dist(3, cw)[ci]
            shape(s, I(x), I(y_top), I(cw), I(0.42),
                  DEEP_PURPLE, DEEP_PURPLE, title, sz=14, color=WHITE, bold=True)
            shape(s, I(x), I(y_top + 0.46), I(cw), I(1.0),
                  PURPLE_L, DEEP_PURPLE, desc, sz=12, color=DARK_TEXT, align=PP_ALIGN.LEFT)
            text(s, I(x + 0.1), I(y_top + 1.48), I(cw - 0.2), I(0.25),
                 cli, sz=10, color=MID_GRAY, bold=True)

    # ── MSR researcher quote ──
    shape(s, I(M), I(6.0), I(U), I(0.85), FAINT_GRAY, LIGHT_GRAY, "", rr=True)
    text_italic(s, I(M + 0.2), I(6.05), I(U - 0.4), I(0.55),
         "\u201CNovel capabilities include adding to workflows for guardrailing "
         "AI-generated policies. The IAM angle is where customers maintain "
         "policies over a portal and program policies by hand or their custom "
         "automation. When this applies to policies you consume, it makes "
         "sense to figure out what to surface.\u201D",
         sz=11, color=DARK_TEXT)
    text(s, I(M + 0.2), I(6.6), I(U - 0.4), I(0.25),
         "\u2014 MSR Researcher", sz=10, color=MID_GRAY, bold=True)


# ═════════════════════════════════════════════════════════════════════════
#   SLIDE 7 — DEMO 2: Formal Verification
# ═════════════════════════════════════════════════════════════════════════

def slide_demo_verification(prs, BL):
    s = prs.slides.add_slide(BL)
    setbg(s, CHARCOAL)

    text(s, I(1.5), I(1.2), I(10.3), I(1.5),
         "LIVE DEMO", sz=72, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    text(s, I(1.5), I(3.0), I(10.3), I(0.6),
         "Policy Intelligence \u2014 Formal Verification with Z3",
         sz=24, color=PURPLE_L, align=PP_ALIGN.CENTER)

    cues = [
        ("regorus diff",       "Are these two policies equivalent? Prove it \u2014 or find the difference."),
        ("regorus gen-tests",  "Auto-generate test inputs that cover every path."),
        ("--cover-line",       "Synthesize an input that reaches a specific line."),
        ("coverage",           "MC/DC condition coverage \u2014 every condition independently toggled."),
    ]
    cy = 4.2
    for label, desc in cues:
        text(s, I(3.0), I(cy), I(2.2), I(0.35),
             label, sz=14, color=GREEN_L, bold=True)
        text(s, I(5.3), I(cy), I(5.5), I(0.35),
             desc, sz=13, color=LIGHT_GRAY)
        cy += 0.45

    text(s, I(1.5), I(6.5), I(10.3), I(0.5),
         "Z3-powered  \u00B7  Works across Rego, Cedar, Azure Policy",
         sz=14, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════
#   SLIDE 8 — Copilot + Policy Intelligence
# ═════════════════════════════════════════════════════════════════════════

def slide_copilot_pi(prs, BL):
    s = prs.slides.add_slide(BL)
    setbg(s, WHITE)

    text(s, I(M), I(0.2), I(U), I(0.55),
         "Copilot + Policy Intelligence", sz=36, color=DEEP_PURPLE, bold=True)

    text(s, I(M), I(0.72), I(U), I(0.35),
         "AI writes the policy \u2014 math proves it\u2019s correct.",
         sz=18, color=MID_GRAY)

    # ── Circular feedback loop layout ──
    #
    #   Developer ──→ Copilot ──→ Policy Intelligence
    #                    ↑                    │
    #                    │         ┌──────────┴──────────┐
    #                    │         ↓                     ↓
    #                    │    ISSUE FOUND           VERIFIED
    #                    │    + counterexample      Safe to deploy
    #                    └────────┘
    #

    # Row 1: Developer → Copilot → PI
    r1y = 1.25
    bh = 1.35

    # Developer box
    shape(s, I(0.4), I(r1y), I(2.3), I(bh), BLUE_L, BLUE_D,
          "Developer\n\n\u201COnly allow reads\nfor interns\u201D",
          sz=14, color=DARK_TEXT)

    # Arrow →
    text(s, I(2.75), I(r1y + bh/2 - 0.2), I(0.4), I(0.4),
         "\u2192", sz=28, color=BLUE_D, bold=True, align=PP_ALIGN.CENTER)

    # Copilot box
    shape(s, I(3.2), I(r1y), I(2.8), I(bh), ORANGE_L, ORANGE_D,
          "Copilot\n\nGenerates Rego\npolicy from intent",
          sz=14, color=DARK_TEXT)

    # Arrow →
    text(s, I(6.05), I(r1y + bh/2 - 0.2), I(0.4), I(0.4),
         "\u2192", sz=28, color=DEEP_PURPLE, bold=True, align=PP_ALIGN.CENTER)

    # Policy Intelligence box
    shape(s, I(6.5), I(r1y), I(3.2), I(bh), PURPLE_L, DEEP_PURPLE,
          "Policy Intelligence\n\nVerifies formal properties:\n"
          "\u2022 No privilege escalation\n"
          "\u2022 No bypasses possible",
          sz=13, color=DARK_TEXT)

    # Arrow down from PI
    text(s, I(7.9), I(r1y + bh - 0.05), I(0.4), I(0.4),
         "\u2193", sz=28, color=DEEP_PURPLE, bold=True, align=PP_ALIGN.CENTER)

    # Row 2: outcomes side-by-side under PI
    r2y = r1y + bh + 0.45

    # VERIFIED (right outcome)
    shape(s, I(9.8), I(r1y), I(2.9), I(bh), GREEN_XL, GREEN_D,
          "VERIFIED\n\nSafe to deploy\n\u2192 Production",
          sz=14, color=GREEN_D, bold=True)
    # Arrow from PI to Verified
    text(s, I(9.75), I(r1y + bh/2 - 0.25), I(0.4), I(0.4),
         "\u2192", sz=22, color=GREEN_D, bold=True, align=PP_ALIGN.CENTER)

    # ISSUE FOUND box
    shape(s, I(6.5), I(r2y), I(3.2), I(1.1), RED_L, RED_D,
          "ISSUE FOUND\n\n"
          "\u201CIntern with role=manager\n"
          "can still write\u201D",
          sz=13, color=RED_D)

    # ── Feedback loop: Issue → back to Copilot ──
    # Horizontal arrow left from Issue box
    loop_y = r2y + 0.35

    # Curved feedback path rendered as: ← ← ← with label
    shape(s, I(3.2), I(r2y), I(3.1), I(0.45), RGBColor(0xFF, 0xEB, 0xEE), RED_D,
          "\u2190\u2190\u2190  Counterexample fed back to Copilot",
          sz=12, color=RED_D, bold=True)

    # Upward arrow into Copilot
    text(s, I(4.4), I(r2y - 0.45), I(0.4), I(0.45),
         "\u2191", sz=26, color=RED_D, bold=True, align=PP_ALIGN.CENTER)

    # Loop iteration label
    shape(s, I(3.2), I(r2y + 0.55), I(3.1), I(0.55), FAINT_GRAY, LIGHT_GRAY,
          "Copilot auto-revises policy\nusing the counterexample\nand re-submits for verification",
          sz=11, color=MID_GRAY, align=PP_ALIGN.LEFT)

    # ── "Iterate until proof passes" label on the loop ──
    shape(s, I(0.4), I(r2y), I(2.6), I(1.1),
          RGBColor(0xFD, 0xF0, 0xE5), ORANGE_D,
          "Feedback Loop\n\n"
          "Repeat until Z3 proof\npasses or developer\nreviews counterexample",
          sz=12, color=ORANGE_D, bold=False)

    # ── Key Insight bar ──
    shape(s, I(M), I(4.75), I(U), I(0.55), FAINT_GRAY, LIGHT_GRAY,
          "Every AI-generated policy gets a mathematical proof before it reaches "
          "production. Copilot iterates automatically using counterexamples "
          "until the proof passes.",
          sz=13, color=DARK_TEXT, align=PP_ALIGN.LEFT)

    # ── Why this matters ──
    text(s, I(M), I(5.6), I(U), I(0.35),
         "Why this matters", sz=18, color=DEEP_PURPLE, bold=True)

    bullets(s, I(M + 0.2), I(5.95), I(U), I(1.5), [
        "\u2013  AI hallucinations caught before deployment \u2014 not after an incident",
        "\u2013  Counterexamples become regression tests automatically",
        "\u2013  Works across Rego, Cedar, and Azure Policy \u2014 same verification loop",
    ], sz=14, color=DARK_TEXT)


# ═════════════════════════════════════════════════════════════════════════
#   SLIDE 9 — PR Review Bot + Regulatory Impact Analysis (combined)
# ═════════════════════════════════════════════════════════════════════════

def slide_pr_review_regulatory(prs, BL):
    s = prs.slides.add_slide(BL)
    setbg(s, WHITE)

    text(s, I(M), I(0.2), I(U), I(0.55),
         "What\u2019s Next \u2014 CI Integration & Compliance",
         sz=34, color=DEEP_PURPLE, bold=True)

    # ── TOP HALF: PR Review Bot ──
    text(s, I(M), I(0.85), I(U), I(0.4),
         "PR Review Bot \u2014 Formal Diff in CI",
         sz=20, color=BLUE_D, bold=True)

    # Flow: PR → CI → Comment
    shape(s, I(0.4), I(1.35), I(2.2), I(1.0), BLUE_L, BLUE_D,
          "Pull Request\n\nModifies policy.rego",
          sz=12, color=DARK_TEXT)

    text(s, I(2.65), I(1.65), I(0.4), I(0.4),
         "\u2192", sz=24, color=BLUE_D, bold=True, align=PP_ALIGN.CENTER)

    shape(s, I(3.1), I(1.35), I(2.6), I(1.0), FAINT_GRAY, BLUE_D,
          "CI Pipeline\n\nruns regorus diff\nv1 vs v2",
          sz=12, color=DARK_TEXT)

    text(s, I(5.75), I(1.65), I(0.4), I(0.4),
         "\u2192", sz=24, color=BLUE_D, bold=True, align=PP_ALIGN.CENTER)

    # PR comment box
    shape(s, I(6.2), I(1.2), I(6.5), I(1.3), FAINT_GRAY, BLUE_D,
          "", sz=12, color=DARK_TEXT, align=PP_ALIGN.LEFT)
    text(s, I(6.35), I(1.25), I(6.2), I(0.25),
         "PR Comment \u2014 Formal Diff Report", sz=12, color=BLUE_D, bold=True)
    text(s, I(6.35), I(1.52), I(6.2), I(0.22),
         "Input #1: role=contractor, resource=internal-docs \u2192 was ALLOW, now DENY",
         sz=10, color=DARK_TEXT)
    text(s, I(6.35), I(1.74), I(6.2), I(0.22),
         "Input #2: role=vendor, action=read \u2192 was DENY, now ALLOW  (wider access)",
         sz=10, color=RED_D)

    # Outcome bullets
    bullets(s, I(6.35), I(2.0), I(6.0), I(0.5), [
        "No behavioral change \u2192 auto-approve",
        "Changes found \u2192 reviewer sees exact divergent inputs",
        "Wider access detected \u2192 mandatory security review",
    ], sz=10, color=MID_GRAY, spacing=Pt(3))

    # ── Divider ──
    shape(s, I(M), I(3.1), I(U), I(0.03), LIGHT_GRAY, None, "", rr=False)

    # ── BOTTOM HALF: Regulatory Impact Analysis ──
    text(s, I(M), I(3.35), I(U), I(0.4),
         "Regulatory Impact Analysis \u2014 Scan All Policies for Compliance Gaps",
         sz=20, color=TEAL_D, bold=True)

    # Flow: Regulation → Formalize → Scan fleet
    shape(s, I(0.4), I(3.95), I(2.4), I(1.1), RGBColor(0xE0, 0xF2, 0xF1), TEAL_D,
          "Regulation\n\n\u201CMFA required for\nall PII access\u201D",
          sz=12, color=DARK_TEXT)

    text(s, I(2.85), I(4.3), I(0.4), I(0.4),
         "\u2192", sz=24, color=TEAL_D, bold=True, align=PP_ALIGN.CENTER)

    shape(s, I(3.3), I(3.95), I(2.6), I(1.1), RGBColor(0xE0, 0xF2, 0xF1), TEAL_D,
          "Formalize & Scan\n\n\u2200 request:\naccesses_pii(r) \u2192 has_mfa(r)",
          sz=12, color=DARK_TEXT)

    text(s, I(5.95), I(4.3), I(0.4), I(0.4),
         "\u2192", sz=24, color=TEAL_D, bold=True, align=PP_ALIGN.CENTER)

    # Results
    shape(s, I(6.4), I(3.95), I(2.5), I(1.1), GREEN_XL, GREEN_D,
          "39 services\nCompliant\n\nMathematical proof\nfor each",
          sz=12, color=GREEN_D)

    shape(s, I(9.2), I(3.95), I(3.5), I(1.1), RED_L, RED_D,
          "3 services \u2014 Gap Found\n\n"
          "billing-api: no MFA on /invoices\n"
          "hr-portal: admin override bypasses PII\n"
          "analytics: logging exposes email",
          sz=11, color=RED_D, align=PP_ALIGN.LEFT)

    # ── Bottom takeaway ──
    shape(s, I(M), I(5.6), I(U), I(0.6), FAINT_GRAY, LIGHT_GRAY,
          "Each result backed by a mathematical proof \u2014 not sampling. "
          "Counterexamples pinpoint the exact non-compliant path.",
          sz=14, color=DARK_TEXT, align=PP_ALIGN.LEFT)

    text(s, I(M), I(6.4), I(U), I(0.5),
         "Both capabilities work across Rego, Cedar, and Azure Policy \u2014 "
         "because they operate on RVM bytecode.",
         sz=14, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════
#   SLIDE 10 — Closing
# ═════════════════════════════════════════════════════════════════════════

def slide_closing(prs, BL):
    s = prs.slides.add_slide(BL)
    setbg(s, GREEN_D)

    text(s, I(1.5), I(1.0), I(10.3), I(1.2),
         "One Engine. Every Policy.\nFormal Analysis Built In.",
         sz=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    text(s, I(1.5), I(2.8), I(10.3), I(0.6),
         "Regorus Virtual Machine",
         sz=32, color=GREEN_L, align=PP_ALIGN.CENTER)

    lines = [
        "Rego  +  Cedar  +  Azure Policy \u2014 one bytecode target",
        "High performance  \u00B7  Memory safe  \u00B7  Portable  \u00B7  8 language bindings",
        "Policy diff  \u00B7  Subsumption  \u00B7  Auto test generation  \u00B7  Coverage",
        "Open Source (MIT)  \u00B7  github.com/microsoft/regorus",
    ]
    for i, line in enumerate(lines):
        text(s, I(1.5), I(3.8 + i * 0.55), I(10.3), I(0.5),
             line, sz=19, color=WHITE, align=PP_ALIGN.CENTER)

    hyperlink_text(s, I(1.5), I(5.8), I(10.3), I(0.5),
                   "Try it: https://anakrish.github.io/rego-virtual-machine-playground/",
                   "https://anakrish.github.io/rego-virtual-machine-playground/",
                   sz=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    hyperlink_text(s, I(1.5), I(6.4), I(10.3), I(0.5),
                   "github.com/microsoft/regorus",
                   "https://github.com/microsoft/regorus",
                   sz=17, color=GREEN_L, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════
#   MAIN
# ═════════════════════════════════════════════════════════════════════════

CONTENT_SLIDES = [
    slide_problem,              # will become slide 1 (moved to front after generation)
    # (CTO memo is slide 2 — preserved from SlidesToReuse)
    # (Regorus overview is slide 3 — generated separately with extracted images)
    slide_jvm_analogy,          # slide 4
    slide_demo_playground,      # slide 5
    slide_cedar_contrast,       # slide 6
    slide_policy_intelligence,  # slide 7
    slide_demo_verification,    # slide 8
    slide_copilot_pi,           # slide 9
    slide_pr_review_regulatory, # slide 10
    slide_closing,              # slide 11
]


def main():
    # Open SlidesToReuse — has 3 slides: CTO memo, Renewed Urgency, Regorus Engine
    prs = Presentation(REUSE)
    prs.slide_width = I(SW)
    prs.slide_height = I(SH)

    # Extract images from Regorus Engine slide (index 2) before deleting
    engine_slide = prs.slides[2]
    image_blobs = {}
    for sh in engine_slide.shapes:
        if sh.shape_type == 13:  # Picture
            image_blobs[sh.name] = sh.image.blob

    # Delete Regorus Engine (slide 3, index 2) then Renewed Urgency (slide 2, index 1)
    # Keep only CTO Memo (slide 1)
    delete_slide(prs, 2)   # remove "Regorus Engine"
    delete_slide(prs, 1)   # remove "Renewed Urgency"
    # Now: [0] = CTO Memo

    # Find Blank layout for our generated slides
    BL = None
    for layout in prs.slide_layouts:
        if layout.name == "Blank":
            BL = layout
            break
    if BL is None:
        BL = prs.slide_layouts[6]

    # Generate Regorus Overview slide (uses extracted images from SlidesToReuse)
    slide_regorus_overview(prs, BL, image_blobs)

    # Add remaining 9 content slides
    for fn in CONTENT_SLIDES:
        fn(prs, BL)
    # Current order: CTO(0), Overview(1), Problem(2), JVM(3), Demo1(4), Cedar(5), PI(6), Demo2(7), CopilotPI(8), PRReg(9), Closing(10)
    # Desired order: Problem(0), CTO(1), Overview(2), JVM(3), ...

    # Move Problem (currently index 2) to front
    move_slide_to_front(prs, 2)
    # Now: Problem(0), CTO(1), Overview(2), JVM(3), Demo1(4), Cedar(5), PI(6), Demo2(7), CopilotPI(8), PRReg(9), Closing(10)

    out = os.path.join(BASE, "RVM-Final.pptx")
    prs.save(out)
    print(f"Generated: {out}")
    print(f"Total slides: {len(prs.slides)}")
    print()
    labels = [
        "Problem (engine proliferation)",
        "CTO Memo (preserved from SlidesToReuse)",
        "Regorus Overview + Adoption",
        "JVM / RVM Analogy",
        "DEMO 1: RVM Playground",
        "A Comparison with Amazon\u2019s Cedar",
        "Policy Intelligence",
        "DEMO 2: Formal Verification",
        "Copilot + Policy Intelligence",
        "PR Review Bot + Regulatory Impact",
        "Closing",
    ]
    for i, label in enumerate(labels, 1):
        print(f"  {i}. {label}")


if __name__ == "__main__":
    main()
