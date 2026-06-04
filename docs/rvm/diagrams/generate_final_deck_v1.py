#!/usr/bin/env python3
"""Generate the final RVM presentation — demo-driven, minimal slides.

Structure (9 slides):
  1-2  Context from SlidesToReuse.pptx (CTO memo, Renewed Urgency) [preserved]
  3    Title — "The Power of RVM"
  4    Problem — engine proliferation (minimal, sets up Demo 1)
  5    DEMO 1 — RVM Playground
  6    Cedar Contrast — two-column, factual, sets up "and we go further"
  7    Policy Intelligence — what bytecode + Z3 unlocks (sets up Demo 2)
  8    DEMO 2 — Formal Verification / Policy Intelligence
  9    Closing

Output: RVM-Final.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

try:
    from pptx.enum.text import MSO_ANCHOR
except ImportError:
    MSO_ANCHOR = None

BASE = os.path.dirname(os.path.abspath(__file__))
REUSE = os.path.join(BASE, "SlidesToReuse.pptx")

# ── Dimensions ───────────────────────────────────────────────────────────
SW, SH = 13.333, 7.5
M = 0.6
U = SW - 2 * M

# ── Palette ──────────────────────────────────────────────────────────────
CHARCOAL   = RGBColor(0x2D, 0x2D, 0x2D)
DARK_TEXT   = RGBColor(0x33, 0x33, 0x33)
MID_GRAY    = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT_GRAY  = RGBColor(0xE0, 0xE0, 0xE0)
FAINT_GRAY  = RGBColor(0xF5, 0xF5, 0xF5)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

GREEN_D   = RGBColor(0x1B, 0x5E, 0x20)
GREEN_M   = RGBColor(0x2E, 0x7D, 0x32)
GREEN_L   = RGBColor(0xC8, 0xE6, 0xC9)
GREEN_XL  = RGBColor(0xE8, 0xF5, 0xE9)
BLUE_D    = RGBColor(0x15, 0x65, 0xC0)
BLUE_L    = RGBColor(0xE3, 0xF2, 0xFD)
ORANGE_D  = RGBColor(0xEF, 0x6C, 0x00)
ORANGE_L  = RGBColor(0xFF, 0xF3, 0xE0)
RED_D     = RGBColor(0xC6, 0x28, 0x28)
RED_L     = RGBColor(0xFF, 0xCD, 0xD2)
PURPLE_D  = RGBColor(0x7B, 0x1F, 0xA2)
PURPLE_L  = RGBColor(0xF3, 0xE5, 0xF5)
TEAL_D    = RGBColor(0x00, 0x69, 0x5C)
TEAL_L    = RGBColor(0xB2, 0xDF, 0xDB)
AZURE     = RGBColor(0x00, 0x78, 0xD4)
DEEP_PURPLE = RGBColor(0x4A, 0x14, 0x8C)


# ── Helpers ──────────────────────────────────────────────────────────────

def I(v): return Inches(v)

def dist(n, w, total=U, start=M):
    if n <= 1: return [start + (total - w) / 2]
    gap = (total - n * w) / (n - 1)
    return [start + i * (w + gap) for i in range(n)]

def setbg(slide, color):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def shape(slide, l, t, w, h, fill, border=None, txt="",
          sz=14, color=DARK_TEXT, bold=False,
          align=PP_ALIGN.CENTER, rr=True):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rr else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(st, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border; sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    tf = sh.text_frame; tf.word_wrap = True; tf.auto_size = None
    tf.margin_left = Pt(8); tf.margin_right = Pt(8)
    tf.margin_top = Pt(4); tf.margin_bottom = Pt(4)
    if MSO_ANCHOR:
        try: tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except: pass
    if txt:
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = txt
        r.font.size = Pt(sz); r.font.color.rgb = color; r.font.bold = bold
    return sh

def text(slide, l, t, w, h, txt, sz=14, color=DARK_TEXT,
         bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = txt
    r.font.size = Pt(sz); r.font.color.rgb = color; r.font.bold = bold
    return tb

def bullets(slide, l, t, w, h, items, sz=14, color=DARK_TEXT,
            spacing=Pt(8), bold=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        r = p.add_run(); r.text = item
        r.font.size = Pt(sz); r.font.color.rgb = color; r.font.bold = bold
    return tb

def delete_slide(prs, index):
    rIdList = prs.slides._sldIdLst
    sldId = list(rIdList)[index]
    rId = sldId.get(qn("r:id"))
    prs.part.drop_rel(rId)
    rIdList.remove(sldId)


# ═════════════════════════════════════════════════════════════════════════
#   SLIDE 3 — Title
# ═════════════════════════════════════════════════════════════════════════

def slide_title(prs, BL):
    s = prs.slides.add_slide(BL)
    setbg(s, GREEN_D)
    text(s, I(1.5), I(1.6), I(10.3), I(1.2),
         "The Power of RVM", sz=54, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(s, I(1.5), I(3.2), I(10.3), I(0.8),
         "Regorus Virtual Machine", sz=36, color=GREEN_L, align=PP_ALIGN.CENTER)
    text(s, I(1.5), I(4.5), I(10.3), I(0.6),
         "One Engine  \u00B7  Every Policy Language  \u00B7  Any Platform",
         sz=22, color=WHITE, align=PP_ALIGN.CENTER)
    text(s, I(1.5), I(6.0), I(10.3), I(0.5),
         "microsoft/regorus  \u00B7  Open Source (MIT)",
         sz=16, color=GREEN_L, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════
#   SLIDE 4 — Problem (minimal, sets up Demo 1)
# ═════════════════════════════════════════════════════════════════════════

def slide_problem(prs, BL):
    s = prs.slides.add_slide(BL)
    setbg(s, WHITE)

    text(s, I(M), I(0.4), I(U), I(0.8),
         "The Problem", sz=44, color=RED_D, bold=True)

    text(s, I(M), I(1.3), I(U), I(0.6),
         "Azure Policy must support multiple policy languages.\n"
         "Across Microsoft, teams maintain separate engines.",
         sz=20, color=MID_GRAY)

    engines = ["OPA / Rego", "Cedar", "Azure Policy", "AKS Admission", "Custom"]
    bw = 2.1
    for name, x in zip(engines, dist(5, bw)):
        shape(s, I(x), I(2.3), I(bw), I(1.0),
              RED_L, RED_D, name, sz=16, color=RED_D, bold=True)

    shape(s, I(M), I(3.7), I(U), I(0.65), FAINT_GRAY, LIGHT_GRAY,
          "N engines  =  N \u00D7 audits  +  N \u00D7 test suites  +  "
          "N \u00D7 pipelines  +  N \u00D7 teams",
          sz=18, color=RED_D, bold=True)

    text(s, I(1.5), I(5.0), I(10.3), I(0.9),
         "What if one engine could run them all?",
         sz=32, color=GREEN_D, bold=True, align=PP_ALIGN.CENTER)

    text(s, I(1.5), I(6.0), I(10.3), I(0.6),
         "Let me show you.",
         sz=22, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════
#   SLIDE 5 — DEMO 1: RVM Playground
# ═════════════════════════════════════════════════════════════════════════

def slide_demo_playground(prs, BL):
    s = prs.slides.add_slide(BL)
    setbg(s, CHARCOAL)

    text(s, I(1.5), I(1.5), I(10.3), I(1.5),
         "LIVE DEMO", sz=72, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    text(s, I(1.5), I(3.3), I(10.3), I(0.6),
         "RVM Playground \u2014 Cross-language policy in your browser",
         sz=22, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    text(s, I(1.5), I(4.5), I(10.3), I(0.5),
         "Write Rego  \u2192  Compile to bytecode  \u2192  Run",
         sz=18, color=GREEN_L, align=PP_ALIGN.CENTER)

    text(s, I(1.5), I(5.1), I(10.3), I(0.5),
         "Switch to Cedar  \u2192  Same bytecode target  \u2192  Same engine",
         sz=18, color=GREEN_L, align=PP_ALIGN.CENTER)

    text(s, I(1.5), I(6.2), I(10.3), I(0.5),
         "https://anakrish.github.io/rego-virtual-machine-playground/",
         sz=16, color=AZURE, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════
#   SLIDE 6 — Cedar Contrast
# ═════════════════════════════════════════════════════════════════════════

def slide_cedar_contrast(prs, BL):
    s = prs.slides.add_slide(BL)
    setbg(s, WHITE)

    text(s, I(M), I(0.25), I(U), I(0.6),
         "RVM vs. Cedar SDK", sz=36, color=GREEN_D, bold=True)

    text(s, I(M), I(0.82), I(U), I(0.4),
         "Cedar is strong. RVM builds on top of what Cedar offers \u2014 and goes further.",
         sz=16, color=MID_GRAY)

    # ── Left column: Cedar SDK ──
    lx, lw = 0.6, 5.6
    shape(s, I(lx), I(1.45), I(lw), I(0.55),
          ORANGE_D, ORANGE_D, "Cedar SDK  (Amazon)", sz=18, color=WHITE, bold=True)

    cedar_items = [
        ("\u2713  Strong formal foundations \u2014 Dafny proofs, 100M DRT/night",    DARK_TEXT, False),
        ("\u2713  Purpose-built for authorization (RBAC + ABAC)",                     DARK_TEXT, False),
        ("\u2713  Amazon Verified Permissions \u2014 managed cloud service",          DARK_TEXT, False),
        ("\u2713  Open source (Apache 2.0), ~5.7K GitHub stars",                     DARK_TEXT, False),
        ("",                                                                          DARK_TEXT, False),
        ("\u2013  Single language \u2014 Cedar only, no Rego or Azure Policy",       RED_D,     False),
        ("\u2013  AST tree-walking evaluator \u2014 no bytecode compilation",        RED_D,     False),
        ("\u2013  No mid-evaluation host callbacks (no HostAwait)",                   RED_D,     False),
        ("\u2013  No suspendable execution, no step-through debugging",              RED_D,     False),
        ("\u2013  Limited expressiveness \u2014 no loops, no user-defined functions", RED_D,     False),
    ]
    cy = 2.15
    for (txt_val, clr, bld) in cedar_items:
        if txt_val == "":
            cy += 0.12
            continue
        text(s, I(lx + 0.15), I(cy), I(lw - 0.2), I(0.28),
             txt_val, sz=12, color=clr, bold=bld)
        cy += 0.3

    # ── Right column: RVM ──
    rx, rw = 6.8, 6.0
    shape(s, I(rx), I(1.45), I(rw), I(0.55),
          GREEN_D, GREEN_D, "Regorus / RVM  (Microsoft)", sz=18, color=WHITE, bold=True)

    rvm_items = [
        ("\u2713  Multi-language: Rego + Cedar + Azure Policy + extensible",         DARK_TEXT),
        ("\u2713  Compiled register-based bytecode \u2014 high performance",         DARK_TEXT),
        ("\u2713  HostAwait: pause mid-eval for external host data",                 DARK_TEXT),
        ("\u2713  Suspendable execution \u2014 breakpoints, step-through debugging", DARK_TEXT),
        ("\u2713  Rust core \u2014 memory safe, no_std, no GC",                      DARK_TEXT),
        ("\u2713  8 language bindings (C, C++, C#, Go, Java, Python, Ruby, JS)",     DARK_TEXT),
        ("\u2713  Compile once, run anywhere \u2014 WASM, cloud, edge, bare metal",  DARK_TEXT),
        ("\u2713  Instruction budget + memory limits per evaluation",                 DARK_TEXT),
        ("\u2713  Open source (MIT), interactive WASM playground",                    DARK_TEXT),
        ("\u2713  Full Rego expressiveness \u2014 loops, functions, partial rules",   DARK_TEXT),
    ]
    ry = 2.15
    for i, (txt_val, clr) in enumerate(rvm_items):
        bg = GREEN_XL if i % 2 == 0 else FAINT_GRAY
        shape(s, I(rx), I(ry), I(rw), I(0.33),
              bg, None, txt_val, sz=12, color=clr, align=PP_ALIGN.LEFT)
        ry += 0.37

    # ── Bottom banner ──
    shape(s, I(M), I(6.1), I(U), I(0.55), GREEN_D, GREEN_D,
          "RVM compiles Cedar to bytecode \u2014 "
          "Cedar users get RVM performance + tooling for free",
          sz=16, color=WHITE, bold=True)


# ═════════════════════════════════════════════════════════════════════════
#   SLIDE 7 — Policy Intelligence (sets up Demo 2)
# ═════════════════════════════════════════════════════════════════════════

def slide_policy_intelligence(prs, BL):
    s = prs.slides.add_slide(BL)
    setbg(s, WHITE)

    text(s, I(M), I(0.25), I(U), I(0.6),
         "Policy Intelligence", sz=36, color=DEEP_PURPLE, bold=True)

    text(s, I(M), I(0.82), I(U), I(0.45),
         "Bytecode + Z3 unlocks formal analysis "
         "across every language RVM supports.",
         sz=17, color=MID_GRAY)

    # ── AWS context bar ──
    shape(s, I(M), I(1.45), I(U), I(0.7), ORANGE_L, ORANGE_D,
          "AWS invested ~8 years building formal analysis (Zelkova, IAM Access Analyzer, Cedar proofs) "
          "\u2014 but only for IAM/Cedar.  None of the capabilities below exist in their stack.",
          sz=13, color=DARK_TEXT, align=PP_ALIGN.LEFT)

    # ── Capability cards ──
    caps = [
        ("Policy Diff",
         "Prove two policies are equivalent\nor find a distinguishing input.",
         "regorus diff",
         DEEP_PURPLE, PURPLE_L),
        ("Policy Subsumption",
         "Prove new policy is a superset\nof old \u2014 safe migration guarantee.",
         "regorus subsumes",
         DEEP_PURPLE, PURPLE_L),
        ("Auto Test Generation",
         "Z3-driven \u2014 generate inputs\nthat cover every code path.",
         "regorus gen-tests",
         DEEP_PURPLE, PURPLE_L),
        ("MC/DC Coverage",
         "Modified condition / decision\ncoverage \u2014 each condition toggled.",
         "--condition-coverage",
         DEEP_PURPLE, PURPLE_L),
        ("Dead Rule Detection",
         "UNSAT path = provably unreachable.\nNo inputs can ever trigger it.",
         "UNSAT analysis",
         DEEP_PURPLE, PURPLE_L),
        ("Input Synthesis",
         "Z3-driven input generation\nwith JSON output. Target any line.",
         "--cover-line / --avoid-line",
         DEEP_PURPLE, PURPLE_L),
    ]

    cw = 3.8
    rows = [caps[:3], caps[3:]]
    for ri, row in enumerate(rows):
        y_top = 2.45 + ri * 2.15
        for ci, (title, desc, cli, accent, bg) in enumerate(row):
            x = dist(3, cw)[ci]
            # Title bar
            shape(s, I(x), I(y_top), I(cw), I(0.45),
                  accent, accent, title, sz=15, color=WHITE, bold=True)
            # Description box
            shape(s, I(x), I(y_top + 0.5), I(cw), I(1.1),
                  bg, accent, desc, sz=12, color=DARK_TEXT, align=PP_ALIGN.LEFT)
            # CLI label
            text(s, I(x + 0.1), I(y_top + 1.6), I(cw - 0.2), I(0.3),
                 cli, sz=10, color=MID_GRAY, bold=True)

    # ── Bottom ──
    text(s, I(1.5), I(6.6), I(10.3), I(0.5),
         "Let me show you.",
         sz=22, color=DEEP_PURPLE, bold=True, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════
#   SLIDE 8 — DEMO 2: Formal Verification
# ═════════════════════════════════════════════════════════════════════════

def slide_demo_verification(prs, BL):
    s = prs.slides.add_slide(BL)
    setbg(s, CHARCOAL)

    text(s, I(1.5), I(1.2), I(10.3), I(1.5),
         "LIVE DEMO", sz=72, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    text(s, I(1.5), I(3.0), I(10.3), I(0.6),
         "Policy Intelligence \u2014 Formal Verification with Z3",
         sz=24, color=PURPLE_L, align=PP_ALIGN.CENTER)

    cues = [
        ("regorus diff",       "Are these two policies equivalent? Prove it or find the difference."),
        ("regorus gen-tests",  "Auto-generate test inputs that cover every path."),
        ("--cover-line",       "Synthesize an input that reaches a specific line."),
        ("coverage",           "MC/DC condition coverage \u2014 every condition independently toggled."),
    ]
    cy = 4.2
    for label, desc in cues:
        text(s, I(3.0), I(cy), I(2.2), I(0.35),
             label, sz=14, color=GREEN_L, bold=True)
        text(s, I(5.3), I(cy), I(5.5), I(0.35),
             desc, sz=13, color=LIGHT_GRAY)
        cy += 0.45

    text(s, I(1.5), I(6.5), I(10.3), I(0.5),
         "Z3-powered  \u00B7  Works across Rego, Cedar, Azure Policy",
         sz=14, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════
#   SLIDE 9 — Closing
# ═════════════════════════════════════════════════════════════════════════

def slide_closing(prs, BL):
    s = prs.slides.add_slide(BL)
    setbg(s, GREEN_D)

    text(s, I(1.5), I(1.0), I(10.3), I(1.2),
         "One Engine. Every Policy.\nFormal Analysis Built In.",
         sz=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    text(s, I(1.5), I(2.8), I(10.3), I(0.6),
         "Regorus Virtual Machine",
         sz=32, color=GREEN_L, align=PP_ALIGN.CENTER)

    lines = [
        "Rego  +  Cedar  +  Azure Policy \u2014 one bytecode target",
        "High performance  \u00B7  Memory safe  \u00B7  Portable  \u00B7  8 language bindings",
        "Policy diff  \u00B7  Subsumption  \u00B7  Auto test generation  \u00B7  Coverage",
        "Open Source (MIT)  \u00B7  github.com/microsoft/regorus",
    ]
    for i, line in enumerate(lines):
        text(s, I(1.5), I(3.8 + i * 0.55), I(10.3), I(0.5),
             line, sz=19, color=WHITE, align=PP_ALIGN.CENTER)

    text(s, I(1.5), I(5.8), I(10.3), I(0.5),
         "Try it: https://anakrish.github.io/rego-virtual-machine-playground/",
         sz=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    text(s, I(1.5), I(6.4), I(10.3), I(0.5),
         "github.com/microsoft/regorus",
         sz=17, color=GREEN_L, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════
#   MAIN
# ═════════════════════════════════════════════════════════════════════════

SLIDE_FNS = [
    slide_title,
    slide_problem,
    slide_demo_playground,
    slide_cedar_contrast,
    slide_policy_intelligence,
    slide_demo_verification,
    slide_closing,
]


def main():
    # Open SlidesToReuse — keeps slides 1-2 (CTO memo, Renewed Urgency)
    prs = Presentation(REUSE)
    prs.slide_width = I(SW)
    prs.slide_height = I(SH)

    # Delete slide 3 (index 2) — outdated Regorus Engine slide
    delete_slide(prs, 2)

    # Find Blank layout
    BL = None
    for layout in prs.slide_layouts:
        if layout.name == "Blank":
            BL = layout
            break
    if BL is None:
        BL = prs.slide_layouts[6]

    # Add our 7 content slides (3-9)
    for fn in SLIDE_FNS:
        fn(prs, BL)

    out = os.path.join(BASE, "RVM-Final.pptx")
    prs.save(out)
    print(f"Generated: {out}")
    print(f"Total slides: {len(prs.slides)}")
    print()
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        for shp in slide.shapes:
            if shp.has_text_frame:
                t = shp.text_frame.text.strip()
                if t:
                    title = t[:60]
                    break
        print(f"  {i:2d}. {title}")


if __name__ == "__main__":
    main()
