#!/usr/bin/env python3
"""Approach A: Merge SlidesToReuse.pptx context slides with generated content.

For each version A-E:
  - Opens SlidesToReuse.pptx (keeps original slides 1-2 with SmartArt/images)
  - Removes slide 3 (old Regorus Engine slide)
  - Adds version-specific RVM content slides after the context
  - Saves as RVM-Merged-{version}.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os, shutil

try:
    from pptx.enum.text import MSO_ANCHOR
except ImportError:
    MSO_ANCHOR = None

BASE = os.path.dirname(os.path.abspath(__file__))
SRC = os.path.join(BASE, "SlidesToReuse.pptx")

# ── Palette ──────────────────────────────────────────────────────────────────
CHARCOAL   = RGBColor(0x2D, 0x2D, 0x2D)
DARK_TEXT  = RGBColor(0x33, 0x33, 0x33)
MID_GRAY   = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
FAINT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

GREEN_D  = RGBColor(0x1B, 0x5E, 0x20)
GREEN_M  = RGBColor(0x2E, 0x7D, 0x32)
GREEN_L  = RGBColor(0xC8, 0xE6, 0xC9)
GREEN_XL = RGBColor(0xE8, 0xF5, 0xE9)
BLUE_D   = RGBColor(0x15, 0x65, 0xC0)
BLUE_L   = RGBColor(0xE3, 0xF2, 0xFD)
ORANGE_D = RGBColor(0xEF, 0x6C, 0x00)
ORANGE_L = RGBColor(0xFF, 0xF3, 0xE0)
RED_D    = RGBColor(0xC6, 0x28, 0x28)
RED_L    = RGBColor(0xFF, 0xCD, 0xD2)
PURPLE_D = RGBColor(0x7B, 0x1F, 0xA2)
PURPLE_L = RGBColor(0xF3, 0xE5, 0xF5)
TEAL_D   = RGBColor(0x00, 0x69, 0x5C)
TEAL_L   = RGBColor(0xB2, 0xDF, 0xDB)
AZURE    = RGBColor(0x00, 0x78, 0xD4)

SW = 13.333
SH = 7.5
M = 0.6
U = SW - 2 * M


# ── Helpers ──────────────────────────────────────────────────────────────────

def I(v):
    return Inches(v)


def dist(n, w, total=U, start=M):
    if n <= 1:
        return [start + (total - w) / 2]
    gap = (total - n * w) / (n - 1)
    return [start + i * (w + gap) for i in range(n)]


def bgfill(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def shape(slide, l, t, w, h, fill, border=None, txt="",
          sz=14, color=DARK_TEXT, bold=False,
          align=PP_ALIGN.CENTER, rr=True):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rr else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(st, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    tf = sh.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    if MSO_ANCHOR:
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass
    if txt:
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.color.rgb = color
        r.font.bold = bold
    return sh


def text(slide, l, t, w, h, txt, sz=14, color=DARK_TEXT,
         bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(sz)
    r.font.color.rgb = color
    r.font.bold = bold
    return tb


def bullets(slide, l, t, w, h, items, sz=14, color=DARK_TEXT,
            spacing=Pt(8), bold=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        r = p.add_run()
        r.text = item
        r.font.size = Pt(sz)
        r.font.color.rgb = color
        r.font.bold = bold
    return tb


def delete_slide(prs, index):
    """Delete slide at given index (0-based)."""
    sldIdLst = prs.slides._sldIdLst
    sldIds = list(sldIdLst)
    sldId = sldIds[index]
    rId = sldId.get(qn("r:id"))
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)


# ── Shared slide builders ────────────────────────────────────────────────────

def add_title(prs, BL):
    s = prs.slides.add_slide(BL)
    bgfill(s, GREEN_D)
    text(s, I(1.5), I(1.6), I(10.3), I(1.2),
         "The Power of RVM", sz=54, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER)
    text(s, I(1.5), I(3.2), I(10.3), I(0.8),
         "Regorus Virtual Machine", sz=36, color=GREEN_L,
         align=PP_ALIGN.CENTER)
    text(s, I(1.5), I(4.5), I(10.3), I(0.6),
         "One Engine  \u00B7  Every Policy Language  \u00B7  Any Platform",
         sz=22, color=WHITE, align=PP_ALIGN.CENTER)
    text(s, I(1.5), I(6.0), I(10.3), I(0.5),
         "microsoft/regorus  \u00B7  Open Source (MIT)",
         sz=16, color=GREEN_L, align=PP_ALIGN.CENTER)


def add_demo(prs, BL):
    s = prs.slides.add_slide(BL)
    bgfill(s, CHARCOAL)
    text(s, I(1.5), I(2.0), I(10.3), I(1.5),
         "LIVE DEMO", sz=72, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER)
    text(s, I(1.5), I(3.8), I(10.3), I(0.6),
         "RVM Playground \u2014 Cross-language policy in your browser",
         sz=22, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    text(s, I(1.5), I(5.0), I(10.3), I(0.5),
         "https://anakrish.github.io/rego-virtual-machine-playground/",
         sz=16, color=AZURE, align=PP_ALIGN.CENTER)


def add_closing(prs, BL):
    s = prs.slides.add_slide(BL)
    bgfill(s, GREEN_D)
    text(s, I(1.5), I(1.5), I(10.3), I(1.2),
         "One Engine. Every Policy. Any Platform.",
         sz=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(s, I(1.5), I(3.0), I(10.3), I(0.8),
         "Regorus Virtual Machine",
         sz=32, color=GREEN_L, align=PP_ALIGN.CENTER)
    for i, item in enumerate([
        "Rego  +  Cedar  +  Azure Policy",
        "High performance  \u00B7  Memory safe  \u00B7  Portable bytecode",
        "8 language bindings  \u00B7  WASM Playground  \u00B7  Open Source (MIT)",
    ]):
        text(s, I(1.5), I(4.2 + i * 0.55), I(10.3), I(0.5),
             item, sz=20, color=WHITE, align=PP_ALIGN.CENTER)
    text(s, I(1.5), I(5.9), I(10.3), I(0.5),
         "Try it: https://anakrish.github.io/"
         "rego-virtual-machine-playground/",
         sz=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(s, I(1.5), I(6.5), I(10.3), I(0.5),
         "github.com/microsoft/regorus",
         sz=17, color=GREEN_L, align=PP_ALIGN.CENTER)


def add_competitive(prs, BL, appendix=False):
    s = prs.slides.add_slide(BL)
    bgfill(s, WHITE)
    if appendix:
        text(s, I(M), I(0.12), I(U), I(0.35),
             "APPENDIX", sz=11, color=MID_GRAY, bold=True)
    text(s, I(M), I(0.3), I(U), I(0.7),
         "Competitive Landscape",
         sz=36, color=GREEN_D, bold=True)
    lx = 0.5
    competitors = [
        ("OPA  (Go)", [
            "+ Mature ecosystem, large community",
            "+ Rego language",
            "\u2013 Go runtime overhead, GC pauses",
            "\u2013 Embeddable via cgo (complex, heavy)",
            "\u2013 Single language \u2014 no Cedar or Azure Policy",
        ]),
        ("Cedar SDK  (Rust)", [
            "+ Strong formal foundations",
            "+ Amazon Verified Permissions \u2014 managed cloud service",
            "\u2013 Single language \u2014 no Rego or Azure Policy",
            "\u2013 No bytecode \u2014 AST tree-walking evaluator",
            "\u2013 No mid-evaluation host callbacks",
        ]),
        ("Custom Engines", [
            "\u2013 One-off implementations per product",
            "\u2013 Per-product maintenance & security burden",
            "\u2013 No shared tooling or analysis",
        ]),
    ]
    cy = 1.2
    for title, items in competitors:
        text(s, I(lx), I(cy), I(5.0), I(0.35),
             title, sz=15, color=RED_D, bold=True)
        cy += 0.38
        for item in items:
            text(s, I(lx + 0.15), I(cy), I(5.2), I(0.25),
                 item, sz=11, color=DARK_TEXT)
            cy += 0.26
        cy += 0.15
    text(s, I(5.8), I(3.3), I(0.8), I(0.5),
         "vs", sz=22, color=ORANGE_D, bold=True,
         align=PP_ALIGN.CENTER)
    rx, rw = 6.8, 5.9
    shape(s, I(rx), I(1.2), I(rw), I(0.5),
          GREEN_D, GREEN_D, "Regorus / RVM",
          sz=20, color=WHITE, bold=True)
    rvm_feats = [
        "Multi-language: Rego + Cedar + Azure Policy + more",
        "Compiled register-based bytecode with short-circuit optimization",
        "Rust core \u2014 memory safe, no_std, no GC, WASM compatible",
        "8 language bindings: C, C++, C#, Go, Java, Python, Ruby, JS",
        "HostAwait: suspendable execution for mid-eval host callbacks",
        "Breakpoints, stepping, debug \u2014 suspendable by design",
        "Portable binary artifacts \u2014 compile once, run anywhere",
        "Interactive WASM playground in browser",
        "Policy intelligence: formal analysis via Z3 (upcoming)",
        "Open source (MIT) \u2014 transparent and auditable",
    ]
    fy = 1.85
    for i, feat in enumerate(rvm_feats):
        bg_c = GREEN_XL if i % 2 == 0 else FAINT_GRAY
        shape(s, I(rx), I(fy + i * 0.44), I(rw), I(0.38),
              bg_c, None, feat,
              sz=11, color=DARK_TEXT, align=PP_ALIGN.LEFT)


# ── Version-specific slide builders ──────────────────────────────────────────

def version_A(prs, BL):
    """Demo-First: Title → Why We Built → DEMO → Delivers → Competitive → Close"""
    add_title(prs, BL)

    # Why We Built This
    s = prs.slides.add_slide(BL)
    bgfill(s, WHITE)
    text(s, I(M), I(0.25), I(U), I(0.6),
         "Why We Built This", sz=36, color=RED_D, bold=True)
    text(s, I(M), I(0.85), I(U), I(0.5),
         "Azure Policy must support multiple policy languages. "
         "Across Microsoft, teams maintain separate engines.",
         sz=16, color=MID_GRAY)
    engines = ["OPA / Rego", "Cedar", "Azure Policy", "AKS Admission", "Custom"]
    bw = 2.1
    for name, x in zip(engines, dist(5, bw)):
        shape(s, I(x), I(1.5), I(bw), I(0.7),
              RED_L, RED_D, name, sz=13, color=RED_D, bold=True)
    text(s, I(M), I(2.4), I(U), I(0.4),
         "N engines = N audits + N test suites + N pipelines + N teams",
         sz=14, color=RED_D, bold=True)
    text(s, I(M), I(2.95), I(U), I(0.3),
         "The JVM solved this for applications. RVM solves it for policy.",
         sz=16, color=GREEN_D, bold=True)

    def flow_compact(slide, y, label, accent, boxes):
        text(slide, I(M), I(y - 0.25), I(3), I(0.3),
             label, sz=14, color=accent, bold=True)
        bw_f, gap = 2.2, 0.35
        total_flow = len(boxes) * bw_f + (len(boxes) - 1) * gap
        sx = (SW - total_flow) / 2
        for i, (txt_val, bg_c) in enumerate(boxes):
            x = sx + i * (bw_f + gap)
            shape(slide, I(x), I(y), I(bw_f), I(0.7),
                  bg_c, accent, txt_val, sz=11, color=DARK_TEXT)
            if i < len(boxes) - 1:
                ax = x + bw_f + (gap - 0.2) / 2
                text(slide, I(ax), I(y + 0.1), I(0.2), I(0.4),
                     "\u2192", sz=18, color=accent, bold=True,
                     align=PP_ALIGN.CENTER)

    flow_compact(s, 3.5, "JVM", ORANGE_D, [
        ("Java \u00B7 Kotlin\nScala \u00B7 Groovy", ORANGE_L),
        ("Compilers", ORANGE_L),
        ("Java Bytecode", RGBColor(0xFF, 0xE0, 0xB2)),
        ("JVM: Any Platform", ORANGE_L),
    ])
    flow_compact(s, 4.6, "RVM", GREEN_D, [
        ("Rego \u00B7 Cedar\nAzure Policy \u00B7 more", GREEN_XL),
        ("Compilers", GREEN_XL),
        ("RVM Bytecode", GREEN_L),
        ("RVM: Any Platform", GREEN_XL),
    ])
    text(s, I(M), I(5.7), I(U), I(0.5),
         "One engine to audit, one engine to optimize, "
         "one engine to deploy \u2014 any new language gets everything for free.",
         sz=15, color=DARK_TEXT, bold=True)

    add_demo(prs, BL)

    # What RVM Delivers
    s = prs.slides.add_slide(BL)
    bgfill(s, WHITE)
    text(s, I(M), I(0.25), I(U), I(0.6),
         "What RVM Delivers", sz=36, color=GREEN_D, bold=True)
    cols = [
        ("Performance", [
            "Compiled register-based bytecode",
            "Short-circuit optimization",
            "No cloning, no boxing overhead",
        ], GREEN_D),
        ("Portability", [
            "Compile once, run anywhere",
            "WASM, cloud, edge, bare metal",
            "8 language bindings",
        ], BLUE_D),
        ("Safety", [
            "Rust core \u2014 memory safe",
            "Instruction budget for CPU",
            "Per-eval memory limits",
            "Sandboxed \u2014 no I/O",
        ], TEAL_D),
        ("Extensibility", [
            "HostAwait: mid-eval callbacks",
            "Suspendable execution",
            "Forward-compatible serialization",
        ], PURPLE_D),
    ]
    cw = 2.8
    for (title, items, accent), x in zip(cols, dist(4, cw)):
        shape(s, I(x), I(1.2), I(cw), I(0.5),
              accent, accent, title, sz=15, color=WHITE, bold=True)
        bullets(s, I(x + 0.1), I(1.85), I(cw - 0.15), I(2.5),
                ["\u2013  " + it for it in items],
                sz=12, color=DARK_TEXT)
    text(s, I(M), I(4.6), I(U), I(0.4),
         "One engine to certify  \u00B7  Reproducible results  "
         "\u00B7  Open source (MIT)  \u00B7  Structured audit logs",
         sz=15, color=DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)

    add_competitive(prs, BL)
    add_closing(prs, BL)


def version_B(prs, BL):
    """Bookend: Title → Problem → DEMO → How It Works → Close"""
    add_title(prs, BL)

    # The Problem
    s = prs.slides.add_slide(BL)
    bgfill(s, WHITE)
    text(s, I(M), I(0.3), I(U), I(0.7),
         "The Problem", sz=40, color=RED_D, bold=True)
    text(s, I(M), I(1.1), I(U), I(0.6),
         "Azure Policy must support multiple policy languages.\n"
         "Across Microsoft, teams maintain separate engines.",
         sz=20, color=MID_GRAY)
    engines = ["OPA / Rego", "Cedar", "Azure Policy", "AKS Admission", "Custom"]
    bw = 2.1
    for name, x in zip(engines, dist(5, bw)):
        shape(s, I(x), I(2.1), I(bw), I(1.0),
              RED_L, RED_D, name, sz=15, color=RED_D, bold=True)
    shape(s, I(M), I(3.5), I(U), I(0.65),
          FAINT_GRAY, LIGHT_GRAY,
          "N engines  =  N \u00D7 audits  +  N \u00D7 test suites  +  "
          "N \u00D7 pipelines  +  N \u00D7 teams",
          sz=18, color=RED_D, bold=True)
    text(s, I(M), I(4.7), I(U), I(0.8),
         "What if one engine could run them all?",
         sz=28, color=GREEN_D, bold=True, align=PP_ALIGN.CENTER)
    text(s, I(M), I(5.6), I(U), I(0.6),
         "Let me show you.",
         sz=22, color=MID_GRAY, align=PP_ALIGN.CENTER)

    add_demo(prs, BL)

    # How It Works
    s = prs.slides.add_slide(BL)
    bgfill(s, WHITE)
    text(s, I(M), I(0.25), I(U), I(0.6),
         "How It Works", sz=36, color=GREEN_D, bold=True)
    text(s, I(M), I(0.85), I(U), I(0.4),
         "What Java did for applications, RVM does for policy.",
         sz=17, color=MID_GRAY)

    def flow_row(slide, y, label, accent, boxes):
        text(slide, I(M), I(y - 0.3), I(3), I(0.3),
             label, sz=14, color=accent, bold=True)
        bw_f, gap = 2.4, 0.4
        total_flow = len(boxes) * bw_f + (len(boxes) - 1) * gap
        sx = (SW - total_flow) / 2
        for i, (txt_val, bg_c) in enumerate(boxes):
            x = sx + i * (bw_f + gap)
            shape(slide, I(x), I(y), I(bw_f), I(0.75),
                  bg_c, accent, txt_val, sz=12, color=DARK_TEXT)
            if i < len(boxes) - 1:
                ax = x + bw_f + (gap - 0.25) / 2
                text(slide, I(ax), I(y + 0.15), I(0.25), I(0.4),
                     "\u2192", sz=18, color=accent, bold=True,
                     align=PP_ALIGN.CENTER)

    flow_row(s, 1.5, "JVM", ORANGE_D, [
        ("Java \u00B7 Kotlin\nScala \u00B7 Groovy", ORANGE_L),
        ("Compilers", ORANGE_L),
        ("Java Bytecode", RGBColor(0xFF, 0xE0, 0xB2)),
        ("JVM: Any Platform", ORANGE_L),
    ])
    flow_row(s, 2.6, "RVM", GREEN_D, [
        ("Rego \u00B7 Cedar\nAzure Policy \u00B7 more", GREEN_XL),
        ("Compilers", GREEN_XL),
        ("RVM Bytecode", GREEN_L),
        ("RVM: Any Platform", GREEN_XL),
    ])
    text(s, I(M), I(3.7), I(U), I(0.4),
         "What this means in practice", sz=18, color=BLUE_D, bold=True)
    impl_cols = [
        ("Security & Compliance", [
            "One engine to audit and certify",
            "Memory safe \u2014 Rust, no_std",
            "Sandboxed \u2014 no file or network I/O",
            "Structured audit logs",
        ], TEAL_D),
        ("Performance & Portability", [
            "Compiled register-based bytecode",
            "Compile once, run anywhere",
            "Cloud, edge, WASM, bare metal",
            "8 language bindings",
        ], GREEN_D),
        ("Unique Capabilities", [
            "HostAwait: pause mid-eval for host data",
            "Instruction budget + memory limits",
            "Breakpoints & step-through debugging",
            "Forward-compatible binary artifacts",
        ], PURPLE_D),
    ]
    cw = 3.7
    for (title, items, accent), x in zip(impl_cols, dist(3, cw)):
        shape(s, I(x), I(4.15), I(cw), I(0.45),
              accent, accent, title, sz=14, color=WHITE, bold=True)
        bullets(s, I(x + 0.1), I(4.7), I(cw - 0.15), I(2.2),
                ["\u2013  " + it for it in items],
                sz=12, color=DARK_TEXT, spacing=Pt(5))

    add_closing(prs, BL)


def version_C(prs, BL):
    """Three-Act: Title → Problem → JVM → DEMO → Under Hood → Competitive → Close"""
    add_title(prs, BL)

    # Problem
    s = prs.slides.add_slide(BL)
    bgfill(s, WHITE)
    text(s, I(M), I(0.3), I(U), I(0.7),
         "The Problem: Policy Engine Proliferation",
         sz=36, color=RED_D, bold=True)
    text(s, I(M), I(1.05), I(U), I(0.5),
         "Azure Policy must support multiple policy languages. "
         "Across Microsoft, teams maintain separate engines.",
         sz=18, color=MID_GRAY)
    engines = ["OPA / Rego", "Cedar", "Azure Policy", "AKS Admission", "Custom"]
    bw = 2.1
    for name, x in zip(engines, dist(5, bw)):
        shape(s, I(x), I(1.8), I(bw), I(1.0),
              RED_L, RED_D, name, sz=15, color=RED_D, bold=True)
    costs = [
        "N separate security audits",
        "N separate testing frameworks",
        "N separate deployment pipelines",
        "N separate teams to maintain",
    ]
    bullets(s, I(M + 0.2), I(3.2), I(5.5), I(2.5), costs,
            sz=17, color=DARK_TEXT, spacing=Pt(10))
    shape(s, I(M), I(5.5), I(U), I(0.65),
          FAINT_GRAY, LIGHT_GRAY,
          "N engines  =  N \u00D7 audits  +  N \u00D7 test suites  +  "
          "N \u00D7 pipelines  +  N \u00D7 teams",
          sz=18, color=RED_D, bold=True)
    text(s, I(M), I(6.3), I(U), I(0.5),
         "Multiplicative cost, inconsistent behavior, fragmented tooling.",
         sz=17, color=DARK_TEXT, bold=True)

    # JVM Analogy
    s = prs.slides.add_slide(BL)
    bgfill(s, WHITE)
    text(s, I(M), I(0.3), I(U), I(0.7),
         "What Java Did for Apps, RVM Does for Policy",
         sz=34, color=GREEN_D, bold=True)

    def flow_row(slide, y, label, accent, boxes):
        text(slide, I(M), I(y - 0.45), I(5), I(0.4),
             label, sz=17, color=accent, bold=True)
        bw_f, gap = 2.4, 0.55
        total = 4 * bw_f + 3 * gap
        sx = (SW - total) / 2
        for i, (txt_val, bg_c) in enumerate(boxes):
            x = sx + i * (bw_f + gap)
            shape(slide, I(x), I(y), I(bw_f), I(0.95),
                  bg_c, accent, txt_val, sz=14, color=DARK_TEXT)
            if i < 3:
                ax = x + bw_f + (gap - 0.3) / 2
                text(slide, I(ax), I(y + 0.2), I(0.3), I(0.5),
                     "\u2192", sz=24, color=accent, bold=True,
                     align=PP_ALIGN.CENTER)

    flow_row(s, 1.5, "JVM  (Applications)", ORANGE_D, [
        ("Java \u00B7 Kotlin\nScala \u00B7 Groovy", ORANGE_L),
        ("Compilers", ORANGE_L),
        ("Java Bytecode", RGBColor(0xFF, 0xE0, 0xB2)),
        ("JVM: Any Platform", ORANGE_L),
    ])
    flow_row(s, 3.3, "RVM  (Policy)", GREEN_D, [
        ("Rego \u00B7 Cedar\nAzure Policy \u00B7 more", GREEN_XL),
        ("Compilers", GREEN_XL),
        ("RVM Bytecode", GREEN_L),
        ("RVM: Any Platform", GREEN_XL),
    ])
    text(s, I(M), I(4.8), I(U), I(0.4),
         "What this gives you", sz=20, color=BLUE_D, bold=True)
    bullets(s, I(M + 0.2), I(5.3), I(U), I(2.0), [
        "One security surface to audit",
        "One engine to optimize",
        "One test framework to maintain",
        "One artifact to deploy",
        "Any new language gets the entire tooling ecosystem for free",
    ], sz=15, color=DARK_TEXT)

    add_demo(prs, BL)

    # Under the Hood
    s = prs.slides.add_slide(BL)
    bgfill(s, WHITE)
    text(s, I(M), I(0.25), I(U), I(0.6),
         "Under the Hood", sz=36, color=GREEN_D, bold=True)
    text(s, I(M), I(0.85), I(U), I(0.4),
         "What makes the playground possible \u2014 and what it means "
         "for production.", sz=16, color=MID_GRAY)
    bc_cols = [
        ("Performance", [
            "Register-based bytecode, minimal dispatch",
            "Short-circuit optimization (PolicyOp)",
            "No cloning, no boxing",
        ], GREEN_D),
        ("Portability", [
            "Compile once, run anywhere",
            "Cloud, edge, WASM, bare metal",
            "Forward-compatible binary serialization",
        ], BLUE_D),
        ("Extensibility", [
            "HostAwait: pause mid-eval for host data",
            "Suspendable execution by design",
            "Breakpoints and step-through debugging",
        ], PURPLE_D),
    ]
    cw = 3.7
    for (title, items, accent), x in zip(bc_cols, dist(3, cw)):
        shape(s, I(x), I(1.4), I(cw), I(0.45),
              accent, accent, title, sz=14, color=WHITE, bold=True)
        bullets(s, I(x + 0.1), I(1.95), I(cw - 0.15), I(1.6),
                ["\u2013  " + it for it in items],
                sz=11, color=DARK_TEXT, spacing=Pt(5))
    text(s, I(M), I(3.7), I(U), I(0.4),
         "Trust Through Safety", sz=20, color=TEAL_D, bold=True)
    safety_cols = [
        ("Memory Safety", [
            "Written in Rust \u2014 no buffer overflows",
            "no_std compatible \u2014 bare metal ready",
            "No unsafe code in evaluation path",
        ]),
        ("Execution Safety", [
            "Configurable memory limits per evaluation",
            "Instruction budget \u2014 bounded CPU",
            "Sandboxed \u2014 no file or network I/O",
        ]),
        ("Auditability", [
            "Structured audit logs \u2014 every decision traceable",
            "Deterministic: same input = same output",
            "Open source (MIT) \u2014 fully transparent",
        ]),
    ]
    for (title, items), x in zip(safety_cols, dist(3, cw)):
        shape(s, I(x), I(4.15), I(cw), I(0.4),
              TEAL_D, TEAL_D, title, sz=13, color=WHITE, bold=True)
        bullets(s, I(x + 0.1), I(4.65), I(cw - 0.15), I(1.5),
                ["\u2013  " + it for it in items],
                sz=11, color=DARK_TEXT, spacing=Pt(4))
    shape(s, I(M), I(6.0), I(U), I(0.55),
          FAINT_GRAY, LIGHT_GRAY,
          "One engine to certify  \u00B7  Reproducible results  "
          "\u00B7  Resource & time bounded  \u00B7  Open source (MIT)",
          sz=14, color=DARK_TEXT, bold=True)

    add_competitive(prs, BL)
    add_closing(prs, BL)


def version_D(prs, BL):
    """Progressive Demo: Title → One Problem → Extended DEMO → Close"""
    add_title(prs, BL)

    # One Problem — minimal
    s = prs.slides.add_slide(BL)
    bgfill(s, WHITE)
    text(s, I(1.5), I(1.5), I(10.3), I(1.0),
         "Azure Policy must support\nmultiple policy languages.",
         sz=36, color=RED_D, bold=True, align=PP_ALIGN.CENTER)
    text(s, I(1.5), I(3.0), I(10.3), I(0.8),
         "Across Microsoft, teams maintain separate engines.\n"
         "Every engine means another audit, another test suite, "
         "another deployment.",
         sz=20, color=MID_GRAY, align=PP_ALIGN.CENTER)
    engines = ["OPA / Rego", "Cedar", "Azure Policy", "AKS Admission", "Custom"]
    bw = 2.1
    for name, x in zip(engines, dist(5, bw)):
        shape(s, I(x), I(4.2), I(bw), I(0.8),
              RED_L, RED_D, name, sz=14, color=RED_D, bold=True)
    text(s, I(1.5), I(5.5), I(10.3), I(0.8),
         "What if one engine could run them all?",
         sz=30, color=GREEN_D, bold=True, align=PP_ALIGN.CENTER)

    # Extended DEMO with cues
    s = prs.slides.add_slide(BL)
    bgfill(s, CHARCOAL)
    text(s, I(1.5), I(0.8), I(10.3), I(1.2),
         "LIVE DEMO", sz=72, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER)
    text(s, I(1.5), I(2.2), I(10.3), I(0.5),
         "RVM Playground \u2014 Cross-language policy in your browser",
         sz=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    cues = [
        ("1. Write", "Rego policy in Monaco editor with syntax highlighting"),
        ("2. Compile", "See RVM bytecode and instruction count"),
        ("3. Run", "Set input data, get instant evaluation results"),
        ("4. Switch", "Same policy logic in Cedar \u2014 same bytecode target"),
        ("5. Debug", "Step through bytecode, inspect registers"),
        ("6. Compare", "Side-by-side assembly for Rego vs Cedar"),
    ]
    cw_d = 5.5
    cx = (SW - cw_d) / 2
    cy = 3.2
    for label, desc in cues:
        text(s, I(cx), I(cy), I(1.2), I(0.35),
             label, sz=14, color=GREEN_L, bold=True)
        text(s, I(cx + 1.2), I(cy), I(cw_d - 1.2), I(0.35),
             desc, sz=13, color=LIGHT_GRAY)
        cy += 0.42
    text(s, I(1.5), I(6.2), I(10.3), I(0.5),
         "https://anakrish.github.io/rego-virtual-machine-playground/",
         sz=14, color=AZURE, align=PP_ALIGN.CENTER)

    # Closing — richer since fewer slides
    s = prs.slides.add_slide(BL)
    bgfill(s, GREEN_D)
    text(s, I(1.5), I(1.2), I(10.3), I(1.0),
         "One Engine. Every Policy. Any Platform.",
         sz=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(s, I(1.5), I(2.5), I(10.3), I(0.8),
         "Regorus Virtual Machine",
         sz=32, color=GREEN_L, align=PP_ALIGN.CENTER)
    for i, item in enumerate([
        "Rego  +  Cedar  +  Azure Policy \u2014 one bytecode target",
        "High performance  \u00B7  Memory safe  \u00B7  Portable",
        "8 language bindings: C, C++, C#, Go, Java, Python, Ruby, JS",
        "WASM Playground  \u00B7  Debugger  \u00B7  Disassembler",
        "Open Source (MIT)  \u00B7  github.com/microsoft/regorus",
    ]):
        text(s, I(1.5), I(3.8 + i * 0.5), I(10.3), I(0.5),
             item, sz=18, color=WHITE, align=PP_ALIGN.CENTER)
    text(s, I(1.5), I(6.3), I(10.3), I(0.5),
         "Try it: https://anakrish.github.io/"
         "rego-virtual-machine-playground/",
         sz=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def version_E(prs, BL):
    """Appendix: Title → Problem+Solution → DEMO → Closing w/ takeaways + 4 appendix"""
    add_title(prs, BL)

    # Problem + Solution combined
    s = prs.slides.add_slide(BL)
    bgfill(s, WHITE)
    lw = 5.8
    text(s, I(M), I(0.3), I(lw), I(0.6),
         "The Problem", sz=32, color=RED_D, bold=True)
    text(s, I(M), I(0.9), I(lw), I(0.5),
         "Azure Policy must support multiple policy languages.\n"
         "Across Microsoft, teams maintain separate engines.",
         sz=15, color=MID_GRAY)
    engines_l = ["OPA / Rego", "Cedar", "Azure Policy"]
    bw_l = 1.7
    for i, name in enumerate(engines_l):
        x = M + i * (bw_l + 0.2)
        shape(s, I(x), I(1.65), I(bw_l), I(0.6),
              RED_L, RED_D, name, sz=12, color=RED_D, bold=True)
    for i, name in enumerate(["AKS Admission", "Custom"]):
        x = M + i * (bw_l + 0.2)
        shape(s, I(x), I(2.35), I(bw_l), I(0.6),
              RED_L, RED_D, name, sz=12, color=RED_D, bold=True)
    text(s, I(M), I(3.15), I(lw), I(0.4),
         "N audits  \u00B7  N test suites  \u00B7  N pipelines  \u00B7  N teams",
         sz=13, color=RED_D, bold=True)
    # Divider
    shape(s, I(6.5), I(0.5), I(0.02), I(6.0), LIGHT_GRAY, None)
    # Solution
    rx = 7.0
    rw = 5.8
    text(s, I(rx), I(0.3), I(rw), I(0.6),
         "The Solution", sz=32, color=GREEN_D, bold=True)
    text(s, I(rx), I(0.9), I(rw), I(0.5),
         "What Java did for applications,\nRVM does for policy.",
         sz=15, color=MID_GRAY)
    flow = [
        ("Rego \u00B7 Cedar\nAzure Policy", GREEN_XL, GREEN_D),
        ("Compilers", GREEN_XL, GREEN_D),
        ("RVM\nBytecode", GREEN_L, GREEN_D),
        ("Any\nPlatform", GREEN_XL, GREEN_D),
    ]
    fw = 1.2
    positions = dist(4, fw, total=rw, start=rx)
    for i, ((txt_val, bg_c, border), x) in enumerate(zip(flow, positions)):
        shape(s, I(x), I(1.65), I(fw), I(0.85),
              bg_c, border, txt_val, sz=10, color=DARK_TEXT)
        if i < 3:
            gap_x = positions[i + 1] - x - fw
            text(s, I(x + fw), I(1.85), I(gap_x), I(0.4),
                 "\u2192", sz=16, color=GREEN_D, bold=True,
                 align=PP_ALIGN.CENTER)
    bullets(s, I(rx), I(2.8), I(rw), I(2.5), [
        "\u2013  One engine to audit and certify",
        "\u2013  Compiled bytecode \u2014 high performance",
        "\u2013  Rust core \u2014 memory safe, no_std",
        "\u2013  8 language bindings \u2014 deploy anywhere",
        "\u2013  HostAwait \u2014 suspend mid-eval for host data",
        "\u2013  WASM playground in browser",
    ], sz=13, color=DARK_TEXT, spacing=Pt(6))
    shape(s, I(M), I(5.8), I(U), I(0.6),
          GREEN_D, GREEN_D,
          "Let me show you \u2014 live, in your browser.",
          sz=20, color=WHITE, bold=True)

    add_demo(prs, BL)

    # Closing with takeaway boxes
    s = prs.slides.add_slide(BL)
    bgfill(s, GREEN_D)
    text(s, I(1.5), I(0.8), I(10.3), I(1.0),
         "One Engine. Every Policy. Any Platform.",
         sz=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(s, I(1.5), I(2.0), I(10.3), I(0.6),
         "Regorus Virtual Machine",
         sz=32, color=GREEN_L, align=PP_ALIGN.CENTER)
    takeaways = [
        ("Multi-Language", "Rego + Cedar +\nAzure Policy \u2192\none bytecode target"),
        ("High Performance", "Compiled register VM\nShort-circuit optimization\nNo GC, no cloning"),
        ("Deploy Anywhere", "Cloud \u00B7 Edge \u00B7 WASM\n8 language bindings\nno_std compatible"),
        ("Safe & Auditable", "Memory safe (Rust)\nInstruction + memory limits\nOpen source (MIT)"),
    ]
    tw = 2.7
    for (title, desc), x in zip(takeaways, dist(4, tw)):
        shape(s, I(x), I(3.0), I(tw), I(0.45),
              GREEN_M, GREEN_M, title, sz=14, color=WHITE, bold=True)
        shape(s, I(x), I(3.55), I(tw), I(1.3),
              RGBColor(0x15, 0x4A, 0x15), None, desc,
              sz=12, color=GREEN_L)
    text(s, I(1.5), I(5.3), I(10.3), I(0.5),
         "Try it: https://anakrish.github.io/"
         "rego-virtual-machine-playground/",
         sz=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(s, I(1.5), I(5.9), I(10.3), I(0.5),
         "github.com/microsoft/regorus",
         sz=17, color=GREEN_L, align=PP_ALIGN.CENTER)

    # ── Appendix slides ──────────────────────────────────────────────────
    add_competitive(prs, BL, appendix=True)

    # Tooling
    s = prs.slides.add_slide(BL)
    bgfill(s, WHITE)
    text(s, I(M), I(0.12), I(U), I(0.35),
         "APPENDIX", sz=11, color=MID_GRAY, bold=True)
    text(s, I(M), I(0.3), I(U), I(0.7),
         "Developer Tooling & Policy Intelligence",
         sz=34, color=GREEN_D, bold=True)
    text(s, I(M), I(0.95), I(U), I(0.4),
         "Tools built once work across every policy language.",
         sz=15, color=MID_GRAY)
    lx_t, lw_t = 0.4, 2.0
    text(s, I(lx_t), I(1.5), I(lw_t), I(0.3),
         "Policy Languages", sz=13, color=BLUE_D,
         bold=True, align=PP_ALIGN.CENTER)
    for i, lang in enumerate(["Rego", "Cedar", "Azure Policy", "Future DSLs"]):
        shape(s, I(lx_t), I(1.9 + i * 0.5), I(lw_t), I(0.4),
              BLUE_L, BLUE_D, lang, sz=12, color=DARK_TEXT)
    text(s, I(2.5), I(2.9), I(0.5), I(0.4),
         "\u2192", sz=20, color=BLUE_D, bold=True,
         align=PP_ALIGN.CENTER)
    shape(s, I(3.1), I(1.7), I(2.4), I(2.7),
          GREEN_D, GREEN_D,
          "RVM\n\nUnified AST/IR\nCompiler\nRuntime",
          sz=14, color=WHITE, bold=True)
    text(s, I(5.6), I(2.9), I(0.5), I(0.4),
         "\u2192", sz=20, color=GREEN_M, bold=True,
         align=PP_ALIGN.CENTER)
    tcw, tg, tx0 = 1.7, 0.15, 6.2
    tool_cols = [
        ("Editor Experience", PURPLE_D, PURPLE_L,
         ["LSP Server", "IntelliSense", "Go-to-Definition", "Autocomplete"]),
        ("Code Quality", ORANGE_D, ORANGE_L,
         ["Linter", "Formatter", "Static Analysis", "Dead Rule Finder"]),
        ("DevOps & CI", TEAL_D, TEAL_L,
         ["Test Framework", "Doc Generator", "CI/CD Gates", "Coverage"]),
    ]
    for ci, (title, hc, ic, items) in enumerate(tool_cols):
        tx = tx0 + ci * (tcw + tg)
        shape(s, I(tx), I(1.7), I(tcw), I(0.4),
              hc, hc, title, sz=11, color=WHITE, bold=True)
        for j, tool in enumerate(items):
            shape(s, I(tx), I(2.18 + j * 0.42), I(tcw), I(0.36),
                  ic, hc, tool, sz=10, color=DARK_TEXT)
    pi_x, pi_w, pi_y = 6.2, 5.5, 4.2
    shape(s, I(pi_x), I(pi_y), I(pi_w), I(0.45),
          RGBColor(0x4A, 0x14, 0x8C), RGBColor(0x4A, 0x14, 0x8C),
          "Policy Intelligence \u2014 Formal Analysis "
          "(separate presentation)",
          sz=14, color=WHITE, bold=True)
    for i, t in enumerate([
        "Z3-powered\nSatisfiability",
        "Conflict &\nRedundancy Detection",
        "Policy\nEquivalence Proofs",
        "Coverage &\nReachability",
    ]):
        iw = pi_w / 4 - 0.06
        ix = pi_x + i * (iw + 0.08)
        shape(s, I(ix), I(pi_y + 0.55), I(iw), I(0.75),
              PURPLE_L, RGBColor(0x4A, 0x14, 0x8C), t,
              sz=11, color=DARK_TEXT)
    text(s, I(M), I(6.2), I(U), I(0.5),
         "Without RVM: N tools \u00D7 M languages.  "
         "With RVM: build each tool once.",
         sz=15, color=DARK_TEXT, bold=True,
         align=PP_ALIGN.CENTER)

    # Safety
    s = prs.slides.add_slide(BL)
    bgfill(s, WHITE)
    text(s, I(M), I(0.12), I(U), I(0.35),
         "APPENDIX", sz=11, color=MID_GRAY, bold=True)
    text(s, I(M), I(0.3), I(U), I(0.7),
         "Trust Through Safety",
         sz=36, color=GREEN_D, bold=True)
    col_w = 3.7
    safety = [
        ("Memory Safety", [
            "Written in Rust \u2014 no buffer overflows",
            "no_std compatible \u2014 bare metal ready",
            "No unsafe code in evaluation path",
        ]),
        ("Execution Safety", [
            "Configurable memory limits per evaluation",
            "Instruction budget \u2014 bounded CPU time",
            "Deterministic: same input = same output",
            "Sandboxed \u2014 no file or network I/O",
        ]),
        ("Auditability", [
            "Full disassembly with source mapping",
            "Breakpoints & step-through debugging",
            "Structured audit logs \u2014 every decision traceable",
            "Forward-compatible binary serialization",
            "Deterministic \u2014 enables formal verification",
            "Open source (MIT) \u2014 fully transparent",
        ]),
    ]
    for (title, items), x in zip(safety, dist(3, col_w)):
        shape(s, I(x), I(1.3), I(col_w), I(0.5),
              GREEN_M, GREEN_M, title, sz=16, color=WHITE, bold=True)
        bullets(s, I(x + 0.15), I(1.95), I(col_w - 0.2), I(3.0),
                ["\u2013  " + it for it in items],
                sz=12, color=DARK_TEXT, spacing=Pt(6))
    text(s, I(M), I(5.2), I(U), I(0.4),
         "Compliance Impact", sz=18, color=BLUE_D,
         bold=True, align=PP_ALIGN.CENTER)
    cbw = 2.7
    for t, x in zip(
        ["One engine to certify", "Reproducible results",
         "Resource & time bounded", "Open source (MIT)"],
        dist(4, cbw),
    ):
        shape(s, I(x), I(5.7), I(cbw), I(0.75),
              BLUE_L, BLUE_D, t, sz=14, color=DARK_TEXT, bold=True)

    # Bytecode
    s = prs.slides.add_slide(BL)
    bgfill(s, WHITE)
    text(s, I(M), I(0.12), I(U), I(0.35),
         "APPENDIX", sz=11, color=MID_GRAY, bold=True)
    text(s, I(M), I(0.3), I(U), I(0.7),
         "Why Bytecode Matters",
         sz=36, color=GREEN_D, bold=True)
    text(s, I(M), I(1.0), I(U), I(0.4),
         "Compiling policy to bytecode unlocks capabilities "
         "you can\u2019t get from a tree-walking interpreter.",
         sz=17, color=MID_GRAY)
    pillars = [
        ("Performance", [
            "Much faster than tree-walking",
            "Register-based, minimal dispatch",
            "No cloning, no boxing",
        ], RED_D),
        ("Portable Artifacts", [
            "Compile once, run anywhere",
            "Ship pre-compiled bundles",
            "Forward-compatible serialization",
        ], GREEN_D),
        ("Tooling Ecosystem", [
            "Disassembly + source mapping",
            "Breakpoints & step-through",
            "Build once, all languages",
        ], AZURE),
        ("Suspendable Execution", [
            "HostAwait: pause mid-eval",
            "Fetch external data on demand",
            "Resume where you left off",
        ], PURPLE_D),
        ("Safety & Governance", [
            "Instruction budget for CPU",
            "Per-evaluation memory limits",
            "Deterministic, one engine to certify",
        ], TEAL_D),
    ]
    pw = 2.3
    for (title, items, accent), x in zip(pillars, dist(5, pw)):
        shape(s, I(x), I(1.7), I(pw), I(0.55),
              accent, accent, title, sz=14, color=WHITE, bold=True)
        bullets(s, I(x + 0.1), I(2.4), I(pw - 0.15), I(2.0),
                ["\u2013  " + it for it in items],
                sz=11, color=DARK_TEXT)
    half = U / 2 - 0.1
    shape(s, I(M), I(4.5), I(half), I(0.55),
          RED_L, None,
          "Without bytecode: interpreter \u00D7 tooling \u00D7 deploy per language",
          sz=12, color=RED_D, bold=True)
    shape(s, I(M + half + 0.2), I(4.5), I(half), I(0.55),
          GREEN_L, None,
          "With RVM: compile any language \u2192 one artifact \u2192 "
          "tooling for free",
          sz=12, color=GREEN_D, bold=True)


# ── Generate all merged versions ─────────────────────────────────────────────

versions = [
    ("A-DemoFirst", version_A),
    ("B-Bookend", version_B),
    ("C-ThreeAct", version_C),
    ("D-ProgressiveDemo", version_D),
    ("E-Appendix", version_E),
]

for name, func in versions:
    prs = Presentation(SRC)
    # Delete slide 3 (old "Regorus Engine" slide, index 2)
    delete_slide(prs, 2)
    BL = prs.slide_layouts[6]  # Blank
    func(prs, BL)
    out = os.path.join(BASE, f"RVM-Merged-{name}.pptx")
    prs.save(out)
    print(f"  {name}: {out}  ({len(prs.slides)} slides)")

print("Done — all merged versions generated.")
