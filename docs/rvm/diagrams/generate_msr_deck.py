#!/usr/bin/env python3
"""Generate the MSR Verification Workshop presentation.

Title: Policy Analysis for Rego and Cedar
Duration: 10 minutes
Style: Academic — clean, content-dense, white backgrounds

Structure (8 slides):
  1    Title
  2    CTO Memo (preserved from SlidesToReuse.pptx)
  3    Why Formal Assurance Matters for Policy
  4    The RVM: One Bytecode, Multiple Languages
  5    The Assurance Spectrum
  6    Technical Approach: Z3 Encoding
  7    Where We Are Today (Results)
  8    Open Challenges & Future Work

Output: RVM-Policy-Analysis-MSR-Workshop.pptx
"""

import os
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

try:
    from pptx.enum.text import MSO_ANCHOR
except ImportError:
    MSO_ANCHOR = None

BASE = os.path.dirname(os.path.abspath(__file__))
REUSE = os.path.join(BASE, "SlidesToReuse.pptx")

# ── Dimensions ───────────────────────────────────────────────────────────
SW, SH = 13.333, 7.5
M = 0.8          # wider margins for academic feel
U = SW - 2 * M

# ── Palette (muted, academic) ───────────────────────────────────────────
BLACK       = RGBColor(0x1A, 0x1A, 0x1A)
DARK_TEXT   = RGBColor(0x2A, 0x2A, 0x2A)
MID_GRAY    = RGBColor(0x60, 0x60, 0x60)
LIGHT_GRAY  = RGBColor(0xD8, 0xD8, 0xD8)
FAINT_GRAY  = RGBColor(0xF2, 0xF2, 0xF2)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

NAVY        = RGBColor(0x0D, 0x47, 0xA1)
NAVY_L      = RGBColor(0xE3, 0xF2, 0xFD)
TEAL        = RGBColor(0x00, 0x69, 0x5C)
TEAL_L      = RGBColor(0xE0, 0xF2, 0xF1)
GREEN_D     = RGBColor(0x1B, 0x5E, 0x20)
GREEN_L     = RGBColor(0xE8, 0xF5, 0xE9)
RED_D       = RGBColor(0xC6, 0x28, 0x28)
RED_L       = RGBColor(0xFF, 0xCD, 0xD2)
ORANGE_D    = RGBColor(0xEF, 0x6C, 0x00)
ORANGE_L    = RGBColor(0xFF, 0xF3, 0xE0)
PURPLE_D    = RGBColor(0x4A, 0x14, 0x8C)
PURPLE_L    = RGBColor(0xF3, 0xE5, 0xF5)
AZURE       = RGBColor(0x00, 0x78, 0xD4)


# ── Helpers ──────────────────────────────────────────────────────────────

def I(v): return Inches(v)

def dist(n, w, total=None, start=None):
    if total is None: total = U
    if start is None: start = M
    if n <= 1: return [start + (total - w) / 2]
    gap = (total - n * w) / (n - 1)
    return [start + i * (w + gap) for i in range(n)]

def setbg(slide, color):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def shape(slide, l, t, w, h, fill, border=None, txt="",
          sz=14, color=DARK_TEXT, bold=False,
          align=PP_ALIGN.CENTER, rr=True):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rr else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(st, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border; sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    tf = sh.text_frame; tf.word_wrap = True; tf.auto_size = None
    tf.margin_left = Pt(8); tf.margin_right = Pt(8)
    tf.margin_top = Pt(4); tf.margin_bottom = Pt(4)
    if MSO_ANCHOR:
        try: tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except: pass
    if txt:
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = txt
        r.font.size = Pt(sz); r.font.color.rgb = color; r.font.bold = bold
    return sh

def text(slide, l, t, w, h, txt, sz=14, color=DARK_TEXT,
         bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = txt
    r.font.size = Pt(sz); r.font.color.rgb = color; r.font.bold = bold
    return tb

def text_italic(slide, l, t, w, h, txt, sz=14, color=DARK_TEXT,
                align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = txt
    r.font.size = Pt(sz); r.font.color.rgb = color; r.font.italic = True
    return tb

def bullets(slide, l, t, w, h, items, sz=14, color=DARK_TEXT,
            spacing=Pt(6), bold=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        r = p.add_run(); r.text = item
        r.font.size = Pt(sz); r.font.color.rgb = color; r.font.bold = bold
    return tb

def hyperlink_text(slide, l, t, w, h, txt, url, sz=14, color=AZURE,
                   bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = txt
    r.font.size = Pt(sz); r.font.color.rgb = color
    r.font.bold = bold; r.font.underline = True
    r.hyperlink.address = url
    return tb

def delete_slide(prs, index):
    rIdList = prs.slides._sldIdLst
    sldId = list(rIdList)[index]
    rId = sldId.get(qn("r:id"))
    prs.part.drop_rel(rId)
    rIdList.remove(sldId)

def move_slide_to_front(prs, src_index):
    rIdList = prs.slides._sldIdLst
    items = list(rIdList)
    target = items[src_index]
    rIdList.remove(target)
    rIdList.insert(0, target)


# ═════════════════════════════════════════════════════════════════════════
#   SLIDE 1 — Title
# ═════════════════════════════════════════════════════════════════════════

def slide_title(prs, BL):
    s = prs.slides.add_slide(BL)
    setbg(s, WHITE)

    text(s, I(M), I(1.1), I(U), I(1.0),
         "Policy Analysis for Rego, Cedar, and Azure Policy",
         sz=42, color=BLACK, bold=True, align=PP_ALIGN.CENTER)

    text(s, I(M), I(2.4), I(U), I(0.6),
         "Formal Verification via Bytecode-Level Z3 Encoding",
         sz=22, color=MID_GRAY, align=PP_ALIGN.CENTER)

    shape(s, I(SW/2 - 2.5), I(3.3), I(5.0), I(0.03),
          LIGHT_GRAY, None, "", rr=False)

    text(s, I(M), I(3.6), I(U), I(0.5),
         "Anand Krishnamoorthi",
         sz=24, color=DARK_TEXT, align=PP_ALIGN.CENTER)

    text(s, I(M), I(4.15), I(U), I(0.4),
         "Azure Policy Team  \u00B7  Microsoft",
         sz=18, color=MID_GRAY, align=PP_ALIGN.CENTER)

    text(s, I(M), I(4.7), I(U), I(0.4),
         "MSR Verification Workshop  \u00B7  2026",
         sz=16, color=MID_GRAY, align=PP_ALIGN.CENTER)

    hyperlink_text(s, I(M), I(5.5), I(U/3 - 0.1), I(0.4),
                   "github.com/microsoft/regorus",
                   "https://github.com/microsoft/regorus",
                   sz=14, color=AZURE, align=PP_ALIGN.CENTER)

    hyperlink_text(s, I(M + U/3), I(5.5), I(U/3 - 0.1), I(0.4),
                   "RVM Playground",
                   "https://anakrish.github.io/rego-virtual-machine-playground/",
                   sz=14, color=AZURE, align=PP_ALIGN.CENTER)

    hyperlink_text(s, I(M + 2*U/3 + 0.1), I(5.5), I(U/3 - 0.1), I(0.4),
                   "RVM Policy Analysis Demo",
                   "https://anakrish.github.io/rvm-policy-analysis/",
                   sz=14, color=AZURE, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════
#   SLIDE 2 — CTO Memo (preserved from SlidesToReuse.pptx)
# ═════════════════════════════════════════════════════════════════════════


# ═════════════════════════════════════════════════════════════════════════
#   SLIDE 3 — Why Formal Assurance Matters for Policy
# ═════════════════════════════════════════════════════════════════════════

def slide_why_assurance(prs, BL):
    s = prs.slides.add_slide(BL)
    setbg(s, WHITE)

    text(s, I(M), I(0.3), I(U), I(0.5),
         "Why Formal Analysis?",
         sz=34, color=BLACK, bold=True)

    shape(s, I(M), I(0.85), I(U), I(0.03), LIGHT_GRAY, None, "", rr=False)

    # Two columns: the problem / current gaps
    lw = U / 2 - 0.2
    rx = M + U / 2 + 0.2
    rw = lw

    text(s, I(M), I(1.1), I(lw), I(0.3),
         "The problem",
         sz=18, color=NAVY, bold=True)

    bullets(s, I(M + 0.1), I(1.5), I(lw - 0.1), I(2.5), [
        "Billions of policy evaluations/day across Azure",
        "1400+ built-in Azure Policy definitions",
        "Rego/OPA is the de facto standard for K8s,\n"
        "service mesh, IAM",
        "Every policy change is a regression risk\n"
        "at cloud scale",
        "LLMs are starting to author policies \u2014\n"
        "faster than humans can review",
    ], sz=13, color=DARK_TEXT)

    text(s, I(rx), I(1.1), I(rw), I(0.3),
         "What we do today (and why it\u2019s not enough)",
         sz=18, color=RED_D, bold=True)

    today_items = [
        ("Manual tests",
         "Hand-written; miss edge cases by definition"),
        ("Code review",
         "Can\u2019t enumerate all input combinations"),
        ("Shadow mode",
         "Needs production traffic, takes weeks,\n"
         "only covers inputs that actually happen"),
    ]
    ty = 1.5
    for (title, desc) in today_items:
        text(s, I(rx + 0.1), I(ty), I(rw - 0.1), I(0.2),
             title, sz=13, color=RED_D, bold=True)
        text(s, I(rx + 0.1), I(ty + 0.22), I(rw - 0.2), I(0.4),
             desc, sz=12, color=MID_GRAY)
        ty += 0.65

    # Bottom question
    shape(s, I(M), I(4.2), I(U), I(0.5), FAINT_GRAY, LIGHT_GRAY,
          "Can we get mathematical guarantees \u2014 "
          "not for a few test inputs, but for all possible inputs?",
          sz=17, color=DARK_TEXT, bold=True, align=PP_ALIGN.LEFT)

    # What kinds of guarantees (folded from Assurance Spectrum)
    text(s, I(M), I(5.0), I(U), I(0.25),
         "What we can answer with Z3",
         sz=16, color=NAVY, bold=True)

    caps = [
        ("Property verification", "\u201CNo input bypasses this deny.\u201D"),
        ("Equivalence", "\u201CDoes v2 behave the same as v1?\u201D"),
        ("Subsumption", "\u201CIs the new policy at least as restrictive?\u201D"),
        ("Test generation", "Minimal inputs covering all paths + MC/DC"),
        ("Gap detection", "Inputs where no rule matches (silent fail-open)"),
    ]
    cx = M
    cw = U / len(caps)
    for (label, desc) in caps:
        text(s, I(cx + 0.05), I(5.3), I(cw - 0.1), I(0.2),
             label, sz=11, color=NAVY, bold=True)
        text(s, I(cx + 0.05), I(5.5), I(cw - 0.1), I(0.4),
             desc, sz=10, color=MID_GRAY)
        cx += cw


# ═════════════════════════════════════════════════════════════════════════
#   SLIDE 4 — The RVM: One Bytecode, Multiple Languages
# ═════════════════════════════════════════════════════════════════════════

def slide_rvm(prs, BL):
    s = prs.slides.add_slide(BL)
    setbg(s, WHITE)

    text(s, I(M), I(0.3), I(U), I(0.5),
         "The RVM: One Bytecode, Multiple Languages",
         sz=34, color=BLACK, bold=True)

    shape(s, I(M), I(0.85), I(U), I(0.03), LIGHT_GRAY, None, "", rr=False)

    # ── Diagram: Languages → Compilers → RVM Bytecode → (Evaluate | Analyze)
    # Three source language boxes funneling into center bytecode,
    # then fanning out to two paths

    # Source language boxes (stacked vertically on left)
    lang_x = 0.5
    lang_w = 2.3
    lang_h = 0.65
    lang_colors = [
        ("Rego (OPA)", NAVY_L, NAVY),
        ("Cedar (Amazon)", TEAL_L, TEAL),
        ("Azure Policy", PURPLE_L, PURPLE_D),
    ]
    lang_ys = [1.2, 2.0, 2.8]
    for i, (lbl, bg, border) in enumerate(lang_colors):
        shape(s, I(lang_x), I(lang_ys[i]), I(lang_w), I(lang_h),
              bg, border, lbl, sz=14, color=border, bold=True)

    # Arrows from languages to compilers
    arr_x1 = lang_x + lang_w + 0.05
    for y in lang_ys:
        text(s, I(arr_x1), I(y + 0.1), I(0.35), I(0.4),
             "\u2192", sz=22, color=MID_GRAY, bold=True, align=PP_ALIGN.CENTER)

    # Compiler box
    comp_x = 3.2
    comp_w = 1.8
    comp_h = 2.25
    shape(s, I(comp_x), I(1.2), I(comp_w), I(comp_h),
          GREEN_L, GREEN_D,
          "Per-Language\nCompilers\n\nAST \u2192 register\nbytecode",
          sz=12, color=DARK_TEXT)

    # Arrow to bytecode
    text(s, I(comp_x + comp_w + 0.05), I(2.05), I(0.35), I(0.4),
         "\u2192", sz=26, color=GREEN_D, bold=True, align=PP_ALIGN.CENTER)

    # Central RVM Bytecode box (prominent)
    bc_x = 5.5
    bc_w = 2.5
    bc_h = 2.25
    shape(s, I(bc_x), I(1.2), I(bc_w), I(bc_h),
          ORANGE_L, ORANGE_D,
          "RVM Bytecode\n\nRegister-based IR\nUntyped values\nStructured guards\nExplicit control flow",
          sz=13, color=DARK_TEXT, bold=False)

    # Two fan-out arrows from bytecode → Evaluate and → Analyze
    arr_x2 = bc_x + bc_w + 0.05

    # Upper arrow → Evaluate
    text(s, I(arr_x2), I(1.55), I(0.5), I(0.4),
         "\u2192", sz=26, color=GREEN_D, bold=True, align=PP_ALIGN.CENTER)
    # Lower arrow → Analyze
    text(s, I(arr_x2), I(2.7), I(0.5), I(0.4),
         "\u2192", sz=26, color=PURPLE_D, bold=True, align=PP_ALIGN.CENTER)

    # Evaluate box (top-right)
    eval_x = 8.6
    eval_w = 2.2
    eval_h = 1.0
    shape(s, I(eval_x), I(1.2), I(eval_w), I(eval_h),
          GREEN_L, GREEN_D,
          "Evaluate\n\nFast execution engine\n(production use)",
          sz=12, color=GREEN_D)

    # Arrow from Evaluate → Output
    text(s, I(eval_x + eval_w + 0.05), I(1.45), I(0.4), I(0.4),
         "\u2192", sz=22, color=GREEN_D, bold=True, align=PP_ALIGN.CENTER)

    shape(s, I(11.2), I(1.2), I(1.6), I(eval_h),
          FAINT_GRAY, GREEN_D,
          "allow / deny\naudit / log",
          sz=12, color=DARK_TEXT, bold=True)

    # Analyze box (bottom-right)
    ana_x = 8.6
    ana_w = 2.2
    ana_h = 1.0
    shape(s, I(ana_x), I(2.5), I(ana_w), I(ana_h),
          PURPLE_L, PURPLE_D,
          "Analyze (Z3)\n\nSymbolic execution\n\u2200-input proofs",
          sz=12, color=PURPLE_D)

    # Arrow from Analyze → Outputs
    text(s, I(ana_x + ana_w + 0.05), I(2.75), I(0.4), I(0.4),
         "\u2192", sz=22, color=PURPLE_D, bold=True, align=PP_ALIGN.CENTER)

    shape(s, I(11.2), I(2.5), I(1.6), I(ana_h),
          FAINT_GRAY, PURPLE_D,
          "proofs / tests\ndiffs / coverage",
          sz=12, color=DARK_TEXT, bold=True)

    # What this buys you
    shape(s, I(M), I(3.75), I(U), I(0.5), NAVY_L, NAVY,
          "Compile once, then both run and analyze the same bytecode. "
          "Works for Rego, Cedar, and Azure Policy without separate toolchains.",
          sz=14, color=NAVY, bold=False, align=PP_ALIGN.LEFT)

    # RVM properties
    text(s, I(M), I(4.5), I(U), I(0.3),
         "RVM properties that enable analysis",
         sz=17, color=DARK_TEXT, bold=True)

    details = [
        ("Register-based (not stack)",
         "Named registers map directly to Z3 variables \u2014 no stack simulation needed"),
        ("Structured guards",
         "AllOfGuard / AnyOfGuard \u2192 Z3 And/Or with short-circuit semantics"),
        ("Untyped values",
         "Values carry dynamic type tags \u2014 modeled as Z3 algebraic datatypes "
         "(String | Number | Bool | Set | Object | Undefined)"),
        ("Explicit control flow",
         "No exceptions, no hidden branches \u2014 complete symbolic execution is tractable"),
    ]
    dy = 4.95
    for (title, desc) in details:
        text(s, I(M + 0.15), I(dy), I(2.9), I(0.25),
             title, sz=13, color=NAVY, bold=True)
        text(s, I(M + 3.2), I(dy), I(U - 3.4), I(0.25),
             desc, sz=12, color=MID_GRAY)
        dy += 0.36


# ═════════════════════════════════════════════════════════════════════════
#   SLIDE 5 — AI + Formal Verification
# ═════════════════════════════════════════════════════════════════════════

def slide_ai_verification(prs, BL):
    s = prs.slides.add_slide(BL)
    setbg(s, WHITE)

    text(s, I(M), I(0.3), I(U), I(0.5),
         "AI + Formal Verification",
         sz=34, color=BLACK, bold=True)

    shape(s, I(M), I(0.85), I(U), I(0.03), LIGHT_GRAY, None, "", rr=False)

    text(s, I(M), I(1.05), I(U), I(0.35),
         "LLMs generate policies fast. How do you know they\u2019re correct?",
         sz=18, color=MID_GRAY)

    # Feedback loop: Intent \u2192 LLM \u2192 Policy \u2192 Z3 \u2192 result
    loop_y = 1.7
    bh = 1.1

    shape(s, I(0.5), I(loop_y), I(2.2), I(bh), FAINT_GRAY, MID_GRAY,
          "Natural Language\nIntent\n\n\u201CDeny VMs without\nencryption tags\u201D",
          sz=12, color=DARK_TEXT)

    text(s, I(2.75), I(loop_y + bh/2 - 0.15), I(0.4), I(0.3),
         "\u2192", sz=24, color=MID_GRAY, bold=True, align=PP_ALIGN.CENTER)

    shape(s, I(3.2), I(loop_y), I(2.3), I(bh), PURPLE_L, PURPLE_D,
          "Copilot / LLM\n\nGenerates Rego,\nCedar, or Azure Policy",
          sz=12, color=PURPLE_D)

    text(s, I(5.55), I(loop_y + bh/2 - 0.15), I(0.4), I(0.3),
         "\u2192", sz=24, color=PURPLE_D, bold=True, align=PP_ALIGN.CENTER)

    shape(s, I(6.0), I(loop_y), I(2.0), I(bh), ORANGE_L, ORANGE_D,
          "Generated\nPolicy\n\n(candidate)",
          sz=12, color=ORANGE_D)

    text(s, I(8.05), I(loop_y + bh/2 - 0.15), I(0.4), I(0.3),
         "\u2192", sz=24, color=ORANGE_D, bold=True, align=PP_ALIGN.CENTER)

    shape(s, I(8.5), I(loop_y), I(2.2), I(bh), NAVY_L, NAVY,
          "RVM Compile\n+ Z3 Verify\n\nProve properties",
          sz=12, color=NAVY)

    text(s, I(10.75), I(loop_y + bh/2 - 0.15), I(0.4), I(0.3),
         "\u2192", sz=24, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    shape(s, I(11.2), I(loop_y), I(1.6), I(bh), GREEN_L, GREEN_D,
          "\u2713 Proof\n\nDeploy with\nconfidence",
          sz=12, color=GREEN_D, bold=True)

    # Counterexample loop
    shape(s, I(3.2), I(loop_y + bh + 0.15), I(7.5), I(0.5),
          RED_L, RED_D,
          "\u2717 Counterexample \u2192 concrete failing input \u2192 "
          "LLM revises \u2192 re-verify (iterate until proof)",
          sz=12, color=RED_D, align=PP_ALIGN.LEFT)

    # Brief notes below
    text(s, I(M), I(3.7), I(U), I(0.3),
         "Why this matters",
         sz=18, color=DARK_TEXT, bold=True)

    notes = [
        "LLMs produce syntactically valid but semantically wrong policies "
        "(overly broad wildcards, missing negation, wrong conditions).",
        "Z3 counterexamples are concrete inputs \u2014 the LLM can use them "
        "directly to understand and fix the mistake.",
        "This gives you a generate \u2192 verify \u2192 fix loop that converges "
        "on provably correct policy.",
    ]
    ny = 4.1
    for note in notes:
        text(s, I(M + 0.1), I(ny), I(U - 0.2), I(0.4),
             note, sz=13, color=MID_GRAY)
        ny += 0.5


# ═════════════════════════════════════════════════════════════════════════
#   SLIDE 6 — The Assurance Spectrum
# ═════════════════════════════════════════════════════════════════════════

def slide_assurance_spectrum(prs, BL):
    s = prs.slides.add_slide(BL)
    setbg(s, WHITE)

    text(s, I(M), I(0.3), I(U), I(0.5),
         "The Assurance Spectrum",
         sz=34, color=BLACK, bold=True)

    text(s, I(M), I(0.8), I(U), I(0.35),
         "What kinds of questions can we answer about authorization policies?",
         sz=17, color=MID_GRAY)

    shape(s, I(M), I(1.2), I(U), I(0.03), LIGHT_GRAY, None, "", rr=False)

    # Spectrum items — left label, right description, color coding
    spectrum = [
        ("Property Verification",
         "\u201CNo input can bypass this deny rule.\u201D\n"
         "Prove \u2200 input: rule(input) \u2260 bypass  (UNSAT = holds).",
         PURPLE_L, PURPLE_D),
        ("Equivalence / Regression",
         "\u201CDoes v2 behave identically to v1?\u201D\n"
         "Find a distinguishing input, or prove equivalence.",
         NAVY_L, NAVY),
        ("Subsumption",
         "\u201CIs the new policy at least as restrictive?\u201D\n"
         "Prove \u2200 input: allow(new) \u2192 allow(old).",
         TEAL_L, TEAL),
        ("MC/DC Coverage",
         "Each boolean sub-condition independently flips the outcome.\n"
         "DO-178C level coverage, auto-generated.",
         GREEN_L, GREEN_D),
        ("Test Generation",
         "Synthesize minimal inputs covering all decision paths.\n"
         "Z3 finds satisfying assignments per path constraint.",
         ORANGE_L, ORANGE_D),
        ("Gap Detection",
         "Find inputs where no rule matches (undefined result).\n"
         "Undefined = potential silent fail-open.",
         RED_L, RED_D),
    ]

    lw = 2.6     # label width
    dw = U - lw - 0.3  # description width
    bh = 0.8
    sy = 1.4

    for i, (label, desc, bg, accent) in enumerate(spectrum):
        y = sy + i * (bh + 0.15)
        shape(s, I(M), I(y), I(lw), I(bh),
              accent, accent, label, sz=14, color=WHITE, bold=True)
        shape(s, I(M + lw + 0.15), I(y), I(dw), I(bh),
              bg, accent, desc, sz=12, color=DARK_TEXT, align=PP_ALIGN.LEFT)

    # Footer
    text(s, I(M), I(sy + 6 * (bh + 0.15) + 0.1), I(U), I(0.35),
         "All six share a single Z3 encoding of the RVM bytecode.",
         sz=14, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════
#   SLIDE 7 — DEMO
# ═════════════════════════════════════════════════════════════════════════

def slide_demo(prs, BL):
    s = prs.slides.add_slide(BL)
    setbg(s, WHITE)

    text(s, I(M), I(0.5), I(U), I(0.8),
         "DEMO",
         sz=72, color=BLACK, bold=True, align=PP_ALIGN.CENTER)

    shape(s, I(SW/2 - 3.0), I(1.5), I(6.0), I(0.03),
          LIGHT_GRAY, None, "", rr=False)

    # ── Demo 1: RVM Playground (quick, smaller)
    d1y = 1.9
    shape(s, I(M), I(d1y), I(U), I(1.3), FAINT_GRAY, LIGHT_GRAY,
          "", sz=1, color=WHITE)

    text(s, I(M + 0.3), I(d1y + 0.15), I(U - 0.6), I(0.3),
         "\u2460  RVM Playground \u2014 Multi-Language Support",
         sz=20, color=NAVY, bold=True)

    text(s, I(M + 0.3), I(d1y + 0.5), I(U - 0.6), I(0.35),
         "Quick look: compile Rego, Cedar, and Azure Policy to the same "
         "RVM bytecode, step through execution",
         sz=15, color=MID_GRAY)

    hyperlink_text(s, I(M + 0.3), I(d1y + 0.9), I(U - 0.6), I(0.3),
                   "anakrish.github.io/rego-virtual-machine-playground",
                   "https://anakrish.github.io/rego-virtual-machine-playground/",
                   sz=14, color=AZURE)

    # ── Demo 2: RVM Policy Analysis (main showcase)
    d2y = 3.5
    shape(s, I(M), I(d2y), I(U), I(2.8), NAVY_L, NAVY,
          "", sz=1, color=WHITE)

    text(s, I(M + 0.3), I(d2y + 0.15), I(U - 0.6), I(0.35),
         "\u2461  RVM Policy Analysis",
         sz=24, color=NAVY, bold=True)

    text(s, I(M + 0.3), I(d2y + 0.55), I(U - 0.6), I(0.35),
         "Semantic diff, subsumption, test generation (MC/DC), property verification",
         sz=15, color=DARK_TEXT)

    text(s, I(M + 0.3), I(d2y + 0.95), I(U - 0.6), I(0.35),
         "~10K lines of Rust (RVM bytecode \u2192 Z3), entirely written with Copilot (Claude Opus). 3rd time lucky.",
         sz=14, color=PURPLE_D)

    hyperlink_text(s, I(M + 0.3), I(d2y + 1.4), I(U - 0.6), I(0.35),
                   "\u25B6  anakrish.github.io/rvm-policy-analysis",
                   "https://anakrish.github.io/rvm-policy-analysis/",
                   sz=18, color=AZURE, bold=True)


# ═════════════════════════════════════════════════════════════════════════
#   SLIDE 7 — Where We Are Today (Results)
# ═════════════════════════════════════════════════════════════════════════

def slide_results(prs, BL):
    s = prs.slides.add_slide(BL)
    setbg(s, WHITE)

    text(s, I(M), I(0.3), I(U), I(0.5),
         "Where We Are Today", sz=34, color=BLACK, bold=True)

    shape(s, I(M), I(0.85), I(U), I(0.03), LIGHT_GRAY, None, "", rr=False)

    # Tool results in a table-like layout
    text(s, I(M), I(1.05), I(U), I(0.35),
         "Implemented capabilities (regorus CLI)",
         sz=18, color=DARK_TEXT, bold=True)

    tools = [
        ("regorus diff",
         "Finds exact distinguishing input between two policy versions",
         "Input: two .rego files + schema  \u2192  Output: concrete input where "
         "results diverge, or \u201Cequivalent\u201D proof"),
        ("regorus subsumes",
         "Proves one policy is strictly at least as restrictive as another",
         "Used for Azure Policy: prove v2 definition is tighter than v1\n"
         "with effect precedence (deny > audit > allow)"),
        ("regorus gen-tests",
         "Z3-driven test generation with MC/DC condition coverage",
         "Generates minimal input set covering all decision paths;\n"
         "--condition-coverage flag enables full MC/DC"),
        ("regorus analyze",
         "Property verification: find input producing a target output",
         "UNSAT = property holds for all inputs (no bypass possible);\n"
         "SAT = counterexample returned as concrete JSON input"),
    ]

    tw1 = 2.5   # tool name
    tw2 = 4.5   # one-liner
    tw3 = U - tw1 - tw2 - 0.4  # detail
    ty = 1.5
    for (name, one_liner, detail) in tools:
        shape(s, I(M), I(ty), I(tw1), I(0.85),
              NAVY, NAVY, name, sz=14, color=WHITE, bold=True)
        shape(s, I(M + tw1 + 0.15), I(ty), I(tw2), I(0.85),
              NAVY_L, NAVY, one_liner, sz=12, color=DARK_TEXT,
              align=PP_ALIGN.LEFT)
        shape(s, I(M + tw1 + tw2 + 0.3), I(ty), I(tw3), I(0.85),
              FAINT_GRAY, LIGHT_GRAY, detail, sz=11, color=MID_GRAY,
              align=PP_ALIGN.LEFT)
        ty += 1.0

    # Concrete result highlight
    shape(s, I(M), I(5.65), I(U), I(0.7), GREEN_L, GREEN_D, "", rr=True)
    text(s, I(M + 0.15), I(5.7), I(U - 0.3), I(0.25),
         "Real Azure Policy result: group governance (139 lines, 5 decision paths)",
         sz=15, color=GREEN_D, bold=True)
    text(s, I(M + 0.15), I(5.98), I(U - 0.3), I(0.3),
         "23 auto-generated tests  \u00B7  100% line coverage  \u00B7  "
         "100% MC/DC coverage  \u00B7  Zero manual effort  \u00B7  "
         "Works across Rego and Azure Policy from the same bytecode",
         sz=13, color=DARK_TEXT)

    text(s, I(M), I(6.6), I(U), I(0.35),
         "Cedar policies also analyzable \u2014 same bytecode, same tooling.",
         sz=13, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════
#   SLIDE 8 — Open Challenges & Future Work
# ═════════════════════════════════════════════════════════════════════════

def slide_challenges(prs, BL):
    s = prs.slides.add_slide(BL)
    setbg(s, WHITE)

    text(s, I(M), I(0.3), I(U), I(0.5),
         "Road Ahead",
         sz=34, color=BLACK, bold=True)

    shape(s, I(M), I(0.85), I(U), I(0.03), LIGHT_GRAY, None, "", rr=False)

    # Context: AWS
    shape(s, I(M), I(1.05), I(U), I(0.9), ORANGE_L, ORANGE_D, "", rr=True)
    text(s, I(M + 0.15), I(1.1), I(U - 0.3), I(0.25),
         "AWS has ~8 years of production formal analysis",
         sz=16, color=ORANGE_D, bold=True)
    text(s, I(M + 0.15), I(1.38), I(U - 0.3), I(0.5),
         "Zelkova (2017) \u2192 IAM Access Analyzer \u2192 Cedar + Verified Permissions. "
         "~1B SMT queries/day. Named enterprise customers. \"Provable Security\" brand.",
         sz=12, color=DARK_TEXT)

    # The real differentiator
    shape(s, I(M), I(2.2), I(U), I(1.1), GREEN_L, GREEN_D, "", rr=True)
    text(s, I(M + 0.15), I(2.25), I(U - 0.3), I(0.25),
         "The hard problem: formal analysis of Rego",
         sz=16, color=GREEN_D, bold=True)
    text(s, I(M + 0.15), I(2.55), I(U - 0.3), I(0.7),
         "Cedar was purpose-built for analysis: no loops, no user-defined functions, no partial rules. "
         "Rego has all of those.\n"
         "Encoding Rego\u2019s expressiveness into Z3 (bounded loop unrolling, function inlining, "
         "partial rule aggregation) is a\n"
         "genuinely difficult translation problem. "
         "Rego/OPA remains the dominant policy language in the Kubernetes ecosystem.",
         sz=13, color=DARK_TEXT)

    # What needs to happen
    text(s, I(M), I(3.6), I(U), I(0.3),
         "Ideas for next steps",
         sz=18, color=NAVY, bold=True)

    items = [
        ("Productize",
         "Prototype (CLI, feature branch) \u2192 production service. "
         "Performance, reliability, documentation."),
        ("CI integration",
         "Formal diff on policy PRs. Equivalent \u2192 auto-approve. "
         "Wider access \u2192 review with counterexample."),
        ("AI verification loop",
         "LLM generates policy \u2192 Z3 verifies or returns counterexample \u2192 "
         "LLM revises \u2192 iterate. Exploratory."),
        ("Scalability",
         "Compositional analysis for large policy sets. "
         "Avoiding full enumeration over deep nesting and comprehensions."),
        ("Publication",
         "The RVM-to-Z3 translation and the multi-language bytecode design "
         "are candidates for a formal methods venue."),
        ("Verified engine",
         "Formal verification of the RVM implementation itself using Verus. "
         "Two layers of trust: verified policies on a verified engine."),
    ]
    iy = 4.0
    for (title, desc) in items:
        text(s, I(M + 0.1), I(iy), I(2.0), I(0.22),
             title, sz=14, color=NAVY, bold=True)
        text(s, I(M + 2.2), I(iy), I(U - 2.4), I(0.35),
             desc, sz=12, color=MID_GRAY)
        iy += 0.48

    text(s, I(M), I(6.9), I(U), I(0.35),
         "Questions and feedback welcome  \u2014  "
         "github.com/microsoft/regorus",
         sz=15, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════
#   MAIN
# ═════════════════════════════════════════════════════════════════════════

CONTENT_SLIDES = [
    slide_title,             # slide 1 (will be moved to front)
    # CTO memo (background) is slide 2 (preserved from SlidesToReuse)
    slide_rvm,               # slide 3
    slide_why_assurance,     # slide 4
    slide_demo,              # slide 5
    slide_challenges,        # slide 6
]


def main():
    # Open SlidesToReuse — has 3 slides: CTO memo, Renewed Urgency, Regorus Engine
    prs = Presentation(REUSE)
    prs.slide_width = I(SW)
    prs.slide_height = I(SH)

    # Delete Regorus Engine (index 2) and Renewed Urgency (index 1)
    # Keep only CTO Memo (index 0)
    delete_slide(prs, 2)
    delete_slide(prs, 1)
    # Now: [0] = CTO Memo

    # Find Blank layout
    BL = None
    for layout in prs.slide_layouts:
        if layout.name == "Blank":
            BL = layout
            break
    if BL is None:
        BL = prs.slide_layouts[6]

    # Add content slides (appended after CTO memo)
    for fn in CONTENT_SLIDES:
        fn(prs, BL)
    # Current: CTO(0), Title(1), RVM(2), Why(3), Demo(4), Challenges(5)
    # Desired: Title(0), CTO(1), RVM(2), Why(3), Demo(4), Challenges(5)

    # Move Title (index 1) to front
    move_slide_to_front(prs, 1)

    out = os.path.join(BASE, "RVM-Policy-Analysis-MSR-Workshop.pptx")
    prs.save(out)
    print(f"Generated: {out}")
    print(f"Total slides: {len(prs.slides)}")
    print()
    labels = [
        "Title",
        "CTO Memo \u2014 Background",
        "The RVM",
        "Why Formal Analysis?",
        "DEMO",
        "Road Ahead",
    ]
    for i, label in enumerate(labels, 1):
        print(f"  {i}. {label}")


if __name__ == "__main__":
    main()
