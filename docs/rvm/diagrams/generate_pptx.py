#!/usr/bin/env python3
"""Generate RVM Power of the Virtual Machine PowerPoint presentation.

Slide deck (10 slides):
  1. Title
  2. The Problem — Policy Fragmentation
  3. The JVM Analogy
  4. Strategic Ecosystem
  5. Why Bytecode Matters
  6. Trust Through Safety
  7. Playground Experience
  8. Competitive Landscape
  9. Developer Tooling & Policy Intelligence
 10. Closing
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

try:
    from pptx.enum.text import MSO_ANCHOR
except ImportError:
    MSO_ANCHOR = None

# ── Dimensions ───────────────────────────────────────────────────────────────
SW = 13.333
SH = 7.5
M = 0.6          # margin
U = SW - 2 * M   # usable width

# ── Palette — intentionally restrained ───────────────────────────────────────
CHARCOAL  = RGBColor(0x2D, 0x2D, 0x2D)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
MID_GRAY  = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
FAINT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

# Accent colours — used sparingly
GREEN_D   = RGBColor(0x1B, 0x5E, 0x20)
GREEN_M   = RGBColor(0x2E, 0x7D, 0x32)
GREEN_L   = RGBColor(0xC8, 0xE6, 0xC9)
GREEN_XL  = RGBColor(0xE8, 0xF5, 0xE9)
BLUE_D    = RGBColor(0x15, 0x65, 0xC0)
BLUE_L    = RGBColor(0xE3, 0xF2, 0xFD)
ORANGE_D  = RGBColor(0xEF, 0x6C, 0x00)
ORANGE_L  = RGBColor(0xFF, 0xF3, 0xE0)
RED_D     = RGBColor(0xC6, 0x28, 0x28)
RED_L     = RGBColor(0xFF, 0xCD, 0xD2)
PURPLE_D  = RGBColor(0x7B, 0x1F, 0xA2)
PURPLE_L  = RGBColor(0xF3, 0xE5, 0xF5)
TEAL_D    = RGBColor(0x00, 0x69, 0x5C)
TEAL_L    = RGBColor(0xB2, 0xDF, 0xDB)
AZURE     = RGBColor(0x00, 0x78, 0xD4)


# ── Helpers ──────────────────────────────────────────────────────────────────

def I(v):
    return Inches(v)


def dist(n, w, total=U, start=M):
    """Evenly space *n* boxes of width *w*."""
    if n <= 1:
        return [start + (total - w) / 2]
    gap = (total - n * w) / (n - 1)
    return [start + i * (w + gap) for i in range(n)]


def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def shape(slide, l, t, w, h, fill, border=None, txt="",
          sz=14, color=DARK_TEXT, bold=False,
          align=PP_ALIGN.CENTER, rr=True):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rr else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(st, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    tf = sh.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    if MSO_ANCHOR:
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass
    if txt:
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.color.rgb = color
        r.font.bold = bold
    return sh


def text(slide, l, t, w, h, txt, sz=14, color=DARK_TEXT,
         bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(sz)
    r.font.color.rgb = color
    r.font.bold = bold
    return tb


def bullets(slide, l, t, w, h, items, sz=14, color=DARK_TEXT,
            spacing=Pt(8), bold=False):
    """Plain-text bullet list using em-dash prefix."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        r = p.add_run()
        r.text = item
        r.font.size = Pt(sz)
        r.font.color.rgb = color
        r.font.bold = bold
    return tb


def line(slide, x1, y1, x2, y2, color=LIGHT_GRAY, width=Pt(1)):
    """Draw a thin connector line."""
    cn = slide.shapes.add_connector(
        1, I(x1), I(y1), I(x2), I(y2))  # straight connector
    cn.line.color.rgb = color
    cn.line.width = width
    return cn


# ── Build ────────────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = I(SW)
prs.slide_height = I(SH)
BL = prs.slide_layouts[6]


# ═══════════════════════════════════════════════════════════════════════════════
# 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, GREEN_D)

text(s, I(1.5), I(1.6), I(10.3), I(1.2),
     "The Power of RVM", sz=54, color=WHITE, bold=True,
     align=PP_ALIGN.CENTER)

text(s, I(1.5), I(3.2), I(10.3), I(0.8),
     "Regorus Virtual Machine", sz=36, color=GREEN_L,
     align=PP_ALIGN.CENTER)

text(s, I(1.5), I(4.5), I(10.3), I(0.6),
     "One Engine  \u00B7  Every Policy Language  \u00B7  Any Platform",
     sz=22, color=WHITE, align=PP_ALIGN.CENTER)

text(s, I(1.5), I(6.0), I(10.3), I(0.5),
     "microsoft/regorus  \u00B7  Open Source (MIT)",
     sz=16, color=GREEN_L, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# 2 — The Problem: Policy Fragmentation
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, WHITE)

text(s, I(M), I(0.3), I(U), I(0.7),
     "The Problem: Policy Fragmentation",
     sz=36, color=RED_D, bold=True)

text(s, I(M), I(1.05), I(U), I(0.5),
     "Every team builds its own policy engine. The cost multiplies.",
     sz=18, color=MID_GRAY)

# Engine boxes — plain, no emoji
engines = ["OPA / Rego", "Cedar", "Azure Policy", "RBAC", "Custom"]
bw = 2.1
for name, x in zip(engines, dist(5, bw)):
    shape(s, I(x), I(1.8), I(bw), I(1.0),
          RED_L, RED_D, name, sz=15, color=RED_D, bold=True)

# Cost items — simple text list, no boxes
costs = [
    "N separate security audits",
    "N separate testing frameworks",
    "N separate deployment pipelines",
    "N separate teams to maintain",
]
bullets(s, I(M + 0.2), I(3.2), I(5.5), I(2.5), costs,
        sz=17, color=DARK_TEXT, spacing=Pt(10))

# Formula — single accent bar
shape(s, I(M), I(5.5), I(U), I(0.65),
      FAINT_GRAY, LIGHT_GRAY,
      "N engines  =  N \u00D7 audits  +  N \u00D7 test suites  +  "
      "N \u00D7 pipelines  +  N \u00D7 teams",
      sz=18, color=RED_D, bold=True)

text(s, I(M), I(6.3), I(U), I(0.5),
     "Multiplicative cost, inconsistent behavior, fragmented tooling.",
     sz=17, color=DARK_TEXT, bold=True)


# ═══════════════════════════════════════════════════════════════════════════════
# 3 — The JVM Analogy
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, WHITE)

text(s, I(M), I(0.3), I(U), I(0.7),
     "What Java Did for Apps, RVM Does for Policy",
     sz=34, color=GREEN_D, bold=True)


def flow_row(slide, y, label, accent, boxes):
    text(slide, I(M), I(y - 0.45), I(5), I(0.4),
         label, sz=17, color=accent, bold=True)
    bw, gap = 2.4, 0.55
    total = 4 * bw + 3 * gap
    sx = (SW - total) / 2
    for i, (txt_val, bg_c) in enumerate(boxes):
        x = sx + i * (bw + gap)
        shape(slide, I(x), I(y), I(bw), I(0.95),
              bg_c, accent, txt_val, sz=14, color=DARK_TEXT)
        if i < 3:
            ax = x + bw + (gap - 0.3) / 2
            text(slide, I(ax), I(y + 0.2), I(0.3), I(0.5),
                 "\u2192", sz=24, color=accent, bold=True,
                 align=PP_ALIGN.CENTER)


flow_row(s, 1.5, "JVM  (Applications)", ORANGE_D, [
    ("Java \u00B7 Kotlin\nScala \u00B7 Groovy", ORANGE_L),
    ("Compilers", ORANGE_L),
    ("Java Bytecode", RGBColor(0xFF, 0xE0, 0xB2)),
    ("JVM: Any Platform", ORANGE_L),
])

flow_row(s, 3.3, "RVM  (Policy)", GREEN_D, [
    ("Rego \u00B7 Cedar\nAzure Policy \u00B7 RBAC", GREEN_XL),
    ("Compilers", GREEN_XL),
    ("RVM Bytecode", GREEN_L),
    ("RVM: Any Platform", GREEN_XL),
])

# Benefits as a clean numbered list — not boxes
text(s, I(M), I(4.8), I(U), I(0.4),
     "What this gives you", sz=20, color=BLUE_D, bold=True)

benefit_list = [
    "One security surface to audit",
    "One engine to optimize",
    "One test framework to maintain",
    "One artifact to deploy",
    "Any new language gets the entire tooling ecosystem for free",
]
bullets(s, I(M + 0.2), I(5.3), I(U), I(2.0), benefit_list,
        sz=15, color=DARK_TEXT)


# ═══════════════════════════════════════════════════════════════════════════════
# 4 — Strategic Ecosystem
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, WHITE)

text(s, I(M), I(0.3), I(U), I(0.7),
     "Strategic Ecosystem: RVM at the Center",
     sz=36, color=GREEN_D, bold=True)

# Left — sources
sx, sw = 0.5, 2.5
text(s, I(sx), I(1.3), I(sw), I(0.4),
     "Policy Sources", sz=15, color=BLUE_D,
     bold=True, align=PP_ALIGN.CENTER)
sources = ["Azure Policy", "Azure RBAC", "AKS Admission",
           "OPA / Rego", "Cedar"]
for i, src in enumerate(sources):
    shape(s, I(sx), I(1.85 + i * 0.68), I(sw), I(0.55),
          BLUE_L, BLUE_D, src, sz=14, color=DARK_TEXT)

# Centre — RVM
rx, rw = 4.3, 4.6
ry, rh = 2.0, 2.5
shape(s, I(rx), I(ry), I(rw), I(rh),
      GREEN_D, GREEN_D,
      "Regorus Virtual Machine\n\n"
      "One Secure, Audited,\nHigh-Performance Engine",
      sz=19, color=WHITE, bold=True)

# Right — targets
tx, tw = 10.2, 2.5
text(s, I(tx), I(1.3), I(tw), I(0.4),
     "Deploy Anywhere", sz=15, color=GREEN_M,
     bold=True, align=PP_ALIGN.CENTER)
targets = ["Cloud Scale", "Edge / IoT", "Browser (WASM)",
           "Confidential Compute"]
for i, tgt in enumerate(targets):
    shape(s, I(tx), I(1.85 + i * 0.68), I(tw), I(0.55),
          GREEN_XL, GREEN_M, tgt, sz=14, color=DARK_TEXT)

# Arrows
text(s, I(3.15), I(2.85), I(1.0), I(0.5),
     "\u2192", sz=26, color=BLUE_D, bold=True,
     align=PP_ALIGN.CENTER)
text(s, I(9.1), I(2.85), I(1.0), I(0.5),
     "\u2192", sz=26, color=GREEN_M, bold=True,
     align=PP_ALIGN.CENTER)
text(s, I(6.2), I(4.7), I(1.0), I(0.5),
     "\u2193", sz=26, color=PURPLE_D, bold=True,
     align=PP_ALIGN.CENTER)

# Language bindings — simple text
text(s, I(3.5), I(5.3), I(6.3), I(0.4),
     "Language Bindings", sz=15, color=PURPLE_D,
     bold=True, align=PP_ALIGN.CENTER)
shape(s, I(3.2), I(5.8), I(6.8), I(0.6),
      PURPLE_L, PURPLE_D,
      "Rust  \u00B7  C  \u00B7  C++  \u00B7  C#  \u00B7  Go  "
      "\u00B7  Java  \u00B7  Python  \u00B7  Ruby  \u00B7  JS",
      sz=15, color=DARK_TEXT)


# ═══════════════════════════════════════════════════════════════════════════════
# 5 — Why Bytecode Matters
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, WHITE)

text(s, I(M), I(0.3), I(U), I(0.7),
     "Why Bytecode Matters",
     sz=36, color=GREEN_D, bold=True)

text(s, I(M), I(1.0), I(U), I(0.4),
     "Compiling policy to bytecode unlocks capabilities "
     "you can\u2019t get from a tree-walking interpreter.",
     sz=17, color=MID_GRAY)

# 5 pillars — coloured header only, items as plain text beneath
pillars = [
    ("Performance", [
        "Much faster than tree-walking",
        "Register-based, minimal dispatch",
        "No cloning, no boxing",
    ], RED_D),
    ("Portable Artifacts", [
        "Compile once, run anywhere",
        "Ship pre-compiled bundles",
        "Forward-compatible serialization",
    ], GREEN_D),
    ("Tooling Ecosystem", [
        "Disassembly + source mapping",
        "Breakpoints & step-through",
        "Build once, all languages",
    ], AZURE),
    ("Suspendable Execution", [
        "HostAwait: pause mid-eval",
        "Fetch external data on demand",
        "Resume where you left off",
    ], PURPLE_D),
    ("Safety & Governance", [
        "Instruction budget for CPU",
        "Per-evaluation memory limits",
        "Deterministic, one engine to certify",
    ], TEAL_D),
]

pw = 2.3
for (title, items, accent), x in zip(pillars, dist(5, pw)):
    # Header bar
    shape(s, I(x), I(1.7), I(pw), I(0.55),
          accent, accent, title,
          sz=14, color=WHITE, bold=True)
    # Items as plain indented text
    bullets(s, I(x + 0.1), I(2.4), I(pw - 0.15), I(2.0),
            ["\u2013  " + it for it in items],
            sz=11, color=DARK_TEXT)

# Contrast bar — two halves
half = U / 2 - 0.1
shape(s, I(M), I(4.5), I(half), I(0.55),
      RED_L, None,
      "Without bytecode: interpreter \u00D7 tooling \u00D7 deploy per language",
      sz=12, color=RED_D, bold=True)
shape(s, I(M + half + 0.2), I(4.5), I(half), I(0.55),
      GREEN_L, None,
      "With RVM: compile any language \u2192 one artifact \u2192 "
      "tooling for free",
      sz=12, color=GREEN_D, bold=True)


# ═══════════════════════════════════════════════════════════════════════════════
# 6 — Trust Through Safety
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, WHITE)

text(s, I(M), I(0.3), I(U), I(0.7),
     "Trust Through Safety",
     sz=36, color=GREEN_D, bold=True)

# Three columns — header + bullet list (no coloured boxes per item)
col_w = 3.7
col_xs = dist(3, col_w)

safety = [
    ("Memory Safety", [
        "Written in Rust \u2014 no buffer overflows",
        "no_std compatible \u2014 bare metal ready",
        "No unsafe code in evaluation path",
    ]),
    ("Execution Safety", [
        "Configurable memory limits per evaluation",
        "Instruction budget \u2014 bounded CPU time",
        "Deterministic: same input = same output",
        "Sandboxed \u2014 no file or network I/O",
    ]),
    ("Auditability", [
        "Full disassembly with source mapping",
        "Breakpoints & step-through debugging",
        "Structured audit logs \u2014 every decision traceable",
        "Forward-compatible binary serialization",
        "Deterministic \u2014 enables formal verification",
        "Open source (MIT) \u2014 fully transparent",
    ]),
]

for (title, items), x in zip(safety, col_xs):
    shape(s, I(x), I(1.3), I(col_w), I(0.5),
          GREEN_M, GREEN_M, title,
          sz=16, color=WHITE, bold=True)
    bullets(s, I(x + 0.15), I(1.95), I(col_w - 0.2), I(3.0),
            ["\u2013  " + it for it in items],
            sz=12, color=DARK_TEXT, spacing=Pt(6))

# Compliance row
text(s, I(M), I(5.2), I(U), I(0.4),
     "Compliance Impact", sz=18, color=BLUE_D,
     bold=True, align=PP_ALIGN.CENTER)

comp = ["One engine to certify", "Reproducible results",
        "Resource & time bounded", "Open source (MIT)"]
bw = 2.7
for t, x in zip(comp, dist(4, bw)):
    shape(s, I(x), I(5.7), I(bw), I(0.75),
          BLUE_L, BLUE_D, t,
          sz=14, color=DARK_TEXT, bold=True)


# ═══════════════════════════════════════════════════════════════════════════════
# 7 — Playground Experience
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, WHITE)

text(s, I(M), I(0.3), I(U), I(0.7),
     "RVM Playground",
     sz=36, color=AZURE, bold=True)

text(s, I(M), I(1.0), I(U), I(0.4),
     "Cross-language policy comparison, entirely in your browser.",
     sz=17, color=MID_GRAY)

# 4 steps — header + body, arrows in between
steps = [
    ("WRITE", "Monaco editor\nSyntax highlighting\nExample library",
     BLUE_L, BLUE_D),
    ("COMPILE", "See bytecode\nInstruction count\nError reporting",
     ORANGE_L, ORANGE_D),
    ("RUN", "Custom input\nSet data\nInstant results",
     GREEN_L, GREEN_M),
    ("COMPARE", "View assembly\nStep through\nSide-by-side",
     PURPLE_L, PURPLE_D),
]

cw = 2.6
positions = dist(4, cw)
for i, ((title, desc, bg_c, accent), x) in enumerate(
        zip(steps, positions)):
    shape(s, I(x), I(1.7), I(cw), I(0.5),
          accent, accent, title,
          sz=16, color=WHITE, bold=True)
    shape(s, I(x), I(2.3), I(cw), I(1.5),
          bg_c, accent, desc,
          sz=14, color=DARK_TEXT)
    if i < 3:
        gap = positions[i + 1] - x - cw
        text(s, I(x + cw), I(2.7), I(gap), I(0.5),
             "\u2192", sz=22, color=MID_GRAY, bold=True,
             align=PP_ALIGN.CENTER)

# Powered by WASM — simple text, no orange box
text(s, I(M), I(4.4), I(U), I(0.5),
     "Powered by WebAssembly \u2014 the full Regorus engine "
     "compiled to WASM, no server required.",
     sz=16, color=DARK_TEXT, bold=True,
     align=PP_ALIGN.CENTER)

text(s, I(M), I(5.1), I(U), I(0.5),
     "https://anakrish.github.io/"
     "rego-virtual-machine-playground/",
     sz=16, color=AZURE, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# 8 — Competitive Landscape
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, WHITE)

text(s, I(M), I(0.3), I(U), I(0.7),
     "Competitive Landscape",
     sz=36, color=GREEN_D, bold=True)

# Left — competitors as titled sections with bullet lists
lx = 0.5

competitors = [
    ("OPA  (Go)", [
        "+ Mature ecosystem, large community",
        "+ Rego language",
        "\u2013 Go runtime overhead, GC pauses",
        "\u2013 Embeddable via cgo (complex, heavy)",
        "\u2013 Single language \u2014 no Cedar or Azure Policy",
    ]),
    ("Cedar SDK  (Rust)", [
        "+ Strong formal foundations",
        "+ Amazon Verified Permissions \u2014 managed cloud service",
        "\u2013 Single language \u2014 no Rego or Azure Policy",
        "\u2013 No bytecode \u2014 AST tree-walking evaluator",
        "\u2013 No mid-evaluation host callbacks",
    ]),
    ("Custom Engines", [
        "\u2013 One-off implementations per product",
        "\u2013 Per-product maintenance & security burden",
        "\u2013 No shared tooling or analysis",
    ]),
]

cy = 1.2
for title, items in competitors:
    text(s, I(lx), I(cy), I(5.0), I(0.35),
         title, sz=15, color=RED_D, bold=True)
    cy += 0.38
    for item in items:
        text(s, I(lx + 0.15), I(cy), I(5.2), I(0.25),
             item, sz=11, color=DARK_TEXT)
        cy += 0.26
    cy += 0.15

# Divider
text(s, I(5.8), I(3.3), I(0.8), I(0.5),
     "vs", sz=22, color=ORANGE_D, bold=True,
     align=PP_ALIGN.CENTER)

# Right — RVM
rx = 6.8
rw = 5.9

shape(s, I(rx), I(1.2), I(rw), I(0.5),
      GREEN_D, GREEN_D, "Regorus / RVM",
      sz=20, color=WHITE, bold=True)

rvm_feats = [
    "Multi-language: Rego + Cedar + Azure Policy + RBAC",
    "Compiled register-based bytecode with short-circuit optimization",
    "Rust core \u2014 memory safe, no_std, no GC, WASM compatible",
    "8 language bindings: C, C++, C#, Go, Java, Python, Ruby, JS",
    "HostAwait: suspendable execution for mid-eval host callbacks",
    "Breakpoints, stepping, debug \u2014 suspendable by design",
    "Portable binary artifacts \u2014 compile once, run anywhere",
    "Interactive WASM playground in browser",
    "Policy intelligence: formal analysis via Z3 (upcoming)",
    "Open source (MIT) \u2014 transparent and auditable",
]

fy = 1.85
for i, feat in enumerate(rvm_feats):
    bg_c = GREEN_XL if i % 2 == 0 else FAINT_GRAY
    shape(s, I(rx), I(fy + i * 0.44), I(rw), I(0.38),
          bg_c, None, feat,
          sz=11, color=DARK_TEXT, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════════════
# 9 — Developer Tooling & Policy Intelligence
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, WHITE)

text(s, I(M), I(0.3), I(U), I(0.7),
     "Developer Tooling & Policy Intelligence",
     sz=34, color=GREEN_D, bold=True)

text(s, I(M), I(0.95), I(U), I(0.4),
     "Tools built once work across every policy language.",
     sz=15, color=MID_GRAY)

# Left — languages
lx, lw = 0.4, 2.0
text(s, I(lx), I(1.5), I(lw), I(0.3),
     "Policy Languages", sz=13, color=BLUE_D,
     bold=True, align=PP_ALIGN.CENTER)

langs = ["Rego", "Cedar", "Azure Policy", "RBAC", "Future DSLs"]
for i, lang in enumerate(langs):
    shape(s, I(lx), I(1.9 + i * 0.5), I(lw), I(0.4),
          BLUE_L, BLUE_D, lang, sz=12, color=DARK_TEXT)

# Arrow in
text(s, I(2.5), I(2.9), I(0.5), I(0.4),
     "\u2192", sz=20, color=BLUE_D, bold=True,
     align=PP_ALIGN.CENTER)

# Centre — RVM
shape(s, I(3.1), I(1.7), I(2.4), I(2.7),
      GREEN_D, GREEN_D,
      "RVM\n\nUnified AST/IR\nCompiler\nRuntime",
      sz=14, color=WHITE, bold=True)

# Arrow out
text(s, I(5.6), I(2.9), I(0.5), I(0.4),
     "\u2192", sz=20, color=GREEN_M, bold=True,
     align=PP_ALIGN.CENTER)

# Right — 3 tooling columns, plain text
tcw = 1.7
tg = 0.15
tx0 = 6.2

tool_cols = [
    ("Editor Experience", PURPLE_D, PURPLE_L,
     ["LSP Server", "IntelliSense", "Go-to-Definition", "Autocomplete"]),
    ("Code Quality", ORANGE_D, ORANGE_L,
     ["Linter", "Formatter", "Static Analysis", "Dead Rule Finder"]),
    ("DevOps & CI", TEAL_D, TEAL_L,
     ["Test Framework", "Doc Generator", "CI/CD Gates", "Coverage"]),
]

for ci, (title, hc, ic, items) in enumerate(tool_cols):
    tx = tx0 + ci * (tcw + tg)
    shape(s, I(tx), I(1.7), I(tcw), I(0.4),
          hc, hc, title, sz=11, color=WHITE, bold=True)
    for j, tool in enumerate(items):
        shape(s, I(tx), I(2.18 + j * 0.42), I(tcw), I(0.36),
              ic, hc, tool, sz=10, color=DARK_TEXT)

# Policy Intelligence
pi_x, pi_w, pi_y = 6.2, 5.5, 4.2
shape(s, I(pi_x), I(pi_y), I(pi_w), I(0.45),
      RGBColor(0x4A, 0x14, 0x8C), RGBColor(0x4A, 0x14, 0x8C),
      "Policy Intelligence  \u2014  Formal Analysis "
      "(separate presentation)",
      sz=14, color=WHITE, bold=True)

intel = [
    "Z3-powered\nSatisfiability",
    "Conflict &\nRedundancy Detection",
    "Policy\nEquivalence Proofs",
    "Coverage &\nReachability",
]
iw = pi_w / 4 - 0.06
for i, t in enumerate(intel):
    ix = pi_x + i * (iw + 0.08)
    shape(s, I(ix), I(pi_y + 0.55), I(iw), I(0.75),
          PURPLE_L, RGBColor(0x4A, 0x14, 0x8C), t,
          sz=11, color=DARK_TEXT)

# Bottom note
text(s, I(M), I(6.2), I(U), I(0.5),
     "Without RVM: N tools \u00D7 M languages.  "
     "With RVM: build each tool once.",
     sz=15, color=DARK_TEXT, bold=True,
     align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# 10 — Closing
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, GREEN_D)

text(s, I(1.5), I(1.5), I(10.3), I(1.2),
     "One Engine. Every Policy. Any Platform.",
     sz=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

text(s, I(1.5), I(3.0), I(10.3), I(0.8),
     "Regorus Virtual Machine",
     sz=32, color=GREEN_L, align=PP_ALIGN.CENTER)

closing = [
    "Rego  +  Cedar  +  Azure Policy  +  RBAC",
    "High performance  \u00B7  Memory safe  \u00B7  Portable bytecode",
    "8 language bindings  \u00B7  WASM Playground  \u00B7  "
    "Open Source (MIT)",
]
for i, item in enumerate(closing):
    text(s, I(1.5), I(4.2 + i * 0.55), I(10.3), I(0.5),
         item, sz=20, color=WHITE, align=PP_ALIGN.CENTER)

text(s, I(1.5), I(5.9), I(10.3), I(0.5),
     "Try it: https://anakrish.github.io/"
     "rego-virtual-machine-playground/",
     sz=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

text(s, I(1.5), I(6.5), I(10.3), I(0.5),
     "github.com/microsoft/regorus",
     sz=17, color=GREEN_L, align=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "RVM-Power-Presentation.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
