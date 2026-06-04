#!/usr/bin/env python3
"""Generate RVM presentation — Version B: Bookend

Slide deck (5 slides):
  1. Title
  2. The Problem  (reframed, condensed)
  3. LIVE DEMO
  4. How It Works  (JVM analogy + key capabilities)
  5. Closing
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

try:
    from pptx.enum.text import MSO_ANCHOR
except ImportError:
    MSO_ANCHOR = None

# ── Dimensions ───────────────────────────────────────────────────────────────
SW = 13.333
SH = 7.5
M = 0.6
U = SW - 2 * M

# ── Palette ──────────────────────────────────────────────────────────────────
CHARCOAL   = RGBColor(0x2D, 0x2D, 0x2D)
DARK_TEXT  = RGBColor(0x33, 0x33, 0x33)
MID_GRAY   = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
FAINT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

GREEN_D  = RGBColor(0x1B, 0x5E, 0x20)
GREEN_M  = RGBColor(0x2E, 0x7D, 0x32)
GREEN_L  = RGBColor(0xC8, 0xE6, 0xC9)
GREEN_XL = RGBColor(0xE8, 0xF5, 0xE9)
BLUE_D   = RGBColor(0x15, 0x65, 0xC0)
BLUE_L   = RGBColor(0xE3, 0xF2, 0xFD)
ORANGE_D = RGBColor(0xEF, 0x6C, 0x00)
ORANGE_L = RGBColor(0xFF, 0xF3, 0xE0)
RED_D    = RGBColor(0xC6, 0x28, 0x28)
RED_L    = RGBColor(0xFF, 0xCD, 0xD2)
PURPLE_D = RGBColor(0x7B, 0x1F, 0xA2)
PURPLE_L = RGBColor(0xF3, 0xE5, 0xF5)
TEAL_D   = RGBColor(0x00, 0x69, 0x5C)
TEAL_L   = RGBColor(0xB2, 0xDF, 0xDB)
AZURE    = RGBColor(0x00, 0x78, 0xD4)


# ── Helpers ──────────────────────────────────────────────────────────────────

def I(v):
    return Inches(v)


def dist(n, w, total=U, start=M):
    if n <= 1:
        return [start + (total - w) / 2]
    gap = (total - n * w) / (n - 1)
    return [start + i * (w + gap) for i in range(n)]


def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def shape(slide, l, t, w, h, fill, border=None, txt="",
          sz=14, color=DARK_TEXT, bold=False,
          align=PP_ALIGN.CENTER, rr=True):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rr else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(st, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    tf = sh.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    if MSO_ANCHOR:
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass
    if txt:
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.color.rgb = color
        r.font.bold = bold
    return sh


def text(slide, l, t, w, h, txt, sz=14, color=DARK_TEXT,
         bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(sz)
    r.font.color.rgb = color
    r.font.bold = bold
    return tb


def bullets(slide, l, t, w, h, items, sz=14, color=DARK_TEXT,
            spacing=Pt(8), bold=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        r = p.add_run()
        r.text = item
        r.font.size = Pt(sz)
        r.font.color.rgb = color
        r.font.bold = bold
    return tb


# ── Build ────────────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = I(SW)
prs.slide_height = I(SH)
BL = prs.slide_layouts[6]


# ═══════════════════════════════════════════════════════════════════════════════
# 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, GREEN_D)

text(s, I(1.5), I(1.6), I(10.3), I(1.2),
     "The Power of RVM", sz=54, color=WHITE, bold=True,
     align=PP_ALIGN.CENTER)
text(s, I(1.5), I(3.2), I(10.3), I(0.8),
     "Regorus Virtual Machine", sz=36, color=GREEN_L,
     align=PP_ALIGN.CENTER)
text(s, I(1.5), I(4.5), I(10.3), I(0.6),
     "One Engine  \u00B7  Every Policy Language  \u00B7  Any Platform",
     sz=22, color=WHITE, align=PP_ALIGN.CENTER)
text(s, I(1.5), I(6.0), I(10.3), I(0.5),
     "microsoft/regorus  \u00B7  Open Source (MIT)",
     sz=16, color=GREEN_L, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# 2 — The Problem  (condensed, punchy)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, WHITE)

text(s, I(M), I(0.3), I(U), I(0.7),
     "The Problem",
     sz=40, color=RED_D, bold=True)

text(s, I(M), I(1.1), I(U), I(0.6),
     "Azure Policy must support multiple policy languages.\n"
     "Across Microsoft, teams maintain separate engines.",
     sz=20, color=MID_GRAY)

# Engine boxes — 5 engines
engines = ["OPA / Rego", "Cedar", "Azure Policy", "AKS Admission", "Custom"]
bw = 2.1
for name, x in zip(engines, dist(5, bw)):
    shape(s, I(x), I(2.1), I(bw), I(1.0),
          RED_L, RED_D, name, sz=15, color=RED_D, bold=True)

# Cost formula
shape(s, I(M), I(3.5), I(U), I(0.65),
      FAINT_GRAY, LIGHT_GRAY,
      "N engines  =  N \u00D7 audits  +  N \u00D7 test suites  +  "
      "N \u00D7 pipelines  +  N \u00D7 teams",
      sz=18, color=RED_D, bold=True)

# Single question
text(s, I(M), I(4.7), I(U), I(0.8),
     "What if one engine could run them all?",
     sz=28, color=GREEN_D, bold=True, align=PP_ALIGN.CENTER)

text(s, I(M), I(5.6), I(U), I(0.6),
     "Let me show you.",
     sz=22, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# 3 — LIVE DEMO
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, CHARCOAL)

text(s, I(1.5), I(2.0), I(10.3), I(1.5),
     "LIVE DEMO", sz=72, color=WHITE, bold=True,
     align=PP_ALIGN.CENTER)
text(s, I(1.5), I(3.8), I(10.3), I(0.6),
     "RVM Playground \u2014 Cross-language policy in your browser",
     sz=22, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
text(s, I(1.5), I(5.0), I(10.3), I(0.5),
     "https://anakrish.github.io/rego-virtual-machine-playground/",
     sz=16, color=AZURE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# 4 — How It Works  (JVM analogy + key capabilities)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, WHITE)

text(s, I(M), I(0.25), I(U), I(0.6),
     "How It Works",
     sz=36, color=GREEN_D, bold=True)

text(s, I(M), I(0.85), I(U), I(0.4),
     "What Java did for applications, RVM does for policy.",
     sz=17, color=MID_GRAY)

# Compact JVM/RVM flow
def flow_row(slide, y, label, accent, boxes):
    text(slide, I(M), I(y - 0.3), I(3), I(0.3),
         label, sz=14, color=accent, bold=True)
    bw_f, gap = 2.4, 0.4
    total_flow = len(boxes) * bw_f + (len(boxes) - 1) * gap
    sx = (SW - total_flow) / 2
    for i, (txt_val, bg_c) in enumerate(boxes):
        x = sx + i * (bw_f + gap)
        shape(slide, I(x), I(y), I(bw_f), I(0.75),
              bg_c, accent, txt_val, sz=12, color=DARK_TEXT)
        if i < len(boxes) - 1:
            ax = x + bw_f + (gap - 0.25) / 2
            text(slide, I(ax), I(y + 0.15), I(0.25), I(0.4),
                 "\u2192", sz=18, color=accent, bold=True,
                 align=PP_ALIGN.CENTER)

flow_row(s, 1.5, "JVM", ORANGE_D, [
    ("Java \u00B7 Kotlin\nScala \u00B7 Groovy", ORANGE_L),
    ("Compilers", ORANGE_L),
    ("Java Bytecode", RGBColor(0xFF, 0xE0, 0xB2)),
    ("JVM: Any Platform", ORANGE_L),
])

flow_row(s, 2.6, "RVM", GREEN_D, [
    ("Rego \u00B7 Cedar\nAzure Policy \u00B7 more", GREEN_XL),
    ("Compilers", GREEN_XL),
    ("RVM Bytecode", GREEN_L),
    ("RVM: Any Platform", GREEN_XL),
])

# Key implications — 3 columns
text(s, I(M), I(3.7), I(U), I(0.4),
     "What this means in practice", sz=18, color=BLUE_D, bold=True)

impl_cols = [
    ("Security & Compliance", [
        "One engine to audit and certify",
        "Memory safe \u2014 Rust, no_std",
        "Sandboxed \u2014 no file or network I/O",
        "Structured audit logs",
    ], TEAL_D),
    ("Performance & Portability", [
        "Compiled register-based bytecode",
        "Compile once, run anywhere",
        "Cloud, edge, WASM, bare metal",
        "8 language bindings",
    ], GREEN_D),
    ("Unique Capabilities", [
        "HostAwait: pause mid-eval for host data",
        "Instruction budget + memory limits",
        "Breakpoints & step-through debugging",
        "Forward-compatible binary artifacts",
    ], PURPLE_D),
]

cw = 3.7
for (title, items, accent), x in zip(impl_cols, dist(3, cw)):
    shape(s, I(x), I(4.15), I(cw), I(0.45),
          accent, accent, title,
          sz=14, color=WHITE, bold=True)
    bullets(s, I(x + 0.1), I(4.7), I(cw - 0.15), I(2.2),
            ["\u2013  " + it for it in items],
            sz=12, color=DARK_TEXT, spacing=Pt(5))


# ═══════════════════════════════════════════════════════════════════════════════
# 5 — Closing
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, GREEN_D)

text(s, I(1.5), I(1.5), I(10.3), I(1.2),
     "One Engine. Every Policy. Any Platform.",
     sz=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
text(s, I(1.5), I(3.0), I(10.3), I(0.8),
     "Regorus Virtual Machine",
     sz=32, color=GREEN_L, align=PP_ALIGN.CENTER)

closing = [
    "Rego  +  Cedar  +  Azure Policy",
    "High performance  \u00B7  Memory safe  \u00B7  Portable bytecode",
    "8 language bindings  \u00B7  WASM Playground  \u00B7  Open Source (MIT)",
]
for i, item in enumerate(closing):
    text(s, I(1.5), I(4.2 + i * 0.55), I(10.3), I(0.5),
         item, sz=20, color=WHITE, align=PP_ALIGN.CENTER)

text(s, I(1.5), I(5.9), I(10.3), I(0.5),
     "Try it: https://anakrish.github.io/"
     "rego-virtual-machine-playground/",
     sz=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
text(s, I(1.5), I(6.5), I(10.3), I(0.5),
     "github.com/microsoft/regorus",
     sz=17, color=GREEN_L, align=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "RVM-Presentation-B-Bookend.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
