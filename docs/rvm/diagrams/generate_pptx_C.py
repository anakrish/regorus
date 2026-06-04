#!/usr/bin/env python3
"""Generate RVM presentation — Version C: Three-Act

Slide deck (7 slides):
  1. Title
  2. The Problem  (Azure Policy needs multiple languages)
  3. The Answer   (JVM analogy)
  4. LIVE DEMO
  5. Under the Hood  (bytecode + safety combined)
  6. Competitive Landscape
  7. Closing
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

try:
    from pptx.enum.text import MSO_ANCHOR
except ImportError:
    MSO_ANCHOR = None

# ── Dimensions ───────────────────────────────────────────────────────────────
SW = 13.333
SH = 7.5
M = 0.6
U = SW - 2 * M

# ── Palette ──────────────────────────────────────────────────────────────────
CHARCOAL   = RGBColor(0x2D, 0x2D, 0x2D)
DARK_TEXT  = RGBColor(0x33, 0x33, 0x33)
MID_GRAY   = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
FAINT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

GREEN_D  = RGBColor(0x1B, 0x5E, 0x20)
GREEN_M  = RGBColor(0x2E, 0x7D, 0x32)
GREEN_L  = RGBColor(0xC8, 0xE6, 0xC9)
GREEN_XL = RGBColor(0xE8, 0xF5, 0xE9)
BLUE_D   = RGBColor(0x15, 0x65, 0xC0)
BLUE_L   = RGBColor(0xE3, 0xF2, 0xFD)
ORANGE_D = RGBColor(0xEF, 0x6C, 0x00)
ORANGE_L = RGBColor(0xFF, 0xF3, 0xE0)
RED_D    = RGBColor(0xC6, 0x28, 0x28)
RED_L    = RGBColor(0xFF, 0xCD, 0xD2)
PURPLE_D = RGBColor(0x7B, 0x1F, 0xA2)
PURPLE_L = RGBColor(0xF3, 0xE5, 0xF5)
TEAL_D   = RGBColor(0x00, 0x69, 0x5C)
TEAL_L   = RGBColor(0xB2, 0xDF, 0xDB)
AZURE    = RGBColor(0x00, 0x78, 0xD4)


# ── Helpers ──────────────────────────────────────────────────────────────────

def I(v):
    return Inches(v)


def dist(n, w, total=U, start=M):
    if n <= 1:
        return [start + (total - w) / 2]
    gap = (total - n * w) / (n - 1)
    return [start + i * (w + gap) for i in range(n)]


def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def shape(slide, l, t, w, h, fill, border=None, txt="",
          sz=14, color=DARK_TEXT, bold=False,
          align=PP_ALIGN.CENTER, rr=True):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rr else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(st, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    tf = sh.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    if MSO_ANCHOR:
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass
    if txt:
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.color.rgb = color
        r.font.bold = bold
    return sh


def text(slide, l, t, w, h, txt, sz=14, color=DARK_TEXT,
         bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(sz)
    r.font.color.rgb = color
    r.font.bold = bold
    return tb


def bullets(slide, l, t, w, h, items, sz=14, color=DARK_TEXT,
            spacing=Pt(8), bold=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        r = p.add_run()
        r.text = item
        r.font.size = Pt(sz)
        r.font.color.rgb = color
        r.font.bold = bold
    return tb


# ── Build ────────────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = I(SW)
prs.slide_height = I(SH)
BL = prs.slide_layouts[6]


# ═══════════════════════════════════════════════════════════════════════════════
# 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, GREEN_D)

text(s, I(1.5), I(1.6), I(10.3), I(1.2),
     "The Power of RVM", sz=54, color=WHITE, bold=True,
     align=PP_ALIGN.CENTER)
text(s, I(1.5), I(3.2), I(10.3), I(0.8),
     "Regorus Virtual Machine", sz=36, color=GREEN_L,
     align=PP_ALIGN.CENTER)
text(s, I(1.5), I(4.5), I(10.3), I(0.6),
     "One Engine  \u00B7  Every Policy Language  \u00B7  Any Platform",
     sz=22, color=WHITE, align=PP_ALIGN.CENTER)
text(s, I(1.5), I(6.0), I(10.3), I(0.5),
     "microsoft/regorus  \u00B7  Open Source (MIT)",
     sz=16, color=GREEN_L, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# 2 — The Problem: Policy Engine Proliferation
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, WHITE)

text(s, I(M), I(0.3), I(U), I(0.7),
     "The Problem: Policy Engine Proliferation",
     sz=36, color=RED_D, bold=True)

text(s, I(M), I(1.05), I(U), I(0.5),
     "Azure Policy must support multiple policy languages. "
     "Across Microsoft, teams maintain separate engines.",
     sz=18, color=MID_GRAY)

engines = ["OPA / Rego", "Cedar", "Azure Policy", "AKS Admission", "Custom"]
bw = 2.1
for name, x in zip(engines, dist(5, bw)):
    shape(s, I(x), I(1.8), I(bw), I(1.0),
          RED_L, RED_D, name, sz=15, color=RED_D, bold=True)

costs = [
    "N separate security audits",
    "N separate testing frameworks",
    "N separate deployment pipelines",
    "N separate teams to maintain",
]
bullets(s, I(M + 0.2), I(3.2), I(5.5), I(2.5), costs,
        sz=17, color=DARK_TEXT, spacing=Pt(10))

shape(s, I(M), I(5.5), I(U), I(0.65),
      FAINT_GRAY, LIGHT_GRAY,
      "N engines  =  N \u00D7 audits  +  N \u00D7 test suites  +  "
      "N \u00D7 pipelines  +  N \u00D7 teams",
      sz=18, color=RED_D, bold=True)

text(s, I(M), I(6.3), I(U), I(0.5),
     "Multiplicative cost, inconsistent behavior, fragmented tooling.",
     sz=17, color=DARK_TEXT, bold=True)


# ═══════════════════════════════════════════════════════════════════════════════
# 3 — The Answer: JVM Analogy
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, WHITE)

text(s, I(M), I(0.3), I(U), I(0.7),
     "What Java Did for Apps, RVM Does for Policy",
     sz=34, color=GREEN_D, bold=True)

def flow_row(slide, y, label, accent, boxes):
    text(slide, I(M), I(y - 0.45), I(5), I(0.4),
         label, sz=17, color=accent, bold=True)
    bw_f, gap = 2.4, 0.55
    total = 4 * bw_f + 3 * gap
    sx = (SW - total) / 2
    for i, (txt_val, bg_c) in enumerate(boxes):
        x = sx + i * (bw_f + gap)
        shape(slide, I(x), I(y), I(bw_f), I(0.95),
              bg_c, accent, txt_val, sz=14, color=DARK_TEXT)
        if i < 3:
            ax = x + bw_f + (gap - 0.3) / 2
            text(slide, I(ax), I(y + 0.2), I(0.3), I(0.5),
                 "\u2192", sz=24, color=accent, bold=True,
                 align=PP_ALIGN.CENTER)

flow_row(s, 1.5, "JVM  (Applications)", ORANGE_D, [
    ("Java \u00B7 Kotlin\nScala \u00B7 Groovy", ORANGE_L),
    ("Compilers", ORANGE_L),
    ("Java Bytecode", RGBColor(0xFF, 0xE0, 0xB2)),
    ("JVM: Any Platform", ORANGE_L),
])

flow_row(s, 3.3, "RVM  (Policy)", GREEN_D, [
    ("Rego \u00B7 Cedar\nAzure Policy \u00B7 more", GREEN_XL),
    ("Compilers", GREEN_XL),
    ("RVM Bytecode", GREEN_L),
    ("RVM: Any Platform", GREEN_XL),
])

text(s, I(M), I(4.8), I(U), I(0.4),
     "What this gives you", sz=20, color=BLUE_D, bold=True)

benefit_list = [
    "One security surface to audit",
    "One engine to optimize",
    "One test framework to maintain",
    "One artifact to deploy",
    "Any new language gets the entire tooling ecosystem for free",
]
bullets(s, I(M + 0.2), I(5.3), I(U), I(2.0), benefit_list,
        sz=15, color=DARK_TEXT)


# ═══════════════════════════════════════════════════════════════════════════════
# 4 — LIVE DEMO
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, CHARCOAL)

text(s, I(1.5), I(2.0), I(10.3), I(1.5),
     "LIVE DEMO", sz=72, color=WHITE, bold=True,
     align=PP_ALIGN.CENTER)
text(s, I(1.5), I(3.8), I(10.3), I(0.6),
     "RVM Playground \u2014 Cross-language policy in your browser",
     sz=22, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
text(s, I(1.5), I(5.0), I(10.3), I(0.5),
     "https://anakrish.github.io/rego-virtual-machine-playground/",
     sz=16, color=AZURE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# 5 — Under the Hood  (bytecode + safety, one slide)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, WHITE)

text(s, I(M), I(0.25), I(U), I(0.6),
     "Under the Hood",
     sz=36, color=GREEN_D, bold=True)

text(s, I(M), I(0.85), I(U), I(0.4),
     "What makes the playground possible \u2014 and what it means "
     "for production.",
     sz=16, color=MID_GRAY)

# Top row: 3 bytecode pillars
bc_cols = [
    ("Performance", [
        "Register-based bytecode, minimal dispatch",
        "Short-circuit optimization (PolicyOp)",
        "No cloning, no boxing",
    ], GREEN_D),
    ("Portability", [
        "Compile once, run anywhere",
        "Cloud, edge, WASM, bare metal",
        "Forward-compatible binary serialization",
    ], BLUE_D),
    ("Extensibility", [
        "HostAwait: pause mid-eval for host data",
        "Suspendable execution by design",
        "Breakpoints and step-through debugging",
    ], PURPLE_D),
]

cw = 3.7
for (title, items, accent), x in zip(bc_cols, dist(3, cw)):
    shape(s, I(x), I(1.4), I(cw), I(0.45),
          accent, accent, title,
          sz=14, color=WHITE, bold=True)
    bullets(s, I(x + 0.1), I(1.95), I(cw - 0.15), I(1.6),
            ["\u2013  " + it for it in items],
            sz=11, color=DARK_TEXT, spacing=Pt(5))

# Bottom row: safety + compliance
text(s, I(M), I(3.7), I(U), I(0.4),
     "Trust Through Safety", sz=20, color=TEAL_D, bold=True)

safety_cols = [
    ("Memory Safety", [
        "Written in Rust \u2014 no buffer overflows",
        "no_std compatible \u2014 bare metal ready",
        "No unsafe code in evaluation path",
    ]),
    ("Execution Safety", [
        "Configurable memory limits per evaluation",
        "Instruction budget \u2014 bounded CPU",
        "Sandboxed \u2014 no file or network I/O",
    ]),
    ("Auditability", [
        "Structured audit logs \u2014 every decision traceable",
        "Deterministic: same input = same output",
        "Open source (MIT) \u2014 fully transparent",
    ]),
]

for (title, items), x in zip(safety_cols, dist(3, cw)):
    shape(s, I(x), I(4.15), I(cw), I(0.4),
          TEAL_D, TEAL_D, title,
          sz=13, color=WHITE, bold=True)
    bullets(s, I(x + 0.1), I(4.65), I(cw - 0.15), I(1.5),
            ["\u2013  " + it for it in items],
            sz=11, color=DARK_TEXT, spacing=Pt(4))

# Compliance bar
shape(s, I(M), I(6.0), I(U), I(0.55),
      FAINT_GRAY, LIGHT_GRAY,
      "One engine to certify  \u00B7  Reproducible results  "
      "\u00B7  Resource & time bounded  \u00B7  Open source (MIT)",
      sz=14, color=DARK_TEXT, bold=True)


# ═══════════════════════════════════════════════════════════════════════════════
# 6 — Competitive Landscape
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, WHITE)

text(s, I(M), I(0.3), I(U), I(0.7),
     "Competitive Landscape",
     sz=36, color=GREEN_D, bold=True)

lx = 0.5
competitors = [
    ("OPA  (Go)", [
        "+ Mature ecosystem, large community",
        "+ Rego language",
        "\u2013 Go runtime overhead, GC pauses",
        "\u2013 Embeddable via cgo (complex, heavy)",
        "\u2013 Single language \u2014 no Cedar or Azure Policy",
    ]),
    ("Cedar SDK  (Rust)", [
        "+ Strong formal foundations",
        "+ Amazon Verified Permissions \u2014 managed cloud service",
        "\u2013 Single language \u2014 no Rego or Azure Policy",
        "\u2013 No bytecode \u2014 AST tree-walking evaluator",
        "\u2013 No mid-evaluation host callbacks",
    ]),
    ("Custom Engines", [
        "\u2013 One-off implementations per product",
        "\u2013 Per-product maintenance & security burden",
        "\u2013 No shared tooling or analysis",
    ]),
]

cy = 1.2
for title, items in competitors:
    text(s, I(lx), I(cy), I(5.0), I(0.35),
         title, sz=15, color=RED_D, bold=True)
    cy += 0.38
    for item in items:
        text(s, I(lx + 0.15), I(cy), I(5.2), I(0.25),
             item, sz=11, color=DARK_TEXT)
        cy += 0.26
    cy += 0.15

text(s, I(5.8), I(3.3), I(0.8), I(0.5),
     "vs", sz=22, color=ORANGE_D, bold=True,
     align=PP_ALIGN.CENTER)

rx, rw = 6.8, 5.9
shape(s, I(rx), I(1.2), I(rw), I(0.5),
      GREEN_D, GREEN_D, "Regorus / RVM",
      sz=20, color=WHITE, bold=True)

rvm_feats = [
    "Multi-language: Rego + Cedar + Azure Policy + more",
    "Compiled register-based bytecode with short-circuit optimization",
    "Rust core \u2014 memory safe, no_std, no GC, WASM compatible",
    "8 language bindings: C, C++, C#, Go, Java, Python, Ruby, JS",
    "HostAwait: suspendable execution for mid-eval host callbacks",
    "Breakpoints, stepping, debug \u2014 suspendable by design",
    "Portable binary artifacts \u2014 compile once, run anywhere",
    "Interactive WASM playground in browser",
    "Policy intelligence: formal analysis via Z3 (upcoming)",
    "Open source (MIT) \u2014 transparent and auditable",
]

fy = 1.85
for i, feat in enumerate(rvm_feats):
    bg_c = GREEN_XL if i % 2 == 0 else FAINT_GRAY
    shape(s, I(rx), I(fy + i * 0.44), I(rw), I(0.38),
          bg_c, None, feat,
          sz=11, color=DARK_TEXT, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════════════
# 7 — Closing
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, GREEN_D)

text(s, I(1.5), I(1.5), I(10.3), I(1.2),
     "One Engine. Every Policy. Any Platform.",
     sz=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
text(s, I(1.5), I(3.0), I(10.3), I(0.8),
     "Regorus Virtual Machine",
     sz=32, color=GREEN_L, align=PP_ALIGN.CENTER)

closing = [
    "Rego  +  Cedar  +  Azure Policy",
    "High performance  \u00B7  Memory safe  \u00B7  Portable bytecode",
    "8 language bindings  \u00B7  WASM Playground  \u00B7  Open Source (MIT)",
]
for i, item in enumerate(closing):
    text(s, I(1.5), I(4.2 + i * 0.55), I(10.3), I(0.5),
         item, sz=20, color=WHITE, align=PP_ALIGN.CENTER)

text(s, I(1.5), I(5.9), I(10.3), I(0.5),
     "Try it: https://anakrish.github.io/"
     "rego-virtual-machine-playground/",
     sz=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
text(s, I(1.5), I(6.5), I(10.3), I(0.5),
     "github.com/microsoft/regorus",
     sz=17, color=GREEN_L, align=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "RVM-Presentation-C-ThreeAct.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
