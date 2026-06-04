#!/usr/bin/env python3
"""Generate RVM presentation — Version D: Progressive Demo

Slide deck (4 slides):
  1. Title
  2. One Problem  (minimal, single powerful statement)
  3. LIVE DEMO   (extended, with narration cues on slide)
  4. Closing
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

try:
    from pptx.enum.text import MSO_ANCHOR
except ImportError:
    MSO_ANCHOR = None

# ── Dimensions ───────────────────────────────────────────────────────────────
SW = 13.333
SH = 7.5
M = 0.6
U = SW - 2 * M

# ── Palette ──────────────────────────────────────────────────────────────────
CHARCOAL   = RGBColor(0x2D, 0x2D, 0x2D)
DARK_TEXT  = RGBColor(0x33, 0x33, 0x33)
MID_GRAY   = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
FAINT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

GREEN_D  = RGBColor(0x1B, 0x5E, 0x20)
GREEN_M  = RGBColor(0x2E, 0x7D, 0x32)
GREEN_L  = RGBColor(0xC8, 0xE6, 0xC9)
GREEN_XL = RGBColor(0xE8, 0xF5, 0xE9)
BLUE_D   = RGBColor(0x15, 0x65, 0xC0)
BLUE_L   = RGBColor(0xE3, 0xF2, 0xFD)
ORANGE_D = RGBColor(0xEF, 0x6C, 0x00)
ORANGE_L = RGBColor(0xFF, 0xF3, 0xE0)
RED_D    = RGBColor(0xC6, 0x28, 0x28)
RED_L    = RGBColor(0xFF, 0xCD, 0xD2)
PURPLE_D = RGBColor(0x7B, 0x1F, 0xA2)
PURPLE_L = RGBColor(0xF3, 0xE5, 0xF5)
TEAL_D   = RGBColor(0x00, 0x69, 0x5C)
TEAL_L   = RGBColor(0xB2, 0xDF, 0xDB)
AZURE    = RGBColor(0x00, 0x78, 0xD4)


# ── Helpers ──────────────────────────────────────────────────────────────────

def I(v):
    return Inches(v)


def dist(n, w, total=U, start=M):
    if n <= 1:
        return [start + (total - w) / 2]
    gap = (total - n * w) / (n - 1)
    return [start + i * (w + gap) for i in range(n)]


def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def shape(slide, l, t, w, h, fill, border=None, txt="",
          sz=14, color=DARK_TEXT, bold=False,
          align=PP_ALIGN.CENTER, rr=True):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rr else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(st, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    tf = sh.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    if MSO_ANCHOR:
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass
    if txt:
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.color.rgb = color
        r.font.bold = bold
    return sh


def text(slide, l, t, w, h, txt, sz=14, color=DARK_TEXT,
         bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(sz)
    r.font.color.rgb = color
    r.font.bold = bold
    return tb


def bullets(slide, l, t, w, h, items, sz=14, color=DARK_TEXT,
            spacing=Pt(8), bold=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        r = p.add_run()
        r.text = item
        r.font.size = Pt(sz)
        r.font.color.rgb = color
        r.font.bold = bold
    return tb


# ── Build ────────────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = I(SW)
prs.slide_height = I(SH)
BL = prs.slide_layouts[6]


# ═══════════════════════════════════════════════════════════════════════════════
# 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, GREEN_D)

text(s, I(1.5), I(1.6), I(10.3), I(1.2),
     "The Power of RVM", sz=54, color=WHITE, bold=True,
     align=PP_ALIGN.CENTER)
text(s, I(1.5), I(3.2), I(10.3), I(0.8),
     "Regorus Virtual Machine", sz=36, color=GREEN_L,
     align=PP_ALIGN.CENTER)
text(s, I(1.5), I(4.5), I(10.3), I(0.6),
     "One Engine  \u00B7  Every Policy Language  \u00B7  Any Platform",
     sz=22, color=WHITE, align=PP_ALIGN.CENTER)
text(s, I(1.5), I(6.0), I(10.3), I(0.5),
     "microsoft/regorus  \u00B7  Open Source (MIT)",
     sz=16, color=GREEN_L, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# 2 — One Problem  (minimal, powerful)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, WHITE)

text(s, I(1.5), I(1.5), I(10.3), I(1.0),
     "Azure Policy must support\nmultiple policy languages.",
     sz=36, color=RED_D, bold=True, align=PP_ALIGN.CENTER)

text(s, I(1.5), I(3.0), I(10.3), I(0.8),
     "Across Microsoft, teams maintain separate engines.\n"
     "Every engine means another audit, another test suite, "
     "another deployment.",
     sz=20, color=MID_GRAY, align=PP_ALIGN.CENTER)

# Minimal engine boxes
engines = ["OPA / Rego", "Cedar", "Azure Policy", "AKS Admission", "Custom"]
bw = 2.1
for name, x in zip(engines, dist(5, bw)):
    shape(s, I(x), I(4.2), I(bw), I(0.8),
          RED_L, RED_D, name, sz=14, color=RED_D, bold=True)

text(s, I(1.5), I(5.5), I(10.3), I(0.8),
     "What if one engine could run them all?",
     sz=30, color=GREEN_D, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# 3 — LIVE DEMO  (extended, with narration cues)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, CHARCOAL)

text(s, I(1.5), I(0.8), I(10.3), I(1.2),
     "LIVE DEMO", sz=72, color=WHITE, bold=True,
     align=PP_ALIGN.CENTER)
text(s, I(1.5), I(2.2), I(10.3), I(0.5),
     "RVM Playground \u2014 Cross-language policy in your browser",
     sz=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Narration cues — what to show during the demo
cues = [
    ("1. Write", "Rego policy in Monaco editor with syntax highlighting"),
    ("2. Compile", "See RVM bytecode and instruction count"),
    ("3. Run", "Set input data, get instant evaluation results"),
    ("4. Switch", "Same policy logic in Cedar \u2014 same bytecode target"),
    ("5. Debug", "Step through bytecode, inspect registers"),
    ("6. Compare", "Side-by-side assembly for Rego vs Cedar"),
]

cw = 5.5
cx = (SW - cw) / 2
cy = 3.2
for label, desc in cues:
    text(s, I(cx), I(cy), I(1.2), I(0.35),
         label, sz=14, color=GREEN_L, bold=True)
    text(s, I(cx + 1.2), I(cy), I(cw - 1.2), I(0.35),
         desc, sz=13, color=LIGHT_GRAY)
    cy += 0.42

text(s, I(1.5), I(6.2), I(10.3), I(0.5),
     "https://anakrish.github.io/rego-virtual-machine-playground/",
     sz=14, color=AZURE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# 4 — Closing
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, GREEN_D)

text(s, I(1.5), I(1.2), I(10.3), I(1.0),
     "One Engine. Every Policy. Any Platform.",
     sz=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
text(s, I(1.5), I(2.5), I(10.3), I(0.8),
     "Regorus Virtual Machine",
     sz=32, color=GREEN_L, align=PP_ALIGN.CENTER)

# More detailed closing since there are fewer slides
closing = [
    "Rego  +  Cedar  +  Azure Policy \u2014 one bytecode target",
    "High performance  \u00B7  Memory safe  \u00B7  Portable",
    "8 language bindings: C, C++, C#, Go, Java, Python, Ruby, JS",
    "WASM Playground  \u00B7  Debugger  \u00B7  Disassembler",
    "Open Source (MIT)  \u00B7  github.com/microsoft/regorus",
]
for i, item in enumerate(closing):
    text(s, I(1.5), I(3.8 + i * 0.5), I(10.3), I(0.5),
         item, sz=18, color=WHITE, align=PP_ALIGN.CENTER)

text(s, I(1.5), I(6.3), I(10.3), I(0.5),
     "Try it: https://anakrish.github.io/"
     "rego-virtual-machine-playground/",
     sz=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "RVM-Presentation-D-ProgressiveDemo.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
