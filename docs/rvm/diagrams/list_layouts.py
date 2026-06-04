#!/usr/bin/env python3
"""List slide layouts in SlidesToReuse.pptx"""
from pptx import Presentation
import os
base = os.path.dirname(os.path.abspath(__file__))
prs = Presentation(os.path.join(base, "SlidesToReuse.pptx"))
for i, layout in enumerate(prs.slide_layouts):
    phs = [p.placeholder_format.idx for p in layout.placeholders]
    name = layout.name
    print(f"  [{i}] {name}  placeholders={phs}")
