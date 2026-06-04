#!/usr/bin/env python3
"""Read and dump contents of SlidesToReuse.pptx"""
from pptx import Presentation
import os

path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                    "SlidesToReuse.pptx")
prs = Presentation(path)
print(f"Slide dimensions: {prs.slide_width/914400:.3f}\" x {prs.slide_height/914400:.3f}\"")
print(f"Total slides: {len(prs.slides)}\n")

for si, slide in enumerate(prs.slides, 1):
    print("=" * 80)
    print(f"SLIDE {si}")
    bg = slide.background
    if bg.fill.type is not None:
        try:
            print(f"  Background: {bg.fill.fore_color.rgb}")
        except Exception:
            print("  Background: (gradient/pattern)")
    print(f"  Shapes: {len(slide.shapes)}\n")

    for sh in slide.shapes:
        stype = sh.shape_type
        print(f"  [{sh.name}]  type={stype}  "
              f"pos=({sh.left/914400:.2f}\",{sh.top/914400:.2f}\")  "
              f"size=({sh.width/914400:.2f}\"x{sh.height/914400:.2f}\")")

        # Fill
        if hasattr(sh, "fill"):
            try:
                if sh.fill.type is not None:
                    print(f"    fill={sh.fill.fore_color.rgb}")
            except Exception:
                pass

        # Line
        if hasattr(sh, "line"):
            try:
                if sh.line.color and sh.line.color.rgb:
                    print(f"    border={sh.line.color.rgb}")
            except Exception:
                pass

        # Text
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                if not para.text.strip():
                    continue
                for run in para.runs:
                    f = run.font
                    parts = []
                    if f.size:
                        parts.append(f"sz={f.size.pt:.0f}")
                    if f.bold:
                        parts.append("bold")
                    if f.italic:
                        parts.append("italic")
                    try:
                        if f.color and f.color.rgb:
                            parts.append(f"color={f.color.rgb}")
                    except Exception:
                        pass
                    if f.name:
                        parts.append(f"font={f.name}")
                    det = ", ".join(parts)
                    txt = run.text.replace("\n", "\\n")
                    print(f"    \"{txt}\"  [{det}]")

        # Image
        if stype == 13:
            print(f"    image: {sh.image.content_type}")

        # Table
        if sh.has_table:
            tbl = sh.table
            print(f"    table: {len(tbl.rows)}r x {len(tbl.columns)}c")
            for ri, row in enumerate(tbl.rows):
                for ci, cell in enumerate(row.cells):
                    t = cell.text.strip()
                    if t:
                        print(f"      [{ri},{ci}] {t[:100]}")

        print()
