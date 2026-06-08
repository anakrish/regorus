#!/usr/bin/env python3
"""
Generate RegorusValue.pptx by rendering mermaid diagrams to high-res PNG
and embedding them in slides with natural (aspect-ratio-preserving) scaling.
"""

import os
import re
import subprocess
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from PIL import Image

BLUE_DARK = RGBColor(0x1A, 0x52, 0x76)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Available area for the image (leave room for title at top)
MARGIN_L = Inches(0.3)
MARGIN_T = Inches(0.75)
MAX_IMG_W = Inches(12.7)
MAX_IMG_H = Inches(6.5)

SOURCES = [
    ("value-compare/src/baseline_visual.md", "Baseline: Arc + BTreeMap \u2014 24 bytes"),
    ("value-compare/src/v1/visual.md",       "v1: HashMap + SmolStr \u2014 24 bytes"),
    ("value-compare/src/v2/visual.md",       "v2: Arc<str> keys \u2014 24 bytes"),
    ("value-compare/src/v3/visual.md",       "v3: HashSet + cached hash \u2014 24 bytes"),
    ("value-compare/src/v4/visual.md",       "v4: ArcStr thin pointer \u2014 16 bytes"),
    ("value-compare/src/v5/visual.md",       "v5: NaN-boxed Value \u2014 8 bytes"),
    ("value-compare/src/v6/visual.md",       "v6: Portable tagged pointer \u2014 8 bytes"),
    ("value-compare/src/v7/visual.md",       "v7: Tagged pointer + schema \u2014 8 bytes"),
    ("value-compare/src/v8/visual.md",       "v8: Flattened Number + schema \u2014 16 bytes"),
    ("value-compare/src/v9/visual.md",       "v9: Arena-allocated, Copy \u2014 16 bytes"),
]


def extract_mermaid(md_path: str) -> str:
    text = Path(md_path).read_text()
    m = re.search(r'```mermaid\s*\n(.*?)```', text, re.DOTALL)
    if not m:
        raise ValueError(f"No mermaid block found in {md_path}")
    return m.group(1).strip()


def render_png(mermaid_code: str, out_png: str, scale: int = 4):
    with tempfile.NamedTemporaryFile(mode='w', suffix='.mmd', delete=False) as f:
        f.write(mermaid_code)
        mmd_path = f.name
    try:
        result = subprocess.run(
            ["mmdc", "-i", mmd_path, "-o", out_png, "-b", "white",
             "-t", "default", "-s", str(scale)],
            capture_output=True, text=True, timeout=60
        )
        if result.returncode != 0:
            print(f"  mmdc stderr: {result.stderr[:500]}")
            raise RuntimeError(f"mmdc failed for {out_png}")
    finally:
        os.unlink(mmd_path)


def fit_image(img_w_px, img_h_px, max_w_emu, max_h_emu):
    """Scale image to fit within max bounds while preserving aspect ratio.
    Returns (width_emu, height_emu)."""
    aspect = img_w_px / img_h_px
    max_aspect = max_w_emu / max_h_emu

    if aspect > max_aspect:
        # Image is wider relative to box -> constrain by width
        w = max_w_emu
        h = int(w / aspect)
    else:
        # Image is taller relative to box -> constrain by height
        h = max_h_emu
        w = int(h * aspect)
    return w, h


def add_title(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12), Inches(0.45))
    tf = tb.text_frame
    tf.word_wrap = True
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(22)
    run.font.color.rgb = BLUE_DARK
    run.font.bold = True


def main():
    workspace = Path(__file__).parent
    png_dir = workspace / "svg_diagrams"
    png_dir.mkdir(exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for md_rel, title in SOURCES:
        md_path = workspace / md_rel
        png_path = png_dir / (md_rel.replace("/", "_").replace(".md", ".png"))

        print(f"Processing {md_rel}...")

        # Extract and render
        mermaid_code = extract_mermaid(str(md_path))
        render_png(mermaid_code, str(png_path))
        print(f"  -> {png_path.name}")

        # Get actual pixel dimensions
        im = Image.open(str(png_path))
        img_w_px, img_h_px = im.size
        im.close()

        # Fit to available area preserving aspect ratio
        w_emu, h_emu = fit_image(img_w_px, img_h_px, int(MAX_IMG_W), int(MAX_IMG_H))

        # Center horizontally, place below title
        left = int(MARGIN_L) + (int(MAX_IMG_W) - w_emu) // 2
        top = int(MARGIN_T) + (int(MAX_IMG_H) - h_emu) // 2

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        add_title(slide, title)
        slide.shapes.add_picture(str(png_path), left, top, w_emu, h_emu)

        print(f"  {img_w_px}x{img_h_px}px -> {w_emu/914400:.1f}x{h_emu/914400:.1f}in")

    out = workspace / "RegorusValue.pptx"
    prs.save(str(out))
    print(f"\nSaved {out} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
