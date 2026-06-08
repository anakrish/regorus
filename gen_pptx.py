#!/usr/bin/env python3
"""Generate RegorusValue.pptx with editable diagram shapes for baseline + v1-v9."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ── Blue Academic palette ──
BLUE_BG     = RGBColor(0xEB, 0xF5, 0xFB)  # group fill
BLUE_INLINE = RGBColor(0xD6, 0xEA, 0xF8)  # inline variant fill
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLUE_DARK   = RGBColor(0x1A, 0x52, 0x76)  # text
BLUE_BORDER = RGBColor(0x24, 0x71, 0xA3)  # stroke
NOTE_BG     = RGBColor(0xFE, 0xF9, 0xE7)  # note fill (yellow)
NOTE_BORDER = RGBColor(0xF1, 0xC4, 0x0F)
NOTE_TEXT   = RGBColor(0x7D, 0x66, 0x08)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

def make_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

# ── Helper functions ──

def add_box(slide, left, top, width, height, text, fill=WHITE, border=BLUE_BORDER,
            text_color=BLUE_DARK, font_size=10, bold=False, italic=False,
            border_width=Pt(1), rounded=False, dash=False):
    """Add a rectangle with text. Returns the shape."""
    if rounded:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = 0.1
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = border_width
    if dash:
        shape.line.dash_style = 4  # dash
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = text_color
    run.font.bold = bold
    run.font.italic = italic
    return shape

def add_group_box(slide, left, top, width, height, title, border=BLUE_BORDER,
                  fill=BLUE_BG, border_width=Pt(2)):
    """Add a group/subgraph rectangle with a title at top."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = border_width
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    tf.paragraphs[0].space_before = Pt(2)
    tf.paragraphs[0].space_after = Pt(0)
    run = tf.paragraphs[0].add_run()
    run.text = title
    run.font.size = Pt(11)
    run.font.color.rgb = BLUE_DARK
    run.font.bold = True
    # anchor text to top
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    return shape

def add_arrow(slide, start_shape, end_shape, label="", dashed=False):
    """Add a connector/arrow between two shapes (approximated with a line + optional label)."""
    # Get centers
    s = start_shape
    e = end_shape
    sx = s.left + s.width
    sy = s.top + s.height // 2
    ex = e.left
    ey = e.top + e.height // 2
    
    connector = slide.shapes.add_connector(1, sx, sy, ex, ey)  # straight
    connector.line.color.rgb = BLUE_BORDER
    connector.line.width = Pt(1.5)
    if dashed:
        connector.line.dash_style = 4

    # Add arrowhead at end
    connector.end_x = ex
    connector.end_y = ey
    
    if label:
        # Add a small text box at midpoint
        mx = (sx + ex) // 2 - Inches(0.3)
        my = min(sy, ey) - Inches(0.2)
        tb = slide.shapes.add_textbox(mx, my, Inches(0.6), Inches(0.2))
        tf = tb.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = label
        run.font.size = Pt(8)
        run.font.color.rgb = BLUE_BORDER
        run.font.bold = True

def add_note_box(slide, left, top, width, height, lines):
    """Add a dashed yellow note box."""
    shape = add_box(slide, left, top, width, height, "",
                    fill=NOTE_BG, border=NOTE_BORDER, text_color=NOTE_TEXT,
                    font_size=9, dash=True, border_width=Pt(1.5))
    tf = shape.text_frame
    tf.paragraphs[0].text = ""
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        run.font.size = Pt(9)
        run.font.color.rgb = NOTE_TEXT
        if i == 0:
            run.font.bold = True
    return shape

def add_title(slide, text, subtitle=""):
    """Add slide title."""
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
    tf = tb.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(24)
    run.font.color.rgb = BLUE_DARK
    run.font.bold = True
    if subtitle:
        p = tf.add_paragraph()
        run2 = p.add_run()
        run2.text = subtitle
        run2.font.size = Pt(14)
        run2.font.color.rgb = BLUE_BORDER

# ── Constants for layout ──
V_LEFT = Inches(0.4)
V_TOP = Inches(1.0)
BOX_W = Inches(1.8)
BOX_H = Inches(0.35)
BOX_GAP = Inches(0.08)
HEAP_LEFT = Inches(4.8)

def variant_boxes(slide, variants, left, top, box_w=BOX_W, box_h=BOX_H, gap=BOX_GAP):
    """Draw a vertical stack of variant boxes. Returns list of shapes."""
    shapes = []
    y = top
    for text, fill in variants:
        s = add_box(slide, left, y, box_w, box_h, text, fill=fill,
                    font_size=9, rounded=(fill == BLUE_INLINE))
        shapes.append(s)
        y += box_h + gap
    return shapes

def heap_boxes(slide, items, left, top, box_w=Inches(2.2), box_h=BOX_H, gap=BOX_GAP):
    """Draw heap target boxes. Returns dict name->shape."""
    shapes = {}
    y = top
    for name, text in items:
        s = add_box(slide, left, y, box_w, box_h, text, fill=WHITE, font_size=9)
        shapes[name] = s
        y += box_h + gap
    return shapes


# ═══════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════

def slide_baseline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_title(slide, "Baseline: Arc + BTreeMap", "enum Value — 24 bytes")
    
    # Value group box
    vg = add_group_box(slide, V_LEFT, V_TOP, Inches(2.6), Inches(3.2), "enum Value — 24 bytes")
    
    variants = [
        ("object — 8B thin ptr", WHITE),
        ("string — 16B fat ptr", WHITE),
        ("array — 8B thin ptr", WHITE),
        ("set — 8B thin ptr", WHITE),
        ("number — 16B nested enum", WHITE),
        ("null · bool · undefined — ≤ 1B", BLUE_INLINE),
    ]
    vshapes = variant_boxes(slide, variants, V_LEFT + Inches(0.3), V_TOP + Inches(0.4), box_w=Inches(2.0))
    
    # Heap targets
    hl = HEAP_LEFT
    htop = V_TOP + Inches(0.2)
    heap = heap_boxes(slide, [
        ("obj", "BTreeMap‹Value, Value› (ordered)"),
        ("str", "str"),
        ("arr", "Vec‹Value›"),
        ("set", "BTreeSet‹Value› (ordered)"),
    ], hl, htop)
    
    # Number group
    ng = add_group_box(slide, hl, htop + Inches(1.8), Inches(2.2), Inches(1.2), "enum Number — 16 bytes")
    n_inline = add_box(slide, hl + Inches(0.1), htop + Inches(2.1), Inches(2.0), BOX_H,
                       "u64 · i64 · f64", fill=BLUE_INLINE, font_size=9, rounded=True)
    n_bigint = add_box(slide, hl + Inches(0.1), htop + Inches(2.55), Inches(2.0), BOX_H,
                       "bigint", fill=WHITE, font_size=9)
    
    bigint_target = add_box(slide, hl + Inches(3.0), htop + Inches(2.55), Inches(1.2), BOX_H,
                            "BigInt", fill=WHITE, font_size=9)
    
    # Arrows
    for v, h, lbl in [(vshapes[0], heap["obj"], "Arc"),
                       (vshapes[1], heap["str"], "Arc"),
                       (vshapes[2], heap["arr"], "Arc"),
                       (vshapes[3], heap["set"], "Arc")]:
        add_arrow(slide, v, h, lbl)
    add_arrow(slide, n_bigint, bigint_target, "Arc")
    
    # Note
    add_note_box(slide, Inches(0.4), Inches(4.5), Inches(3.5), Inches(0.7),
                 ["24B = 8B tag + 16B largest payload",
                  "String (fat ptr) and Number tie at 16B"])

def slide_v1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "v1: HashMap + SmolStr", "enum Value — 24 bytes")
    
    vg = add_group_box(slide, V_LEFT, V_TOP, Inches(2.6), Inches(3.2), "enum Value — 24 bytes")
    variants = [
        ("object — 8B thin ptr", WHITE),
        ("string — 16B fat ptr", WHITE),
        ("array — 8B thin ptr", WHITE),
        ("set — 8B thin ptr", WHITE),
        ("number — 16B nested enum", WHITE),
        ("null · bool · undefined — ≤ 1B", BLUE_INLINE),
    ]
    vshapes = variant_boxes(slide, variants, V_LEFT + Inches(0.3), V_TOP + Inches(0.4), box_w=Inches(2.0))
    
    hl = HEAP_LEFT
    htop = V_TOP
    
    # ObjectMap group
    omg = add_group_box(slide, hl, htop, Inches(2.6), Inches(1.5),
                        "enum ObjectMap (unordered)")
    sk_box = add_group_box(slide, hl + Inches(0.1), htop + Inches(0.4), Inches(2.4), Inches(0.45),
                           "StringKeyed", fill=BLUE_BG, border_width=Pt(1))
    add_box(slide, hl + Inches(0.2), htop + Inches(0.6), Inches(2.2), BOX_H,
            "HashMap‹SmolStr, Value›", fill=WHITE, font_size=9)
    mx_box = add_group_box(slide, hl + Inches(0.1), htop + Inches(0.95), Inches(2.4), Inches(0.45),
                           "Mixed", fill=BLUE_BG, border_width=Pt(1))
    add_box(slide, hl + Inches(0.2), htop + Inches(1.15), Inches(2.2), BOX_H,
            "HashMap‹Value, Value›", fill=WHITE, font_size=9)
    
    heap = heap_boxes(slide, [
        ("str", "str"),
        ("arr", "Vec‹Value›"),
        ("set", "BTreeSet‹Value› (ordered)"),
    ], hl, htop + Inches(1.7))
    
    # Number
    ng = add_group_box(slide, hl, htop + Inches(3.0), Inches(2.2), Inches(1.2), "enum Number — 16 bytes")
    add_box(slide, hl + Inches(0.1), htop + Inches(3.3), Inches(2.0), BOX_H,
            "u64 · i64 · f64", fill=BLUE_INLINE, font_size=9, rounded=True)
    add_box(slide, hl + Inches(0.1), htop + Inches(3.75), Inches(2.0), BOX_H,
            "bigint → Arc BigInt", fill=WHITE, font_size=9)
    
    # SmolStr
    smg = add_group_box(slide, Inches(8.0), htop, Inches(2.4), Inches(1.2), "SmolStr — 24 bytes")
    add_box(slide, Inches(8.1), htop + Inches(0.4), Inches(2.2), BOX_H,
            "≤ 23 chars: inline", fill=BLUE_INLINE, font_size=9, rounded=True)
    add_box(slide, Inches(8.1), htop + Inches(0.85), Inches(2.2), BOX_H,
            "> 23 chars → Arc str", fill=WHITE, font_size=9)
    
    add_note_box(slide, Inches(0.4), Inches(4.5), Inches(3.5), Inches(0.7),
                 ["24B = 8B tag + 16B largest payload",
                  "Arc‹str› fat pointer drives the size"])

def slide_v2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "v2: Arc‹str› keys (no SmolStr)", "enum Value — 24 bytes")
    
    vg = add_group_box(slide, V_LEFT, V_TOP, Inches(2.6), Inches(3.2), "enum Value — 24 bytes")
    variants = [
        ("object — 8B thin ptr", WHITE),
        ("string — 16B fat ptr", WHITE),
        ("array — 8B thin ptr", WHITE),
        ("set — 8B thin ptr", WHITE),
        ("number — 16B nested enum", WHITE),
        ("null · bool · undefined — ≤ 1B", BLUE_INLINE),
    ]
    vshapes = variant_boxes(slide, variants, V_LEFT + Inches(0.3), V_TOP + Inches(0.4), box_w=Inches(2.0))
    
    hl = HEAP_LEFT
    htop = V_TOP
    
    # ObjectMap
    omg = add_group_box(slide, hl, htop, Inches(2.6), Inches(1.5), "enum ObjectMap (unordered)")
    add_group_box(slide, hl + Inches(0.1), htop + Inches(0.4), Inches(2.4), Inches(0.45),
                  "StringKeyed", fill=BLUE_BG, border_width=Pt(1))
    add_box(slide, hl + Inches(0.2), htop + Inches(0.6), Inches(2.2), BOX_H,
            "HashMap‹Arc‹str›, Value›", fill=WHITE, font_size=9)
    add_group_box(slide, hl + Inches(0.1), htop + Inches(0.95), Inches(2.4), Inches(0.45),
                  "Mixed", fill=BLUE_BG, border_width=Pt(1))
    add_box(slide, hl + Inches(0.2), htop + Inches(1.15), Inches(2.2), BOX_H,
            "HashMap‹Value, Value›", fill=WHITE, font_size=9)
    
    heap = heap_boxes(slide, [
        ("str", "str"),
        ("arr", "Vec‹Value›"),
        ("set", "BTreeSet‹Value› (ordered)"),
    ], hl, htop + Inches(1.7))
    
    # Number
    ng = add_group_box(slide, hl, htop + Inches(3.0), Inches(2.2), Inches(1.2), "enum Number — 16 bytes")
    add_box(slide, hl + Inches(0.1), htop + Inches(3.3), Inches(2.0), BOX_H,
            "u64 · i64 · f64", fill=BLUE_INLINE, font_size=9, rounded=True)
    add_box(slide, hl + Inches(0.1), htop + Inches(3.75), Inches(2.0), BOX_H,
            "bigint → Arc BigInt", fill=WHITE, font_size=9)
    
    # Interner
    ig = add_group_box(slide, Inches(8.0), htop, Inches(2.8), Inches(0.8), "Interner")
    add_box(slide, Inches(8.1), htop + Inches(0.4), Inches(2.6), BOX_H,
            "HashMap‹u64, Arc‹str››", fill=WHITE, font_size=9)
    
    add_note_box(slide, Inches(0.4), Inches(4.5), Inches(3.5), Inches(0.7),
                 ["24B = 8B tag + 16B largest payload",
                  "Arc‹str› fat pointer drives the size"])

def slide_v3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "v3: HashSet + cached hash", "enum Value — 24 bytes")
    
    vg = add_group_box(slide, V_LEFT, V_TOP, Inches(2.6), Inches(3.2), "enum Value — 24 bytes")
    variants = [
        ("object — 8B thin ptr", WHITE),
        ("string — 16B fat ptr", WHITE),
        ("array — 8B thin ptr", WHITE),
        ("set — 8B thin ptr", WHITE),
        ("number — 16B nested enum", WHITE),
        ("null · bool · undefined — ≤ 1B", BLUE_INLINE),
    ]
    variant_boxes(slide, variants, V_LEFT + Inches(0.3), V_TOP + Inches(0.4), box_w=Inches(2.0))
    
    hl = HEAP_LEFT
    htop = V_TOP
    
    # ObjectMap struct
    omg = add_group_box(slide, hl, htop, Inches(2.8), Inches(1.8), "struct ObjectMap (unordered)")
    add_box(slide, hl + Inches(0.1), htop + Inches(0.4), Inches(2.6), BOX_H,
            "strings: HashMap‹Arc‹str›, Value›", fill=WHITE, font_size=9)
    add_box(slide, hl + Inches(0.1), htop + Inches(0.85), Inches(2.6), BOX_H,
            "other: Option‹HashMap‹Value, Value››", fill=WHITE, font_size=9)
    add_box(slide, hl + Inches(0.1), htop + Inches(1.3), Inches(2.6), BOX_H,
            "cached_hash: u64", fill=BLUE_INLINE, font_size=9, rounded=True)
    
    heap = heap_boxes(slide, [
        ("str", "str"),
        ("arr", "Vec‹Value›"),
        ("set", "HashSet‹Value› (unordered)"),
    ], hl, htop + Inches(2.0))
    
    # Number
    ng = add_group_box(slide, hl, htop + Inches(3.3), Inches(2.2), Inches(1.2), "enum Number — 16 bytes")
    add_box(slide, hl + Inches(0.1), htop + Inches(3.6), Inches(2.0), BOX_H,
            "u64 · i64 · f64", fill=BLUE_INLINE, font_size=9, rounded=True)
    add_box(slide, hl + Inches(0.1), htop + Inches(4.05), Inches(2.0), BOX_H,
            "bigint → Arc BigInt", fill=WHITE, font_size=9)
    
    # Interner
    ig = add_group_box(slide, Inches(8.2), htop, Inches(2.8), Inches(0.8), "Interner (reuses v2)")
    add_box(slide, Inches(8.3), htop + Inches(0.4), Inches(2.6), BOX_H,
            "HashMap‹u64, Arc‹str››", fill=WHITE, font_size=9)
    
    add_note_box(slide, Inches(0.4), Inches(4.5), Inches(3.5), Inches(0.7),
                 ["24B = 8B tag + 16B largest payload",
                  "Arc‹str› fat pointer drives the size"])

def slide_v4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "v4: ArcStr thin pointer", "enum Value — 16 bytes")
    
    vg = add_group_box(slide, V_LEFT, V_TOP, Inches(2.6), Inches(3.2), "enum Value — 16 bytes")
    variants = [
        ("object — 8B thin ptr", WHITE),
        ("string: ArcStr — 8B thin ptr", WHITE),
        ("array — 8B thin ptr", WHITE),
        ("set — 8B thin ptr", WHITE),
        ("number — 8B (flattened)", WHITE),
        ("null · bool · undefined — ≤ 1B", BLUE_INLINE),
    ]
    variant_boxes(slide, variants, V_LEFT + Inches(0.3), V_TOP + Inches(0.4), box_w=Inches(2.0))
    
    hl = HEAP_LEFT
    htop = V_TOP
    
    # ObjectMap
    omg = add_group_box(slide, hl, htop, Inches(2.8), Inches(1.8), "struct ObjectMap (unordered)")
    add_box(slide, hl + Inches(0.1), htop + Inches(0.4), Inches(2.6), BOX_H,
            "strings: HashMap‹ArcStr, Value›", fill=WHITE, font_size=9)
    add_box(slide, hl + Inches(0.1), htop + Inches(0.85), Inches(2.6), BOX_H,
            "other: Option‹HashMap‹Value, Value››", fill=WHITE, font_size=9)
    add_box(slide, hl + Inches(0.1), htop + Inches(1.3), Inches(2.6), BOX_H,
            "cached_hash: u64", fill=BLUE_INLINE, font_size=9, rounded=True)
    
    heap = heap_boxes(slide, [
        ("arr", "Vec‹Value›"),
        ("set", "HashSet‹Value› (unordered)"),
    ], hl, htop + Inches(2.0))
    
    # Number
    ng = add_group_box(slide, hl, htop + Inches(3.0), Inches(2.2), Inches(1.2), "enum Number — 16 bytes")
    add_box(slide, hl + Inches(0.1), htop + Inches(3.3), Inches(2.0), BOX_H,
            "u64 · i64 · f64", fill=BLUE_INLINE, font_size=9, rounded=True)
    add_box(slide, hl + Inches(0.1), htop + Inches(3.75), Inches(2.0), BOX_H,
            "bigint → Arc BigInt", fill=WHITE, font_size=9)
    
    # ArcStr detail
    asg = add_group_box(slide, Inches(8.2), htop, Inches(2.4), Inches(0.8), "ArcStr — 8 bytes")
    add_box(slide, Inches(8.3), htop + Inches(0.4), Inches(2.2), BOX_H,
            "thin pointer to heap string", fill=BLUE_INLINE, font_size=9, rounded=True)
    
    # Interner
    ig = add_group_box(slide, Inches(8.2), htop + Inches(1.0), Inches(2.6), Inches(0.8), "Interner")
    add_box(slide, Inches(8.3), htop + Inches(1.4), Inches(2.4), BOX_H,
            "HashSet‹ArcStr› (thread-local)", fill=WHITE, font_size=9)
    
    add_note_box(slide, Inches(0.4), Inches(4.5), Inches(4.0), Inches(0.8),
                 ["16B = 8B tag + 8B payload",
                  "ArcStr thin ptr eliminates fat pointer;",
                  "all variants now ≤ 8B"])

def slide_v5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "v5: NaN-boxed Value", "struct Value — 8 bytes")
    
    # Value box (simple)
    vg = add_group_box(slide, V_LEFT, V_TOP, Inches(1.8), Inches(0.8), "struct Value — 8 bytes")
    add_box(slide, V_LEFT + Inches(0.1), V_TOP + Inches(0.4), Inches(1.6), BOX_H,
            "bits: u64", fill=BLUE_INLINE, font_size=9)
    
    # NaN Boxing group
    nbl = Inches(0.3)
    nbt = V_TOP + Inches(1.2)
    nbg = add_group_box(slide, nbl, nbt, Inches(4.0), Inches(3.5),
                        "IEEE 754 NaN Boxing")
    
    add_box(slide, nbl + Inches(0.1), nbt + Inches(0.4), Inches(3.8), BOX_H,
            "not all-1 → Float (raw f64 inline)", fill=BLUE_INLINE, font_size=9, rounded=True)
    
    tg = add_group_box(slide, nbl + Inches(0.1), nbt + Inches(0.9), Inches(3.8), Inches(2.4),
                       "Tagged: upper 16 bits", fill=BLUE_BG, border_width=Pt(1))
    
    tag_items = [
        ("0xFFF8 → Null · False · True · Undef", BLUE_INLINE),
        ("0xFFFA → HeapNumber ptr", WHITE),
        ("0xFFFC → String (ArcStr ptr)", WHITE),
        ("0xFFFD → Array ptr", WHITE),
        ("0xFFFE → Object ptr", WHITE),
        ("0xFFFF → Set ptr", WHITE),
    ]
    variant_boxes(slide, tag_items, nbl + Inches(0.2), nbt + Inches(1.2),
                  box_w=Inches(3.6), box_h=Inches(0.33))
    
    # HeapNumber
    hl = Inches(4.8)
    hng = add_group_box(slide, hl, V_TOP, Inches(2.0), Inches(1.5), "enum HeapNumber")
    heap_items = [("Int(i64)", WHITE), ("UInt(u64)", WHITE), ("BigInt → Arc", WHITE)]
    y = V_TOP + Inches(0.4)
    for text, fill in heap_items:
        add_box(slide, hl + Inches(0.1), y, Inches(1.8), BOX_H, text, fill=fill, font_size=9)
        y += BOX_H + BOX_GAP
    
    # ObjectMap
    omg = add_group_box(slide, hl, V_TOP + Inches(1.7), Inches(2.8), Inches(1.8),
                        "struct ObjectMap (unordered)")
    add_box(slide, hl + Inches(0.1), V_TOP + Inches(2.1), Inches(2.6), BOX_H,
            "strings: HashMap‹ArcStr, Value›", fill=WHITE, font_size=9)
    add_box(slide, hl + Inches(0.1), V_TOP + Inches(2.55), Inches(2.6), BOX_H,
            "other: Option‹HashMap‹Value, Value››", fill=WHITE, font_size=9)
    add_box(slide, hl + Inches(0.1), V_TOP + Inches(3.0), Inches(2.6), BOX_H,
            "cached_hash: u64", fill=BLUE_INLINE, font_size=9, rounded=True)
    
    # Targets
    heap_boxes(slide, [
        ("str", "str"),
        ("arr", "Vec‹Value›"),
        ("set", "HashSet‹Value› (unordered)"),
    ], hl + Inches(3.2), V_TOP)
    
    # Interner
    ig = add_group_box(slide, hl, V_TOP + Inches(3.7), Inches(2.6), Inches(0.8), "Interner")
    add_box(slide, hl + Inches(0.1), V_TOP + Inches(4.1), Inches(2.4), BOX_H,
            "HashSet‹ArcStr› (thread-local)", fill=WHITE, font_size=9)
    
    add_note_box(slide, Inches(0.3), Inches(5.0), Inches(4.5), Inches(0.8),
                 ["8B = single u64",
                  "NaN boxing steals tag bits from f64 quiet-NaN space;",
                  "pointers limited to 48 bits"])

def slide_v6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "v6: Portable tagged pointer", "struct Value — 8 bytes")
    
    # Value box
    vg = add_group_box(slide, V_LEFT, V_TOP, Inches(1.8), Inches(0.8), "struct Value — 8 bytes")
    add_box(slide, V_LEFT + Inches(0.1), V_TOP + Inches(0.4), Inches(1.6), BOX_H,
            "bits: usize", fill=BLUE_INLINE, font_size=9)
    
    # Tags group
    tl = Inches(0.3)
    tt = V_TOP + Inches(1.2)
    add_group_box(slide, tl, tt, Inches(3.2), Inches(2.2),
                  "Low 3-bit tags (portable)")
    tag_items = [
        ("0b001 → u61 unsigned inline", BLUE_INLINE),
        ("0b010 → u61 negative inline", BLUE_INLINE),
        ("0b011 → Null · False · True · Undef", BLUE_INLINE),
        ("0b000 → Arc‹HeapValue›", WHITE),
    ]
    variant_boxes(slide, tag_items, tl + Inches(0.1), tt + Inches(0.4),
                  box_w=Inches(3.0), box_h=Inches(0.38))
    
    # HeapValue
    hl = Inches(4.0)
    hvg = add_group_box(slide, hl, V_TOP, Inches(2.8), Inches(3.8), "enum HeapValue")
    hv_items = [
        ("Float(f64)", WHITE), ("BigInt → Arc", WHITE),
        ("LargeUInt(u64)", WHITE), ("LargeInt(i64)", WHITE),
        ("String(ArcStr)", WHITE), ("Array(Vec‹Value›)", WHITE),
        ("Object(ObjectMap)", WHITE), ("Set(HashSet‹Value›) — unordered", WHITE),
    ]
    variant_boxes(slide, hv_items, hl + Inches(0.1), V_TOP + Inches(0.4),
                  box_w=Inches(2.6), box_h=Inches(0.36))
    
    # ObjectMap
    omg = add_group_box(slide, Inches(7.3), V_TOP, Inches(2.8), Inches(1.8),
                        "struct ObjectMap (unordered)")
    add_box(slide, Inches(7.4), V_TOP + Inches(0.4), Inches(2.6), BOX_H,
            "strings: HashMap‹ArcStr, Value›", fill=WHITE, font_size=9)
    add_box(slide, Inches(7.4), V_TOP + Inches(0.85), Inches(2.6), BOX_H,
            "other: Option‹HashMap‹Value, Value››", fill=WHITE, font_size=9)
    add_box(slide, Inches(7.4), V_TOP + Inches(1.3), Inches(2.6), BOX_H,
            "cached_hash: u64", fill=BLUE_INLINE, font_size=9, rounded=True)
    
    # ArcStr
    add_group_box(slide, Inches(7.3), V_TOP + Inches(2.0), Inches(2.4), Inches(0.8), "ArcStr — 8 bytes")
    add_box(slide, Inches(7.4), V_TOP + Inches(2.4), Inches(2.2), BOX_H,
            "thin pointer to heap string", fill=BLUE_INLINE, font_size=9, rounded=True)
    
    # Interner
    add_group_box(slide, Inches(7.3), V_TOP + Inches(3.0), Inches(2.6), Inches(0.8), "Interner")
    add_box(slide, Inches(7.4), V_TOP + Inches(3.4), Inches(2.4), BOX_H,
            "HashSet‹ArcStr› (thread-local)", fill=WHITE, font_size=9)
    
    add_note_box(slide, Inches(0.3), Inches(5.0), Inches(4.5), Inches(0.8),
                 ["8B = single usize",
                  "Low 3 alignment bits of pointers used as tag;",
                  "portable — no VA-width assumptions"])

def slide_v7(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "v7: Tagged pointer + schema-shared objects", "struct Value — 8 bytes")
    
    # Value box
    vg = add_group_box(slide, V_LEFT, V_TOP, Inches(1.8), Inches(0.8), "struct Value — 8 bytes")
    add_box(slide, V_LEFT + Inches(0.1), V_TOP + Inches(0.4), Inches(1.6), BOX_H,
            "bits: usize", fill=BLUE_INLINE, font_size=9)
    
    # Tags
    tl = Inches(0.3)
    tt = V_TOP + Inches(1.2)
    add_group_box(slide, tl, tt, Inches(3.2), Inches(2.2), "Low 3-bit tags (portable)")
    tag_items = [
        ("0b001 → u61 unsigned inline", BLUE_INLINE),
        ("0b010 → u61 negative inline", BLUE_INLINE),
        ("0b011 → Null · False · True · Undef", BLUE_INLINE),
        ("0b000 → Arc‹HeapValue›", WHITE),
    ]
    variant_boxes(slide, tag_items, tl + Inches(0.1), tt + Inches(0.4),
                  box_w=Inches(3.0), box_h=Inches(0.38))
    
    # HeapValue
    hl = Inches(3.8)
    hvg = add_group_box(slide, hl, V_TOP, Inches(2.5), Inches(3.8), "enum HeapValue")
    hv_items = [
        ("Float(f64)", WHITE), ("BigInt → Arc", WHITE),
        ("LargeUInt(u64)", WHITE), ("LargeInt(i64)", WHITE),
        ("String(ArcStr)", WHITE), ("Array(Vec‹Value›)", WHITE),
        ("Object(ObjectMap)", WHITE), ("Set(HashSet) — unordered", WHITE),
    ]
    variant_boxes(slide, hv_items, hl + Inches(0.1), V_TOP + Inches(0.4),
                  box_w=Inches(2.3), box_h=Inches(0.36))
    
    # ObjectMap with Compact + Map
    oml = Inches(6.6)
    omg = add_group_box(slide, oml, V_TOP, Inches(3.5), Inches(4.2), "struct ObjectMap")
    
    # Compact sub
    cg = add_group_box(slide, oml + Inches(0.1), V_TOP + Inches(0.4), Inches(3.3), Inches(1.6),
                       "Compact (sorted via Schema)", fill=BLUE_BG, border_width=Pt(1))
    add_box(slide, oml + Inches(0.2), V_TOP + Inches(0.8), Inches(3.1), BOX_H,
            "schema: Arc‹Schema›", fill=WHITE, font_size=9)
    add_box(slide, oml + Inches(0.2), V_TOP + Inches(1.25), Inches(3.1), BOX_H,
            "values: Box‹[Value]›", fill=WHITE, font_size=9)
    add_box(slide, oml + Inches(0.2), V_TOP + Inches(1.7), Inches(3.1), BOX_H,
            "cached_hash: u64", fill=BLUE_INLINE, font_size=9, rounded=True)
    
    # Map fallback sub
    mg = add_group_box(slide, oml + Inches(0.1), V_TOP + Inches(2.1), Inches(3.3), Inches(1.8),
                       "Map — fallback (unordered)", fill=BLUE_BG, border_width=Pt(1))
    add_box(slide, oml + Inches(0.2), V_TOP + Inches(2.5), Inches(3.1), BOX_H,
            "strings: HashMap‹ArcStr, Value›", fill=WHITE, font_size=9)
    add_box(slide, oml + Inches(0.2), V_TOP + Inches(2.95), Inches(3.1), BOX_H,
            "other: Option‹HashMap‹Value, Value››", fill=WHITE, font_size=9)
    add_box(slide, oml + Inches(0.2), V_TOP + Inches(3.4), Inches(3.1), BOX_H,
            "cached_hash: u64", fill=BLUE_INLINE, font_size=9, rounded=True)
    
    # Schema
    sl = Inches(6.6)
    st = V_TOP + Inches(4.4)
    sg = add_group_box(slide, sl, st, Inches(3.0), Inches(1.2), "Schema — shared via Arc")
    add_box(slide, sl + Inches(0.1), st + Inches(0.4), Inches(2.8), BOX_H,
            "keys: Arc‹[ArcStr]›", fill=WHITE, font_size=9)
    add_box(slide, sl + Inches(0.1), st + Inches(0.85), Inches(2.8), BOX_H,
            "lookup: HashMap‹ArcStr, u32›", fill=WHITE, font_size=9)
    
    # Interner + Schema Cache
    add_group_box(slide, Inches(0.3), V_TOP + Inches(3.6), Inches(2.6), Inches(0.8), "Interner")
    add_box(slide, Inches(0.4), V_TOP + Inches(4.0), Inches(2.4), BOX_H,
            "HashSet‹ArcStr› (thread-local)", fill=WHITE, font_size=9)
    
    add_group_box(slide, Inches(0.3), V_TOP + Inches(4.6), Inches(3.2), Inches(0.8),
                  "Schema Cache (thread-local)")
    add_box(slide, Inches(0.4), V_TOP + Inches(5.0), Inches(3.0), BOX_H,
            "HashMap‹Vec‹ArcStr›, Arc‹Schema››", fill=WHITE, font_size=9)
    
    add_note_box(slide, Inches(0.3), Inches(6.0), Inches(4.5), Inches(0.8),
                 ["8B = single usize",
                  "Same tagged pointer as v6;",
                  "schema-shared objects reduce heap, not Value size"])

def slide_v8(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "v8: Flattened Number + schema objects", "enum Value — 16 bytes")
    
    vg = add_group_box(slide, V_LEFT, V_TOP, Inches(2.6), Inches(4.4), "enum Value — 16 bytes")
    variants = [
        ("object — 8B thin ptr", WHITE),
        ("string: ArcStr — 8B thin ptr", WHITE),
        ("array — 8B thin ptr", WHITE),
        ("set — 8B thin ptr", WHITE),
        ("uint(u64) — 8B", BLUE_INLINE),
        ("int(i64) — 8B", BLUE_INLINE),
        ("float(f64) — 8B", BLUE_INLINE),
        ("bigint — 8B thin ptr", WHITE),
        ("null · bool · undefined — ≤ 1B", BLUE_INLINE),
    ]
    variant_boxes(slide, variants, V_LEFT + Inches(0.3), V_TOP + Inches(0.4), box_w=Inches(2.0))
    
    # ObjectMap with Compact + Map
    oml = Inches(4.0)
    omg = add_group_box(slide, oml, V_TOP, Inches(3.5), Inches(4.0), "struct ObjectMap")
    
    cg = add_group_box(slide, oml + Inches(0.1), V_TOP + Inches(0.4), Inches(3.3), Inches(1.5),
                       "Compact (sorted via Schema)", fill=BLUE_BG, border_width=Pt(1))
    add_box(slide, oml + Inches(0.2), V_TOP + Inches(0.8), Inches(3.1), BOX_H,
            "schema: Arc‹Schema›", fill=WHITE, font_size=9)
    add_box(slide, oml + Inches(0.2), V_TOP + Inches(1.25), Inches(3.1), BOX_H,
            "values: Box‹[Value]›", fill=WHITE, font_size=9)
    add_box(slide, oml + Inches(0.2), V_TOP + Inches(1.6), Inches(3.1), BOX_H,
            "cached_hash: u64", fill=BLUE_INLINE, font_size=9, rounded=True)
    
    mg = add_group_box(slide, oml + Inches(0.1), V_TOP + Inches(2.0), Inches(3.3), Inches(1.8),
                       "Map — fallback (unordered)", fill=BLUE_BG, border_width=Pt(1))
    add_box(slide, oml + Inches(0.2), V_TOP + Inches(2.4), Inches(3.1), BOX_H,
            "strings: HashMap‹ArcStr, Value›", fill=WHITE, font_size=9)
    add_box(slide, oml + Inches(0.2), V_TOP + Inches(2.85), Inches(3.1), BOX_H,
            "other: Option‹HashMap‹Value, Value››", fill=WHITE, font_size=9)
    add_box(slide, oml + Inches(0.2), V_TOP + Inches(3.3), Inches(3.1), BOX_H,
            "cached_hash: u64", fill=BLUE_INLINE, font_size=9, rounded=True)
    
    # Schema + Cache + ArcStr
    sr = Inches(8.0)
    sg = add_group_box(slide, sr, V_TOP, Inches(3.0), Inches(1.2), "Schema — shared via Arc")
    add_box(slide, sr + Inches(0.1), V_TOP + Inches(0.4), Inches(2.8), BOX_H,
            "keys: Arc‹[ArcStr]›", fill=WHITE, font_size=9)
    add_box(slide, sr + Inches(0.1), V_TOP + Inches(0.85), Inches(2.8), BOX_H,
            "lookup: HashMap‹ArcStr, u32›", fill=WHITE, font_size=9)
    
    add_group_box(slide, sr, V_TOP + Inches(1.4), Inches(3.2), Inches(0.8),
                  "Schema Cache (thread-local)")
    add_box(slide, sr + Inches(0.1), V_TOP + Inches(1.8), Inches(3.0), BOX_H,
            "HashMap‹Vec‹ArcStr›, Arc‹Schema››", fill=WHITE, font_size=9)
    
    add_group_box(slide, sr, V_TOP + Inches(2.4), Inches(2.4), Inches(0.8), "ArcStr — 8 bytes")
    add_box(slide, sr + Inches(0.1), V_TOP + Inches(2.8), Inches(2.2), BOX_H,
            "thin pointer to heap string", fill=BLUE_INLINE, font_size=9, rounded=True)
    
    # Heap targets
    heap_boxes(slide, [
        ("arr", "Vec‹Value›"),
        ("set", "HashSet‹Value› (unordered)"),
        ("bigint", "BigInt"),
    ], sr, V_TOP + Inches(3.4))
    
    add_note_box(slide, Inches(0.4), Inches(5.7), Inches(4.0), Inches(0.8),
                 ["16B = 8B tag + 8B payload",
                  "Flattened numbers + ArcStr thin ptr;",
                  "all variants ≤ 8B"])

def slide_v9(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "v9: Arena-allocated, Copy Value", "enum Value‹'a› — 16 bytes, Copy")
    
    vg = add_group_box(slide, V_LEFT, V_TOP, Inches(2.8), Inches(5.0),
                       "enum Value‹'a› — 16 bytes, Copy")
    ag = add_group_box(slide, V_LEFT + Inches(0.1), V_TOP + Inches(0.4), Inches(2.6), Inches(4.0),
                       "Arena variants", fill=BLUE_BG, border_width=Pt(1))
    arena_items = [
        ("object — 8B ref", WHITE),
        ("string: &'a ArenaStr — 8B ref", WHITE),
        ("array: &'a ArenaArray — 8B ref", WHITE),
        ("set: &'a ArenaSet (unordered) — 8B ref", WHITE),
        ("uint(u64) — 8B", BLUE_INLINE),
        ("int(i64) — 8B", BLUE_INLINE),
        ("float(f64) — 8B", BLUE_INLINE),
        ("bigint: &'a BigInt — 8B ref", WHITE),
        ("null · bool · undefined — ≤ 1B", BLUE_INLINE),
    ]
    variant_boxes(slide, arena_items, V_LEFT + Inches(0.2), V_TOP + Inches(0.7),
                  box_w=Inches(2.4))
    add_box(slide, V_LEFT + Inches(0.1), V_TOP + Inches(4.5), Inches(2.6), BOX_H,
            "Ext(&'a ExtValue) — 8B ref", fill=WHITE, font_size=9)
    
    # ObjectMap
    oml = Inches(3.6)
    omg = add_group_box(slide, oml, V_TOP, Inches(3.8), Inches(4.2), "struct ObjectMap‹'a›")
    
    cg = add_group_box(slide, oml + Inches(0.1), V_TOP + Inches(0.4), Inches(3.6), Inches(1.7),
                       "Compact (sorted via Schema)", fill=BLUE_BG, border_width=Pt(1))
    add_box(slide, oml + Inches(0.2), V_TOP + Inches(0.8), Inches(3.4), BOX_H,
            "schema: Arc‹Schema›", fill=WHITE, font_size=9)
    add_box(slide, oml + Inches(0.2), V_TOP + Inches(1.2), Inches(3.4), BOX_H,
            "values: &'a [Value‹'a›]", fill=WHITE, font_size=9)
    add_box(slide, oml + Inches(0.2), V_TOP + Inches(1.6), Inches(3.4), BOX_H,
            "arena_keys: &'a [&'a ArenaStr]", fill=WHITE, font_size=9)
    add_box(slide, oml + Inches(0.2), V_TOP + Inches(1.95), Inches(3.4), BOX_H - Inches(0.05),
            "cached_hash: u64", fill=BLUE_INLINE, font_size=9, rounded=True)
    
    mg = add_group_box(slide, oml + Inches(0.1), V_TOP + Inches(2.2), Inches(3.6), Inches(1.8),
                       "Map — fallback (unordered)", fill=BLUE_BG, border_width=Pt(1))
    add_box(slide, oml + Inches(0.2), V_TOP + Inches(2.6), Inches(3.4), BOX_H,
            "strings: HashMap‹&'a str, Value, _, &'a Bump›", fill=WHITE, font_size=8)
    add_box(slide, oml + Inches(0.2), V_TOP + Inches(3.05), Inches(3.4), BOX_H,
            "others: Option‹HashMap‹Value, Value, _, &'a Bump››", fill=WHITE, font_size=8)
    add_box(slide, oml + Inches(0.2), V_TOP + Inches(3.5), Inches(3.4), BOX_H,
            "cached_hash: u64", fill=BLUE_INLINE, font_size=9, rounded=True)
    
    # Right column
    sr = Inches(7.8)
    
    # Schema
    sg = add_group_box(slide, sr, V_TOP, Inches(3.2), Inches(1.2), "Schema — shared via Arc")
    add_box(slide, sr + Inches(0.1), V_TOP + Inches(0.4), Inches(3.0), BOX_H,
            "keys: Arc‹[Arc‹str›]›", fill=WHITE, font_size=9)
    add_box(slide, sr + Inches(0.1), V_TOP + Inches(0.85), Inches(3.0), BOX_H,
            "lookup: HashMap‹Arc‹str›, u32›", fill=WHITE, font_size=9)
    
    # Schema Cache
    add_group_box(slide, sr, V_TOP + Inches(1.4), Inches(3.4), Inches(0.8),
                  "Schema Cache (thread-local)")
    add_box(slide, sr + Inches(0.1), V_TOP + Inches(1.8), Inches(3.2), BOX_H,
            "HashMap‹Vec‹Arc‹str››, Arc‹Schema››", fill=WHITE, font_size=8)
    
    # StringInterner
    add_group_box(slide, sr, V_TOP + Inches(2.4), Inches(3.0), Inches(0.8),
                  "StringInterner‹'a› (arena-backed)")
    add_box(slide, sr + Inches(0.1), V_TOP + Inches(2.8), Inches(2.8), BOX_H,
            "HashSet‹&'a str›", fill=WHITE, font_size=9)
    
    # Bump arena
    add_group_box(slide, sr, V_TOP + Inches(3.4), Inches(3.2), Inches(1.2), "bumpalo::Bump arena")
    add_box(slide, sr + Inches(0.1), V_TOP + Inches(3.8), Inches(3.0), Inches(0.6),
            "all Value data lives here\nbulk deallocation\nzero-cost clone (Copy)",
            fill=BLUE_INLINE, font_size=8, rounded=True)
    
    add_note_box(slide, Inches(0.4), Inches(6.2), Inches(4.5), Inches(0.8),
                 ["16B = 8B tag + 8B payload",
                  "Arena refs (thin ptrs) + niche optimization;",
                  "Copy — zero-cost clone"])


# ═══════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════

def main():
    prs = make_prs()
    
    slide_baseline(prs)
    slide_v1(prs)
    slide_v2(prs)
    slide_v3(prs)
    slide_v4(prs)
    slide_v5(prs)
    slide_v6(prs)
    slide_v7(prs)
    slide_v8(prs)
    slide_v9(prs)
    
    out = "RegorusValue.pptx"
    prs.save(out)
    print(f"Saved {out} with {len(prs.slides)} slides")

if __name__ == "__main__":
    main()
