#!/usr/bin/env python3
"""
Generate RegorusValue.pptx with fully editable native PowerPoint shapes + arrows.
All shapes, connectors, and text are editable.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# ── Blue Academic palette ──
C_BG       = RGBColor(0xEB, 0xF5, 0xFB)
C_INLINE   = RGBColor(0xD6, 0xEA, 0xF8)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT     = RGBColor(0x1A, 0x52, 0x76)
C_BORDER   = RGBColor(0x24, 0x71, 0xA3)
C_NOTE_BG  = RGBColor(0xFE, 0xF9, 0xE7)
C_NOTE_BD  = RGBColor(0xF1, 0xC4, 0x0F)
C_NOTE_TX  = RGBColor(0x7D, 0x66, 0x08)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Helpers ──

class Box:
    """Track shape position for arrow drawing."""
    def __init__(self, shape, left, top, width, height):
        self.shape = shape
        self.left = left
        self.top = top
        self.width = width
        self.height = height

    @property
    def right(self): return self.left + self.width
    @property
    def bottom(self): return self.top + self.height
    @property
    def cx(self): return self.left + self.width // 2
    @property
    def cy(self): return self.top + self.height // 2
    @property
    def right_mid(self): return (self.right, self.cy)
    @property
    def left_mid(self): return (self.left, self.cy)
    @property
    def bottom_mid(self): return (self.cx, self.bottom)
    @property
    def top_mid(self): return (self.cx, self.top)


def _set_text(shape, text, font_size=10, color=C_TEXT, bold=False, italic=False,
              align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p0 = tf.paragraphs[0]
    p0.alignment = align
    p0.space_before = Pt(1)
    p0.space_after = Pt(1)
    run = p0.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return tf


def _multiline(shape, lines, font_size=9, color=C_TEXT, bold_first=False, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        if bold_first and i == 0:
            run.font.bold = True


def add_box(slide, l, t, w, h, text, fill=C_WHITE, border=C_BORDER,
            font_size=10, bold=False, rounded=False, border_w=Pt(1),
            dash=False, text_color=C_TEXT):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shp_type, l, t, w, h)
    if rounded:
        shape.adjustments[0] = 0.08
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = border_w
    if dash:
        shape.line.dash_style = 4
    _set_text(shape, text, font_size=font_size, color=text_color, bold=bold)
    return Box(shape, l, t, w, h)


def add_group(slide, l, t, w, h, title, fill=C_BG, border_w=Pt(2)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = C_BORDER
    shape.line.width = border_w
    _set_text(shape, title, font_size=11, bold=True, align=PP_ALIGN.LEFT)
    # Anchor text to top
    shape.text_frame.paragraphs[0].space_before = Pt(2)
    return Box(shape, l, t, w, h)


def add_note(slide, l, t, w, h, lines):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_NOTE_BG
    shape.line.color.rgb = C_NOTE_BD
    shape.line.width = Pt(1.5)
    shape.line.dash_style = 4
    _multiline(shape, lines, font_size=9, color=C_NOTE_TX, bold_first=True)
    return Box(shape, l, t, w, h)


def add_arrow(slide, x1, y1, x2, y2, label="", dashed=False):
    """Draw a line with arrowhead from (x1,y1) to (x2,y2), with optional label."""
    # Use a freeform-ish approach: add_connector with XML arrowhead
    cxn = slide.shapes.add_connector(1, x1, y1, x2, y2)  # MSO_CONNECTOR.STRAIGHT
    cxn.line.color.rgb = C_BORDER
    cxn.line.width = Pt(1.5)
    if dashed:
        cxn.line.dash_style = 4

    # Add arrowhead at end via XML
    ln = cxn._element.find(qn('a:ln'))
    if ln is None:
        spPr = cxn._element.find(qn('p:spPr'))
        if spPr is None:
            spPr = cxn._element.find(qn('p:cxnSp'))
        ln = cxn._element.xpath('.//a:ln')[0] if cxn._element.xpath('.//a:ln') else None

    if ln is not None:
        tail = etree.SubElement(ln, qn('a:tailEnd'))
        tail.set('type', 'triangle')
        tail.set('w', 'med')
        tail.set('len', 'med')

    if label:
        # Place label near midpoint
        mx = (x1 + x2) // 2
        my = (y1 + y2) // 2 - Inches(0.18)
        tb = slide.shapes.add_textbox(mx - Inches(0.3), my, Inches(0.7), Inches(0.2))
        tf = tb.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = label
        run.font.size = Pt(8)
        run.font.color.rgb = C_BORDER
        run.font.bold = True


def arrow_right(slide, src, dst, label="", dashed=False):
    """Arrow from right edge of src to left edge of dst."""
    x1, y1 = src.right, src.cy
    x2, y2 = dst.left, dst.cy
    add_arrow(slide, x1, y1, x2, y2, label, dashed)


def arrow_down(slide, src, dst, label="", dashed=False):
    """Arrow from bottom of src to top of dst."""
    x1, y1 = src.cx, src.bottom
    x2, y2 = dst.cx, dst.top
    add_arrow(slide, x1, y1, x2, y2, label, dashed)


def add_title(slide, title, subtitle=""):
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12), Inches(0.5))
    tf = tb.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = title
    run.font.size = Pt(24)
    run.font.color.rgb = C_TEXT
    run.font.bold = True
    if subtitle:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = subtitle
        r.font.size = Pt(13)
        r.font.color.rgb = C_BORDER


# ── Layout constants ──
# Slide is 13.333" x 7.5" (widescreen). Use full width with 3 columns.
I = Inches
BW = I(2.4)   # box width
BH = I(0.36)  # box height
BG = I(0.07)  # box gap

# Column positions for 3-column layout
COL1_L = I(0.4)          # Left column (Value enum)
COL2_L = I(4.0)          # Middle column (heap targets / ObjectMap)
COL3_L = I(8.5)          # Right column (extras: Schema, Interner, ArcStr)
TOP = I(0.9)             # Below title


def vstack(slide, items, left, top, bw=BW, bh=BH, bg=BG):
    """Draw vertical stack of (text, fill) boxes. Returns list of Box."""
    boxes = []
    y = top
    for text, fill in items:
        b = add_box(slide, left, y, bw, bh, text, fill=fill, font_size=10,
                    rounded=(fill == C_INLINE))
        boxes.append(b)
        y += bh + bg
    return boxes


# Each slide -- full-width 3-column layout (13.333in wide, 7.5in tall)
# COL1 ~0.4in, COL2 ~4.0in, COL3 ~8.5in

def slide_baseline(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Baseline: Arc + BTreeMap", "enum Value -- 24 bytes")

    # Col1: Value group
    VW, VH = I(3.2), I(3.6)
    add_group(s, COL1_L, TOP, VW, VH, "enum Value -- 24 bytes")
    vb = vstack(s, [
        ("object -- 8B thin ptr", C_WHITE),
        ("string -- 16B fat ptr", C_WHITE),
        ("array -- 8B thin ptr", C_WHITE),
        ("set -- 8B thin ptr", C_WHITE),
        ("number -- 16B nested enum", C_WHITE),
        ("null / bool / undefined", C_INLINE),
    ], COL1_L + I(0.3), TOP + I(0.45), bw=I(2.6))

    # Col2: Heap targets + Number
    obj_t = add_box(s, COL2_L, TOP + I(0.2),  I(3.5), BH, "BTreeMap<Value, Value> (ordered)", font_size=10)
    str_t = add_box(s, COL2_L, TOP + I(0.7),  I(3.5), BH, "str", font_size=10)
    arr_t = add_box(s, COL2_L, TOP + I(1.2),  I(3.5), BH, "Vec<Value>", font_size=10)
    set_t = add_box(s, COL2_L, TOP + I(1.7),  I(3.5), BH, "BTreeSet<Value> (ordered)", font_size=10)

    NT = TOP + I(2.4)
    add_group(s, COL2_L, NT, I(3.5), I(1.3), "enum Number -- 16 bytes")
    n_inl = add_box(s, COL2_L+I(0.15), NT+I(0.4), I(3.2), BH, "u64 / i64 / f64",
                    fill=C_INLINE, font_size=10, rounded=True)
    n_big = add_box(s, COL2_L+I(0.15), NT+I(0.85), I(3.2), BH, "bigint", font_size=10)

    # Col3: BigInt target
    big_t = add_box(s, COL3_L, NT+I(0.85), I(2.0), BH, "BigInt", font_size=10)

    # Arrows
    arrow_right(s, vb[0], obj_t, "Arc")
    arrow_right(s, vb[1], str_t, "Arc")
    arrow_right(s, vb[2], arr_t, "Arc")
    arrow_right(s, vb[3], set_t, "Arc")
    arrow_right(s, n_big, big_t, "Arc")
    add_arrow(s, vb[4].right, vb[4].cy, COL2_L, NT + I(0.65), dashed=True)

    # Note
    add_note(s, COL1_L, TOP + I(4.2), I(5.5), I(0.7),
             ["24B = 8B tag + 16B largest payload",
              "String (fat ptr) and Number tie at 16B"])


def slide_v1(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "v1: HashMap + SmolStr", "enum Value -- 24 bytes")

    VW = I(3.2)
    add_group(s, COL1_L, TOP, VW, I(3.6), "enum Value -- 24 bytes")
    vb = vstack(s, [
        ("object -- 8B thin ptr", C_WHITE),
        ("string -- 16B fat ptr", C_WHITE),
        ("array -- 8B thin ptr", C_WHITE),
        ("set -- 8B thin ptr", C_WHITE),
        ("number -- 16B nested enum", C_WHITE),
        ("null / bool / undefined", C_INLINE),
    ], COL1_L+I(0.3), TOP+I(0.45), bw=I(2.6))

    # Col2: ObjectMap + heap targets
    add_group(s, COL2_L, TOP, I(3.8), I(2.1), "enum ObjectMap (unordered)")
    add_group(s, COL2_L+I(0.1), TOP+I(0.4), I(3.6), I(0.65), "StringKeyed", fill=C_BG, border_w=Pt(1))
    sk = add_box(s, COL2_L+I(0.15), TOP+I(0.6), I(3.5), BH, "HashMap<SmolStr, Value>", font_size=10)
    add_group(s, COL2_L+I(0.1), TOP+I(1.15), I(3.6), I(0.65), "Mixed", fill=C_BG, border_w=Pt(1))
    add_box(s, COL2_L+I(0.15), TOP+I(1.35), I(3.5), BH, "HashMap<Value, Value>", font_size=10)

    str_t = add_box(s, COL2_L, TOP+I(2.3), I(3.5), BH, "str", font_size=10)
    arr_t = add_box(s, COL2_L, TOP+I(2.8), I(3.5), BH, "Vec<Value>", font_size=10)
    set_t = add_box(s, COL2_L, TOP+I(3.3), I(3.5), BH, "BTreeSet<Value> (ordered)", font_size=10)

    # Number
    NT = TOP + I(3.9)
    add_group(s, COL2_L, NT, I(3.5), I(1.3), "enum Number -- 16 bytes")
    add_box(s, COL2_L+I(0.15), NT+I(0.4), I(3.2), BH, "u64 / i64 / f64",
            fill=C_INLINE, font_size=10, rounded=True)
    add_box(s, COL2_L+I(0.15), NT+I(0.85), I(3.2), BH, "bigint -> Arc BigInt", font_size=10)

    # Col3: SmolStr
    add_group(s, COL3_L, TOP, I(4.0), I(1.3), "SmolStr -- 24 bytes")
    add_box(s, COL3_L+I(0.15), TOP+I(0.4), I(3.7), BH, "<= 23 chars: inline",
            fill=C_INLINE, font_size=10, rounded=True)
    add_box(s, COL3_L+I(0.15), TOP+I(0.85), I(3.7), BH, "> 23 chars -> Arc str", font_size=10)

    # Arrows
    arrow_right(s, vb[0], Box(None, COL2_L, TOP+I(1.0), 0, 0), "Arc")
    arrow_right(s, vb[1], str_t, "Arc")
    arrow_right(s, vb[2], arr_t, "Arc")
    arrow_right(s, vb[3], set_t, "Arc")
    add_arrow(s, vb[4].right, vb[4].cy, COL2_L, NT+I(0.6), dashed=True)
    arrow_right(s, sk, Box(None, COL3_L, TOP+I(0.6), 0, 0), "", dashed=True)

    add_note(s, COL1_L, TOP+I(5.5), I(5.5), I(0.7),
             ["24B = 8B tag + 16B largest payload",
              "Arc<str> fat pointer drives the size"])


def slide_v2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "v2: Arc<str> keys (no SmolStr)", "enum Value -- 24 bytes")

    VW = I(3.2)
    add_group(s, COL1_L, TOP, VW, I(3.6), "enum Value -- 24 bytes")
    vb = vstack(s, [
        ("object -- 8B thin ptr", C_WHITE),
        ("string -- 16B fat ptr", C_WHITE),
        ("array -- 8B thin ptr", C_WHITE),
        ("set -- 8B thin ptr", C_WHITE),
        ("number -- 16B nested enum", C_WHITE),
        ("null / bool / undefined", C_INLINE),
    ], COL1_L+I(0.3), TOP+I(0.45), bw=I(2.6))

    # Col2: ObjectMap
    add_group(s, COL2_L, TOP, I(3.8), I(2.1), "enum ObjectMap (unordered)")
    add_group(s, COL2_L+I(0.1), TOP+I(0.4), I(3.6), I(0.65), "StringKeyed", fill=C_BG, border_w=Pt(1))
    add_box(s, COL2_L+I(0.15), TOP+I(0.6), I(3.5), BH, "HashMap<Arc<str>, Value>", font_size=10)
    add_group(s, COL2_L+I(0.1), TOP+I(1.15), I(3.6), I(0.65), "Mixed", fill=C_BG, border_w=Pt(1))
    add_box(s, COL2_L+I(0.15), TOP+I(1.35), I(3.5), BH, "HashMap<Value, Value>", font_size=10)

    str_t = add_box(s, COL2_L, TOP+I(2.3), I(3.5), BH, "str", font_size=10)
    arr_t = add_box(s, COL2_L, TOP+I(2.8), I(3.5), BH, "Vec<Value>", font_size=10)
    set_t = add_box(s, COL2_L, TOP+I(3.3), I(3.5), BH, "BTreeSet<Value> (ordered)", font_size=10)

    NT = TOP + I(3.9)
    add_group(s, COL2_L, NT, I(3.5), I(1.3), "enum Number -- 16 bytes")
    add_box(s, COL2_L+I(0.15), NT+I(0.4), I(3.2), BH, "u64 / i64 / f64",
            fill=C_INLINE, font_size=10, rounded=True)
    add_box(s, COL2_L+I(0.15), NT+I(0.85), I(3.2), BH, "bigint -> Arc BigInt", font_size=10)

    # Col3: Interner
    add_group(s, COL3_L, TOP, I(4.0), I(1.0), "Interner")
    add_box(s, COL3_L+I(0.15), TOP+I(0.45), I(3.7), BH, "HashMap<u64, Arc<str>>", font_size=10)

    arrow_right(s, vb[0], Box(None, COL2_L, TOP+I(1.0), 0, 0), "Arc")
    arrow_right(s, vb[1], str_t, "Arc")
    arrow_right(s, vb[2], arr_t, "Arc")
    arrow_right(s, vb[3], set_t, "Arc")
    add_arrow(s, vb[4].right, vb[4].cy, COL2_L, NT+I(0.6), dashed=True)

    add_note(s, COL1_L, TOP+I(5.5), I(5.5), I(0.7),
             ["24B = 8B tag + 16B largest payload",
              "Arc<str> fat pointer drives the size"])


def slide_v3(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "v3: HashSet + cached hash", "enum Value -- 24 bytes")

    VW = I(3.2)
    add_group(s, COL1_L, TOP, VW, I(3.6), "enum Value -- 24 bytes")
    vb = vstack(s, [
        ("object -- 8B thin ptr", C_WHITE),
        ("string -- 16B fat ptr", C_WHITE),
        ("array -- 8B thin ptr", C_WHITE),
        ("set -- 8B thin ptr", C_WHITE),
        ("number -- 16B nested enum", C_WHITE),
        ("null / bool / undefined", C_INLINE),
    ], COL1_L+I(0.3), TOP+I(0.45), bw=I(2.6))

    # Col2: struct ObjectMap + heap
    add_group(s, COL2_L, TOP, I(3.8), I(1.9), "struct ObjectMap (unordered)")
    add_box(s, COL2_L+I(0.15), TOP+I(0.4), I(3.5), BH, "strings: HashMap<Arc<str>, Value>", font_size=10)
    add_box(s, COL2_L+I(0.15), TOP+I(0.9), I(3.5), BH, "other: Option<HashMap<Value, Value>>", font_size=9)
    add_box(s, COL2_L+I(0.15), TOP+I(1.35), I(3.5), BH, "cached_hash: u64",
            fill=C_INLINE, font_size=10, rounded=True)

    str_t = add_box(s, COL2_L, TOP+I(2.2),  I(3.5), BH, "str", font_size=10)
    arr_t = add_box(s, COL2_L, TOP+I(2.7),  I(3.5), BH, "Vec<Value>", font_size=10)
    set_t = add_box(s, COL2_L, TOP+I(3.2),  I(3.5), BH, "HashSet<Value> (unordered)", font_size=10)

    NT = TOP + I(3.8)
    add_group(s, COL2_L, NT, I(3.5), I(1.3), "enum Number -- 16 bytes")
    add_box(s, COL2_L+I(0.15), NT+I(0.4), I(3.2), BH, "u64 / i64 / f64",
            fill=C_INLINE, font_size=10, rounded=True)
    add_box(s, COL2_L+I(0.15), NT+I(0.85), I(3.2), BH, "bigint -> Arc BigInt", font_size=10)

    # Col3: Interner
    add_group(s, COL3_L, TOP, I(4.0), I(1.0), "Interner (reuses v2)")
    add_box(s, COL3_L+I(0.15), TOP+I(0.45), I(3.7), BH, "HashMap<u64, Arc<str>>", font_size=10)

    arrow_right(s, vb[0], Box(None, COL2_L, TOP+I(0.8), 0, 0), "Arc")
    arrow_right(s, vb[1], str_t, "Arc")
    arrow_right(s, vb[2], arr_t, "Arc")
    arrow_right(s, vb[3], set_t, "Arc")
    add_arrow(s, vb[4].right, vb[4].cy, COL2_L, NT+I(0.6), dashed=True)

    add_note(s, COL1_L, TOP+I(5.5), I(5.5), I(0.7),
             ["24B = 8B tag + 16B largest payload",
              "Arc<str> fat pointer drives the size"])


def slide_v4(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "v4: ArcStr thin pointer", "enum Value -- 16 bytes")

    VW = I(3.2)
    add_group(s, COL1_L, TOP, VW, I(3.6), "enum Value -- 16 bytes")
    vb = vstack(s, [
        ("object -- 8B thin ptr", C_WHITE),
        ("string: ArcStr -- 8B thin ptr", C_WHITE),
        ("array -- 8B thin ptr", C_WHITE),
        ("set -- 8B thin ptr", C_WHITE),
        ("number -- 8B (flattened)", C_WHITE),
        ("null / bool / undefined", C_INLINE),
    ], COL1_L+I(0.3), TOP+I(0.45), bw=I(2.6))

    # Col2: ObjectMap + heap
    add_group(s, COL2_L, TOP, I(3.8), I(1.9), "struct ObjectMap (unordered)")
    add_box(s, COL2_L+I(0.15), TOP+I(0.4), I(3.5), BH, "strings: HashMap<ArcStr, Value>", font_size=10)
    add_box(s, COL2_L+I(0.15), TOP+I(0.9), I(3.5), BH, "other: Option<HashMap<Value, Value>>", font_size=9)
    add_box(s, COL2_L+I(0.15), TOP+I(1.35), I(3.5), BH, "cached_hash: u64",
            fill=C_INLINE, font_size=10, rounded=True)

    arr_t = add_box(s, COL2_L, TOP+I(2.2), I(3.5), BH, "Vec<Value>", font_size=10)
    set_t = add_box(s, COL2_L, TOP+I(2.7), I(3.5), BH, "HashSet<Value> (unordered)", font_size=10)

    NT = TOP + I(3.3)
    add_group(s, COL2_L, NT, I(3.5), I(1.3), "enum Number -- 16 bytes")
    add_box(s, COL2_L+I(0.15), NT+I(0.4), I(3.2), BH, "u64 / i64 / f64",
            fill=C_INLINE, font_size=10, rounded=True)
    add_box(s, COL2_L+I(0.15), NT+I(0.85), I(3.2), BH, "bigint -> Arc BigInt", font_size=10)

    # Col3: ArcStr + Interner
    add_group(s, COL3_L, TOP, I(4.0), I(1.0), "ArcStr -- 8 bytes")
    add_box(s, COL3_L+I(0.15), TOP+I(0.45), I(3.7), BH, "thin pointer to heap string",
            fill=C_INLINE, font_size=10, rounded=True)

    add_group(s, COL3_L, TOP+I(1.3), I(4.0), I(1.0), "Interner")
    add_box(s, COL3_L+I(0.15), TOP+I(1.75), I(3.7), BH, "HashSet<ArcStr> (thread-local)", font_size=10)

    arrow_right(s, vb[0], Box(None, COL2_L, TOP+I(0.8), 0, 0), "Arc")
    arrow_right(s, vb[2], arr_t, "Arc")
    arrow_right(s, vb[3], set_t, "Arc")
    add_arrow(s, vb[1].right, vb[1].cy, COL3_L, TOP+I(0.45)+BH//2, "", dashed=True)
    add_arrow(s, vb[4].right, vb[4].cy, COL2_L, NT+I(0.6), dashed=True)

    add_note(s, COL1_L, TOP+I(5.0), I(6.0), I(0.8),
             ["16B = 8B tag + 8B payload",
              "ArcStr thin ptr eliminates fat pointer; all variants now <= 8B"])


def slide_v5(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "v5: NaN-boxed Value", "struct Value -- 8 bytes")

    # Col1: Value + NaN Boxing
    add_group(s, COL1_L, TOP, I(2.2), I(0.9), "struct Value -- 8 bytes")
    add_box(s, COL1_L+I(0.15), TOP+I(0.45), I(1.9), BH, "bits: u64",
            fill=C_INLINE, font_size=10)

    NL = COL1_L
    NT_ = TOP + I(1.3)
    add_group(s, NL, NT_, I(4.2), I(4.0), "IEEE 754 NaN Boxing")
    add_box(s, NL+I(0.15), NT_+I(0.4), I(3.9), BH+I(0.05),
            "not all-1 -> Float (raw f64 inline)", fill=C_INLINE, font_size=10, rounded=True)

    add_group(s, NL+I(0.1), NT_+I(0.9), I(4.0), I(2.9), "Tagged: upper 16 bits",
              fill=C_BG, border_w=Pt(1))
    tags = vstack(s, [
        ("0xFFF8 -> Null / False / True / Undef", C_INLINE),
        ("0xFFFA -> HeapNumber ptr", C_WHITE),
        ("0xFFFC -> String (ArcStr ptr)", C_WHITE),
        ("0xFFFD -> Array ptr", C_WHITE),
        ("0xFFFE -> Object ptr", C_WHITE),
        ("0xFFFF -> Set ptr", C_WHITE),
    ], NL+I(0.2), NT_+I(1.2), bw=I(3.8), bh=I(0.36))

    # Col2: HeapNumber + ObjectMap
    HL = I(5.0)
    add_group(s, HL, TOP, I(2.5), I(1.7), "enum HeapNumber")
    vstack(s, [
        ("Int(i64)", C_WHITE), ("UInt(u64)", C_WHITE), ("BigInt -> Arc", C_WHITE),
    ], HL+I(0.15), TOP+I(0.4), bw=I(2.2))

    add_group(s, HL, TOP+I(2.0), I(3.5), I(1.9), "struct ObjectMap (unordered)")
    add_box(s, HL+I(0.15), TOP+I(2.4), I(3.2), BH, "strings: HashMap<ArcStr, Value>", font_size=10)
    add_box(s, HL+I(0.15), TOP+I(2.9), I(3.2), BH, "other: Option<HashMap<V, V>>", font_size=9)
    add_box(s, HL+I(0.15), TOP+I(3.35), I(3.2), BH, "cached_hash: u64",
            fill=C_INLINE, font_size=10, rounded=True)

    # Col3: Heap targets + Interner
    str_t = add_box(s, COL3_L, TOP+I(0.2), I(3.5), BH, "str", font_size=10)
    arr_t = add_box(s, COL3_L, TOP+I(0.7), I(3.5), BH, "Vec<Value>", font_size=10)
    set_t = add_box(s, COL3_L, TOP+I(1.2), I(3.5), BH, "HashSet<Value> (unordered)", font_size=10)

    add_group(s, COL3_L, TOP+I(2.0), I(3.8), I(1.0), "Interner")
    add_box(s, COL3_L+I(0.15), TOP+I(2.45), I(3.5), BH, "HashSet<ArcStr> (thread-local)", font_size=10)

    # Arrows from tags
    arrow_right(s, tags[1], Box(None, HL, TOP+I(0.8), 0, 0), "Arc")
    arrow_right(s, tags[2], str_t, "ArcStr")
    arrow_right(s, tags[3], arr_t, "Arc")
    arrow_right(s, tags[4], Box(None, HL, TOP+I(2.7), 0, 0), "Arc")
    arrow_right(s, tags[5], set_t, "Arc")

    add_note(s, COL1_L, TOP+I(5.6), I(6.0), I(0.8),
             ["8B = single u64 -- NaN boxing steals tag bits from f64 quiet-NaN space;",
              "pointers limited to 48 bits"])


def slide_v6(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "v6: Portable tagged pointer", "struct Value -- 8 bytes")

    # Col1: Value + Tags
    add_group(s, COL1_L, TOP, I(2.2), I(0.9), "struct Value -- 8 bytes")
    add_box(s, COL1_L+I(0.15), TOP+I(0.45), I(1.9), BH, "bits: usize", fill=C_INLINE, font_size=10)

    TT = TOP + I(1.3)
    add_group(s, COL1_L, TT, I(3.5), I(2.3), "Low 3-bit tags (portable)")
    tags = vstack(s, [
        ("0b001 -> u61 unsigned inline", C_INLINE),
        ("0b010 -> u61 negative inline", C_INLINE),
        ("0b011 -> Null / False / True / Undef", C_INLINE),
        ("0b000 -> Arc<HeapValue>", C_WHITE),
    ], COL1_L+I(0.15), TT+I(0.4), bw=I(3.2), bh=I(0.4))

    # Col2: HeapValue
    HL = I(4.5)
    add_group(s, HL, TOP, I(3.2), I(4.0), "enum HeapValue")
    hv = vstack(s, [
        ("Float(f64)", C_WHITE),
        ("BigInt -> Arc", C_WHITE),
        ("LargeUInt(u64)", C_WHITE),
        ("LargeInt(i64)", C_WHITE),
        ("String(ArcStr)", C_WHITE),
        ("Array(Vec<Value>)", C_WHITE),
        ("Object(ObjectMap)", C_WHITE),
        ("Set(HashSet<Value>) -- unordered", C_WHITE),
    ], HL+I(0.15), TOP+I(0.4), bw=I(2.9), bh=I(0.4))

    # Col3: ObjectMap + ArcStr + Interner
    OL = COL3_L
    add_group(s, OL, TOP, I(4.0), I(1.9), "struct ObjectMap (unordered)")
    add_box(s, OL+I(0.15), TOP+I(0.4), I(3.7), BH, "strings: HashMap<ArcStr, Value>", font_size=10)
    add_box(s, OL+I(0.15), TOP+I(0.9), I(3.7), BH, "other: Option<HashMap<V, V>>", font_size=9)
    add_box(s, OL+I(0.15), TOP+I(1.35), I(3.7), BH, "cached_hash: u64",
            fill=C_INLINE, font_size=10, rounded=True)

    add_group(s, OL, TOP+I(2.2), I(3.5), I(1.0), "ArcStr -- 8 bytes")
    add_box(s, OL+I(0.15), TOP+I(2.65), I(3.2), BH, "thin pointer to heap string",
            fill=C_INLINE, font_size=10, rounded=True)

    add_group(s, OL, TOP+I(3.5), I(3.8), I(1.0), "Interner")
    add_box(s, OL+I(0.15), TOP+I(3.95), I(3.5), BH, "HashSet<ArcStr> (thread-local)", font_size=10)

    arrow_right(s, tags[3], Box(None, HL, hv[0].cy, 0, 0), "Arc")
    arrow_right(s, hv[6], Box(None, OL, TOP+I(0.9), 0, 0), "", dashed=True)
    arrow_right(s, hv[4], Box(None, OL, TOP+I(2.65)+BH//2, 0, 0), "", dashed=True)

    add_note(s, COL1_L, TOP+I(5.0), I(6.5), I(0.8),
             ["8B = single usize -- low 3 alignment bits of pointers used as tag;",
              "portable -- no VA-width assumptions"])


def slide_v7(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "v7: Tagged pointer + schema-shared objects", "struct Value -- 8 bytes")

    # Col1: Value + Tags
    add_group(s, COL1_L, TOP, I(2.2), I(0.9), "struct Value -- 8 bytes")
    add_box(s, COL1_L+I(0.15), TOP+I(0.45), I(1.9), BH, "bits: usize", fill=C_INLINE, font_size=10)

    TT = TOP + I(1.2)
    add_group(s, COL1_L, TT, I(3.2), I(2.1), "Low 3-bit tags (portable)")
    tags = vstack(s, [
        ("0b001 -> u61 unsigned", C_INLINE),
        ("0b010 -> u61 negative", C_INLINE),
        ("0b011 -> Null / Bool / Undef", C_INLINE),
        ("0b000 -> Arc<HeapValue>", C_WHITE),
    ], COL1_L+I(0.15), TT+I(0.35), bw=I(2.9), bh=I(0.37))

    # HeapValue (between col1 and col2)
    HL = I(4.0)
    add_group(s, HL, TOP, I(2.8), I(3.8), "enum HeapValue")
    hv = vstack(s, [
        ("Float(f64)", C_WHITE), ("BigInt -> Arc", C_WHITE),
        ("LargeUInt(u64)", C_WHITE), ("LargeInt(i64)", C_WHITE),
        ("String(ArcStr)", C_WHITE), ("Array(Vec<Value>)", C_WHITE),
        ("Object(ObjectMap)", C_WHITE), ("Set(HashSet) -- unord.", C_WHITE),
    ], HL+I(0.15), TOP+I(0.4), bw=I(2.5), bh=I(0.37))

    # Col3: ObjectMap with Compact + Map
    OL = I(7.2)
    add_group(s, OL, TOP, I(5.5), I(4.5), "struct ObjectMap")

    add_group(s, OL+I(0.1), TOP+I(0.4), I(5.3), I(1.7),
              "Compact (sorted via Schema)", fill=C_BG, border_w=Pt(1))
    add_box(s, OL+I(0.2), TOP+I(0.75), I(5.1), BH, "schema: Arc<Schema>", font_size=10)
    add_box(s, OL+I(0.2), TOP+I(1.2), I(5.1), BH, "values: Box<[Value]>", font_size=10)
    add_box(s, OL+I(0.2), TOP+I(1.65), I(5.1), BH, "cached_hash: u64",
            fill=C_INLINE, font_size=10, rounded=True)

    add_group(s, OL+I(0.1), TOP+I(2.2), I(5.3), I(2.0),
              "Map -- fallback (unordered)", fill=C_BG, border_w=Pt(1))
    add_box(s, OL+I(0.2), TOP+I(2.55), I(5.1), BH, "strings: HashMap<ArcStr, Value>", font_size=10)
    add_box(s, OL+I(0.2), TOP+I(3.0), I(5.1), BH, "other: Option<HashMap<Value, Value>>", font_size=10)
    add_box(s, OL+I(0.2), TOP+I(3.45), I(5.1), BH, "cached_hash: u64",
            fill=C_INLINE, font_size=10, rounded=True)

    # Bottom row: Schema + Interner + Schema Cache
    BT = TOP + I(4.8)
    add_group(s, OL, BT, I(3.5), I(1.2), "Schema -- shared via Arc")
    add_box(s, OL+I(0.15), BT+I(0.4), I(3.2), BH, "keys: Arc<[ArcStr]>", font_size=10)
    add_box(s, OL+I(0.15), BT+I(0.85), I(3.2), BH, "lookup: HashMap<ArcStr, u32>", font_size=9)

    add_group(s, COL1_L, TOP+I(3.6), I(3.2), I(0.85), "Interner")
    add_box(s, COL1_L+I(0.15), TOP+I(4.0), I(2.9), BH, "HashSet<ArcStr> (thread-local)", font_size=9)

    add_group(s, COL1_L, TOP+I(4.7), I(3.8), I(0.85), "Schema Cache (thread-local)")
    add_box(s, COL1_L+I(0.15), TOP+I(5.1), I(3.5), BH, "HashMap<Vec<ArcStr>, Arc<Schema>>", font_size=9)

    # Arrows
    arrow_right(s, tags[3], Box(None, HL, hv[0].cy, 0, 0), "Arc")
    arrow_right(s, hv[6], Box(None, OL, TOP+I(2.8), 0, 0), "", dashed=True)

    add_note(s, COL1_L, TOP+I(5.8), I(6.5), I(0.7),
             ["8B = single usize -- same tagged pointer as v6;",
              "schema-shared objects reduce heap, not Value size"])


def slide_v8(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "v8: Flattened Number + schema objects", "enum Value -- 16 bytes")

    # Col1: Value (tall -- 9 variants)
    VW = I(3.2)
    add_group(s, COL1_L, TOP, VW, I(4.8), "enum Value -- 16 bytes")
    vb = vstack(s, [
        ("object -- 8B thin ptr", C_WHITE),
        ("string: ArcStr -- 8B thin ptr", C_WHITE),
        ("array -- 8B thin ptr", C_WHITE),
        ("set -- 8B thin ptr", C_WHITE),
        ("uint(u64) -- 8B", C_INLINE),
        ("int(i64) -- 8B", C_INLINE),
        ("float(f64) -- 8B", C_INLINE),
        ("bigint -- 8B thin ptr", C_WHITE),
        ("null / bool / undefined", C_INLINE),
    ], COL1_L+I(0.3), TOP+I(0.45), bw=I(2.6))

    # Col2: ObjectMap with Compact + Map
    OL = COL2_L
    add_group(s, OL, TOP, I(4.2), I(4.2), "struct ObjectMap")

    add_group(s, OL+I(0.1), TOP+I(0.4), I(4.0), I(1.6),
              "Compact (sorted via Schema)", fill=C_BG, border_w=Pt(1))
    add_box(s, OL+I(0.2), TOP+I(0.75), I(3.8), BH, "schema: Arc<Schema>", font_size=10)
    add_box(s, OL+I(0.2), TOP+I(1.2), I(3.8), BH, "values: Box<[Value]>", font_size=10)
    add_box(s, OL+I(0.2), TOP+I(1.6), I(3.8), BH, "cached_hash: u64",
            fill=C_INLINE, font_size=10, rounded=True)

    add_group(s, OL+I(0.1), TOP+I(2.1), I(4.0), I(1.9),
              "Map -- fallback (unordered)", fill=C_BG, border_w=Pt(1))
    add_box(s, OL+I(0.2), TOP+I(2.45), I(3.8), BH, "strings: HashMap<ArcStr, Value>", font_size=10)
    add_box(s, OL+I(0.2), TOP+I(2.9), I(3.8), BH, "other: Option<HashMap<Value, Value>>", font_size=9)
    add_box(s, OL+I(0.2), TOP+I(3.35), I(3.8), BH, "cached_hash: u64",
            fill=C_INLINE, font_size=10, rounded=True)

    # Col3: Schema + Cache + ArcStr + heap targets
    SL = COL3_L
    add_group(s, SL, TOP, I(4.0), I(1.2), "Schema -- shared via Arc")
    add_box(s, SL+I(0.15), TOP+I(0.4), I(3.7), BH, "keys: Arc<[ArcStr]>", font_size=10)
    add_box(s, SL+I(0.15), TOP+I(0.85), I(3.7), BH, "lookup: HashMap<ArcStr, u32>", font_size=10)

    add_group(s, SL, TOP+I(1.5), I(4.2), I(0.9), "Schema Cache (thread-local)")
    add_box(s, SL+I(0.15), TOP+I(1.9), I(3.9), BH, "HashMap<Vec<ArcStr>, Arc<Schema>>", font_size=9)

    add_group(s, SL, TOP+I(2.7), I(3.5), I(0.9), "ArcStr -- 8 bytes")
    add_box(s, SL+I(0.15), TOP+I(3.1), I(3.2), BH, "thin pointer to heap string",
            fill=C_INLINE, font_size=10, rounded=True)

    arr_t = add_box(s, SL, TOP+I(3.9), I(3.5), BH, "Vec<Value>", font_size=10)
    set_t = add_box(s, SL, TOP+I(4.4), I(3.5), BH, "HashSet<Value> (unordered)", font_size=10)
    big_t = add_box(s, SL, TOP+I(4.9), I(2.0), BH, "BigInt", font_size=10)

    # Arrows
    arrow_right(s, vb[0], Box(None, OL, TOP+I(2.7), 0, 0), "Arc")
    arrow_right(s, vb[2], arr_t, "Arc")
    arrow_right(s, vb[3], set_t, "Arc")
    arrow_right(s, vb[7], big_t, "Arc")
    add_arrow(s, vb[1].right, vb[1].cy, SL, TOP+I(3.1)+BH//2, "", dashed=True)

    add_note(s, COL1_L, TOP+I(5.6), I(6.0), I(0.8),
             ["16B = 8B tag + 8B payload",
              "Flattened numbers + ArcStr thin ptr; all variants <= 8B"])


def slide_v9(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "v9: Arena-allocated, Copy Value", "enum Value<'a> -- 16 bytes, Copy")

    # Col1: Value (tall)
    VW = I(3.2)
    add_group(s, COL1_L, TOP, VW, I(5.4), "enum Value<'a> -- 16 bytes, Copy")
    add_group(s, COL1_L+I(0.1), TOP+I(0.4), I(3.0), I(4.4), "Arena variants",
              fill=C_BG, border_w=Pt(1))
    vb = vstack(s, [
        ("object -- 8B ref", C_WHITE),
        ("string: &'a ArenaStr -- 8B", C_WHITE),
        ("array: &'a ArenaArray -- 8B", C_WHITE),
        ("set: &'a ArenaSet (unord.) -- 8B", C_WHITE),
        ("uint(u64) -- 8B", C_INLINE),
        ("int(i64) -- 8B", C_INLINE),
        ("float(f64) -- 8B", C_INLINE),
        ("bigint: &'a BigInt -- 8B", C_WHITE),
        ("null / bool / undefined", C_INLINE),
    ], COL1_L+I(0.2), TOP+I(0.7), bw=I(2.8))
    add_box(s, COL1_L+I(0.1), TOP+I(4.9), I(3.0), BH,
            "Ext(&'a ExtValue) -- 8B ref", font_size=10)

    # Col2: ObjectMap
    OL = COL2_L
    add_group(s, OL, TOP, I(4.2), I(4.5), "struct ObjectMap<'a>")

    add_group(s, OL+I(0.1), TOP+I(0.4), I(4.0), I(1.8),
              "Compact (sorted via Schema)", fill=C_BG, border_w=Pt(1))
    add_box(s, OL+I(0.2), TOP+I(0.75), I(3.8), BH, "schema: Arc<Schema>", font_size=10)
    add_box(s, OL+I(0.2), TOP+I(1.15), I(3.8), BH, "values: &'a [Value<'a>]", font_size=10)
    add_box(s, OL+I(0.2), TOP+I(1.55), I(3.8), BH, "arena_keys: &'a [&'a ArenaStr]", font_size=10)
    add_box(s, OL+I(0.2), TOP+I(1.92), I(3.8), BH-I(0.03), "cached_hash: u64",
            fill=C_INLINE, font_size=10, rounded=True)

    add_group(s, OL+I(0.1), TOP+I(2.3), I(4.0), I(2.0),
              "Map -- fallback (unordered)", fill=C_BG, border_w=Pt(1))
    add_box(s, OL+I(0.2), TOP+I(2.65), I(3.8), BH,
            "strings: HashMap<&'a str, Value, _, &'a Bump>", font_size=9)
    add_box(s, OL+I(0.2), TOP+I(3.1), I(3.8), BH,
            "others: Option<HashMap<V,V,_,&'a Bump>>", font_size=9)
    add_box(s, OL+I(0.2), TOP+I(3.55), I(3.8), BH, "cached_hash: u64",
            fill=C_INLINE, font_size=10, rounded=True)

    # Col3: Schema + Cache + Interner + Bump
    SL = COL3_L
    add_group(s, SL, TOP, I(4.0), I(1.2), "Schema -- shared via Arc")
    add_box(s, SL+I(0.15), TOP+I(0.4), I(3.7), BH, "keys: Arc<[Arc<str>]>", font_size=10)
    add_box(s, SL+I(0.15), TOP+I(0.85), I(3.7), BH, "lookup: HashMap<Arc<str>, u32>", font_size=10)

    add_group(s, SL, TOP+I(1.5), I(4.2), I(0.9), "Schema Cache (thread-local)")
    add_box(s, SL+I(0.15), TOP+I(1.9), I(3.9), BH,
            "HashMap<Vec<Arc<str>>, Arc<Schema>>", font_size=9)

    add_group(s, SL, TOP+I(2.7), I(3.8), I(0.9), "StringInterner<'a> (arena-backed)")
    add_box(s, SL+I(0.15), TOP+I(3.1), I(3.5), BH, "HashSet<&'a str>", font_size=10)

    add_group(s, SL, TOP+I(3.9), I(4.0), I(1.4), "bumpalo::Bump arena")
    add_box(s, SL+I(0.15), TOP+I(4.3), I(3.7), I(0.7),
            "all Value data lives here\nbulk deallocation\nzero-cost clone (Copy)",
            fill=C_INLINE, font_size=10, rounded=True)

    # Arrows
    arrow_right(s, vb[0], Box(None, OL, TOP+I(2.8), 0, 0), "&'a")

    add_note(s, COL1_L, TOP+I(5.7), I(6.0), I(0.8),
             ["16B = 8B tag + 8B payload",
              "Arena refs (thin ptrs) + niche optimization; Copy -- zero-cost clone"])


# main
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_baseline(prs)
    slide_v1(prs)
    slide_v2(prs)
    slide_v3(prs)
    slide_v4(prs)
    slide_v5(prs)
    slide_v6(prs)
    slide_v7(prs)
    slide_v8(prs)
    slide_v9(prs)

    out = "RegorusValue.pptx"
    prs.save(out)
    print(f"Saved {out} with {len(prs.slides)} slides")

if __name__ == "__main__":
    main()
