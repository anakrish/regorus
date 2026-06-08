#!/usr/bin/env python3
"""
Generate RegorusValue.pptx by:
1. Extracting mermaid blocks from visual.md files
2. Rendering them to SVG via mmdc
3. Embedding SVGs into PowerPoint slides
"""

import os
import re
import subprocess
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BLUE_DARK = RGBColor(0x1A, 0x52, 0x76)
BLUE_BORDER = RGBColor(0x24, 0x71, 0xA3)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

SOURCES = [
    ("value-compare/src/baseline_visual.md", "Baseline: Arc + BTreeMap — 24 bytes"),
    ("value-compare/src/v1/visual.md",       "v1: HashMap + SmolStr — 24 bytes"),
    ("value-compare/src/v2/visual.md",       "v2: Arc‹str› keys — 24 bytes"),
    ("value-compare/src/v3/visual.md",       "v3: HashSet + cached hash — 24 bytes"),
    ("value-compare/src/v4/visual.md",       "v4: ArcStr thin pointer — 16 bytes"),
    ("value-compare/src/v5/visual.md",       "v5: NaN-boxed Value — 8 bytes"),
    ("value-compare/src/v6/visual.md",       "v6: Portable tagged pointer — 8 bytes"),
    ("value-compare/src/v7/visual.md",       "v7: Tagged pointer + schema — 8 bytes"),
    ("value-compare/src/v8/visual.md",       "v8: Flattened Number + schema — 16 bytes"),
    ("value-compare/src/v9/visual.md",       "v9: Arena-allocated, Copy — 16 bytes"),
]


def extract_mermaid(md_path: str) -> str:
    """Extract the mermaid code block from a markdown file."""
    text = Path(md_path).read_text()
    m = re.search(r'```mermaid\s*\n(.*?)```', text, re.DOTALL)
    if not m:
        raise ValueError(f"No mermaid block found in {md_path}")
    return m.group(1).strip()


def render_svg(mermaid_code: str, out_svg: str):
    """Render mermaid code to SVG using mmdc."""
    with tempfile.NamedTemporaryFile(mode='w', suffix='.mmd', delete=False) as f:
        f.write(mermaid_code)
        mmd_path = f.name
    try:
        result = subprocess.run(
            ["mmdc", "-i", mmd_path, "-o", out_svg, "-b", "transparent", "-t", "default"],
            capture_output=True, text=True, timeout=60
        )
        if result.returncode != 0:
            print(f"  mmdc stderr: {result.stderr[:500]}")
            raise RuntimeError(f"mmdc failed for {out_svg}")
    finally:
        os.unlink(mmd_path)


def render_png(mermaid_code: str, out_png: str, scale: int = 4):
    """Render mermaid code to high-res PNG using mmdc."""
    with tempfile.NamedTemporaryFile(mode='w', suffix='.mmd', delete=False) as f:
        f.write(mermaid_code)
        mmd_path = f.name
    try:
        result = subprocess.run(
            ["mmdc", "-i", mmd_path, "-o", out_png, "-b", "white",
             "-t", "default", "-s", str(scale)],
            capture_output=True, text=True, timeout=60
        )
        if result.returncode != 0:
            print(f"  mmdc stderr: {result.stderr[:500]}")
            raise RuntimeError(f"mmdc failed for {out_png}")
    finally:
        os.unlink(mmd_path)


def add_title(slide, text):
    """Add a slide title."""
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12), Inches(0.45))
    tf = tb.text_frame
    tf.word_wrap = True
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(22)
    run.font.color.rgb = BLUE_DARK
    run.font.bold = True


def main():
    workspace = Path(__file__).parent
    svg_dir = workspace / "svg_diagrams"
    svg_dir.mkdir(exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for md_rel, title in SOURCES:
        md_path = workspace / md_rel
        svg_path = svg_dir / (md_rel.replace("/", "_").replace(".md", ".svg"))
        png_path = svg_dir / (md_rel.replace("/", "_").replace(".md", ".png"))

        print(f"Processing {md_rel}...")

        # Extract mermaid
        mermaid_code = extract_mermaid(str(md_path))

        # Render to SVG (for editability — user can open in PowerPoint 365)
        render_svg(mermaid_code, str(svg_path))
        print(f"  → {svg_path.name}")

        # Render to high-res PNG (for PPTX embedding)
        render_png(mermaid_code, str(png_path))
        print(f"  → {png_path.name}")

        # Add slide with PNG
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        add_title(slide, title)

        slide.shapes.add_picture(
            str(png_path),
            Inches(0.3), Inches(0.7),
            Inches(12.7), Inches(6.5)
        )

    out = workspace / "RegorusValue.pptx"
    prs.save(str(out))
    print(f"\nSaved {out} with {len(prs.slides)} slides")
    print(f"SVGs also available in {svg_dir}/")


if __name__ == "__main__":
    main()
